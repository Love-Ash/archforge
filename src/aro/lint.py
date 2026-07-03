# -*- coding: utf-8 -*-
"""
aro(아로): 빌드된 .pptx를 받아 반복 결함을 기계로 차단하는 한글 특화 품질 린터.

존재 이유: 규칙을 문서에 적어둬도 만들 때마다 같은 실수가 반복된다. 기계로 검사 가능한 결함은
사람 눈이 아니라 이 린터가 막는다. 한글 덱이 조용히 깨지는 지점(라틴 전용 폰트로의 폴백, CJK 자간,
소형 글자)과 AI 생성 덱 특유의 티(긴 대시, 상투 카피, 레이아웃 재탕)를 빌드 산출물에서 잡는다.
임의 pptx를 받는 최후 방어선이라 상속·슬롯·autofit·표·그룹 같은 우회 경로까지 덮는다(적대 감사 반영).

ERROR = 배포 차단:
  E1 CJK 텍스트를 렌더하는 폰트 슬롯(a:ea 우선, 없으면 a:latin)이 라틴 전용(IBM Plex Mono 등)이거나 그 굵기 변형
  E2 대시류 문자가 렌더 텍스트에 포함(em/en/figure dash + 수학 마이너스 U+2212·전각 하이픈 U+FF0D·박스 罫線 U+2500 등)
  E3 실효 폰트(autofit fontScale·문단 상속 반영)가 HARD_MIN(기본 5.0pt) 미만 = 판독 불가
  E4 연속 CJK(2자 이상)에 유의미 양수 트래킹(spc>50 = 0.5pt) = 자간 벌어짐
WARN = 참고(배포는 통과):
  W1 본문급 프레임(넓고 글자수 많음)의 실효 폰트가 BODY_MIN(기본 9.0pt) 미만 = 본문이 작다
  W8 좁은 프레임의 소형 CJK(실효 [HARD_MIN, SMALL_MIN=7.5)pt) = 기기 목업·카드 내부 텍스트 판독 위험
     (E3 판독불가 하한 위, W1 본문급 아래의 회색지대. 폰 목업 안 서브텍스트를 프리플라이트에 가시화)
  W6 레이아웃 골격 재탕: 도형 bbox 시그니처 코사인>0.90 클러스터 4장+(렌더 불요, 항상 검사)
  W7 이미지 위 텍스트 저대비 ratio<2.5 (--render <pages> 지정 시에만, 렌더 PNG 픽셀 대조)
  W9 accent 세로바 3개+를 리스트 마커로 반복 = 색으로 구조 잡기(클로드 티). 커넥터·제로폭 커버
  W10 직접 그린 도식(단면도 등 장식 텍스처)이 여러 페이지에서 거의 동일 반복 = 재탕인지 의도인지 눈 확정
  W11 AI 티 카피: 버즈워드(좁은 사전, 전 페이지)·뻔한 오프닝 상투구(p1~3만)
  W12 푸터 baseline 어긋남: 지배 baseline(표지 제외, 0.05in 버킷 중앙값) 대비 0.03~0.25in 벗어난 페이지
  W13 PPT 자체 그림자·글로·3D 효과 2개+(올드 티. 자식 없는 빈 effectLst 는 안 셈)
  W14 타이틀 다수(3+, 절반+)가 서술형 명사구 = 액션타이틀 아님. --ghost 로 고스트덱(타이틀 나열)도 출력
  W15 텍스트끼리 겹침 추정(실효 글리프 bbox 근사, 교집>작은쪽 45%) = 가림·충돌, 페이지당 최대 2건
  W16 화면 밖 넘침: 텍스트=실효 글리프 box 경계 밖 0.15in+, 그림=잉크 bbox(알파 트림) 0.12in+
      (풀블리드 70%+ 제외). 장식 도형의 모서리 블리드는 표준 기법이라 검사 대상 아님(렌더 실측 기각)
  W17 텍스트가 비배경 그림의 잉크 경계에 걸침(글리프의 25~75%만 안) = 잘려 보임. 완전 위는 W7 소관
  W5 폰트 크기를 run·문단 어디에서도 못 찾음(상속 미해석 = 조용한 스킵 방지)

W15~W17 기하 강건성(적대 검증 12건 실측 재현 후 수정, 2026-07-03): 그룹 off/chOff 아핀 변환,
wrap=none(word_wrap=False, python-pptx add_textbox 기본) 한 줄 실폭, autofit 퍼센트 문자열
('62.5%')·lnSpcReduction, 문단별 정렬, 빈 문단 endParaRPr 크기, 회전 도형(축정렬 확장 bbox,
회전 텍스트는 스킵), 그림 srcRect 크롭·flipH/V·P모드 tRNS 투명, W17은 사진-캡션 사이 솔리드
카드(z순서)를 억제. 남은 한계: 플레이스홀더의 레이아웃 lstStyle 정렬 상속은 좌정렬로 후퇴.

사용: aro <built.pptx> [--hard-min 5.0] [--body-min 9.0] [--strict] [--render <pages>] [--ghost] [--json]
  --strict: WARN도 exit 1로 취급. --render: W7 활성화. --ghost: 타이틀 나열(수평 논리 눈검수).
반환: ERROR가 하나라도 있으면 exit 1, 아니면 0.
"""
import os
import sys
import argparse
import colorsys
from collections import Counter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_IN = 914400
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# 한글 글리프가 없어 한글을 주면 폴백되는 라틴 전용 폰트(접두 매칭이라 굵기 변형도 커버)
LATIN_ONLY_FONTS = (
    "ibm plex mono", "courier new", "consolas", "cascadia mono", "cascadia code",
    "roboto mono", "jetbrains mono", "fira code", "source code pro", "pt mono",
    "dejavu sans mono", "dejavu sans", "menlo", "monaco", "sf mono", "space mono",
    "nanumgothiccoding", "switzer",
)

# 긴 대시류 + 전각 하이픈마이너스 + 2·3배 em + 박스 수평선(코드포인트로 구성: 소스에 대시 문자 없음)
LONG_DASHES = {chr(c) for c in (0x2012, 0x2013, 0x2014, 0x2015, 0x2212, 0xFF0D, 0x2E3A, 0x2E3B, 0x2500)}


def is_cjk(ch):
    o = ord(ch)
    return (
        0xAC00 <= o <= 0xD7A3        # 한글 음절
        or 0x1100 <= o <= 0x11FF     # 한글 자모
        or 0x3130 <= o <= 0x318F     # 호환 자모
        or 0x3040 <= o <= 0x30FF     # 가나
        or 0x4E00 <= o <= 0x9FFF     # 한자
    )


def has_cjk(text):
    return any(is_cjk(c) for c in text)


def run_fonts(run):
    """run rPr의 a:latin/a:ea/a:cs typeface. run.font.name은 a:latin만 보므로 XML을 직접 읽는다."""
    out = {}
    rPr = run._r.find(NS + "rPr")
    if rPr is not None:
        for slot in ("latin", "ea", "cs"):
            el = rPr.find(NS + slot)
            if el is not None:
                tf = el.get("typeface")
                if tf:
                    out[slot] = tf
    return out


def run_track(run):
    """run 의 트래킹(spc, 1/100 pt 단위 정수). 없으면 None."""
    rPr = run._r.find(NS + "rPr")
    if rPr is None:
        return None
    spc = rPr.get("spc")
    return int(spc) if spc is not None else None


def _pct_attr(v, default):
    """OOXML 퍼센트 유니언 타입: '62500'(1/1000 %)과 '62.5%'(문자열형) 둘 다 유효
    (ST_TextFontScalePercentOrPercentString). int() 단독은 후자에서 ValueError로 죽어
    blanket except가 1.0으로 삼켰다(적대 검증 실측 2026-07-03)."""
    if not v:
        return default
    v = v.strip()
    if v.endswith("%"):
        return float(v[:-1]) / 100.0
    return int(v) / 100000.0


def frame_autofit(tf):
    """(fontScale, lnSpcReduction) 비율 쌍. normAutofit 없으면 (1.0, 0.0)."""
    try:
        bodyPr = tf._txBody.find(NS + "bodyPr")
        if bodyPr is not None:
            na = bodyPr.find(NS + "normAutofit")
            if na is not None:
                return (_pct_attr(na.get("fontScale"), 1.0),
                        _pct_attr(na.get("lnSpcReduction"), 0.0))
    except Exception:
        pass
    return 1.0, 0.0


def frame_font_scale(tf):
    """텍스트프레임 autofit fontScale(비율). E3 등 기존 소비자 호환 유지."""
    return frame_autofit(tf)[0]


def collect_frames(shapes):
    """(text_frame, width_emu) 리스트. 그룹 재귀 + 네이티브 표 셀 포함."""
    out = []
    for sp in shapes:
        try:
            st = sp.shape_type
        except Exception:
            st = None
        if st == MSO_SHAPE_TYPE.GROUP:
            out += collect_frames(sp.shapes)
            continue
        if getattr(sp, "has_table", False):
            tbl = sp.table
            ncol = len(tbl.columns) or 1
            for row in tbl.rows:
                for cell in row.cells:
                    out.append((cell.text_frame, (sp.width or 0) // ncol))
            continue
        if sp.has_text_frame:
            out.append((sp.text_frame, sp.width or 0))
    return out


def iter_shapes(shapes):
    """그룹 재귀로 모든 도형을 평탄하게 순회."""
    for sp in shapes:
        try:
            if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
                for inner in iter_shapes(sp.shapes):
                    yield inner
                continue
        except Exception:
            pass
        yield sp


def _is_pic(sp):
    try:
        return sp.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


# W15~W17 기하 소비자용 절대좌표 순회. 그룹 자식의 raw left/top은 그룹 chOff 좌표계라
# 그룹이 이동·리사이즈된 pptx(off!=chOff desync, PowerPoint 드래그 표준 동작)에서
# 슬라이드 좌표와 어긋난다(적대 검증 실측 2026-07-03). off/ext vs chOff/chExt 아핀을
# 합성해 (도형, z순서, 절대 xfrm 함수 계수)로 순회한다. xf=(ax,bx,ay,by): abs = a*raw + b (EMU).
def iter_shapes_geo(shapes, xf=(1.0, 0.0, 1.0, 0.0), _z=None):
    if _z is None:
        _z = [0]
    for sp in shapes:
        try:
            is_grp = sp.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_grp = False
        if is_grp:
            sub = xf
            try:
                x = sp._element.find(NS + "grpSpPr").find(NS + "xfrm")
                off, ext = x.find(NS + "off"), x.find(NS + "ext")
                cho, che = x.find(NS + "chOff"), x.find(NS + "chExt")
                ox, oy = int(off.get("x")), int(off.get("y"))
                ew, eh = int(ext.get("cx")), int(ext.get("cy"))
                cx, cy = int(cho.get("x")), int(cho.get("y"))
                cw_, ch_ = int(che.get("cx")) or ew, int(che.get("cy")) or eh
                sx, sy = ew / float(cw_), eh / float(ch_)
                # 그룹 내부 좌표 g: abs_g = ox + (g - cx)*sx 를 부모 xf에 합성
                ax, bx, ay, by = xf
                sub = (ax * sx, ax * (ox - cx * sx) + bx,
                       ay * sy, ay * (oy - cy * sy) + by)
            except Exception:
                pass
            for inner in iter_shapes_geo(sp.shapes, sub, _z):
                yield inner
            continue
        z = _z[0]
        _z[0] += 1
        yield sp, z, xf


def _geo_rect(sp, xf):
    """xf 적용한 절대 bbox(in). 회전 도형은 중심 고정 축정렬 bbox로 확장. 실패 시 None."""
    try:
        L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
    except Exception:
        return None
    if None in (L, T, Wd, Ht):
        return None
    ax, bx, ay, by = xf
    x = (ax * L + bx) / EMU_PER_IN
    y = (ay * T + by) / EMU_PER_IN
    w = ax * Wd / EMU_PER_IN
    h = ay * Ht / EMU_PER_IN
    rot = 0.0
    try:
        rot = float(sp.rotation or 0.0)
    except Exception:
        pass
    if rot % 360.0:
        import math
        r = math.radians(rot)
        w2 = abs(w * math.cos(r)) + abs(h * math.sin(r))
        h2 = abs(w * math.sin(r)) + abs(h * math.cos(r))
        x, y, w, h = x + (w - w2) / 2, y + (h - h2) / 2, w2, h2
    return x, y, w, h, (rot % 360.0 != 0.0)


def _shape_fill_hex(sp):
    """단색 채움 hex(6자리 대문자) 또는 None. python-pptx 접근자라 네임스페이스·커넥터 내부를 알아서 처리."""
    try:
        f = sp.fill
        if f.type == 1:   # MSO_FILL_TYPE.SOLID
            c = f.fore_color
            if c.type is not None and c.rgb is not None:
                return str(c.rgb).upper()
    except Exception:
        pass
    return None


def _shape_line_hex(sp):
    """선/윤곽 색 hex 또는 None. 커넥터(cxnSp)의 line 색까지 포함(실덱의 커넥터 세로바가 이 경로로 통과했던 실측)."""
    try:
        c = sp.line.color
        if c.type is not None and c.rgb is not None:
            return str(c.rgb).upper()
    except Exception:
        pass
    return None


def _is_accent(hexc):
    """의미 강조색 판별: HSV 채도>=0.55 그리고 명도 0.18~0.78. 배경·괘선·본문·저채도 보조색은 배제."""
    if not hexc or len(hexc) < 6:
        return False
    try:
        r = int(hexc[0:2], 16) / 255.0
        g = int(hexc[2:4], 16) / 255.0
        b = int(hexc[4:6], 16) / 255.0
    except Exception:
        return False
    mx, mn = max(r, g, b), min(r, g, b)
    sv = 0.0 if mx == 0 else (mx - mn) / mx
    _, l, _ = colorsys.rgb_to_hls(r, g, b)
    return sv >= 0.55 and 0.18 <= l <= 0.78


def accent_vbars_check(slide, si, sw, sh, warns):
    """W9: accent 세로바를 리스트 마커로 반복해 항목 구조를 색으로 잡는 AI 생성물 티(실덱 실측).
    적대 감사(2026-07-02) 반영: 색은 sp.line/sp.fill 접근자로 읽어 네임스페이스·커넥터를 커버,
    세로바는 제로폭 커넥터(w~0)를 명시 포함, 오른쪽 인접 텍스트로 리스트 마커임을 확증."""
    bars, texts = [], []
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht):
            continue
        x, y, w, h = L / EMU_PER_IN, T / EMU_PER_IN, Wd / EMU_PER_IN, Ht / EMU_PER_IN
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            texts.append((x, y, w, h))
        hexc = _shape_line_hex(sp) or _shape_fill_hex(sp)
        if not _is_accent(hexc):
            continue
        if 0.2 <= h <= 1.0 and (w < 0.05 or h > 3 * w):   # 세로 막대(제로폭 커넥터 포함)
            bars.append((x, y, w, h, hexc))
    if len(bars) < 3:
        return
    hues = {b[4] for b in bars}
    if len(hues) != 1:                       # 다색은 데이터 인코딩(범례)이라 정당
        return
    xs = [b[0] for b in bars]
    if max(xs) - min(xs) > 0.15:             # 세로 정렬 스택만(가로 분산=차트·구분선 제외)
        return
    rt = 0
    for bx, by, bw, bh, _hx in bars:
        for tx, ty, tw, th in texts:
            if tx > bx and (tx - (bx + bw)) < 0.6 and not (ty > by + bh or ty + th < by):
                rt += 1
                break
    if rt >= len(bars) - 1:
        warns.append((si, "W9", "accent 세로바 %d개로 항목 구조 반복(구조는 괘선·여백·활자로, accent는 점 하나로)" % len(bars),
                      "x=%.2fin hue=%s" % (min(xs), next(iter(hues)))))


def _fill_tokens(slide, sw, sh):
    """슬라이드의 단색 채움 도형을 (색, 24분할 위치·크기) 토큰 멀티셋으로. 풀블리드 배경은 제외."""
    t = Counter()
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:
            continue
        fh = _shape_fill_hex(sp)
        if fh is None:
            continue
        t[(fh, round(L / sw * 24), round(T / sh * 24), round(Wd / sw * 24), round(Ht / sh * 24))] += 1
    return t


def theme_ea_font(prs):
    """테마 minorFont(본문)의 a:ea typeface. 런에 폰트가 없을 때 CJK가 실제로 폴백되는 슬롯.
    반환: 폰트명 / ""(빈 슬롯 = Windows에서 Malgun 폴백) / None(테마 파싱 실패)."""
    import re as _re
    try:
        for part in prs.part.package.iter_parts():
            pn = str(part.partname)
            if "/theme/" in pn and pn.endswith(".xml"):
                m = _re.search(rb'<a:minorFont>.*?<a:ea[^>]*typeface="([^"]*)"', part.blob, _re.S)
                if m:
                    return m.group(1).decode("utf-8", "replace")
    except Exception:
        pass
    return None


# AI 티 카피: 명백한 클리셰만(좁은 사전 = 오탐 억제). 맥락상 정당할 수 있는 일반어는 안 넣는다.
BUZZWORDS = (
    "시너지", "패러다임", "게임체인저", "게임 체인저", "혁신을 가속", "가치를 극대화",
    "미래를 선도", "새로운 지평", "무한한 가능성", "홀리스틱", "엔드투엔드", "엔드 투 엔드",
    "초격차", "글로벌 리더로 도약", "위대한 여정",
    "synergy", "paradigm shift", "game-changer", "game changer", "cutting-edge",
    "state-of-the-art", "seamless", "revolutionize", "leverage synerg", "holistic",
    "unlock the potential", "empower",
)
STALE_OPENINGS = (
    "오늘날", "급변하는", "4차 산업혁명 시대", "바야흐로", "현대 사회에서",
    "디지털 전환의 시대", "디지털 전환의 물결", "알아보겠습니다", "살펴보겠습니다",
    "in today's", "in the rapidly changing", "in this presentation",
)


def copy_cliche_check(page_texts, warns):
    """W11: AI 티 카피. 버즈워드는 전 페이지, 뻔한 오프닝은 도입부(p1~3)만."""
    for si in sorted(page_texts):
        blob = " ".join(page_texts[si])
        low = blob.lower()
        hits = sorted({b for b in BUZZWORDS if b.lower() in low})
        if hits:
            warns.append((si, "W11", "AI 티 버즈워드 %d종(카피 재작성 검토)" % len(hits),
                          ", ".join(hits[:5])))
        if si <= 3:
            op = sorted({o for o in STALE_OPENINGS if o.lower() in low})
            if op:
                warns.append((si, "W11", "뻔한 오프닝 상투구(표지·도입은 자기 목소리로)",
                              ", ".join(op[:5])))


def footer_top(slide, sw, sh):
    """슬라이드 하단 밴드(y>0.88H)의 최하단 텍스트 top(in). 푸터 없으면 None."""
    best = None
    for sp in iter_shapes(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        try:
            if not sp.text_frame.text.strip():
                continue
            t = sp.top
        except Exception:
            continue
        if t is None or t <= 0.88 * sh:
            continue
        ti = t / EMU_PER_IN
        if best is None or ti > best:
            best = ti
    return best


def footer_check(foot_tops, warns):
    """W12: 페이지간 푸터 baseline 어긋남. 50덱 실측(2026-07-02)에서 절대 편차 방식은 표지
    크레딧·하단 캡션을 푸터로 오인해 17덱 오탐 → 표지(p1) 제외 + 지배 baseline(0.05in 양자화
    최빈, 3페이지+)을 하우스 푸터로 보고, 거기서 '살짝'(0.03~0.25in) 어긋난 페이지만 지목.
    0.25in 넘게 다른 건 캡션·디바이더 등 다른 요소로 보고 무시(존재 여부도 안 본다)."""
    tops = [(si, t) for si, t in foot_tops.items() if t is not None and si > 1]
    if len(tops) < 4:
        return
    q = Counter(round(t / 0.05) for _si, t in tops)
    qv, cnt = q.most_common(1)[0]
    if cnt < 3:
        return
    # 하우스 기준 = 지배 버킷의 실제 값 중앙값(버킷 중심을 쓰면 7.08 다수가 7.10과 딱 0.02
    # 차이로 부동소수점 경계 오탐, 50덱 재스캔서 확인)
    bvals = sorted(t for _si, t in tops if round(t / 0.05) == qv)
    base = bvals[len(bvals) // 2]
    off = [(si, t) for si, t in tops if 0.03 < abs(t - base) <= 0.25]
    if off:
        ex = " ".join("p%d=%.2f" % (si, t) for si, t in off[:4])
        warns.append((0, "W12", "푸터 baseline 어긋남: 하우스 %.2fin 대비 살짝 다른 페이지 %d개(정렬 실수 의심)" % (base, len(off)),
                      ex))


_EFFECT_TAGS = tuple(NS + t for t in ("outerShdw", "innerShdw", "glow", "reflection"))
_3D_TAGS = tuple(NS + t for t in ("sp3d", "scene3d"))


def effects_count(slide):
    """슬라이드의 실효 PPT 효과(그림자·글로·3D) 수와 종류. 일부 생성기는 상속 차단용으로
    자식 없는 빈 effectLst만 남기므로 실제 효과 자식이 있어야만 센다(빈 요소 오탐 방지)."""
    n = 0
    kinds = set()
    for sp in iter_shapes(slide.shapes):
        spPr = getattr(sp._element, "spPr", None)
        if spPr is None:
            continue
        eff = spPr.find(NS + "effectLst")
        if eff is not None:
            for ch in eff:
                if ch.tag in _EFFECT_TAGS:
                    n += 1
                    kinds.add(ch.tag.split("}")[1])
        for tag in _3D_TAGS:
            if spPr.find(tag) is not None:
                n += 1
                kinds.add(tag.split("}")[1])
    return n, kinds


def effects_check_deck(per_page, warns):
    """W13: 덱 단위 1회 집계(페이지별 반복 발화는 소음: 50덱 코퍼스 실측).
    의도적 네온·글로 스타일일 수 있어 WARN, 판단은 눈."""
    hits = [(si, n, kinds) for si, (n, kinds) in per_page.items() if n >= 2]
    if not hits:
        return
    total = sum(n for _si, n, _k in hits)
    kinds = sorted(set().union(*[k for _si, _n, k in hits]))
    pages = ",".join("p%d" % si for si, _n, _k in hits[:6])
    warns.append((0, "W13", "PPT 자체 효과 %d개 / %d개 페이지(그림자·글로·3D = 올드 티 의심, 의도한 스타일이면 무시)" % (total, len(hits)),
                  "%s | %s" % (pages, ",".join(kinds))))


# W15 텍스트 겹침: 생성 덱에서 가장 흔한 결함 축(수정 라운드마다 요소가 포개진다)인데
# 프레임 bbox 는 넉넉히 잡는 관례라 못 쓰고, 실효 글리프 폭을 근사해 잡는다.
_W_CJK, _W_LAT, _W_SP = 0.96, 0.52, 0.28   # 글자폭/폰트크기 비 근사(보수적: 오탐 억제)


def _glyph_w(s, size_pt):
    w = 0.0
    for ch in s:
        if ch == " ":
            w += _W_SP
        elif is_cjk(ch) or ord(ch) > 0x2E80:
            w += _W_CJK
        else:
            w += _W_LAT
    return w * size_pt / 72.0


def _empty_para_pt(para, default_pt):
    """빈 문단(스페이서)의 실효 크기: endParaRPr/defRPr sz 우선(4pt 스페이서를 12pt로
    통산해 팬텀 높이를 만들던 것의 교정, 적대 검증 실측)."""
    try:
        if para.font.size is not None:
            return para.font.size.pt
    except Exception:
        pass
    try:
        epr = para._p.find(NS + "endParaRPr")
        if epr is not None and epr.get("sz"):
            return int(epr.get("sz")) / 100.0
    except Exception:
        pass
    return default_pt


def _text_glyph_boxes(slide, default_pt=12.0):
    """문단 단위 실효 글리프 bbox(in) 근사. 반환 [(x0,y0,x1,y1,대표텍스트,max_pt,frame_id)].
    폭은 run별 실제 크기로 합산, 행간은 line_spacing 실값(없으면 1.2)에 autofit
    lnSpcReduction 반영. 문단별 정렬(명시 없으면 프레임 첫 명시값, 그것도 없으면 좌)로
    x를 잡아 혼합 정렬 프레임의 오배치를 줄인다. wrap=none(word_wrap=False, python-pptx
    add_textbox 기본)은 한 줄로 프레임 밖까지 뻗으므로 랩 접기 없이 실폭 그대로.
    회전 프레임은 추정 무효라 스킵. 그룹은 iter_shapes_geo로 절대좌표 변환.
    실덱 렌더 대조 + 적대 검증 실측 재현으로 캘리브레이션.
    알려진 한계: 플레이스홀더가 레이아웃 lstStyle로 정렬을 상속하는 경우는 좌정렬로 후퇴."""
    import math
    out = []
    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        fx, fy, fw, fh, rotated = geo
        if rotated:
            continue
        fw = max(fw, 0.05)
        tf = sp.text_frame
        scale, lnred = frame_autofit(tf)
        wrap = tf.word_wrap is not False   # None(속성 없음)=OOXML 기본 square=랩
        frame_align = None
        for para in tf.paragraphs:
            if para.alignment is not None:
                frame_align = para.alignment
                break
        paras = []   # (pw, pmx, ptxt, factor, n, align)
        for para in tf.paragraphs:
            pw, pmx, ptxt = 0.0, 0.0, ""
            for r in para.runs:
                t = r.text
                if not t:
                    continue
                sz = r.font.size.pt if r.font.size is not None else \
                    (para.font.size.pt if para.font.size is not None else default_pt)
                sz *= scale
                pw += _glyph_w(t, sz)
                ptxt += t
                if sz > pmx:
                    pmx = sz
            if not ptxt.strip():
                paras.append((0.0, _empty_para_pt(para, default_pt) * scale, "", 1.2, 1, None))
                continue
            ls = para.line_spacing
            if ls is None:
                factor = 1.2
            elif isinstance(ls, float):
                factor = ls
            else:
                try:
                    factor = ls.pt / pmx if pmx else 1.2
                except Exception:
                    factor = 1.2
            n = max(1, math.ceil(pw / (fw * 1.04))) if wrap else 1
            al = para.alignment if para.alignment is not None else frame_align
            paras.append((pw, pmx, ptxt, factor, n, al))
        gh_total = sum(n * pmx * max(f, 0.95) * (1.0 - lnred) / 72.0
                       for (pw, pmx, ptxt, f, n, al) in paras)
        if gh_total <= 0:
            continue
        va = str(tf.vertical_anchor) if tf.vertical_anchor is not None else ""
        if "MIDDLE" in va:
            cy = fy + max(0.0, (fh - gh_total) / 2)
        elif "BOTTOM" in va:
            cy = fy + max(0.0, fh - gh_total)
        else:
            cy = fy
        fid = id(sp)
        for (pw, pmx, ptxt, factor, n, al) in paras:
            ph = n * pmx * max(factor, 0.95) * (1.0 - lnred) / 72.0
            if ptxt.strip():
                gw = pw if not wrap else min(pw, fw)
                a = str(al) if al is not None else ""
                if "CENTER" in a:
                    x0 = fx + (fw - gw) / 2
                elif "RIGHT" in a:
                    x0 = fx + fw - gw
                else:
                    x0 = fx
                out.append((x0, cy, x0 + gw, cy + ph, ptxt[:24], pmx, fid))
            cy += ph
    return out


def text_overlap_check(slide, si, warns):
    """W15: 서로 다른 텍스트 프레임의 실효 글리프 영역이 유의미하게 포개짐(가림·충돌).
    근사 기반이라 WARN. 교집 면적이 작은 쪽의 45% 초과일 때만, 페이지당 최대 2건.
    임계 45%는 렌더 대조 실측: 30~35%대는 전부 오탐(빅넘버 아래 타이틀·2단 사이 침범 추정),
    60%+는 전부 실겹침(차트가 엑서빗 라벨 침범, 캡션 칩이 범례 위 등).
    의도적 레이어는 제외: 동일 텍스트 에코(잔상 타이포 연출), 1~2자 대형 글자(드롭캡·장 번호 I·II 류)."""
    boxes = _text_glyph_boxes(slide)
    hits = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            if a[6] == b[6]:      # 같은 프레임의 문단끼리는 스택이라 제외
                continue
            if a[4].strip() == b[4].strip():
                continue
            if any(len(x[4].strip()) <= 2 and x[5] >= 28 for x in (a, b)):
                continue
            ix = min(a[2], b[2]) - max(a[0], b[0])
            iy = min(a[3], b[3]) - max(a[1], b[1])
            if ix <= 0.02 or iy <= 0.02:
                continue
            area = ix * iy
            amin = min((a[2] - a[0]) * (a[3] - a[1]), (b[2] - b[0]) * (b[3] - b[1]))
            if amin > 0 and area > 0.45 * amin:
                hits.append((area / amin, a[4], b[4]))
    for frac, ta, tb in sorted(hits, reverse=True)[:2]:
        warns.append((si, "W15", "텍스트끼리 겹침 추정 %.0f%%(가림·충돌, 렌더 확인)" % (frac * 100),
                      "%r ~ %r" % (ta, tb)))


def _pic_boxes(slide, sw_in, sh_in):
    """비배경 그림의 실효 잉크 bbox(in)와 z순서. 반환 [(x0,y0,x1,y1,z)].
    슬라이드의 70%+를 덮는 풀블리드·mesh 배경은 제외. 투명 PNG(matplotlib 차트 등)는
    프레임 bbox가 잉크보다 훨씬 커서 오탐 원천 → 알파 불투명 bbox로 트림.
    적대 검증 실측 반영(2026-07-03): P모드+tRNS는 RGBA 변환 후 트림, srcRect 크롭은
    보이는 소스 창으로 좁혀 매핑, flipH/flipV는 창 안에서 미러, 회전 그림은 축정렬
    확장 bbox를 쓰되 잉크 트림은 무효라 건너뜀, 그룹 자식은 절대좌표 변환."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, rotated = geo
        if w * h >= 0.7 * sw_in * sh_in:
            continue
        if not rotated:
            try:
                from PIL import Image
                import io as _io
                im = Image.open(_io.BytesIO(sp.image.blob))
                if im.mode == "P" and "transparency" in im.info:
                    im = im.convert("RGBA")
                if "A" in im.getbands():
                    bb = im.getchannel("A").point(lambda a: 255 if a > 16 else 0).getbbox()
                    if bb is None:
                        continue   # 전면 투명 = 잉크 없음
                    iw, ih = im.size
                    wl, wt = iw * float(sp.crop_left or 0), ih * float(sp.crop_top or 0)
                    wr = iw * (1.0 - float(sp.crop_right or 0))
                    wb = ih * (1.0 - float(sp.crop_bottom or 0))
                    l, t, r, b = bb
                    l, r = max(l, wl), min(r, wr)
                    t, b = max(t, wt), min(b, wb)
                    if l >= r or t >= b:
                        continue   # 크롭 창 안에 잉크 없음
                    try:
                        x2 = sp._element.spPr.find(NS + "xfrm")
                        if x2 is not None and x2.get("flipH") == "1":
                            l, r = wl + (wr - r), wl + (wr - l)
                        if x2 is not None and x2.get("flipV") == "1":
                            t, b = wt + (wb - b), wt + (wb - t)
                    except Exception:
                        pass
                    ww, wh = (wr - wl) or 1.0, (wb - wt) or 1.0
                    x, y, w, h = (x + w * (l - wl) / ww, y + h * (t - wt) / wh,
                                  w * (r - l) / ww, h * (b - t) / wh)
            except Exception:
                pass
        out.append((x, y, x + w, y + h, z))
    return out


def overflow_check(slide, si, sw_in, sh_in, warns):
    """W16: 화면 밖 넘침. 프레임 bbox 기준은 넉넉한 프레임 관례 때문에 오탐이 커서 기각했었지만
    W15의 실효 글리프 bbox가 그 반론을 해소했다(2026-07-03): 텍스트는 실제 글자 영역이 경계를
    뚫을 때만. 비텍스트는 그림(잉크 bbox 트림)만 본다: 장식 도형(글로우 원 등)을 모서리로
    흘리는 블리드는 표준 기법이라 결함이 아님(코퍼스 렌더 실측으로 도형 검사는 기각)."""
    TOL_T, TOL_S = 0.15, 0.12
    hits = []
    for (x0, y0, x1, y1, rep, _mx, _fid) in _text_glyph_boxes(slide):
        over = max(-x0, -y0, x1 - sw_in, y1 - sh_in)
        if over > TOL_T:
            hits.append((over, "텍스트 %r" % rep))
    for (px0, py0, px1, py1, _z) in _pic_boxes(slide, sw_in, sh_in):
        over = max(-px0, -py0, px1 - sw_in, py1 - sh_in)
        if over > TOL_S:
            hits.append((over, "그림 %.1fx%.1fin" % (px1 - px0, py1 - py0)))
    for over, what in sorted(hits, reverse=True)[:2]:
        warns.append((si, "W16", "화면 밖 넘침 %.2fin(잘림, 렌더 확인)" % over, what))


def _occluder_boxes(slide, sw_in, sh_in):
    """그림 위에 얹힌 솔리드 채움 도형(카드·패널) bbox와 z. 사진 위 카드에 캡션이 앉은
    합법 레이아웃이 W17에 걸리던 것의 억제용(적대 검증 실측)."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            continue
        try:
            if sp.fill.type is None or "SOLID" not in str(sp.fill.type):
                continue
        except Exception:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, _rot = geo
        if w * h >= 0.9 * sw_in * sh_in:
            continue
        out.append((x, y, x + w, y + h, z))
    return out


def text_image_straddle_check(slide, si, sw_in, sh_in, warns):
    """W17: 텍스트가 비배경 그림의 잉크 경계에 걸침(글리프의 25~75%만 이미지 안) = 반은 위·
    반은 밖이라 잘리거나 배경이 갈라져 보임. 완전 위(오버레이 캡션)는 W7 대비 게이트 소관,
    완전 밖은 무관. 1제곱인치 미만(아이콘·로고급) 그림 무시. 사진과 텍스트 사이 z에 솔리드
    카드가 끼어 텍스트 영역의 90%+를 받치고 있으면 걸침이 아니라 카드 위 캡션이라 제외.
    페이지당 최대 2건."""
    pics = _pic_boxes(slide, sw_in, sh_in)
    if not pics:
        return
    occl = _occluder_boxes(slide, sw_in, sh_in)
    hits = []
    for (x0, y0, x1, y1, rep, _mx, _fid) in _text_glyph_boxes(slide):
        if len(rep.strip()) < 3:
            continue
        ta = (x1 - x0) * (y1 - y0)
        if ta <= 0:
            continue
        for (px0, py0, px1, py1, pz) in pics:
            if (px1 - px0) * (py1 - py0) < 1.0:
                continue
            ix = min(x1, px1) - max(x0, px0)
            iy = min(y1, py1) - max(y0, py0)
            if ix <= 0 or iy <= 0:
                continue
            frac = (ix * iy) / ta
            if not (0.25 <= frac <= 0.75):
                continue
            carded = False
            for (ox0, oy0, ox1, oy1, oz) in occl:
                if oz <= pz:
                    continue
                cx = min(x1, ox1) - max(x0, ox0)
                cy2 = min(y1, oy1) - max(y0, oy0)
                if cx > 0 and cy2 > 0 and cx * cy2 >= 0.9 * ta:
                    carded = True
                    break
            if not carded:
                hits.append((frac, rep))
    for frac, rep in sorted(hits, reverse=True)[:2]:
        warns.append((si, "W17", "텍스트가 이미지 경계에 걸침(안 %.0f%%, 잘려 보임 의심, 렌더 확인)" % (frac * 100),
                      "%r" % rep))


def action_title_check(titles, warns):
    """W14: 타이틀 다수가 서술형 명사구(시장 현황·경쟁 분석류) = 액션타이틀 아님(MBB: 제목만
    읽어도 주장이 흐르게). 종결어미 휴리스틱이라 주장으로 오판하는 미탐은 있어도(바다 등 '다'
    끝 명사) 오탐은 좁게: 한글 타이틀 3개+가 명사구이고 전체의 절반 이상일 때만 덱 단위 1회."""
    entries = []
    for si in sorted(titles):
        txt = " ".join(titles[si][1]).strip()
        chars = [c for c in txt if not c.isspace()]
        cjk_n = sum(1 for c in chars if is_cjk(c))
        # 한글 3자 미만·한글 비중 30% 미만은 제외: 빅스탯 숫자("300만+ 18.3%")·짧은 브랜드명을
        # 타이틀로 오인하던 것(50덱 스캔 실측)과 영문 타이틀을 걸러낸다.
        if len(chars) < 4 or cjk_n < 3 or cjk_n < 0.3 * len(chars):
            continue
        core = txt.rstrip(" ?!.…”’")
        claim = ("?" in txt or "!" in txt
                 or core.endswith(("다", "까", "요", "자", "죠", "함", "임")))
        entries.append((si, txt, claim))
    nominal = [(si, t) for si, t, c in entries if not c]
    if len(nominal) >= 3 and len(nominal) * 2 >= len(entries):
        ex = " ".join("p%d'%s'" % (si, t[:14]) for si, t in nominal[:4])
        warns.append((0, "W14", "타이틀 %d/%d개가 서술형 명사구(read 덱은 액션타이틀로, 에디토리얼·작품 소개 헤드라인은 무시)"
                      % (len(nominal), len(entries)), ex))
    return entries


def _diagram_clone_marks(inter):
    """두 페이지 공유 채움도형 멀티셋에서 '장식 텍스처 클론'(작은 점·절리 등, 24분할 1x1 이하) 수를 센다.
    카드형(중간 블록)·표(전폭 밴드)는 W6 소관이라 세지 않아 3열 비교카드 오탐을 피한다(적대 감사 반영)."""
    marks = area = 0
    for (fh, gx, gy, gw, gh), c in inter.items():
        if gw <= 1 and gh <= 1:
            marks += c
        if gw >= 1 and gh >= 1:
            area += gw * gh * c
    if marks >= 8 and area / (24.0 * 24.0) >= 0.06:
        return marks
    return 0


def slide_layout_sig(slide, sw, sh, gw=6, gh=4):
    """슬라이드 요소 배치를 gw x gh 그리드 점유 벡터로. 풀블리드 배경은 제외.
    텍스트=1 도형=0.5 이미지=2 가중. 어디에 무엇이 있는지(골격) 비교용."""
    sig = [0.0] * (gw * gh)
    n = 0
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:      # 풀블리드 배경/이미지 제외
            continue
        cx = (L + Wd / 2) / sw; cy = (T + Ht / 2) / sh
        if not (0 <= cx <= 1.0 and 0 <= cy <= 1.0):
            continue
        gc = min(gw - 1, max(0, int(cx * gw))); gr = min(gh - 1, max(0, int(cy * gh)))
        wgt = 2.0 if _is_pic(sp) else (1.0 if getattr(sp, "has_text_frame", False) else 0.5)
        sig[gr * gw + gc] += wgt
        n += 1
    return sig, n


def _cosv(a, b):
    da = sum(x * x for x in a) ** 0.5; db = sum(x * x for x in b) ** 0.5
    if da == 0 or db == 0:
        return 0.0
    return sum(x * y for x, y in zip(a, b)) / (da * db)


def _luma(rgb):
    def lin(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * lin(rgb[0]) + 0.7152 * lin(rgb[1]) + 0.0722 * lin(rgb[2])


def _run_rgb(run):
    try:
        c = run.font.color
        if c is not None and c.type is not None and c.rgb is not None:
            v = c.rgb
            return (v[0], v[1], v[2])
    except Exception:
        pass
    return None


def contrast_check(slide, si, sw, sh, render_dir, warns):
    """이미지 위 텍스트 저대비 감지(렌더 PNG 근사). picture와 겹치는 텍스트 프레임만, 슬라이드당 1회."""
    import glob
    from PIL import Image
    pics = []
    for sp in iter_shapes(slide.shapes):
        if _is_pic(sp):
            try:
                if None not in (sp.left, sp.top, sp.width, sp.height):
                    pics.append((sp.left, sp.top, sp.width, sp.height))
            except Exception:
                pass
    if not pics:
        return 0
    cand = glob.glob(os.path.join(render_dir, "p%02d.png" % si))
    if not cand:
        return 0
    try:
        im = Image.open(cand[0]).convert("RGB"); px = im.load(); PW, PH = im.size
    except Exception:
        return 0
    for sp in iter_shapes(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht):
            continue
        over = any(not (L + Wd <= p0 or L >= p0 + pw or T + Ht <= q0 or T >= q0 + ph)
                   for p0, q0, pw, ph in pics)
        if not over:
            continue
        rgbs = [_run_rgb(r) for para in sp.text_frame.paragraphs for r in para.runs if r.text.strip()]
        rgbs = [c for c in rgbs if c]
        if not rgbs:
            continue
        txt_rgb = rgbs[0]
        x0 = max(0, int(L / sw * PW)); y0 = max(0, int(T / sh * PH))
        x1 = min(PW, int((L + Wd) / sw * PW)); y1 = min(PH, int((T + Ht) / sh * PH))
        if x1 <= x0 or y1 <= y0:
            continue
        sx = max(1, (x1 - x0) // 24); sy = max(1, (y1 - y0) // 24)
        lumas = sorted(_luma(px[x, y]) for x in range(x0, x1, sx) for y in range(y0, y1, sy))
        if not lumas:
            continue
        L_txt = _luma(txt_rgb)
        # 배경 휘도는 텍스트색 반대 극단의 분위수로: 최악-국소 대비를 재고(WCAG는 최악값 기준)
        # 텍스트 잉크가 섞인 평균 오염을 피한다(밝은 글자면 어두운 15퍼센타일, 어두운 글자면 밝은 85퍼센타일).
        if L_txt >= 0.5:
            L_bg = lumas[int(len(lumas) * 0.15)]
        else:
            L_bg = lumas[min(len(lumas) - 1, int(len(lumas) * 0.85))]
        hi = max(L_bg, L_txt); lo = min(L_bg, L_txt)
        ratio = (hi + 0.05) / (lo + 0.05)
        if ratio < 2.5:
            warns.append((si, "W7", "이미지 위 텍스트 대비 낮음 %.1f:1 (스크림·색 보강 필요)" % ratio,
                          "text=%r" % sp.text_frame.text[:20]))
            return 1
    return 1   # 이 페이지의 렌더 PNG를 찾아 검사했음(W7 발화 여부와 무관, 파일명 규약 매치 확인용)


def lint(path, hard_min=5.0, body_min=9.0, small_min=7.5, render_dir=None, ghost=None):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    errors, warns = [], []
    sigs = []
    toks = {}
    page_texts = {}
    foot_tops = {}
    fx_pp = {}
    titles = {}
    thm_ea = theme_ea_font(prs)

    render_png_hits = 0
    for si, slide in enumerate(prs.slides, 1):
        sigs.append(slide_layout_sig(slide, sw, sh))
        toks[si] = _fill_tokens(slide, sw, sh)
        if render_dir:
            try:
                render_png_hits += contrast_check(slide, si, sw, sh, render_dir, warns)
            except Exception as e:
                print("W7 skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            accent_vbars_check(slide, si, sw, sh, warns)
        except Exception as e:
            print("W9 skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            foot_tops[si] = footer_top(slide, sw, sh)
            fx_pp[si] = effects_count(slide)
        except Exception as e:
            print("W12/W13 skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            text_overlap_check(slide, si, warns)
        except Exception as e:
            print("W15 skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            overflow_check(slide, si, sw / EMU_PER_IN, sh / EMU_PER_IN, warns)
            text_image_straddle_check(slide, si, sw / EMU_PER_IN, sh / EMU_PER_IN, warns)
        except Exception as e:
            print("W16/W17 skipped p%02d: %s" % (si, e), file=sys.stderr)
        for tf, w_emu in collect_frames(slide.shapes):
            fw_in = w_emu / EMU_PER_IN
            scale = frame_font_scale(tf)
            for para in tf.paragraphs:
                ptext = "".join(r.text for r in para.runs)
                para_size = para.font.size
                for run in para.runs:
                    t = run.text
                    if not t:
                        continue
                    page_texts.setdefault(si, []).append(t)
                    fonts = run_fonts(run)
                    cjk_font = fonts.get("ea") or fonts.get("latin")

                    # E1: CJK 렌더 슬롯이 라틴 전용(접두 매칭 → 굵기 변형 포함).
                    # 런에 폰트가 아예 없으면(raw python-pptx 전형) 테마 a:ea 가 실효 슬롯이다:
                    # 그 슬롯 판정까지 내려가야 '폰트를 안 준' Malgun 폴백을 잡는다(2026-07-02 감사 사각).
                    if has_cjk(t):
                        cf = (cjk_font or "").strip().lower()
                        if cf and any(cf.startswith(x) for x in LATIN_ONLY_FONTS):
                            errors.append((si, "E1", "CJK 텍스트가 라틴 전용 폰트로 렌더됨(폴백·자간 깨짐)",
                                           "font=%r text=%r" % (cjk_font, t[:24])))
                        elif not cf:
                            th = (thm_ea or "").strip().lower()
                            if thm_ea is None or th == "":
                                errors.append((si, "E1", "CJK 런에 폰트 미지정 + 테마 a:ea 빈 슬롯(Malgun 폴백 확정)",
                                               "text=%r" % t[:24]))
                            elif any(th.startswith(x) for x in LATIN_ONLY_FONTS):
                                errors.append((si, "E1", "CJK 런에 폰트 미지정 + 테마 a:ea 가 라틴 전용(폴백)",
                                               "theme=%r text=%r" % (thm_ea, t[:24])))

                    # E2: 긴 대시류
                    bad = [c for c in t if c in LONG_DASHES]
                    if bad:
                        errors.append((si, "E2", "긴 대시류 문자 포함",
                                       "cp=%s text=%r" % (",".join("U+%04X" % ord(c) for c in sorted(set(bad))), t[:24])))

                    # E3 / W1 / W5: 실효 폰트 크기(run → 문단 상속, autofit 반영)
                    rs = run.font.size if run.font.size is not None else para_size
                    if rs is not None:
                        eff = rs.pt * scale
                        if eff >= 18 and t.strip():
                            cur = titles.get(si)
                            if cur is None or eff > cur[0] + 0.1:
                                titles[si] = (eff, [t])
                            elif abs(eff - cur[0]) <= 0.1:
                                cur[1].append(t)
                        if eff < hard_min:
                            note = "" if scale == 1.0 else " (명목 %.1f * autofit %.2f)" % (rs.pt, scale)
                            errors.append((si, "E3", "실효 폰트 %.1fpt < 하한 %.1fpt(판독 불가)%s"
                                           % (eff, hard_min, note), "text=%r" % t[:24]))
                        elif eff < body_min and fw_in > 4.0 and len(ptext) >= 40:
                            warns.append((si, "W1", "본문급 프레임 글자 %.1fpt < 권장 %.1fpt(출처·캡션이면 무시)"
                                          % (eff, body_min),
                                          "w=%.1fin len=%d text=%r" % (fw_in, len(ptext), ptext[:24])))
                        elif eff < small_min and has_cjk(t) and fw_in <= 4.0:
                            # 좁은 프레임(<=4in)만: 넓은 프레임의 소형 한글은 목업·카드 내부가 아니라
                            # 캡션·주석일 여지가 커 메시지(목업 추정)와 어긋난다(공개 위생 감사 반영).
                            warns.append((si, "W8", "소형 CJK %.1fpt < %.1fpt(목업·카드 내부 추정, 판독 위험, 캡션이면 무시)"
                                          % (eff, small_min),
                                          "w=%.1fin text=%r" % (fw_in, t[:24])))
                    else:
                        warns.append((si, "W5", "폰트 크기를 run·문단에서 못 찾음(상속 미해석)",
                                      "text=%r" % t[:24]))

                    # E4: 연속 CJK 2자 이상 + 유의미 양수 트래킹
                    if sum(1 for c in t if is_cjk(c)) >= 2:
                        tr = run_track(run)
                        if tr is not None and tr > 50:
                            errors.append((si, "E4", "CJK run 에 양수 트래킹 %d(0.5pt 초과, 자간 벌어짐)" % tr,
                                           "text=%r" % t[:24]))

    # --render 를 켰는데 p##.png 규약에 맞는 렌더를 한 장도 못 찾으면 W7 이 조용히 0건 수행된다.
    # 폴더에 png 는 있는데(다른 이름) 매치가 0이면 파일명 규약을 stderr 로 안내(silent no-op 방지).
    if render_dir and render_png_hits == 0:
        anypng = glob.glob(os.path.join(render_dir, "*.png"))
        if anypng:
            print("W7 note: %s 에 p01.png·p02.png 형식 렌더가 없어 이미지 대비 검사를 건너뜀 "
                  "(현재 파일: %s ...)" % (render_dir, os.path.basename(anypng[0])), file=sys.stderr)

    # W6: 레이아웃 골격 재탕. 성긴 슬라이드(디바이더·커버 등 비영 셀<4)는 제외하고, 콘텐츠
    # 슬라이드 중 한 장이 유사(>0.90)한 다른 슬라이드가 3장 이상(같은 골격 4장+)이면 경고.
    # 덱 길이에 불변(전체 쌍 비율은 큰 덱에서 국소 재탕을 희석해 놓치므로 최대 클러스터로 판정).
    content = [(i + 1, sig) for i, (sig, n) in enumerate(sigs) if n >= 3]
    if len(content) >= 4:
        adj = {p: [] for p, _ in content}
        for a in range(len(content)):
            for b in range(a + 1, len(content)):
                pa, sa = content[a]; pb, sb = content[b]
                sim = _cosv(sa, sb)
                if sim > 0.90:
                    adj[pa].append((pb, sim)); adj[pb].append((pa, sim))
        worst_p, worst = max(adj.items(), key=lambda kv: len(kv[1]))
        if len(worst) >= 3:
            ex = " ".join("p%d~p%d(%.2f)" % (worst_p, b, s) for b, s in sorted(worst, key=lambda x: -x[1])[:4])
            warns.append((0, "W6", "레이아웃 골격이 %d개 페이지에서 반복됨(클로드 티 의심, 성긴 디바이더 제외)" % (len(worst) + 1),
                          "예 " + ex))

    # W11 카피 클리셰 · W12 푸터 정렬 · W13 효과 · W14 액션타이틀 (전부 덱 단위)
    try:
        copy_cliche_check(page_texts, warns)
        footer_check(foot_tops, warns)
        effects_check_deck(fx_pp, warns)
        action_title_check(titles, warns)
        if ghost is not None:
            # ghost 는 W14 한글 필터와 무관하게 수집한 전 타이틀(18pt+): 영문·숫자 덱에서도
            # 제목만 읽어 수평 논리를 검토하는 용도라, 필터 후 entries 가 아니라 raw titles 에서 뽑는다.
            ghost.extend((si, " ".join(titles[si][1]).strip()) for si in sorted(titles))
    except Exception as e:
        print("W11/W12/W13/W14 skipped: %s" % e, file=sys.stderr)

    # W10: 직접 그린 도식(단면도 등) 클론 재탕. 장식 텍스처(marks) 경로로 특화해 3열 카드형은
    # W6에 맡기고 여기선 세지 않는다(적대 감사 2026-07-02: 카드형 subblk 오탐 회피). W6가 놓친
    # 0.90 미달·단일 쌍(단면 도식 반복 실측)을 pptx 도형만으로 메운다.
    cadj = {p: [] for p in toks}
    nums_t = sorted(toks)
    for ia in range(len(nums_t)):
        for ib in range(ia + 1, len(nums_t)):
            pa, pb = nums_t[ia], nums_t[ib]
            if _diagram_clone_marks(toks[pa] & toks[pb]):
                cadj[pa].append(pb); cadj[pb].append(pa)
    if cadj:
        cw, cwl = max(cadj.items(), key=lambda kv: len(kv[1]))
        if len(cwl) >= 1:
            grp = sorted({cw, *cwl})
            warns.append((0, "W10", "직접 그린 도식이 %d개 페이지에서 거의 동일 반복(의도된 교육 시퀀스인지 게으른 재탕인지 눈으로 확정)" % len(grp),
                          "페이지 " + ",".join("p%d" % p for p in grp)))

    return errors, warns


def main():
    # 한국어 메시지·argparse 도움말이 non-UTF-8 stdout(파이프·cp949/cp1252)에서 UnicodeEncodeError
    # 로 죽지 않게 파서 생성보다 먼저 재설정한다(--help 는 parse_args 중에 출력되므로 순서가 중요).
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass
    ap = argparse.ArgumentParser(prog="aro",
                                 description="빌드된 .pptx를 배포 전에 기계로 검사하는 한글 특화 품질 린터")
    ap.add_argument("pptx")
    ap.add_argument("--hard-min", type=float, default=5.0)
    ap.add_argument("--body-min", type=float, default=9.0)
    ap.add_argument("--strict", action="store_true")
    ap.add_argument("--small-min", type=float, default=7.5, help="W8 좁은 프레임 소형 CJK 상한(pt)")
    ap.add_argument("--render", default=None, help="렌더 PNG 폴더(p01.png·p02.png 형식) 지정 시 이미지 위 텍스트 대비(W7) 검사 활성화")
    ap.add_argument("--ghost", action="store_true", help="고스트덱(페이지별 타이틀만 나열) 출력: 제목만 읽어 주장이 흐르는지 수평 논리 눈검수")
    ap.add_argument("--json", action="store_true", help="기계 판독용 JSON 출력(에이전트·CI 연동)")
    a = ap.parse_args()

    if not os.path.exists(a.pptx):
        print("aro: 파일을 찾을 수 없습니다: %s" % a.pptx, file=sys.stderr)
        sys.exit(2)

    ghost = [] if (a.ghost or a.json) else None
    try:
        errors, warns = lint(a.pptx, a.hard_min, a.body_min, a.small_min, render_dir=a.render, ghost=ghost)
    except Exception as e:
        print("aro: pptx 를 열 수 없습니다(유효한 .pptx 인지 확인): %s (%s)"
              % (a.pptx, type(e).__name__), file=sys.stderr)
        sys.exit(2)

    if a.json:
        import json
        doc = {
            "file": a.pptx,
            "errors": [{"page": si, "code": c, "message": m, "detail": d} for si, c, m, d in errors],
            "warnings": [{"page": si, "code": c, "message": m, "detail": d} for si, c, m, d in warns],
            "ghost": [{"page": si, "title": t} for si, t in (ghost or [])],
            "summary": {"error_count": len(errors), "warn_count": len(warns),
                        "pass": not errors and not (a.strict and warns)},
        }
        print(json.dumps(doc, ensure_ascii=False, indent=2))
        sys.exit(1 if (errors or (a.strict and warns)) else 0)

    print("=== ARO LINT: %s ===" % a.pptx)
    if ghost:
        print("--- ghost deck (제목만 읽기: 주장이 이야기로 흐르는가) ---")
        for si, txt in ghost:
            print("  p%02d  %s" % (si, txt[:60]))
    if not errors and not warns:
        print("clean: ERROR 0, WARN 0")
    for si, code, msg, detail in errors:
        print("  ERROR p%02d [%s] %s | %s" % (si, code, msg, detail))
    for si, code, msg, detail in warns:
        print("  WARN  p%02d [%s] %s | %s" % (si, code, msg, detail))
    print("--- ERROR %d, WARN %d ---" % (len(errors), len(warns)))

    if errors or (a.strict and warns):
        sys.exit(1)
    sys.exit(0)


if __name__ == "__main__":
    main()
