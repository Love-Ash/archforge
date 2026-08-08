# -*- coding: utf-8 -*-
"""
Archforge: a Hangul-specialized quality linter that takes a built .pptx and mechanically
blocks recurring defects.

Reason to exist: even when rules are written down, the same mistakes repeat every time a deck
is made. Defects that a machine can check should be blocked by this linter, not by human eyes.
It catches the points where Hangul decks silently break (fallback to Latin-only fonts, CJK
tracking, undersized text) and the tells specific to AI-generated decks (long dashes, cliche
copy, recycled layouts) in the build output. Because it is the last line of defense for an
arbitrary pptx, it also covers bypass paths such as inheritance, slots, autofit, tables, and
groups (reflecting adversarial audit findings).

ERROR = blocks deployment:
  E1 The effective render font for Hangul text is Latin-only (no Hangul glyphs) = a silent
     Malgun fallback.
     The effective font is resolved with the measured-by-render COM model (2026-07-10,
     docs/CALIBRATION.md):
     run a:ea > lstStyle inheritance chain (shape -> layout -> master, measured in probe 6)
     > theme ea
     (title family uses majorFont, everything else uses minorFont; if not empty it takes
     priority over run a:latin)
     > (only when the theme ea slot is empty) run a:latin > OS fallback (Malgun).
     The trigger is Hangul-only: kana- or Hanja-only runs are not judged using Hangul coverage
     knowledge (0.2.1).
  E2 A dash-class character appears in the rendered text (em/en/figure dash plus math minus
     U+2212, fullwidth hyphen-minus U+FF0D, box-drawing horizontal line U+2500, etc.).
     The axis of judgment is function, not character: an en dash passes if both neighboring
     tokens are numeric (a range), passes only when directly adjacent if just one side is
     numeric, and is blocked for a spaced-out parenthetical or a word-to-word join. U+2212
     passes when immediately followed by a digit.
     Context is read from the whole paragraph, not the run (to prevent false positives from
     run splitting). --strict blocks everything with no exceptions.
  E3 The effective font (reflecting autofit fontScale, paragraph inheritance, and the
     placeholder inheritance chain) is below HARD_MIN (default 5.0pt) = unreadable.
  E4 Two or more consecutive Hangul characters carry meaningfully positive tracking (spc>50 =
     0.5pt) = tracking has spread apart.
     Hangul-only: kana tracking is normal practice in Japanese design, so it cannot be declared
     a defect outright (0.2.1).
WARN = informational (deployment still passes):
  W1 The effective font of a body-level frame (wide, high character count) is below BODY_MIN
     (default 9.0pt) = body text is too small
  W8 Small CJK text in a narrow frame (effective size in [HARD_MIN, SMALL_MIN=7.5)pt) = risk
     of unreadable text inside device mockups or cards
     (the gray zone above the E3 unreadable floor and below the W1 body-level threshold.
     Surfaces sub-text inside phone mockups during preflight)
  W6 Recycled layout skeleton: shape bbox signature cosine similarity > 0.90, cluster of 4+
     pages (no render needed, always checked)
  W7 Low-contrast text over an image, ratio<2.5 (only when --render <pages> is given, compared
     pixel-by-pixel against the rendered PNG)
  W9 3+ accent-colored vertical bars repeated as list markers = using color to build structure
     (a Claude tell). Covers connectors and zero-width bars
  W10 A hand-drawn diagram (e.g. cross-section, decorative texture) repeated almost identically
      across multiple pages = flags whether it is recycling or intentional for human judgment
  W11 AI-tell copy: buzzwords (narrow dictionary, all pages), cliche stock openings (p1-3 only)
  W12 Footer baseline misalignment: pages deviating 0.03-0.25in from the dominant baseline
      (cover page excluded, median of 0.05in buckets)
  W13 2+ native PPT shadow/glow/3D effects (an old-fashioned tell. An empty effectLst with no
      children is not counted)
  W14 A majority of titles (3+, half or more) are descriptive noun phrases = not action titles.
      --ghost also prints the ghost deck (title listing)
  W15 Estimated text-on-text overlap (approximate effective glyph bbox, intersection > 45% of
      the smaller box) = occlusion or collision, at most 2 findings per page
  W16 Off-canvas overflow: text = 0.15in+ outside the effective glyph box boundary, images =
      0.12in+ outside the ink bbox (alpha-trimmed)
      (full-bleed 70%+ excluded). Corner bleed on decorative shapes is standard technique and
      is not checked (rejected after measured-by-render testing)
  W17 Text straddles the ink boundary of a non-background image (only 25-75% of the glyph is
      inside) = looks cropped. Fully on top of the image is W7's jurisdiction
  W5 Font size found nowhere in the run, paragraph, or inheritance chain (only when the whole
     chain is silent)
  W18 Some region could not be checked due to corrupted or non-standard attributes: the result
      may be incomplete (surfaces the region a guard swallowed in the output contract, not just
      stderr; promoted to exit 1 under --strict)

W15-W17 geometry robustness (fixed after reproducing 12 findings from adversarial
verification, 2026-07-03): group off/chOff affine transform, wrap=none (word_wrap=False, the
python-pptx add_textbox default) single-line actual width, autofit percentage strings
('62.5%') and lnSpcReduction, per-paragraph alignment, empty-paragraph endParaRPr size, rotated
shapes (axis-aligned expanded bbox, rotated text is skipped), picture srcRect crop, flipH/V, P
mode tRNS transparency, and W17 suppresses a solid card (by z order) between a photo and its
caption. Remaining limitation: a placeholder's alignment inherited from layout lstStyle falls
back to left alignment.

Usage: archforge <built.pptx> [--hard-min 5.0] [--body-min 9.0] [--strict] [--render <pages>]
              [--ghost] [--json] [--skip CODES] [--w6-sim 0.90] [--w6-cluster 3]
  --strict: WARN also causes exit 1 + lifts the E2 numeric-context exception. --render:
  enables W7.
  --ghost: lists titles (for eyeballing horizontal logic). --skip: suppresses the given codes
  (e.g. --skip W14,W6 for an editorial deck).
Returns: exit 1 if there is any ERROR, otherwise 0.

Subcommand: archforge skill [--install [DIR]] [--path]
  Prints the bundled agent skill pack (SKILL.md), or installs it to DIR (default
  ./.claude/skills).
"""
import os
import re
import sys
import glob
import math
import argparse
import colorsys
from collections import Counter, namedtuple
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from .messages import M, set_lang, get_lang
    from .findings import Finding, shape_loc
    from .rules import RULES, ALL_CODES, PROFILES, DEFAULT_PROFILE, severity
    from . import config as _config
    from . import reporters as _reporters
except ImportError:   # fallback for standalone file execution (python lint.py)
    from messages import M, set_lang, get_lang
    from findings import Finding, shape_loc
    from rules import RULES, ALL_CODES, PROFILES, DEFAULT_PROFILE, severity
    import config as _config
    import reporters as _reporters

try:
    from .geometry import (_pct_attr,
                     frame_autofit,
                     frame_font_scale,
                     _group_xf,
                     collect_frames,
                     iter_shapes,
                     _is_pic,
                     iter_shapes_geo,
                     _geo_rect)
except ImportError:   # standalone execution
    from geometry import (_pct_attr,
                    frame_autofit,
                    frame_font_scale,
                    _group_xf,
                    collect_frames,
                    iter_shapes,
                    _is_pic,
                    iter_shapes_geo,
                    _geo_rect)
try:
    from .resolution import (_theme_fonts_from_blob,
                     theme_fonts_by_master,
                     theme_ea_by_master,
                     theme_ea_font,
                     _sz_from_defrpr,
                     _lst_defrpr,
                     _lst_sz_pt,
                     StyleResolver, SizeResolver)
except ImportError:   # standalone execution
    from resolution import (_theme_fonts_from_blob,
                    theme_fonts_by_master,
                    theme_ea_by_master,
                    theme_ea_font,
                    _sz_from_defrpr,
                    _lst_defrpr,
                    _lst_sz_pt,
                    StyleResolver, SizeResolver)
try:
    from .ooxml import EMU_PER_IN, NS, NS_P
    from .colors import (_shape_fill_hex, _shape_line_hex, _is_accent,
                         _theme_colors_from_blob, theme_colors_by_master,
                         _cosv, _luma, _run_rgb, _resolve_run_rgb,
                         _COLOR_UNKNOWN)
except ImportError:   # standalone execution
    from ooxml import EMU_PER_IN, NS, NS_P
    from colors import (_shape_fill_hex, _shape_line_hex, _is_accent,
                        _theme_colors_from_blob, theme_colors_by_master,
                        _cosv, _luma, _run_rgb, _resolve_run_rgb,
                        _COLOR_UNKNOWN)

# Rule metadata, profiles, and the full code list have moved to the rules.py registry
# (0.4.0 decomposition).
# PROFILES/ALL_CODES continue to be exposed from this module via the import above (backward
# compatibility).


# Per-run Unicode script detection moved to scripts.py (0.7 decomposition: parsing layer).
# Re-exported here so existing callers (jl.is_hangul etc.) and this module's own detectors
# keep working unchanged. _geometry_unsupported keeps its underscore alias for internal use.
try:
    from .scripts import (is_hangul, has_hangul, is_kana, is_hanja, is_cjk, has_cjk,
                          geometry_unsupported as _geometry_unsupported)
except ImportError:
    from scripts import (is_hangul, has_hangul, is_kana, is_hanja, is_cjk, has_cjk,
                         geometry_unsupported as _geometry_unsupported)

# E1 font kernel and E2 dash kernel moved to fonts.py / dashes.py (0.8, #5
# decomposition: resolution/detection kernels as pure modules). Re-exported here so
# existing callers (jl.e1_violation, jl.dash_violations, deck_lint's star-import)
# keep working unchanged.
try:
    from .fonts import (LATIN_ONLY_FONTS, KOREAN_CAPABLE_EXCEPTIONS,
                        JP_CN_CAPABLE_PREFIXES, run_fonts, para_fonts, run_track,
                        _text_point_attr, _UNIVERSAL_PER_PT, is_latin_only_font,
                        e1_violation, resolve_font_tokens)
    from .dashes import (LONG_DASHES, _EN_DASH, _MINUS, _is_digit_ch,
                         _dash_neighbor, dash_violations)
except ImportError:
    from fonts import (LATIN_ONLY_FONTS, KOREAN_CAPABLE_EXCEPTIONS,
                       JP_CN_CAPABLE_PREFIXES, run_fonts, para_fonts, run_track,
                       _text_point_attr, _UNIVERSAL_PER_PT, is_latin_only_font,
                       e1_violation, resolve_font_tokens)
    from dashes import (LONG_DASHES, _EN_DASH, _MINUS, _is_digit_ch,
                        _dash_neighbor, dash_violations)


class _FldRun:
    """Adapter for a:fld (auto fields such as slide number, date). Since CT_TextField also has
    an rPr+t structure per the schema, PowerPoint renders it with the same rules as a normal
    run, but python-pptx's para.runs only returns a:r, so field text was a blind spot for
    E1/E3/E4 (carried over from the fourth review, 0.5.0). Exposes only the minimal interface
    the checking code uses: ._r for run_fonts/run_track, .font.size for size."""
    __slots__ = ("_r", "text")

    class _Pt:
        __slots__ = ("pt",)

        def __init__(self, pt):
            self.pt = pt

    def __init__(self, fld):
        self._r = fld
        t = fld.find(NS + "t")
        self.text = (t.text or "") if t is not None else ""

    @property
    def font(self):
        return self   # only the .size access is used

    @property
    def size(self):
        try:
            v = self._r.find(NS + "rPr").get("sz")
        except Exception:
            return None
        if not v:
            return None
        try:
            return self._Pt(int(v) / 100.0)
        except (TypeError, ValueError):
            return None
















# Absolute-coordinate traversal for W15-W17 geometry consumers. A group child's raw left/top
# is in the group's chOff coordinate space, which drifts from slide coordinates in a pptx
# where the group has been moved or resized (off!=chOff desync, standard behavior when
# dragging in PowerPoint) (measured in adversarial verification, 2026-07-03). Composes the
# off/ext vs chOff/chExt affine and yields (shape, z-order, absolute xfrm function
# coefficients). xf=(ax,bx,ay,by): abs = a*raw + b (EMU).










def accent_vbars_check(slide, si, sw, sh, warns):
    """W9: an AI-generated-deck tell where accent-colored vertical bars are repeated as list
    markers, using color to build item structure (measured in real decks). Reflects the
    2026-07-02 adversarial audit: colors are read via the sp.line/sp.fill accessors to cover
    namespaces and connectors, vertical bars explicitly include zero-width connectors (w~0),
    and adjacent text to the right confirms it is really a list marker."""
    bars, texts = [], []
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht):
            continue
        x, y, w, h = L / EMU_PER_IN, T / EMU_PER_IN, Wd / EMU_PER_IN, Ht / EMU_PER_IN
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            texts.append((x, y, w, h))
        hexc = _shape_line_hex(sp) or _shape_fill_hex(sp)
        if not _is_accent(hexc):
            continue
        if 0.2 <= h <= 1.0 and (w < 0.05 or h > 3 * w):   # vertical bar (including
                                                           # zero-width connectors)
            bars.append((x, y, w, h, hexc))
    if len(bars) < 3:
        return
    hues = {b[4] for b in bars}
    if len(hues) != 1:                       # multiple colors is legitimate data encoding
                                              # (a legend)
        return
    xs = [b[0] for b in bars]
    if max(xs) - min(xs) > 0.15:             # vertically aligned stack only (horizontal
                                              # spread = chart/divider, excluded)
        return
    rt = 0
    for bx, by, bw, bh, _hx in bars:
        for tx, ty, tw, th in texts:
            if tx > bx and (tx - (bx + bw)) < 0.6 and not (ty > by + bh or ty + th < by):
                rt += 1
                break
    if rt >= len(bars) - 1:
        warns.append(Finding(si, "W9", "w9", (len(bars),),
                         "x=%.2fin hue=%s" % (min(xs), next(iter(hues)))))


def _fill_tokens(slide, sw, sh):
    """Turns the slide's solid-fill shapes into a multiset of (color, 24-grid position/size)
    tokens. Full-bleed backgrounds are excluded."""
    t = Counter()
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:
            continue
        fh = _shape_fill_hex(sp)
        if fh is None:
            continue
        t[(fh, round(L / sw * 24), round(T / sh * 24), round(Wd / sw * 24), round(Ht / sh * 24))] += 1
    return t
























# AI-tell copy: only obvious cliches (a narrow dictionary = suppresses false positives).
# General words that could be legitimate in context are not included.
BUZZWORDS = (
    "시너지", "패러다임", "게임체인저", "게임 체인저", "혁신을 가속", "가치를 극대화",
    "미래를 선도", "새로운 지평", "무한한 가능성", "홀리스틱", "엔드투엔드", "엔드 투 엔드",
    "초격차", "글로벌 리더로 도약", "위대한 여정",
    "synergy", "paradigm shift", "game-changer", "game changer", "cutting-edge",
    "state-of-the-art", "seamless", "revolutionize", "leverage synerg", "holistic",
    "unlock the potential", "empower",
)
STALE_OPENINGS = (
    "오늘날", "급변하는", "4차 산업혁명 시대", "바야흐로", "현대 사회에서",
    "디지털 전환의 시대", "디지털 전환의 물결", "알아보겠습니다", "살펴보겠습니다",
    "in today's", "in the rapidly changing", "in this presentation",
)


def copy_cliche_check(page_texts, warns):
    """W11: AI-tell copy. Buzzwords are checked on every page, cliche openings only in the
    intro (p1-3)."""
    for si in sorted(page_texts):
        blob = " ".join(page_texts[si])
        low = blob.lower()
        hits = sorted({b for b in BUZZWORDS if b.lower() in low})
        if hits:
            warns.append(Finding(si, "W11", "w11_buzz", (len(hits),), ", ".join(hits[:5])))
        if si <= 3:
            op = sorted({o for o in STALE_OPENINGS if o.lower() in low})
            if op:
                warns.append(Finding(si, "W11", "w11_open", (), ", ".join(op[:5])))


def footer_top(slide, sw, sh):
    """The top (in) of the bottommost text in the slide's bottom band (y>0.88H). None if there
    is no footer."""
    best = None
    for sp in iter_shapes(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        try:
            if not sp.text_frame.text.strip():
                continue
            t = sp.top
        except Exception:
            continue
        if t is None or t <= 0.88 * sh:
            continue
        ti = t / EMU_PER_IN
        if best is None or ti > best:
            best = ti
    return best


def footer_check(foot_tops, warns):
    """W12: footer baseline misalignment across pages. In a 50-deck measurement (2026-07-02),
    the absolute-deviation approach mistook cover-page credits and bottom captions for footers,
    producing false positives in 17 decks. Instead this excludes the cover page (p1), treats
    the dominant baseline (mode of the 0.05in quantization buckets, 3+ pages) as the house
    footer, and flags only pages that deviate "slightly" (0.03-0.25in) from it. Anything off
    by more than 0.25in is assumed to be a different element such as a caption or divider and
    is ignored (its existence is not even checked)."""
    tops = [(si, t) for si, t in foot_tops.items() if t is not None and si > 1]
    if len(tops) < 4:
        return
    q = Counter(round(t / 0.05) for _si, t in tops)
    qv, cnt = q.most_common(1)[0]
    if cnt < 3:
        return
    # House baseline = the median of the actual values in the dominant bucket (using the
    # bucket center instead would make a majority of 7.08 values sit exactly 0.02 away from
    # 7.10, a floating-point boundary false positive confirmed by rescanning the 50 decks)
    bvals = sorted(t for _si, t in tops if round(t / 0.05) == qv)
    base = bvals[len(bvals) // 2]
    off = [(si, t) for si, t in tops if 0.03 < abs(t - base) <= 0.25]
    if off:
        ex = " ".join("p%d=%.2f" % (si, t) for si, t in off[:4])
        warns.append(Finding(0, "W12", "w12", (base, len(off)), ex))


_EFFECT_TAGS = tuple(NS + t for t in ("outerShdw", "innerShdw", "glow", "reflection"))
_3D_TAGS = tuple(NS + t for t in ("sp3d", "scene3d"))


def effects_count(slide):
    """The count and kinds of the slide's effective PPT effects (shadow, glow, 3D). Some
    generators leave a childless empty effectLst purely to block inheritance, so this only
    counts elements that actually have effect children (prevents an empty-element false
    positive)."""
    n = 0
    kinds = set()
    for sp in iter_shapes(slide.shapes):
        spPr = getattr(sp._element, "spPr", None)
        if spPr is None:
            continue
        eff = spPr.find(NS + "effectLst")
        if eff is not None:
            for ch in eff:
                if ch.tag in _EFFECT_TAGS:
                    n += 1
                    kinds.add(ch.tag.split("}")[1])
        for tag in _3D_TAGS:
            if spPr.find(tag) is not None:
                n += 1
                kinds.add(tag.split("}")[1])
    return n, kinds


def effects_check_deck(per_page, warns):
    """W13: aggregated once per deck (firing repeatedly per page is noise: measured on the
    50-deck corpus). Could be an intentional neon/glow style, so this is a WARN and the
    judgment is left to human eyes."""
    hits = [(si, n, kinds) for si, (n, kinds) in per_page.items() if n >= 2]
    if not hits:
        return
    total = sum(n for _si, n, _k in hits)
    kinds = sorted(set().union(*[k for _si, _n, k in hits]))
    pages = ",".join("p%d" % si for si, _n, _k in hits[:6])
    warns.append(Finding(0, "W13", "w13", (total, len(hits)),
                         "%s | %s" % (pages, ",".join(kinds))))


# W15 text overlap: the most common defect axis in generated decks (elements pile up with
# every revision round), but the frame bbox is drawn generously by convention and can't be
# used, so this approximates the effective glyph width instead.
_W_CJK, _W_LAT, _W_SP = 0.96, 0.52, 0.28   # character-width/font-size ratio approximation
                                           # (conservative: suppresses false positives)

# Effective glyph bbox (in) per paragraph. Turned a magic index tuple into named fields
# (external review, 2026-07-10).
# sp (owning shape) is for the loc payload of W15-W17 findings (0.5.0); the coordinates are
# already the group's absolute coordinates.
GlyphBox = namedtuple("GlyphBox", "x0 y0 x1 y1 rep max_pt frame_id sp cell para field")
GlyphBox.__new__.__defaults__ = (None, None, None, False)


def _glyph_w(s, size_pt):
    w = 0.0
    for ch in s:
        if ch == " ":
            w += _W_SP
        elif is_cjk(ch) or ord(ch) > 0x2E80:
            w += _W_CJK
        else:
            w += _W_LAT
    return w * size_pt / 72.0


def _empty_para_pt(para, default_pt):
    """The effective size of an empty paragraph (a spacer): prioritizes endParaRPr/defRPr sz
    (fixes a phantom-height issue where a 4pt spacer was counted as 12pt, measured in
    adversarial verification)."""
    try:
        if para.font.size is not None:
            return para.font.size.pt
    except Exception:
        pass
    try:
        epr = para._p.find(NS + "endParaRPr")
        if epr is not None and epr.get("sz"):
            return int(epr.get("sz")) / 100.0
    except Exception:
        pass
    return default_pt


def _text_glyph_boxes(slide, default_pt=12.0, skipped=None, styler=None):
    """Approximates the effective glyph bbox (in) per paragraph. Returns
    [(x0,y0,x1,y1,representative text,max_pt,frame_id)].
    Width is summed from each run's actual size; line height reflects the actual
    line_spacing value (1.2 if absent) combined with autofit lnSpcReduction. x is placed
    using per-paragraph alignment (the frame's first explicit value if unset, otherwise
    left), which reduces misplacement in frames with mixed alignment. wrap=none
    (word_wrap=False, the python-pptx add_textbox default) extends past the frame in a
    single line, so no wrap folding is applied and the actual width is used as-is. A rotated
    frame is skipped since estimation would be invalid. Groups are converted to absolute
    coordinates via iter_shapes_geo.
    Calibrated against real-deck render comparisons plus reproduced adversarial-verification
    measurements.
    Script layer (0.2.1): vertical writing (bodyPr@vert) and frames containing RTL or
    complex-shaping scripts are skipped, since glyph-width approximation is meaningless for
    them, and if a skipped Counter is passed it tallies these and surfaces them via W18.
    0.3.1 (third external review, P0): if given a styler (StyleResolver), a run with no
    explicit size is resolved through the same inheritance chain as E3 (fixes an
    inconsistency where two different effective-style models existed in a single document).
    Native tables compute each cell's rectangle by accumulating column widths and row
    heights, so cell text is included too.
    Known limitations: a placeholder that inherits alignment from layout lstStyle falls back
    to left alignment, and skipping rotated frames is excluded from the W18 tally since it is
    standard decorative practice."""
    import math
    out = []

    def emit_frame(tframe, fx, fy, fw, fh, fid, owner_sp, cell=None, sx=1.0, sy=1.0):
        try:
            bodyPr = tframe._txBody.find(NS + "bodyPr")
            vert = bodyPr.get("vert") if bodyPr is not None else None
            if vert not in (None, "horz"):
                if skipped is not None:
                    skipped["vertical_text"] += 1
                return
        except Exception:
            bodyPr = None
        # Text-frame insets (0.6.0, external review): glyphs start inside the frame, not
        # at its edge. OOXML defaults are lIns/rIns 91440 EMU (0.1in) and tIns/bIns 45720
        # (0.05in), the same order of magnitude as W16's 0.15in tolerance, so ignoring
        # them shifted every glyph box left/up and overstated usable width. Insets live
        # in shape-local units, so group scale (sx/sy) applies.
        li, ri_, ti, bi = 91440, 91440, 45720, 45720
        try:
            if bodyPr is not None:
                li = int(bodyPr.get("lIns", li))
                ri_ = int(bodyPr.get("rIns", ri_))
                ti = int(bodyPr.get("tIns", ti))
                bi = int(bodyPr.get("bIns", bi))
        except Exception:
            pass
        fx += sx * li / EMU_PER_IN
        fy += sy * ti / EMU_PER_IN
        fw = max(fw - sx * (li + ri_) / EMU_PER_IN, 0.0)
        fh = max(fh - sy * (ti + bi) / EMU_PER_IN, 0.0)
        try:
            # All a:t descendants, not para.runs: field text (a:fld) participates in
            # geometry, so it must participate in the complex-script screen too. A
            # field-only Arabic frame used to bypass this check and get measured with
            # the Latin/CJK width model without any W18 (0.6.1, external review).
            frame_text = "".join(t.text or "" for t in tframe._txBody.iter(NS + "t"))
            if _geometry_unsupported(frame_text):
                if skipped is not None:
                    skipped["complex_script"] += 1
                return
        except Exception:
            pass
        fw2 = max(fw, 0.05)
        scale, lnred = frame_autofit(tframe)
        wrap = tframe.word_wrap is not False   # None (no attribute) = OOXML default
                                               # square = wrap
        frame_align = None
        for para in tframe.paragraphs:
            if para.alignment is not None:
                frame_align = para.alignment
                break
        paras = []   # (line_w, pmx, ptxt, factor, n, align, p_idx, field_only)
        for p_idx, para in enumerate(tframe.paragraphs):
            pmx, ptxt = 0.0, ""
            saw_field, saw_run_text = False, False
            segs = [0.0]   # widths of a:br-separated visual lines
            # Document-order walk over a:r / a:fld / a:br (0.6.0, external review): fld
            # text occupies real width and an explicit line break starts a new visual
            # line. Before this, a br-split sentence was measured as one overlong line
            # (width overstated, height understated). Falls back to para.runs on any
            # structural surprise.
            items = []
            try:
                runs_l = list(para.runs)
                r_seen = 0
                for child in para._p:
                    tag = child.tag
                    if tag == NS + "r":
                        if r_seen < len(runs_l):
                            items.append(runs_l[r_seen])
                            r_seen += 1
                    elif tag == NS + "br":
                        items.append(None)
                    elif tag == NS + "fld":
                        items.append(_FldRun(child))
                if r_seen != len(runs_l):
                    items = list(runs_l)
            except Exception:
                items = list(para.runs)
            for r in items:
                if r is None:   # a:br: next visual line
                    segs.append(0.0)
                    ptxt += " "
                    continue
                t = r.text
                if not t:
                    continue
                if isinstance(r, _FldRun):
                    saw_field = True
                else:
                    saw_run_text = True
                if r.font.size is not None:
                    sz = r.font.size.pt
                elif para.font.size is not None:
                    sz = para.font.size.pt
                else:
                    sz = None
                    if styler is not None:
                        try:
                            sz = styler.resolve(tframe, owner_sp, slide, getattr(para, "level", 0))
                        except Exception:
                            sz = None
                    if sz is None:
                        sz = default_pt
                sz *= scale
                segs[-1] += _glyph_w(t, sz)
                ptxt += t
                if sz > pmx:
                    pmx = sz
            if not ptxt.strip():
                paras.append((0.0, _empty_para_pt(para, default_pt) * scale, "", 1.2, 1,
                              None, p_idx, False))
                continue
            ls = para.line_spacing
            if ls is None:
                factor = 1.2
            elif isinstance(ls, float):
                factor = ls
            else:
                try:
                    factor = ls.pt / pmx if pmx else 1.2
                except Exception:
                    factor = 1.2
            if wrap:
                n = sum(max(1, math.ceil(w / (fw2 * 1.04))) for w in segs)
            else:
                n = len(segs)
            line_w = max(segs)
            al = para.alignment if para.alignment is not None else frame_align
            paras.append((line_w, pmx, ptxt, factor, n, al, p_idx,
                          saw_field and not saw_run_text))
        gh_total = sum(n * pmx * max(f, 0.95) * (1.0 - lnred) / 72.0
                       for (pw, pmx, ptxt, f, n, al, _pi, _fo) in paras)
        if gh_total <= 0:
            return
        va = str(tframe.vertical_anchor) if tframe.vertical_anchor is not None else ""
        if "MIDDLE" in va:
            cy = fy + max(0.0, (fh - gh_total) / 2)
        elif "BOTTOM" in va:
            cy = fy + max(0.0, fh - gh_total)
        else:
            cy = fy
        for (pw, pmx, ptxt, factor, n, al, p_idx, field_only) in paras:
            ph = n * pmx * max(factor, 0.95) * (1.0 - lnred) / 72.0
            if ptxt.strip():
                gw = pw if not wrap else min(pw, fw2)
                a = str(al) if al is not None else ""
                if "CENTER" in a:
                    x0 = fx + (fw2 - gw) / 2
                elif "RIGHT" in a:
                    x0 = fx + fw2 - gw
                else:
                    x0 = fx
                out.append(GlyphBox(x0, cy, x0 + gw, cy + ph, ptxt[:24], pmx, fid,
                                    owner_sp, cell, p_idx, field_only))
            cy += ph

    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "has_table", False):
            # Native table cells (third external review, P0: tables in auto-generated decks
            # were a geometry blind spot).
            # Cell rectangle = table origin + accumulated column widths/row heights (xf
            # scaling applied to EMU).
            geo = _geo_rect(sp, xf)
            if geo is None or geo[4]:
                continue
            tx, ty = geo[0], geo[1]
            ax = xf[0]
            ay = xf[2]
            try:
                tbl = sp.table
                col_w = [(c.width or 0) for c in tbl.columns]
                row_h = [(r.height or 0) for r in tbl.rows]
            except Exception:
                continue
            cy_off = 0
            for ri, row in enumerate(tbl.rows):
                cx_off = 0
                for ci, cell in enumerate(row.cells):
                    cw = col_w[ci] if ci < len(col_w) else 0
                    rh = row_h[ri] if ri < len(row_h) else 0
                    try:
                        # Merged regions (0.6.0, external review): continuation cells are
                        # covered by their origin; the origin's rectangle spans the merged
                        # column widths and row heights instead of a single grid cell.
                        if cell.is_spanned:
                            cx_off += cw
                            continue
                        span_w = sum(col_w[ci:ci + max(1, cell.span_width)]) or cw
                        span_h = sum(row_h[ri:ri + max(1, cell.span_height)]) or rh
                    except Exception:
                        span_w, span_h = cw, rh
                    try:
                        ctf = cell.text_frame
                    except Exception:
                        cx_off += cw
                        continue
                    emit_frame(ctf,
                               tx + ax * cx_off / EMU_PER_IN,
                               ty + ay * cy_off / EMU_PER_IN,
                               ax * span_w / EMU_PER_IN,
                               ay * span_h / EMU_PER_IN,
                               id(cell._tc), sp, cell=(ri, ci), sx=ax, sy=ay)
                    cx_off += cw
                cy_off += row_h[ri] if ri < len(row_h) else 0
            continue
        if not getattr(sp, "has_text_frame", False):
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        fx, fy, fw, fh, rotated = geo
        if rotated:
            continue
        emit_frame(sp.text_frame, fx, fy, fw, fh, id(sp), sp, sx=xf[0], sy=xf[2])
    return out


def text_overlap_check(slide, si, warns, boxes: Optional[List[GlyphBox]] = None):
    """W15: the effective glyph regions of two different text frames overlap meaningfully
    (occlusion/collision). This is approximation-based, hence WARN. Fires only when the
    intersection area exceeds 45% of the smaller box, at most 2 findings per page.
    The 45% threshold is measured against render comparisons: cases in the 30-35% range are
    all false positives (estimated as a title under a big number, or bleed-over between two
    columns), while 60%+ are all real overlaps (e.g. a chart overrunning an exhibit label, a
    caption chip sitting on a legend).
    Intentional layering is excluded: an echo of identical text (an afterimage typography
    effect), and 1-2 character oversized glyphs (drop caps, chapter numerals like I/II).
    If boxes is given, it is used as-is without recomputation (a once-per-slide computation
    cache)."""
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    hits = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            if a.frame_id == b.frame_id:      # paragraphs of the same frame are stacked,
                                               # so excluded
                continue
            if a.rep.strip() == b.rep.strip():
                continue
            if any(len(x.rep.strip()) <= 2 and x.max_pt >= 28 for x in (a, b)):
                continue
            ix = min(a.x1, b.x1) - max(a.x0, b.x0)
            iy = min(a.y1, b.y1) - max(a.y0, b.y0)
            if ix <= 0.02 or iy <= 0.02:
                continue
            area = ix * iy
            amin = min((a.x1 - a.x0) * (a.y1 - a.y0), (b.x1 - b.x0) * (b.y1 - b.y0))
            if amin > 0 and area > 0.45 * amin:
                hits.append((area / amin, a, b))
    # Sort key is frac only: GlyphBox's sp field isn't comparable, so a full tuple comparison
    # would fail (0.5.0)
    for frac, a, b in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        # loc: not the frame's raw bbox, but the effective glyph bbox (absolute, in) actually
        # used to judge the overlap, as-is.
        # related carries the counterpart frame so an agent can pinpoint which pair to move
        # (0.5.0).
        loc = shape_loc(a.sp, bbox=[a.x0, a.y0, a.x1 - a.x0, a.y1 - a.y0], cell=a.cell,
                        paragraph=a.para, field=a.field) or {}
        rel = shape_loc(b.sp, bbox=[b.x0, b.y0, b.x1 - b.x0, b.y1 - b.y0], cell=b.cell,
                        paragraph=b.para, field=b.field)
        if rel:
            loc["related"] = rel
        warns.append(Finding(si, "W15", "w15", (frac * 100,), "%r ~ %r" % (a.rep, b.rep),
                             data={"confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False},
                             loc=loc or None))


def _pic_boxes(slide, sw_in, sh_in, skipped=None):
    """The effective ink bbox (in) and z-order of non-background pictures. Returns
    [(x0,y0,x1,y1,z)].
    Full-bleed or mesh backgrounds covering 70%+ of the slide are excluded. A transparent PNG
    (e.g. a matplotlib chart) has a frame bbox much larger than its ink, which is a source of
    false positives, so it is trimmed to the alpha-opaque bbox.
    Performance budget (0.4.0, third review): for images over 25MP, alpha trimming is skipped,
    the frame bbox is used as-is, and this is disclosed via the skipped counter (prevents
    decode blowup on large batches).
    Reflects adversarial-verification measurements (2026-07-03): P mode + tRNS is converted
    to RGBA before trimming, srcRect crop is mapped by narrowing to the visible source
    window, flipH/flipV mirrors within that window, a rotated picture uses the axis-aligned
    expanded bbox but skips ink trimming since it would be invalid, and group children are
    converted to absolute coordinates."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, rotated = geo
        if w * h >= 0.7 * sw_in * sh_in:
            continue
        if not rotated:
            try:
                from PIL import Image
                import io as _io
                im = Image.open(_io.BytesIO(sp.image.blob))
                if im.width * im.height > 25_000_000:
                    if skipped is not None:
                        skipped["image_decode_budget"] += 1
                    raise RuntimeError("image decode budget")   # the except below keeps the
                                                                 # frame bbox
                if im.mode == "P" and "transparency" in im.info:
                    im = im.convert("RGBA")
                if "A" in im.getbands():
                    bb = im.getchannel("A").point(lambda a: 255 if a > 16 else 0).getbbox()
                    if bb is None:
                        continue   # fully transparent = no ink
                    iw, ih = im.size
                    wl, wt = iw * float(sp.crop_left or 0), ih * float(sp.crop_top or 0)
                    wr = iw * (1.0 - float(sp.crop_right or 0))
                    wb = ih * (1.0 - float(sp.crop_bottom or 0))
                    l, t, r, b = bb
                    l, r = max(l, wl), min(r, wr)
                    t, b = max(t, wt), min(b, wb)
                    if l >= r or t >= b:
                        continue   # no ink within the crop window
                    try:
                        x2 = sp._element.spPr.find(NS + "xfrm")
                        if x2 is not None and x2.get("flipH") == "1":
                            l, r = wl + (wr - r), wl + (wr - l)
                        if x2 is not None and x2.get("flipV") == "1":
                            t, b = wt + (wb - b), wt + (wb - t)
                    except Exception:
                        pass
                    ww, wh = (wr - wl) or 1.0, (wb - wt) or 1.0
                    x, y, w, h = (x + w * (l - wl) / ww, y + h * (t - wt) / wh,
                                  w * (r - l) / ww, h * (b - t) / wh)
            except Exception as ex:
                # The frame bbox is kept as the fallback either way, but a real decode
                # failure is no longer silent (0.6.0, external review: "nothing dies
                # silently" only held for the budget path). The budget path already
                # tallied itself above.
                if skipped is not None and str(ex) != "image decode budget":
                    skipped["image_decode"] += 1
        out.append((x, y, x + w, y + h, z, sp))
    return out


def overflow_check(slide, si, sw_in, sh_in, warns,
                   boxes: Optional[List[GlyphBox]] = None, pics: Optional[list] = None):
    """W16: off-canvas overflow. Using the frame bbox as the criterion was previously
    rejected because of the large false-positive rate from the generous-frame convention, but
    W15's effective glyph bbox resolved that objection (2026-07-03): text fires only when the
    actual character area breaches the boundary. For non-text, only pictures (ink bbox,
    trimmed) are checked: bleed where a decorative shape (e.g. a glow circle) spills off a
    corner is standard technique and not a defect (checking shapes was rejected after
    corpus render measurements)."""
    TOL_T, TOL_S = 0.15, 0.12
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    if pics is None:
        pics = _pic_boxes(slide, sw_in, sh_in)
    hits = []
    for gb in boxes:
        over = max(-gb.x0, -gb.y0, gb.x1 - sw_in, gb.y1 - sh_in)
        if over > TOL_T:
            hits.append((over, M("w16_text") % gb.rep, "t|%r" % gb.rep,
                         shape_loc(gb.sp, bbox=[gb.x0, gb.y0, gb.x1 - gb.x0, gb.y1 - gb.y0],
                                   cell=gb.cell, paragraph=gb.para, field=gb.field)))
    for (px0, py0, px1, py1, _z, psp) in pics:
        over = max(-px0, -py0, px1 - sw_in, py1 - sh_in)
        if over > TOL_S:
            hits.append((over, M("w16_pic") % (px1 - px0, py1 - py0),
                         "p|%.1fx%.1f" % (px1 - px0, py1 - py0),
                         shape_loc(psp, bbox=[px0, py0, px1 - px0, py1 - py0])))
    # Sort key is over only: a loc dict isn't comparable, so a full tuple comparison would
    # fail (0.5.0)
    for over, what, fpk, loc in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        # fp_key: detail (what) is a locale-dependent string, so it's excluded from the
        # baseline fingerprint (fourth review)
        warns.append(Finding(si, "W16", "w16", (over,), what, fp_key=fpk, loc=loc,
                             data={"kind": "text" if fpk.startswith("t|") else "picture",
                                   "confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False}))


def _occluder_boxes(slide, sw_in, sh_in):
    """The bbox and z of solid-fill shapes (cards, panels) sitting on top of a picture. Used
    to suppress a legitimate layout, a caption card over a photo, that was falsely caught by
    W17 (measured in adversarial verification)."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            continue
        try:
            if sp.fill.type is None or "SOLID" not in str(sp.fill.type):
                continue
        except Exception:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, _rot = geo
        if w * h >= 0.9 * sw_in * sh_in:
            continue
        out.append((x, y, x + w, y + h, z))
    return out


def text_image_straddle_check(slide, si, sw_in, sh_in, warns,
                              boxes: Optional[List[GlyphBox]] = None, pics: Optional[list] = None):
    """W17: text straddles the ink boundary of a non-background picture (only 25-75% of the
    glyph is inside the image) = half on, half off, so it looks cropped or the background
    looks split. Fully on top (an overlay caption) is the contrast-gate's (W7's) jurisdiction,
    and fully off is irrelevant. Pictures under 1 square inch (icon/logo scale) are ignored.
    If a solid card sits in z-order between the photo and the text and backs 90%+ of the text
    area, this is excluded as a caption on a card rather than a straddle.
    At most 2 findings per page."""
    if pics is None:
        pics = _pic_boxes(slide, sw_in, sh_in)
    if not pics:
        return
    occl = _occluder_boxes(slide, sw_in, sh_in)
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    hits = []
    for gb in boxes:
        rep = gb.rep
        if len(rep.strip()) < 3:
            continue
        ta = (gb.x1 - gb.x0) * (gb.y1 - gb.y0)
        if ta <= 0:
            continue
        for (px0, py0, px1, py1, pz, psp) in pics:
            if (px1 - px0) * (py1 - py0) < 1.0:
                continue
            ix = min(gb.x1, px1) - max(gb.x0, px0)
            iy = min(gb.y1, py1) - max(gb.y0, py0)
            if ix <= 0 or iy <= 0:
                continue
            frac = (ix * iy) / ta
            if not (0.25 <= frac <= 0.75):
                continue
            carded = False
            for (ox0, oy0, ox1, oy1, oz) in occl:
                if oz <= pz:
                    continue
                cx = min(gb.x1, ox1) - max(gb.x0, ox0)
                cy2 = min(gb.y1, oy1) - max(gb.y0, oy0)
                if cx > 0 and cy2 > 0 and cx * cy2 >= 0.9 * ta:
                    carded = True
                    break
            if not carded:
                hits.append((frac, gb, (px0, py0, px1, py1, psp)))
    # Sort key is frac only (sp field isn't comparable, 0.5.0)
    for frac, gb, pic in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        loc = shape_loc(gb.sp, bbox=[gb.x0, gb.y0, gb.x1 - gb.x0, gb.y1 - gb.y0],
                        cell=gb.cell, paragraph=gb.para, field=gb.field) or {}
        rel = shape_loc(pic[4], bbox=[pic[0], pic[1], pic[2] - pic[0], pic[3] - pic[1]])
        if rel:
            loc["related"] = rel
        warns.append(Finding(si, "W17", "w17", (frac * 100,), "%r" % gb.rep, loc=loc or None,
                             data={"confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False}))


# English finite-verb forms that mark a claim-style title. Prefer false negatives over
# false positives: copulas/auxiliaries and ambiguous noun/verb homographs are omitted.
# Measured against the issue fixtures (noun-phrase decks fire; verb/number+unit decks do not).
_EN_CLAIM_VERBS = frozenset({
    "grows", "grew", "rises", "rose", "fall", "falls", "fell", "jumps", "jumped",
    "drives", "drove", "reaches", "reached", "beats", "misses", "missed",
    "expands", "expanded", "cuts", "gains", "gained", "loses", "lost",
    "improves", "improved", "declines", "declined", "surges", "surged",
    "drops", "dropped", "hits", "tops", "topped", "leads", "led", "lags",
    "lagged", "remains", "remained", "continues", "continued", "accelerates",
    "accelerated", "slows", "slowed", "climbs", "climbed", "slides", "slid",
    "boosts", "boosted", "lifts", "lifted", "shrinks", "shrank", "narrows",
    "narrowed", "widens", "widened", "outpaces", "outpaced", "outperforms",
    "outperformed", "underperforms", "underperformed", "adds", "added",
    "opens", "opened", "closes", "closed", "launches", "launched", "ships",
    "shipped", "wins", "won", "fails", "failed", "holds", "held", "keeps",
    "kept", "sets", "set", "breaks", "broke", "posts", "posted", "reports",
    "reported", "delivers", "delivered", "fuels", "fueled", "fuelled",
})



# Structural English titles that should not enter the W14 eligibility pool.
# Cover/section/closing slides are noun phrases by nature; counting them would
# fire W14 on any three-slide deck with a title page and divider.
_EN_STRUCTURAL_TITLES = frozenset({
    "cover", "title", "title page", "agenda", "introduction",
    "overview", "appendix", "thank you", "thanks", "q&a", "q and a", "qa",
    "closing", "contents", "table of contents",
    "toc", "outline", "references", "bibliography", "the end", "end",
    "summary", "background", "agenda overview",
})


def _en_title_key(title: str) -> str:
    key = re.sub(r"\s+", " ", title.strip().lower())
    key = re.sub(r"[^a-z0-9 &]+", "", key).strip()
    return key


def _en_eligible_title(title: str, chars, latin_n: int, cjk_n: int) -> bool:
    """Whether an English title is substance enough for the W14 majority pool.

    Hangul eligibility is a content share test; English needs more than "is Latin"
    so cover/divider/closing slides do not flood the pool on clean decks.
    """
    if cjk_n >= 3:
        return False
    if latin_n < 4 or latin_n < 0.6 * len(chars):
        return False
    if _en_title_key(title) in _EN_STRUCTURAL_TITLES:
        return False
    words = re.findall(r"[A-Za-z]{2,}", title)
    # Need either multi-word substance or enough letters that it is not a label.
    if len(words) >= 2 and sum(len(w) for w in words) >= 10:
        return True
    if len(words) >= 3:
        return True
    return False


def _en_is_claim(title):
    """English W14 claim signal: finite verb from a measured allowlist, or number+unit.
    Bare noun phrases return False. Prefer false negatives on editorial headlines."""
    if "?" in title or "!" in title:
        return True
    # number + unit (%, pp/bp, x multiplier, currency); keeps the Korean contract's spirit
    if re.search(
        r"(?:[$€£]\s*[0-9]|[0-9][0-9,.]*\s*(?:%|pp|bp|[xX]|million|billion|thousand|percent))",
        title,
    ):
        return True
    words = re.findall(r"[A-Za-z]+", title)
    return any(w.lower() in _EN_CLAIM_VERBS for w in words)


# Cover / divider / closing slides in Korean decks. A structural slide is not part of the
# deck's argument, so counting it inflates W14 on both sides of the ratio. Mirrors
# _EN_STRUCTURAL_TITLES, which landed for English in #9.
_KO_STRUCTURAL_TITLES = frozenset({
    "표지", "목차", "차례", "부록", "별첨", "감사합니다", "감사", "질의응답", "질의 응답",
    "마무리", "맺음말", "끝", "참고문헌", "참고 문헌", "출처", "회사소개", "회사 소개",
    "들어가며", "시작하며", "목차 및 구성", "발표를 마칩니다", "이상입니다",
    "경청해 주셔서 감사합니다",
})

# Korean nouns whose last syllable is also a sentence ending. The ending test cannot tell
# 개요 (an overview) from 해요 (a predicate), so the collisions are listed rather than
# guessed at, the way the rule already carried a note about 바다.
#
# The list is incomplete on purpose and safe to leave that way: -자 is a productive agent
# suffix, so no hand-written list keeps up with it, and a word that is absent simply keeps
# the previous reading (counted as a claim, rule stays quiet). Every residual miss
# therefore lands on the false-negative side, which is the side W14 is supposed to err on.
# Words that can also be predicates in a title (가요, 포함, 약함, 이자) are deliberately
# left out: listing them would move errors to the firing side.
_KO_NOUN_FINALS = frozenset({
    "개요", "필요", "중요", "주요", "수요", "소요", "강요", "동요",
    "투자자", "사용자", "소비자", "참가자", "참여자", "가입자", "이용자", "구독자",
    "시청자", "개발자", "경영자", "창업자", "실무자", "담당자", "책임자", "관리자",
    "신청자", "응답자", "협력자", "수혜자", "지원자", "후보자", "방문자", "구매자",
    "판매자", "근로자", "노동자", "종사자", "대상자", "경쟁자", "설계자", "운영자",
    "제작자", "독자", "저자", "기자", "학자", "환자",
    "투자", "숫자", "글자", "문자", "의자", "모자", "상자", "전자", "한자",
    "게임", "모임", "책임", "쓰임",
    "바다", "소다",
})

# Sentence-final interrogatives written without a question mark. Recognising these moves a
# title out of the nominal count, which lowers the firing rate, so they are safe to add.
_KO_INTERROGATIVE_SUFFIX = ("인가", "는가", "은가", "던가", "을까")

_KO_SENTENCE_ENDINGS = ("다", "까", "요", "자", "죠", "함", "임")
_KO_NUMERIC_CLAIM = re.compile(
    r"[0-9][0-9,.]*\s*(%|배|억|조|만|천|pp|bp|x|X|원|건|명|개)")


def _ko_title_key(title):
    return re.sub(r"\s+", " ", title.strip()).rstrip(" ?!.…”’")


def _ko_is_claim(title):
    """Hangul W14 claim signal: a question, a number+unit, a sentence-final interrogative,
    or a sentence ending that is not the tail of a known noun.

    The number+unit branch is there because a title like "매출 3배 성장" is a claim-style
    headline even though it ends in a noun (external review, 2026-07-10)."""
    if "?" in title or "!" in title:
        return True
    if _KO_NUMERIC_CLAIM.search(title):
        return True
    core = _ko_title_key(title)
    if core.endswith(_KO_INTERROGATIVE_SUFFIX):
        return True
    if not core.endswith(_KO_SENTENCE_ENDINGS):
        return False
    return core.split(" ")[-1] not in _KO_NOUN_FINALS


def action_title_check(titles, warns):
    """W14: a majority of titles are descriptive noun phrases (e.g. "Market Overview,"
    "Competitive Analysis") = not action titles (the MBB idea that reading only the titles
    should carry the argument). Hangul titles use a sentence-ending heuristic minus the
    nouns that collide with it (see `_ko_is_claim`); English titles use a finite-verb
    allowlist or number+unit (see `_en_is_claim`). Structural titles are not eligible in
    either language (cover, agenda, appendix, closing / 표지, 목차, 부록, 감사합니다), so a
    clean deck with a title page and divider does not false-fire. Prefer false negatives
    on ambiguous verb/noun forms. Fires once per deck only when 3+ eligible titles are
    noun phrases and they make up at least half of eligible titles. Gate: full profile
    (editorial skips)."""
    entries = []
    for si in sorted(titles):
        txt = " ".join(titles[si][1]).strip()
        chars = [c for c in txt if not c.isspace()]
        if len(chars) < 4:
            continue
        cjk_n = sum(1 for c in chars if is_cjk(c))
        latin_n = sum(1 for c in chars if ("A" <= c <= "Z") or ("a" <= c <= "z"))
        # Hangul path: excludes titles with fewer than 3 Hangul characters or under 30%
        # Hangul share (big-stat numbers / short brand names; measured in the 50-deck scan).
        if cjk_n >= 3 and cjk_n >= 0.3 * len(chars):
            if _ko_title_key(txt) in _KO_STRUCTURAL_TITLES:
                continue
            entries.append((si, txt, _ko_is_claim(txt)))
            continue
        # English path: content-eligible Latin titles only (not structural covers).
        if _en_eligible_title(txt, chars, latin_n, cjk_n):
            entries.append((si, txt, _en_is_claim(txt)))
    nominal = [(si, t) for si, t, c in entries if not c]
    if len(nominal) >= 3 and len(nominal) * 2 >= len(entries):
        ex = " ".join("p%d'%s'" % (si, t[:14]) for si, t in nominal[:4])
        warns.append(Finding(0, "W14", "w14", (len(nominal), len(entries)), ex))
    return entries


def _diagram_clone_marks(inter):
    """Counts "decorative texture clones" (small dots, joints, etc., 1x1 or smaller on the
    24-grid) in the multiset of fill shapes shared between two pages.
    Card-shaped (mid-size block) and table (full-width band) elements are W6's jurisdiction
    and are not counted here, avoiding a false positive on three-column comparison cards
    (reflecting the adversarial audit)."""
    marks = area = 0
    for (fh, gx, gy, gw, gh), c in inter.items():
        if gw <= 1 and gh <= 1:
            marks += c
        if gw >= 1 and gh >= 1:
            area += gw * gh * c
    if marks >= 8 and area / (24.0 * 24.0) >= 0.06:
        return marks
    return 0


def slide_layout_sig(slide, sw, sh, gw=6, gh=4):
    """Turns slide element placement into a gw x gh grid occupancy vector. Full-bleed
    backgrounds are excluded.
    Weights: text=1, shape=0.5, image=2. Used to compare "what is where" (the skeleton)."""
    sig = [0.0] * (gw * gh)
    n = 0
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:      # excludes full-bleed backgrounds/images
            continue
        cx = (L + Wd / 2) / sw; cy = (T + Ht / 2) / sh
        if not (0 <= cx <= 1.0 and 0 <= cy <= 1.0):
            continue
        gc = min(gw - 1, max(0, int(cx * gw))); gr = min(gh - 1, max(0, int(cy * gh)))
        wgt = 2.0 if _is_pic(sp) else (1.0 if getattr(sp, "has_text_frame", False) else 0.5)
        sig[gr * gw + gc] += wgt
        n += 1
    return sig, n








# Sentinel: an explicit color exists but the decoder cannot resolve it (hslClr, scrgbClr,
# sysClr, prstClr, tint/shade transforms...). Falling through to inherited colors produced
# W7 false positives, e.g. an explicit white hslClr run judged with an inherited black
# (0.6.1, external review): an unknown explicit color must stop resolution, not be skipped.




def contrast_check(slide, si, sw, sh, render_dir, warns, styler=None, thm_colors=None,
                   skipped=None):
    """Detects low-contrast text over an image (approximated from the rendered PNG). Only
    text frames overlapping a picture, once per slide.
    Returns: "no_pics" (nothing to check) / "no_png" (there is a picture but no conventional
    render = incomplete) / "ok" (checked).
    Coordinates are absolute, including the group transform (third review P1: fixes pictures
    and text inside a group being misaligned when using raw coordinates)."""
    from PIL import Image
    pics = []
    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if _is_pic(sp):
            geo = _geo_rect(sp, xf)
            if geo is not None:
                x, y, w, h, _rot = geo
                pics.append((x * EMU_PER_IN, y * EMU_PER_IN, w * EMU_PER_IN, h * EMU_PER_IN))
    if not pics:
        return "no_pics"
    cand = glob.glob(os.path.join(render_dir, "p%02d.png" % si))
    if not cand:
        return "no_png"
    try:
        im = Image.open(cand[0]).convert("RGB"); px = im.load(); PW, PH = im.size
    except Exception:
        return "no_png"
    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        gx, gy, gw_, gh_, _rot = geo
        L, T, Wd, Ht = gx * EMU_PER_IN, gy * EMU_PER_IN, gw_ * EMU_PER_IN, gh_ * EMU_PER_IN
        over = any(not (L + Wd <= p0 or L >= p0 + pw or T + Ht <= q0 or T >= q0 + ph)
                   for p0, q0, pw, ph in pics)
        if not over:
            continue
        rgbs = [_resolve_run_rgb(r, para, sp.text_frame, sp, slide, styler, thm_colors)
                for para in sp.text_frame.paragraphs for r in para.runs if r.text.strip()]
        if any(c is _COLOR_UNKNOWN for c in rgbs):
            # An explicit color we cannot decode: judging with an inherited color instead
            # produced false positives (0.6.1). Abstain on this frame and surface it.
            if skipped is not None:
                skipped["w7_color_unknown"] += 1
            continue
        rgbs = [c for c in rgbs if c]
        if not rgbs:
            continue
        txt_rgb = rgbs[0]
        x0 = max(0, int(L / sw * PW)); y0 = max(0, int(T / sh * PH))
        x1 = min(PW, int((L + Wd) / sw * PW)); y1 = min(PH, int((T + Ht) / sh * PH))
        if x1 <= x0 or y1 <= y0:
            continue
        sx = max(1, (x1 - x0) // 24); sy = max(1, (y1 - y0) // 24)
        lumas = sorted(_luma(px[x, y]) for x in range(x0, x1, sx) for y in range(y0, y1, sy))
        if not lumas:
            continue
        L_txt = _luma(txt_rgb)
        # Background luma is taken from the quantile at the opposite extreme of the text
        # color: this measures worst-case local contrast (WCAG is based on the worst case)
        # and avoids mean contamination from mixed-in text ink (dark 15th percentile for light
        # text, light 85th percentile for dark text).
        if L_txt >= 0.5:
            L_bg = lumas[int(len(lumas) * 0.15)]
        else:
            L_bg = lumas[min(len(lumas) - 1, int(len(lumas) * 0.85))]
        hi = max(L_bg, L_txt); lo = min(L_bg, L_txt)
        ratio = (hi + 0.05) / (lo + 0.05)
        if ratio < 2.5:
            warns.append(Finding(si, "W7", "w7", (ratio,),
                                 "text=%r" % sp.text_frame.text[:20],
                                 data={"confidence": "measured",
                                       "evidence_source": "render",
                                       "render_confirmed": True},
                                 loc=shape_loc(sp, bbox=[gx, gy, gw_, gh_])))
            return "ok"
    return "ok"   # found and checked the render PNG for this page (regardless of whether
                  # W7 actually fired)


def collect_wireframe(path):
    """Per-slide shape geometry for the annotated HTML report (0.8.x): slide size in
    inches and, per slide, every shape's absolute bbox with a coarse kind
    (text/picture/other) and a short text sample. Read-only and best-effort: a shape
    whose geometry cannot be resolved is simply omitted (the report draws what it can;
    findings carry their own bboxes independently)."""
    prs = Presentation(path)
    sw, sh = prs.slide_width / EMU_PER_IN, prs.slide_height / EMU_PER_IN
    slides = []
    for slide in prs.slides:
        shapes = []
        for sp, _z, xf in iter_shapes_geo(slide.shapes):
            geo = _geo_rect(sp, xf)
            if geo is None:
                continue
            x, y, w, h, _rot = geo
            if _is_pic(sp):
                kind = "picture"
            elif getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
                kind = "text"
            else:
                kind = "other"
            sample = ""
            try:
                if kind == "text":
                    sample = sp.text_frame.text.strip().splitlines()[0][:40]
            except Exception:
                pass
            shapes.append({"x": x, "y": y, "w": w, "h": h, "kind": kind,
                           "text": sample})
        slides.append(shapes)
    return {"sw": sw, "sh": sh, "slides": slides}


def _zip_preflight(path):
    """A cheap decompression-bomb screen before python-pptx parses the package (0.6.0,
    external review: the linter takes untrusted input, and the only budgets so far were
    per-image decode and pair caps). Entry count, total uncompressed size, and per-entry
    compression ratio are bounded. A violation raises ValueError (a hostile or broken
    file is a usage error, not a deck finding); budgets sit far above any real deck."""
    import zipfile
    try:
        with zipfile.ZipFile(path) as z:
            infos = z.infolist()
    except (OSError, zipfile.BadZipFile) as e:
        raise ValueError("not a readable pptx package: %s" % e)
    if len(infos) > 20000:
        raise ValueError("zip preflight: too many entries (%d > 20000)" % len(infos))
    total = sum(i.file_size for i in infos)
    if total > 2_000_000_000:
        raise ValueError("zip preflight: uncompressed size %d exceeds the 2GB budget" % total)
    for i in infos:
        if i.file_size > 10_000_000 and i.compress_size \
                and i.file_size / i.compress_size > 200:
            raise ValueError("zip preflight: suspicious compression ratio in %r" % i.filename)


def lint(path, hard_min=5.0, body_min=9.0, small_min=7.5, render_dir=None, ghost=None,
         strict=False, w6_sim=0.90, w6_min_cluster=3, profile=DEFAULT_PROFILE):
    """profile is an execution policy (third external review P0: applied at the engine stage,
    not as a CLI post-filter).
    Excluded rules are not run at all, so there is no O(S^2) comparison cost either, internal
    failures of an excluded rule cannot leak into W18, and library callers can use profiles
    too.

    0.4.0 breaking change: the default profile is now core (objective defects only). To also
    get AI-tell and house-style rules (E2, W6, W9-W14), pass profile="full" explicitly. Fixes
    a first-impression problem where a first-time user's no-option run exited 1 over normal
    punctuation (external strategy review, confirmed by the user)."""
    if profile not in PROFILES:
        # Fixes a typo'd profile silently behaving like full (an empty exclusion set)
        # (fourth review)
        raise ValueError("unknown profile %r (choices: %s)" % (profile, ", ".join(sorted(PROFILES))))
    _zip_preflight(path)
    # A package that passes the zip preflight can still be un-parseable (a truncated or
    # attribute-corrupted part). python-pptx raises lxml/opc errors that are not
    # ValueError; the library contract is a controlled ValueError the CLI maps to a clean
    # exit, never a raw traceback (0.7.1, external review: the fuzzer must see no other
    # exception class escape).
    try:
        prs = Presentation(path)
        sw, sh = prs.slide_width, prs.slide_height
    except ValueError:
        raise
    except Exception as e:
        raise ValueError("could not open %r as a valid .pptx: %s: %s"
                         % (path, type(e).__name__, e))
    errors, warns = [], []
    sigs = []
    toks = {}
    page_texts = {}
    foot_tops = {}
    fx_pp = {}
    titles = {}
    excl = PROFILES.get(profile, frozenset())
    fonts_map = theme_fonts_by_master(prs)
    colors_map = theme_colors_by_master(prs) if render_dir else {}
    colors_default = next((v for v in colors_map.values() if v is not None), None)
    thm_default = next((v for v in fonts_map.values() if v is not None), None)
    deck_skipped = Counter()   # deck-level unable-to-check (used to fire W18 p00)
    if render_dir and not os.path.isdir(render_dir):
        # The path where W7 silently ran zero checks if the --render folder didn't exist
        # (third review P0): surfaced as incomplete instead
        deck_skipped["render_dir_missing"] += 1
        print(M("note_render_dir_missing") % render_dir, file=sys.stderr)
    theme_fails = sum(1 for v in fonts_map.values() if v is None)
    if theme_fails:
        # Parse failure (None) and a confirmed empty slot ("") fall back to the same
        # fallback assumption in the E1 branch.
        # This is the point where the distinction disappears, so it's surfaced via W18
        # (adversarial panel, second re-check: prevents silent collapse).
        deck_skipped["theme_parse"] = theme_fails
        print(M("note_theme_parse"), file=sys.stderr)
    styler = StyleResolver(prs)

    render_png_hits = 0
    for si, slide in enumerate(prs.slides, 1):
        # theme font slots of the master this slide uses (prevents E1 false firing in
        # multi-master decks)
        thm_fonts = thm_default
        thm_colors = colors_default
        try:
            pn = str(slide.slide_layout.slide_master.part.partname)
            thm_fonts = fonts_map.get(pn, thm_default)
            thm_colors = colors_map.get(pn, colors_default)
        except Exception:
            pass
        skipped = Counter()   # W18: tallies unable-to-check regions (surfaces silent
                               # degradation in the JSON)
        try:
            slide_part = str(slide.part.partname)   # for the part field of a finding
                                                     # location
        except Exception:
            slide_part = None
        # sig and toks are each guarded separately: wrapping them in one try would, if sig
        # succeeds but toks fails, cause the except's re-append to push sigs out of order and
        # misalign W6 page numbers (reproduced and measured in the adversarial panel).
        # Collection and checking for a rule excluded by the profile is not run at all
        # (third review P0).
        if "W6" not in excl:
            try:
                sig = slide_layout_sig(slide, sw, sh)
            except Exception as e:
                sig = ([0.0] * 24, 0)   # preserves position (sigs is indexed by page order)
                skipped["w6_sig"] += 1
                print("W6 sig skipped p%02d: %s" % (si, e), file=sys.stderr)
            sigs.append(sig)
        if "W10" not in excl:
            try:
                toks[si] = _fill_tokens(slide, sw, sh)
            except Exception as e:
                toks[si] = Counter()
                skipped["w10_tokens"] += 1
                print("W10 tokens skipped p%02d: %s" % (si, e), file=sys.stderr)
        if render_dir:
            try:
                r7 = contrast_check(slide, si, sw, sh, render_dir, warns,
                                    styler=styler, thm_colors=thm_colors, skipped=skipped)
                if r7 == "ok":
                    render_png_hits += 1
                elif r7 == "no_png":
                    # a page with a picture but no conventional render = W7 not run for this
                    # page (incomplete)
                    skipped["w7_no_render"] += 1
            except Exception as e:
                skipped["w7"] += 1
                print("W7 skipped p%02d: %s" % (si, e), file=sys.stderr)
        if "W9" not in excl:
            try:
                accent_vbars_check(slide, si, sw, sh, warns)
            except Exception as e:
                skipped["w9"] += 1
                print("W9 skipped p%02d: %s" % (si, e), file=sys.stderr)
        if "W12" not in excl or "W13" not in excl:
            try:
                foot_tops[si] = footer_top(slide, sw, sh)
                fx_pp[si] = effects_count(slide)
            except Exception as e:
                skipped["w12_w13"] += 1
                print("W12/W13 skipped p%02d: %s" % (si, e), file=sys.stderr)
        # Base geometry data is computed only once per slide (previously W15-W17 each
        # recomputed it independently: text boxes 3 times, PIL picture decode twice. Fixes
        # a large-deck performance issue, external review 2026-07-10)
        sw_in, sh_in = sw / EMU_PER_IN, sh / EMU_PER_IN
        try:
            tboxes = _text_glyph_boxes(slide, skipped=skipped, styler=styler)
        except Exception as e:
            tboxes = None
            skipped["glyph_boxes"] += 1
            print("W15/W16/W17 glyph boxes skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            pboxes = _pic_boxes(slide, sw_in, sh_in, skipped=skipped)
        except Exception as e:
            pboxes = None
            skipped["pic_boxes"] += 1
            print("W16/W17 pic boxes skipped p%02d: %s" % (si, e), file=sys.stderr)
        # Only the axis that failed falls back to an empty list; the axis that is still alive
        # keeps being checked: fixes a shared gate where a tboxes failure was also silencing
        # the unrelated picture W16 check (2 cases reproduced and measured in the adversarial
        # panel).
        if tboxes is not None:
            try:
                text_overlap_check(slide, si, warns, boxes=tboxes)
            except Exception as e:
                skipped["w15"] += 1
                print("W15 skipped p%02d: %s" % (si, e), file=sys.stderr)
        try:
            overflow_check(slide, si, sw_in, sh_in, warns,
                           boxes=tboxes if tboxes is not None else [],
                           pics=pboxes if pboxes is not None else [])
            if tboxes is not None and pboxes is not None:
                text_image_straddle_check(slide, si, sw_in, sh_in, warns, boxes=tboxes, pics=pboxes)
        except Exception as e:
            skipped["w16_w17"] += 1
            print("W16/W17 skipped p%02d: %s" % (si, e), file=sys.stderr)
        # Core-line gates (E1-E4, W1/W5/W8). The guard is per run: a per-frame guard let one
        # run's garbage attribute swallow a real violation in a neighboring run of the same
        # frame, producing a false pass (reproduced and measured in the adversarial panel,
        # 2026-07-10). Swallowed regions are surfaced in the JSON via W18.
        try:
            frames = collect_frames(slide.shapes)
        except Exception as e:
            frames = []
            skipped["frames"] += 1
            print("E1-E4 frames skipped p%02d: %s" % (si, e), file=sys.stderr)
        for tf, w_emu, owner_sp, cell_rc, sp_xf in frames:
            try:
                fw_in = w_emu / EMU_PER_IN
                scale = frame_font_scale(tf)
                paragraphs = list(tf.paragraphs)
                frame_fam = styler.ph_family_of(owner_sp)
            except Exception as e:
                skipped["frame"] += 1
                print("E1-E4 skipped p%02d frame: %s" % (si, e), file=sys.stderr)
                continue
            for pi, para in enumerate(paragraphs):
                try:
                    runs = list(para.runs)
                    # In-document-order items (run_like, index into para.runs or None,
                    # is_fld):
                    # a:fld (an auto field) is rendered with the same rPr as a normal run, so
                    # it must pass the same gates, and a:br is treated as a single line-break
                    # character in E2 context/offsets (carried over from the fourth review,
                    # 0.5.0). If the a:r count doesn't match, falls back to the previous runs
                    # path.
                    items = []
                    try:
                        r_seen = 0
                        for child in para._p:
                            tag = child.tag
                            if tag == NS + "r":
                                if r_seen < len(runs):
                                    items.append((runs[r_seen], r_seen, False))
                                    r_seen += 1
                            elif tag == NS + "br":
                                items.append((None, None, False))
                            elif tag == NS + "fld":
                                items.append((_FldRun(child), None, True))
                        if r_seen != len(runs):
                            items = [(r, i, False) for i, r in enumerate(runs)]
                    except Exception:
                        items = [(r, i, False) for i, r in enumerate(runs)]
                    run_offs = []
                    pos = 0
                    pieces = []
                    for run_like, _ri, _isf in items:
                        piece = "\n" if run_like is None else run_like.text
                        run_offs.append(pos)
                        pos += len(piece)
                        pieces.append(piece)
                    ptext = "".join(pieces)
                    p_fonts = para_fonts(para)   # paragraph defRPr: ranks right after run
                                                 # rPr (probe 7)
                    try:
                        para_size = para.font.size
                    except Exception:
                        para_size = None
                        skipped["para_size"] += 1
                except Exception as e:
                    skipped["para"] += 1
                    print("E1-E4 skipped p%02d para: %s" % (si, e), file=sys.stderr)
                    continue
                for ii, (run, ri, is_fld) in enumerate(items):
                    try:
                        if run is None:   # a:br: contributes context only, not itself checked
                            continue
                        t = run.text
                        if not t:
                            continue
                        if not is_fld:
                            page_texts.setdefault(si, []).append(t)
                        lvl = getattr(para, "level", 0)

                        # E1: effective render font judgment (measured-by-render model, see
                        # e1_violation).
                        # Per-script judgment: Hangul uses the full measured model; kana and
                        # Hanja only fire on fonts with no CJK at all (Inter, mono families)
                        # and pass JP/SC subset fonts (third panel).
                        # A slot missing from run rPr is filled in from the lstStyle
                        # inheritance chain (measured in probe 6: the master lstStyle's a:ea
                        # is actually inherited into rendering); the title family uses the
                        # theme majorFont ea.
                        script = None
                        if has_hangul(t):
                            script = "hangul"
                        elif any(is_kana(c) or is_hanja(c) for c in t):
                            script = "cjk_other"
                        if script:
                            fonts = run_fonts(run)
                            for slot in ("ea", "latin"):
                                if slot not in fonts and slot in p_fonts:
                                    fonts[slot] = p_fonts[slot]   # paragraph defRPr
                                                                  # (probe 7: beats lstStyle)
                            for slot in ("ea", "latin"):
                                if slot not in fonts:
                                    try:
                                        inh = styler.resolve_font(tf, owner_sp, slide, lvl, slot)
                                    except Exception:
                                        inh = None
                                    if inh:
                                        fonts[slot] = inh
                            fonts = resolve_font_tokens(fonts, thm_fonts)
                            eff_thm_ea = (thm_fonts or {}).get("mj-ea" if frame_fam == "title" else "mn-ea")
                            v = e1_violation(t, fonts, eff_thm_ea, script)
                            if v is not None:
                                errors.append(Finding(si, "E1", v[0], (), v[1],
                                                      data=v[2] if len(v) > 2 else None,
                                                      loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))

                        # E2: long-dash class. Context is read from the whole paragraph
                        # (ptext) but only this run's span is reported: prevents a false
                        # positive on a '2020'/'-2021' range split across run boundaries
                        # (confirmed in the adversarial panel)
                        bad = [] if "E2" in excl else \
                            dash_violations(ptext, strict=strict,
                                            span=(run_offs[ii], run_offs[ii] + len(t)))
                        if bad:
                            cps = ["U+%04X" % ord(c) for c in sorted(set(bad))]
                            errors.append(Finding(si, "E2", "e2", (),
                                                  "cp=%s text=%r" % (",".join(cps), t[:24]),
                                                  data={"characters": cps,
                                                        "function": "sentence_punctuation",
                                                        "strict": bool(strict)},
                                                  loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))

                        # E3 / W1 / W5: effective font size (run -> paragraph -> placeholder
                        # inheritance chain, reflecting autofit).
                        # size_src is used to judge title-collection eligibility: a
                        # defaultTextStyle fallback (18pt) is valid for gating purposes but is
                        # not an intended title size, and would flood ghost/W14 (regression
                        # fix).
                        size_src = "explicit"
                        if run.font.size is not None:
                            base_pt = run.font.size.pt
                        elif para_size is not None:
                            base_pt = para_size.pt
                        else:
                            try:
                                base_pt, size_src = styler.resolve_size(tf, owner_sp, slide, lvl)
                            except Exception:
                                base_pt, size_src = None, None
                        if base_pt is not None:
                            eff = base_pt * scale
                            # Title eligibility: explicit size only, or a title-family
                            # placeholder. If an inherited size from lstStyle or the master
                            # bodyStyle exceeds 18pt, body prose gets swept into ghost/W14
                            # (third adversarial panel: reproduced and measured with a 20pt
                            # body using its own lstStyle)
                            # An auto field (e.g. slide number) is not a title even if its
                            # size is large (0.5.0)
                            if eff >= 18 and t.strip() and not is_fld \
                                    and (size_src == "explicit" or frame_fam == "title"):
                                # If a title placeholder exists, it is the title regardless
                                # of size: fixes a 60pt KPI big-number pushing out the real
                                # 26pt title (third review)
                                is_title_ph = frame_fam == "title"
                                cur = titles.get(si)
                                if cur is None or (is_title_ph and not cur[2]) \
                                        or (is_title_ph == cur[2] and eff > cur[0] + 0.1):
                                    titles[si] = (eff, [t], is_title_ph)
                                elif is_title_ph == cur[2] and abs(eff - cur[0]) <= 0.1:
                                    cur[1].append(t)
                            if eff < hard_min:
                                note = "" if scale == 1.0 else M("e3_note") % (base_pt, scale)
                                errors.append(Finding(si, "E3", "e3", (eff, hard_min, note), "text=%r" % t[:24],
                                                      data={"nominal_pt": base_pt,
                                                            "autofit_scale": scale,
                                                            "size_source": size_src},
                                                      loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))
                            elif eff < body_min and fw_in > 4.0 and len(ptext) >= 40:
                                warns.append(Finding(si, "W1", "w1", (eff, body_min),
                                                     "w=%.1fin len=%d text=%r" % (fw_in, len(ptext), ptext[:24]),
                                                     loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))
                            elif eff < small_min and has_cjk(t) and fw_in <= 4.0:
                                # Narrow frames only (<=4in): small Hangul in a wide frame is
                                # more likely to be a caption or annotation rather than being
                                # inside a mockup or card, which would conflict with the
                                # message (mockup assumption) (reflects the public hygiene
                                # audit).
                                warns.append(Finding(si, "W8", "w8", (eff, small_min),
                                                     "w=%.1fin text=%r" % (fw_in, t[:24]),
                                                     loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))
                        else:
                            warns.append(Finding(si, "W5", "w5", (), "text=%r" % t[:24],
                                                 loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))

                        # E4: 2+ consecutive Hangul/Hanja characters + meaningfully positive
                        # tracking. A run mixed with kana is excluded as Japanese (spreading
                        # kana tracking is normal design practice). 0.6.1 (external review):
                        # the run must actually contain Hangul. Tracking on a Hanja-only run
                        # is legitimate convention in Chinese typography, and flagging it as
                        # a universal ERROR contradicted the "other scripts are never falsely
                        # flagged" scope promise. Hanja still counts toward the consecutive
                        # requirement when mixed with Hangul (Korean names, legal terms).
                        if not any(is_kana(c) for c in t) \
                                and any(is_hangul(c) for c in t) \
                                and sum(1 for c in t if is_hangul(c) or is_hanja(c)) >= 2:
                            tr = run_track(run)
                            if tr is not None and tr > 50:
                                # OOXML spc is hundredths of a point: state the unit
                                # explicitly (0.8, external review: raw "tracking": 200
                                # left the consumer guessing pt vs raw)
                                errors.append(Finding(si, "E4", "e4", (tr,), "text=%r" % t[:24],
                                                      data={"tracking_raw_hundredths_pt": tr,
                                                            "tracking_pt": tr / 100.0},
                                                      loc=shape_loc(owner_sp, paragraph=pi, run=ri, part=slide_part, cell=cell_rc, xf=sp_xf, field=is_fld)))
                    except Exception as e:
                        skipped["run"] += 1
                        print("E1-E4 skipped p%02d run: %s" % (si, e), file=sys.stderr)

        # W18: surfaces regions a guard swallowed on this page into the output contract
        # (JSON, text).
        # If it only stayed in stderr, a CI that only looks at the exit code and JSON summary
        # would misread an incomplete check as a pass (confirmed in the adversarial panel,
        # 2026-07-10). Promoted to exit 1 under --strict.
        if skipped:
            det = ", ".join("%s=%d" % (k, v) for k, v in sorted(skipped.items()))
            warns.append(Finding(si, "W18", "w18_page", (), det))

    # If there are zero matches for the naming convention (p01.png) but PNGs with other names
    # exist, hint at the naming convention.
    # Incompleteness itself is surfaced via W18/incomplete by the per-page w7_no_render
    # counter (0.3.1).
    if render_dir and os.path.isdir(render_dir) and render_png_hits == 0:
        anypng = glob.glob(os.path.join(render_dir, "*.png"))
        if anypng:
            print(M("note_render_naming") % (render_dir, os.path.basename(anypng[0])), file=sys.stderr)

    # W6: recycled layout skeleton. Sparse slides (dividers, covers, etc. with fewer than 4
    # non-empty cells) are excluded, and a warning fires if a content slide has w6_min_cluster
    # or more other slides similar to it (>w6_sim).
    # Invariant to deck length (an overall pair ratio would dilute and miss local recycling in
    # a large deck, so this judges by the largest cluster instead).
    # Thresholds are CLI-tunable (--w6-sim/--w6-cluster): a house with strong intentional
    # template consistency can tighten them to suppress this (fixes a genre-blindness
    # complaint, external review 2026-07-10; total suppression is --skip W6).
    try:
        content = [] if "W6" in excl else \
            [(i + 1, sig) for i, (sig, n) in enumerate(sigs) if n >= 3]
        if len(content) > 200:
            # Performance budget (0.4.0): a cap on the O(S^2) pairwise comparison. The
            # truncation is disclosed.
            deck_skipped["w6_capped"] += 1
            content = content[:200]
        if len(content) >= w6_min_cluster + 1:
            adj = {p: [] for p, _ in content}
            for a in range(len(content)):
                for b in range(a + 1, len(content)):
                    pa, sa = content[a]; pb, sb = content[b]
                    sim = _cosv(sa, sb)
                    if sim > w6_sim:
                        adj[pa].append((pb, sim)); adj[pb].append((pa, sim))
            worst_p, worst = max(adj.items(), key=lambda kv: len(kv[1]))
            if len(worst) >= w6_min_cluster:
                ex = " ".join("p%d~p%d(%.2f)" % (worst_p, b, s) for b, s in sorted(worst, key=lambda x: -x[1])[:4])
                # fp_key from the skeleton signature itself, not the page list: page
                # numbers in the fingerprint broke baseline suppression on slide
                # insertion, the exact failure fingerprint v2 exists to prevent
                # (0.6.0, external verification finding)
                sig_map = dict(content)
                fpk = "w6|" + ",".join("%.2f" % v for v in sig_map.get(worst_p, ()))
                warns.append(Finding(0, "W6", "w6", (len(worst) + 1,), M("w6_detail") % ex,
                                     fp_key=fpk))
    except Exception as e:
        deck_skipped["w6"] += 1
        print("W6 skipped: %s" % e, file=sys.stderr)

    # W11 copy cliches, W12 footer alignment, W13 effects, W14 action titles (all deck-level,
    # profile-gated)
    try:
        if "W11" not in excl:
            copy_cliche_check(page_texts, warns)
        if "W12" not in excl:
            footer_check(foot_tops, warns)
        if "W13" not in excl:
            effects_check_deck(fx_pp, warns)
        if "W14" not in excl:
            action_title_check(titles, warns)
        if ghost is not None:
            # ghost is every title (18pt+) collected regardless of W14's Hangul filter: since
            # it is meant for reviewing horizontal logic by reading only titles even in an
            # English/numeric deck, it is pulled from the raw titles rather than the filtered
            # entries.
            ghost.extend((si, " ".join(titles[si][1]).strip()) for si in sorted(titles))
    except Exception as e:
        deck_skipped["w11_w14"] += 1
        print("W11/W12/W13/W14 skipped: %s" % e, file=sys.stderr)

    # W10: clone recycling of a hand-drawn diagram (e.g. cross-section). Specialized to the
    # decorative-texture (marks) path so that three-column card layouts are left to W6 and
    # not counted here (adversarial audit 2026-07-02: avoids a card-shaped subblock false
    # positive). Fills the gap W6 misses, sub-0.90 similarity single pairs (measured on
    # repeated cross-section diagrams), using pptx shapes alone.
    try:
        cadj = {} if "W10" in excl else {p: [] for p in toks}
        nums_t = sorted(cadj)
        if len(nums_t) > 200:
            deck_skipped["w10_capped"] += 1
            nums_t = nums_t[:200]
        for ia in range(len(nums_t)):
            for ib in range(ia + 1, len(nums_t)):
                pa, pb = nums_t[ia], nums_t[ib]
                if _diagram_clone_marks(toks[pa] & toks[pb]):
                    cadj[pa].append(pb); cadj[pb].append(pa)
        if cadj:
            cw, cwl = max(cadj.items(), key=lambda kv: len(kv[1]))
            if len(cwl) >= 1:
                grp = sorted({cw, *cwl})
                pages_key = ",".join("p%d" % p for p in grp)
                # fp_key from the cloned diagram's shared fill tokens, not the page list
                # (page-independent fingerprint contract, 0.6.0)
                try:
                    inter = set(toks[grp[0]])
                    for p in grp[1:]:
                        inter &= set(toks[p])
                    fpk = "w10|" + ",".join(sorted(str(t) for t in inter))[:120]
                except Exception:
                    fpk = "w10|n=%d" % len(grp)
                warns.append(Finding(0, "W10", "w10", (len(grp),),
                                     M("w10_detail") % pages_key, fp_key=fpk))
    except Exception as e:
        deck_skipped["w10"] += 1
        print("W10 skipped: %s" % e, file=sys.stderr)

    # Deck-level W18: also surfaces unable-to-check outcomes from deck-level checks (W6, W10,
    # W11-W14) and theme parsing into the output contract. 0.2.0 only tallied the geometry
    # and core gates, contradicting the documentation's "surfaces everything" claim
    # (confirmed in the second external re-check: fixes the partial W18 implementation).
    if deck_skipped:
        det = ", ".join("%s=%d" % (k, v) for k, v in sorted(deck_skipped.items()))
        warns.append(Finding(0, "W18", "w18_deck", (), det))

    return errors, warns


def _skill_res():
    """The SKILL.md skill pack bundled with the package (importlib.resources, py3.9+).
    Multi-argument joinpath requires 3.11+, so this uses a chain of / instead, which is safe
    down to the 3.9 zip loader."""
    from importlib import resources
    return resources.files("archforge") / "skills" / "archforge-pptx-lint" / "SKILL.md"


def skill_main(argv=None):
    """`archforge skill`: prints the bundled skill pack or installs it into an agent skill
    folder.
    Fixes a distribution gap where the skill pack wasn't shipped to pip users (external
    review, 2026-07-10): SKILL.md is now bundled in the wheel and can be fetched directly via
    this subcommand."""
    ap = argparse.ArgumentParser(prog="archforge skill", description=M("skill_desc"))
    ap.add_argument("--install", nargs="?", const="", metavar="DIR", help=M("help_skill_install"))
    ap.add_argument("--path", action="store_true", help=M("help_skill_path"))
    # Program-wide flag: the value was already applied by main()'s prescan, so this just
    # accepts it (without this, `archforge skill --lang ko` would die as unrecognized: measured
    # in the third adversarial panel)
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    src = _skill_res()
    if a.path:
        print(str(src))
        return 0
    text = src.read_text(encoding="utf-8")
    if a.install is not None:
        root = a.install or os.path.join(".claude", "skills")
        dst_dir = os.path.join(root, "archforge-pptx-lint")
        os.makedirs(dst_dir, exist_ok=True)
        dst = os.path.join(dst_dir, "SKILL.md")
        with open(dst, "w", encoding="utf-8", newline="\n") as f:
            f.write(text)
        print(M("skill_installed") % dst)
        return 0
    sys.stdout.write(text)
    return 0


def _timeout_reexec(argv):
    """--timeout SECONDS: run the whole invocation in a child process with a wall clock,
    so a hostile or pathological deck cannot hang CI (0.6.x, external review: a real
    resource bound needs process isolation, and signal.alarm is POSIX-only). Re-execs
    sys.argv minus --timeout with a sentinel env var so the child does not recurse;
    subprocess timeout is portable to Windows. Returns an exit code, or None if no
    timeout was requested or this is already the child."""
    if os.environ.get("_ARCHFORGE_TIMEOUT_CHILD") == "1":
        return None
    secs = None
    rest = []
    it = iter(range(len(argv)))
    i = 0
    while i < len(argv):
        tok = argv[i]
        if tok == "--timeout" and i + 1 < len(argv):
            secs = argv[i + 1]; i += 2; continue
        if tok.startswith("--timeout="):
            secs = tok.split("=", 1)[1]; i += 1; continue
        rest.append(tok); i += 1
    if secs is None:
        return None
    try:
        secs_f = float(secs)
        # inf passes `> 0` and would silently disable the timeout; nan fails every
        # comparison. Require a finite positive value (0.7.1, external review).
        if not (math.isfinite(secs_f) and secs_f > 0):
            raise ValueError
    except ValueError:
        print(M("err_config") % ("--timeout must be a finite positive number of seconds, "
                                 "got %r" % secs), file=sys.stderr)
        return 2
    import subprocess
    env = dict(os.environ)
    env["_ARCHFORGE_TIMEOUT_CHILD"] = "1"
    try:
        return subprocess.call([sys.executable, "-m", "archforge"] + rest,
                               env=env, timeout=secs_f)
    except subprocess.TimeoutExpired:
        print(M("err_timeout") % secs_f, file=sys.stderr)
        return 124


def main():
    # Reconfigured before parser creation so Korean messages and argparse help don't die with
    # UnicodeEncodeError on non-UTF-8 stdout (pipes, cp949/cp1252); --help is printed during
    # parse_args, so the order here matters.
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass
    argv = sys.argv[1:]
    rc = _timeout_reexec(argv)
    if rc is not None:
        sys.exit(rc)
    # --lang must be finalized before the --help string, so this prescans before the parser
    # is created.
    # When given multiple times, the last value wins per argparse convention (third
    # adversarial panel: fixes a first-match-wins bug).
    lang_arg = None
    for i, tok in enumerate(argv):
        if tok == "--lang" and i + 1 < len(argv):
            lang_arg = argv[i + 1]
        elif tok.startswith("--lang="):
            lang_arg = tok.split("=", 1)[1]
    if lang_arg:
        set_lang(lang_arg)
    # Detecting the skill subcommand looks past leading --lang-style flags:
    # fixes `archforge --lang ko skill` being misinterpreted as "lint a file named skill"
    # (third panel)
    rest = list(argv)
    while rest and (rest[0] == "--lang" or rest[0].startswith("--lang=")):
        rest = rest[2:] if rest[0] == "--lang" else rest[1:]
    if rest and rest[0] == "skill":
        # Warns about the conflict with wanting to lint a file literally named "skill"
        # (adversarial panel finding: prevents silent misbehavior). To lint that file, call it
        # by path, e.g. `archforge ./skill`.
        if os.path.exists("skill"):
            print(M("skill_conflict"), file=sys.stderr)
        sys.exit(skill_main(rest[1:]))
    if rest and rest[0] in ("scan", "demo", "rules", "explain", "baseline", "fix"):
        if os.path.exists(rest[0]) and os.path.isfile(rest[0]):
            print(M("subcmd_conflict") % (rest[0], rest[0]), file=sys.stderr)
        dispatch = {"scan": scan_main, "demo": demo_main,
                    "rules": rules_main, "explain": explain_main,
                    "baseline": baseline_main, "fix": fix_main}
        sys.exit(dispatch[rest[0]](rest[1:]))
    if rest and rest[0] == "lint":
        # Explicit alias for single-file mode (`archforge lint deck.pptx`), so scripts
        # can be unambiguous about the subcommand (0.6.1). Drops just the token; the
        # leading --lang prefix (already prescanned) is preserved for parsing.
        if os.path.exists("lint") and os.path.isfile("lint"):
            print(M("subcmd_conflict") % ("lint", "lint"), file=sys.stderr)
        argv = argv[:len(argv) - len(rest)] + rest[1:]
    ap = argparse.ArgumentParser(prog="archforge", description=M("prog_desc"))
    ap.add_argument("pptx")
    ap.add_argument("--render", default=None, help=M("help_render"))
    ap.add_argument("--write-baseline", default=None, metavar="PATH", help=M("help_write_baseline"))
    _add_common_flags(ap)
    a = ap.parse_args(argv)

    # Validate output-path parents up front so a bad --sarif/--junit/--write-baseline
    # target is a controlled exit 2, not a traceback after linting (0.7.1).
    err = _validate_cli_globals(a)
    if err:
        print(err, file=sys.stderr)
        sys.exit(2)

    try:
        res = _lint_one(a.pptx, a)
    except UsageError as e:
        print(str(e), file=sys.stderr)
        sys.exit(2)
    if res is None:   # --write-baseline: exit after recording
        sys.exit(0)

    if a.sarif:
        import json
        sarif_doc = _reporters.build_sarif(a.pptx, res["errors"], res["warns"])
        with open(a.sarif, "w", encoding="utf-8", newline="\n") as f:
            json.dump(sarif_doc, f, ensure_ascii=False, indent=2)

    if a.junit:
        xml_text = _reporters.build_junit_multi(
            [(a.pptx, res["errors"], res["warns"],
              set(res["summary"]["skipped_codes"]),
              res["summary"]["policy"], None)])
        with open(a.junit, "w", encoding="utf-8", newline="\n") as f:
            f.write(xml_text)

    if a.html:
        _write_html_report(a.html, [(a.pptx, res, None, a.render)])

    if a.json:
        import json
        doc = _reporters.build_json_doc(a.pptx, res["errors"], res["warns"],
                                        res["ghost"], res["summary"],
                                        schema=res["schema"],
                                        capabilities=res["capabilities"],
                                        abstentions=res["abstentions"],
                                        invocation=res["invocation"],
                                        rules_split=res["rules_split"])
        print(json.dumps(doc, ensure_ascii=False, indent=2))
        sys.exit(1 if res["fail"] else 0)

    for line in _reporters.render_text(a.pptx, res["errors"], res["warns"], res["ghost"],
                                       res["profile"], res["profile_excl"], res["skip"],
                                       config_path=res["cfg_path"],
                                       baseline_suppressed=res["baseline_suppressed"],
                                       baseline_path=res["baseline_path"]):
        print(line)

    sys.exit(1 if res["fail"] else 0)


# Maps a W18 skip-reason key (machine string, stable) to the rules it prevented and the
# capability it degrades (0.7 schema 2.0). Reasons not listed default to no affected rules
# and the "meta" capability. Keys mirror the Counter keys written at each guard.
_REASON_RULES = {
    "vertical_text": (["W15", "W16", "W17"], "geometry"),
    "complex_script": (["W15", "W16", "W17"], "geometry"),
    "glyph_boxes": (["W15", "W16", "W17"], "geometry"),
    "pic_boxes": (["W16", "W17"], "geometry"),
    "w15": (["W15"], "geometry"),
    "w16_w17": (["W16", "W17"], "geometry"),
    "image_decode": (["W16", "W17"], "geometry"),
    "image_decode_budget": (["W16", "W17"], "geometry"),
    "frames": (["E1", "E3", "E4", "W1", "W5", "W8"], "typography"),
    "frame": (["E1", "E3", "E4", "W1", "W5", "W8"], "typography"),
    "para": (["E1", "E2", "E3", "E4"], "typography"),
    "para_size": (["E3"], "typography"),
    "run": (["E1", "E2", "E3", "E4"], "typography"),
    "w7": (["W7"], "render"),
    "w7_no_render": (["W7"], "render"),
    "w7_color_unknown": (["W7"], "render"),
    "render_dir_missing": (["W7"], "render"),
    "w9": (["W9"], "structure"),
    "w6_sig": (["W6"], "structure"),
    "w6": (["W6"], "structure"),
    "w6_capped": (["W6"], "structure"),
    "w10_tokens": (["W10"], "structure"),
    "w10": (["W10"], "structure"),
    "w10_capped": (["W10"], "structure"),
    "w11_w14": (["W11", "W12", "W13", "W14"], "structure"),
    "w12_w13": (["W12", "W13"], "structure"),
    "theme_parse": (["E1"], "typography"),
}

# Every skip-reason key a detector can emit must be registered above, so a structural
# abstention never lands as ([], "meta") with structure still reported "complete"
# (0.7.1, external review P0). test_reason_registry_covers_all_keys enforces this.
KNOWN_REASON_KEYS = frozenset(_REASON_RULES)


def _capabilities_and_abstentions(warns, render_requested):
    """Turns the W18 findings into a structured capabilities map and abstentions list
    (0.7 schema 2.0). W18 detail is a machine-key Counter string ('vertical_text=1, ...'),
    so parsing it back is deterministic. Verdict is untouched; this is a richer view of the
    same incompleteness signal."""
    abstentions = []
    degraded = set()
    for f in warns:
        if f.code != "W18":
            continue
        for part in (f.detail or "").split(","):
            part = part.strip()
            if "=" not in part:
                continue
            key, _, cnt = part.partition("=")
            key = key.strip()
            try:
                count = int(cnt)
            except ValueError:
                count = 1
            rules, cap = _REASON_RULES.get(key, ([], "meta"))
            degraded.add(cap)
            abstentions.append({"reason": key, "page": f.page, "count": count,
                                "affected_rules": rules})
    caps = {}
    caps["typography"] = "partial" if "typography" in degraded else "complete"
    caps["geometry"] = "partial" if "geometry" in degraded else "complete"
    caps["structure"] = "partial" if "structure" in degraded else "complete"
    caps["render_contrast"] = ("partial" if "render" in degraded
                               else ("complete" if render_requested else "not_requested"))
    # An unregistered reason (should not happen: enforced by test) still surfaces here
    # rather than vanishing into a "complete" verdict.
    if "meta" in degraded:
        caps["meta"] = "partial"
    return caps, abstentions


def _check_out_dir(path):
    """Error string if an output path's parent directory is missing, so a bad
    --sarif/--junit/--write-baseline target is a controlled exit 2 rather than a
    traceback mid-run (0.7.1, external review)."""
    if not path:
        return None
    d = os.path.dirname(os.path.abspath(path))
    if not os.path.isdir(d):
        return M("err_out_path") % path
    return None


def _html_thumbs(render_dir, n_slides, max_w=640):
    """Optional render thumbnails for the HTML report: p01.png-style files downscaled
    to JPEG bytes. Empty dict when no render dir; failures skip that page (the
    wireframe is the fallback)."""
    out = {}
    if not render_dir or not os.path.isdir(render_dir):
        return out
    try:
        from PIL import Image
        import io as _io
    except ImportError:
        return out
    for si in range(1, n_slides + 1):
        p = os.path.join(render_dir, "p%02d.png" % si)
        if not os.path.exists(p):
            continue
        try:
            im = Image.open(p).convert("RGB")
            if im.width > max_w:
                im = im.resize((max_w, int(im.height * max_w / im.width)))
            buf = _io.BytesIO()
            im.save(buf, "JPEG", quality=72)
            out[si] = buf.getvalue()
        except Exception:
            continue
    return out


def _write_html_report(out_path, entries):
    """entries = [(path, res_or_None, err_or_None, render_dir)] -> one HTML file."""
    items = []
    for p, res, e, rdir in entries:
        if res is None:
            items.append((p, [], [], {}, {"sw": 13.333, "sh": 7.5, "slides": []},
                          {}, e))
            continue
        try:
            wf = collect_wireframe(p)
        except Exception:
            wf = {"sw": 13.333, "sh": 7.5, "slides": []}
        thumbs = _html_thumbs(rdir, len(wf["slides"]))
        items.append((p, res["errors"], res["warns"], res["summary"], wf, thumbs, None))
    with open(out_path, "w", encoding="utf-8", newline="\n") as f:
        f.write(_reporters.build_html_multi(items))


def _validate_cli_globals(a):
    """Validation of CLI-supplied values that are identical for every file in a scan.
    Returns an error string (caller exits 2) or None. Deck-config-supplied values stay
    per-file (a bad config next to one deck is that file's problem, not the batch's)."""
    if a.config and a.no_config:
        return M("err_config") % "--config conflicts with --no-config"
    if a.config and not os.path.exists(a.config):
        return M("err_config") % a.config
    for v, name in ((a.hard_min, "hard_min"), (a.body_min, "body_min"),
                    (a.small_min, "small_min")):
        if v is not None and (not math.isfinite(v) or v <= 0):
            return M("err_config") % ("threshold out of range: %s=%r" % (name, v))
    if a.w6_sim is not None and (not math.isfinite(a.w6_sim) or not (0 < a.w6_sim <= 1)):
        return M("err_config") % ("threshold out of range: w6_sim=%r" % a.w6_sim)
    if a.w6_cluster is not None and a.w6_cluster < 1:
        return M("err_config") % ("threshold out of range: w6_cluster=%r" % a.w6_cluster)
    if a.skip:
        codes_ = {c.strip().upper() for c in a.skip.split(",") if c.strip()}
        bad = sorted(c for c in codes_ if not c.startswith("W"))
        if bad:
            return M("err_skip_e") % ",".join(bad)
        unknown = sorted(c for c in codes_ if c not in ALL_CODES)
        if unknown:
            return M("err_skip_unknown") % ",".join(unknown)
        if "W18" in codes_:
            return M("err_skip_w18")
    for outp in (getattr(a, "sarif", None), getattr(a, "junit", None),
                 getattr(a, "html", None), getattr(a, "write_baseline", None)):
        err = _check_out_dir(outp)
        if err:
            return err
    return None


class UsageError(Exception):
    """A per-file usage/config error (missing file, bad config, invalid pptx, bad flags).

    Single-file mode prints it and exits 2, preserving the existing contract. scan mode
    converts it into a per-file result and keeps scanning: one broken deck must not kill
    the batch or swallow the aggregate report (0.6.0, external review P0)."""


def _pkg_version():
    try:
        from importlib.metadata import version
        return version("archforge")
    except Exception:
        return "unknown"


def _add_common_flags(ap):
    """Flags shared by single-file mode and scan mode (prevents duplicate-definition drift,
    0.5.0)."""
    ap.add_argument("--version", action="version", version="archforge " + _pkg_version())
    ap.add_argument("--hard-min", type=float, default=None, help=M("help_hard_min"))
    ap.add_argument("--body-min", type=float, default=None, help=M("help_body_min"))
    ap.add_argument("--strict", action="store_true", help=M("help_strict"))
    # --strict split into orthogonal policies (0.6.0, external review): failing every
    # advisory warning, failing on incomplete checks, and lifting E2's numeric exemptions
    # are three different decisions. --strict remains as the union for compatibility.
    ap.add_argument("--fail-on-warning", action="store_true", help=M("help_fail_on_warning"))
    ap.add_argument("--fail-incomplete", action="store_true", help=M("help_fail_incomplete"))
    ap.add_argument("--e2-no-exemptions", action="store_true", help=M("help_e2_no_exemptions"))
    ap.add_argument("--small-min", type=float, default=None, help=M("help_small_min"))
    ap.add_argument("--ghost", action="store_true", help=M("help_ghost"))
    ap.add_argument("--json", action="store_true", help=M("help_json"))
    ap.add_argument("--schema", default="1.0", choices=("1.0", "2.0", "1", "2"),
                    help=M("help_schema"))
    ap.add_argument("--skip", default=None, metavar="CODES", help=M("help_skip"))
    ap.add_argument("--profile", default=None, choices=sorted(PROFILES), help=M("help_profile"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    ap.add_argument("--w6-sim", type=float, default=None, help=M("help_w6_sim"))
    ap.add_argument("--w6-cluster", type=int, default=None, help=M("help_w6_cluster"))
    ap.add_argument("--config", default=None, metavar="PATH", help=M("help_config"))
    ap.add_argument("--no-config", action="store_true", help=M("help_no_config"))
    ap.add_argument("--sarif", default=None, metavar="PATH", help=M("help_sarif"))
    ap.add_argument("--junit", default=None, metavar="PATH", help=M("help_junit"))
    ap.add_argument("--html", default=None, metavar="PATH", help=M("help_html"))
    ap.add_argument("--timeout", default=None, metavar="SECONDS", help=M("help_timeout"))
    ap.add_argument("--baseline", default=None, metavar="PATH", help=M("help_baseline"))


def _lint_one(path, a):
    """One file's config resolution -> check -> filter -> summary. Shared by main (single
    mode) and scan_main.
    Usage errors (missing file, bad config, invalid pptx, bad flags) raise UsageError:
    single-file mode turns that into exit 2 (unchanged contract), scan mode turns it into
    a per-file result and keeps going (0.6.0). All flag/config validation happens BEFORE
    the lint run and before --write-baseline recording (0.6.0: a typo'd --skip used to
    record a baseline as if nothing were wrong). If --write-baseline is set, returns None
    after recording. Returned dict: errors/warns/ghost/summary/fail/profile/profile_excl/
    skip/cfg_path/baseline_suppressed/baseline_path."""
    if not os.path.exists(path):
        raise UsageError(M("err_notfound") % path)

    # Explicit --config combined with --no-config used to silently drop the explicit
    # config; that contradiction is now an error (0.6.0, external review).
    if a.config and a.no_config:
        raise UsageError(M("err_config") % "--config conflicts with --no-config")

    # Config file (.archforge.json/.yml): CLI flags always win (0.4.0).
    # Trust boundary (fourth review): since a config file in the deck folder could weaken the
    # gate, the applied config path is recorded in the output contract (JSON summary.config,
    # a text footnote) and can be turned off with --no-config.
    cfg = {}
    cfg_path = None if a.no_config else _config.find_config(path, a.config)
    if a.config and not a.no_config and cfg_path is None:
        raise UsageError(M("err_config") % a.config)
    if cfg_path:
        try:
            cfg, cfg_warns = _config.load_config(cfg_path)
            for wmsg in cfg_warns:
                print("archforge: %s (%s)" % (wmsg, cfg_path), file=sys.stderr)
        except Exception as e:
            raise UsageError(M("err_config") % ("%s (%s)" % (cfg_path, e)))

    def pick(cli_val, cfg_key, default):
        if cli_val is not None:
            return cli_val
        return cfg.get(cfg_key, default)

    try:
        hard_min = float(pick(a.hard_min, "hard_min", 5.0))
        body_min = float(pick(a.body_min, "body_min", 9.0))
        small_min = float(pick(a.small_min, "small_min", 7.5))
        w6_sim = float(pick(a.w6_sim, "w6_sim", 0.90))
        w6_cluster = int(pick(a.w6_cluster, "w6_cluster", 3))
    except (TypeError, ValueError) as e:
        raise UsageError(M("err_config") % ("threshold: %s" % e))
    # Range validation: closes a bypass where --hard-min 0 silently disabled E3 (the
    # unreadable-text block) (formerly X1). NaN slips through ordinary comparisons
    # (NaN <= 0 is False), which re-opened the same bypass via --hard-min nan or a bare
    # NaN literal in .archforge.json (json.load accepts it), so finiteness is checked
    # explicitly (0.6.0, external verification finding).
    if not all(math.isfinite(v) for v in (hard_min, body_min, small_min, w6_sim)) \
            or hard_min <= 0 or body_min <= 0 or small_min <= 0 \
            or not (0 < w6_sim <= 1) or w6_cluster < 1:
        raise UsageError(M("err_config") % (
            "threshold out of range (hard_min/body_min/small_min > 0, "
            "0 < w6_sim <= 1, w6_cluster >= 1, all finite)"))
    profile = pick(a.profile, "profile", DEFAULT_PROFILE)
    if profile not in PROFILES:
        raise UsageError(M("err_config") % ("profile=%r" % profile))
    lang_final = pick(a.lang, "lang", None)
    if lang_final:
        set_lang(lang_final)
    baseline_path = pick(a.baseline, "baseline", None)
    skip_raw = pick(a.skip, "skip", "")
    if isinstance(skip_raw, list):
        skip_raw = ",".join(str(c) for c in skip_raw)

    # --skip validation happens BEFORE the run and before baseline recording (0.6.0).
    # --skip is WARN-only: silently swallowing even E codes would be a footgun that turns
    # off the deployment-blocking gate without a trace (second external re-check). Unknown
    # codes are typos that would make CI look normal (third review P1), and W18 is an
    # incompleteness signal, not something to suppress (third review P1).
    skip = {c.strip().upper() for c in skip_raw.split(",") if c.strip()}
    bad_skip = sorted(c for c in skip if not c.startswith("W"))
    if bad_skip:
        raise UsageError(M("err_skip_e") % ",".join(bad_skip))
    unknown_skip = sorted(c for c in skip if c not in ALL_CODES)
    if unknown_skip:
        raise UsageError(M("err_skip_unknown") % ",".join(unknown_skip))
    if "W18" in skip:
        raise UsageError(M("err_skip_w18"))

    # --strict = the union of the three orthogonal policies (compatibility alias)
    fail_on_warning = a.strict or a.fail_on_warning
    fail_incomplete = a.strict or a.fail_incomplete
    e2_no_exemptions = a.strict or a.e2_no_exemptions

    ghost = [] if (a.ghost or a.json) else None
    try:
        errors, warns = lint(path, hard_min, body_min, small_min, render_dir=a.render,
                             ghost=ghost, strict=e2_no_exemptions, w6_sim=w6_sim,
                             w6_min_cluster=w6_cluster, profile=profile)
    except Exception as e:
        # ValueError carries an intentional diagnosis (zip preflight budgets); other
        # exception types stay name-only to avoid echoing arbitrary parser internals
        reason = type(e).__name__
        if isinstance(e, ValueError) and str(e):
            reason = "%s: %s" % (reason, e)
        raise UsageError(M("err_open") % (path, reason))

    # Baseline recording mode: saves current violations (excluding the W18 incompleteness
    # signal) as a fingerprint and exits. Runs only after all validation above.
    if getattr(a, "write_baseline", None):
        n = _config.write_baseline(a.write_baseline,
                                   [f for f in list(errors) + list(warns) if f.code != "W18"],
                                   profile=profile, lang=get_lang(),
                                   thresholds={"hard_min": hard_min, "body_min": body_min,
                                               "small_min": small_min, "w6_sim": w6_sim,
                                               "w6_cluster": w6_cluster},
                                   artifact=_config.deck_artifact(path))
        print(M("baseline_written") % (n, a.write_baseline))
        return None

    # Incompleteness is determined before filtering: even if W18 is --skip'd, the
    # machine-readable signal remains
    has_w18 = any(w[1] == "W18" for w in warns)
    baseline_suppressed = 0
    if baseline_path:
        try:
            known = _config.load_baseline(baseline_path)
        except Exception as e:
            raise UsageError(M("err_config") % ("baseline %s (%s)" % (baseline_path, e)))
        # Recorded run conditions are checked, not just stored (0.6.0, external review):
        # a baseline made under a different profile or tool version suppresses different
        # things than the reader expects. Warning, not error: baselines are beta.
        meta = _config.load_baseline_meta(baseline_path)
        if meta.get("profile") not in (None, "", profile):
            print(M("note_baseline_meta") % ("profile", meta.get("profile"), profile),
                  file=sys.stderr)
        rec_v = str(meta.get("tool_version") or "")
        cur_v = _pkg_version()
        if rec_v and cur_v != "unknown" and rec_v.split(".")[:2] != cur_v.split(".")[:2]:
            print(M("note_baseline_meta") % ("tool_version", rec_v, cur_v), file=sys.stderr)
        rec_thr = meta.get("threshold_hash")
        if rec_thr:
            import hashlib as _hl
            cur_thr = {"hard_min": hard_min, "body_min": body_min, "small_min": small_min,
                       "w6_sim": w6_sim, "w6_cluster": w6_cluster}
            cur_hash = _hl.sha1((",".join("%s=%r" % (k, cur_thr[k]) for k in sorted(cur_thr)))
                                .encode("utf-8")).hexdigest()[:12]
            if cur_hash != rec_thr:
                print(M("note_baseline_meta") % ("thresholds", rec_thr, cur_hash),
                      file=sys.stderr)
        # Artifact identity (0.8): a baseline written from one deck applied to a
        # different one still suppresses shared fingerprints; the file basename is the
        # identity signal that survives regeneration, so a mismatch is surfaced.
        rec_art = meta.get("artifact") or {}
        rec_name = rec_art.get("file_name")
        if rec_name and rec_name != os.path.basename(path):
            print(M("note_baseline_meta") % ("artifact", rec_name,
                                             os.path.basename(path)), file=sys.stderr)
        errors, s1 = _config.apply_baseline(errors, known)
        warns, s2 = _config.apply_baseline(warns, known)
        baseline_suppressed = s1 + s2
    # Profile exclusions were already not run at the engine stage (0.3.1). Only --skip is
    # filtered here; the applied skip is recorded in the JSON summary to leave a trace.
    profile_excl = PROFILES[profile]
    excluded = skip | profile_excl
    if skip:
        warns = [w for w in warns if w[1] not in skip]

    # Per-rule severity overrides from the config (policy-layer rules only; validated in
    # config.load_config). "off" drops the findings, "warning"/"error" move them between
    # the two lists. Applied before counting so summary/policy see the effective levels;
    # recorded in the summary so a demoted E2 is never invisible (0.8.x, external audit).
    sev_over = cfg.get("severity") or {}
    if sev_over:
        moved_to_warn = [f for f in errors if sev_over.get(f.code) == "warning"]
        moved_to_err = [f for f in warns if sev_over.get(f.code) == "error"]
        errors = [f for f in errors
                  if sev_over.get(f.code) not in ("warning", "off")] + moved_to_err
        warns = [f for f in warns
                 if sev_over.get(f.code) not in ("error", "off")] + moved_to_warn

    schema = "2.0" if str(getattr(a, "schema", "1.0")) in ("2", "2.0") else "1.0"
    caps, abstentions = _capabilities_and_abstentions(warns, bool(a.render))
    # schema 2.0 invocation + rule accounting: skipped_codes mixed profile exclusion and
    # --skip, which mean different things to a consumer (0.7.1, external review section 5).
    invocation = {"profile": profile,
                  "policy": {"fail_on_warning": bool(fail_on_warning),
                             "fail_incomplete": bool(fail_incomplete),
                             "e2_no_exemptions": bool(e2_no_exemptions)},
                  "config": cfg_path,
                  "thresholds": {"hard_min": hard_min, "body_min": body_min,
                                 "small_min": small_min, "w6_sim": w6_sim,
                                 "w6_cluster": w6_cluster}}
    rules_split = {"executed": sorted(ALL_CODES - excluded),
                   "profile_excluded": sorted(profile_excl),
                   "user_suppressed": sorted(skip)}
    fail = bool(errors or (fail_on_warning and warns) or (fail_incomplete and has_w18))
    summary = {"error_count": len(errors), "warn_count": len(warns),
               "pass": not fail,
               # The active failure policy travels with the verdict (0.6.1, external
               # review): identical counts can pass or fail depending on flags, and a
               # JSON consumer could not tell why. Gate on summary.pass.
               "policy": {"fail_on_warning": bool(fail_on_warning),
                          "fail_incomplete": bool(fail_incomplete),
                          "e2_no_exemptions": bool(e2_no_exemptions)},
               "incomplete": has_w18,
               "profile": profile,
               "skipped_codes": sorted(excluded),
               "baseline_suppressed": baseline_suppressed,
               "config": cfg_path}   # always makes visible which config adjusted the gate
                                     # (trust boundary)
    if sev_over:
        summary["severity_overrides"] = dict(sorted(sev_over.items()))

    return {"errors": errors, "warns": warns, "ghost": ghost, "summary": summary,
            "fail": fail,
            "profile": profile, "profile_excl": profile_excl, "skip": skip,
            "cfg_path": cfg_path, "baseline_suppressed": baseline_suppressed,
            "baseline_path": baseline_path,
            "schema": schema, "capabilities": caps, "abstentions": abstentions,
            "invocation": invocation, "rules_split": rules_split}


def _expand_scan_paths(patterns):
    """Expands scan arguments: a directory recurses for .pptx files, a glob pattern is
    globbed, and anything else is taken as a literal file path.
    PowerPoint lock files (~$*.pptx) are excluded, and duplicates are removed while
    preserving order. Returns (files, per_pattern_counts): match counts are tracked per
    input so a typo'd second pattern cannot hide behind a first pattern that matched
    (0.6.1, external review P0: zero-match was only detected for the whole set)."""
    out = []
    counts = []
    for pat in patterns:
        n0 = len(out)
        if os.path.isdir(pat):
            for root, dirs, files in os.walk(pat):
                # os.walk's directory order is filesystem-dependent; sorting both lists
                # makes aggregate output deterministic across machines (0.6.0, external
                # review: matters for snapshot tests and CI diffs)
                dirs.sort()
                for fn in sorted(files):
                    if fn.lower().endswith(".pptx") and not fn.startswith("~$"):
                        out.append(os.path.join(root, fn))
        elif any(ch in pat for ch in "*?["):
            for p in sorted(glob.glob(pat, recursive=True)):
                if p.lower().endswith(".pptx") and not os.path.basename(p).startswith("~$"):
                    out.append(p)
        else:
            # A literal path counts as matched here; a missing file becomes a per-file
            # error entry downstream (visible, not silent)
            out.append(pat)
        counts.append((pat, len(out) - n0))
    seen, uniq = set(), []
    for p in out:
        k = os.path.normcase(os.path.abspath(p))
        if k not in seen:
            seen.add(k)
            uniq.append(p)
    return uniq, counts


def scan_main(argv=None):
    """`archforge scan PATHS...`: lints multiple files, directories, and globs in one run
    (0.5.0, for CI/pre-commit use). Per-file judgment goes through the same path as single
    mode (_lint_one). Batch failure semantics (0.6.0, external review P0): a broken or
    misconfigured file becomes a per-file "error" entry and the scan continues; it never
    aborts the batch or swallows the aggregate report. Exit is 1 if any file fails or
    errors. Zero matches is not a silent pass; it exits 2 (prevents a CI footgun)."""
    ap = argparse.ArgumentParser(prog="archforge scan", description=M("scan_desc"))
    ap.add_argument("paths", nargs="+", help=M("help_scan_paths"))
    ap.add_argument("--allow-empty-pattern", action="store_true",
                    help=M("help_allow_empty_pattern"))
    _add_common_flags(ap)
    # Flags exclusive to single-file mode are not supported: only defaults for _lint_one
    # compatibility are seeded here
    # (--render doesn't make sense in scan since the page render folder differs per deck)
    ap.set_defaults(render=None, write_baseline=None)
    a = ap.parse_args(argv)

    # Global usage errors are not per-file errors (0.6.1, external review): a bad CLI
    # flag must exit 2 up front, not degrade into N identical per-file entries.
    err = _validate_cli_globals(a)
    if err:
        print(err, file=sys.stderr)
        return 2

    files, pattern_counts = _expand_scan_paths(a.paths)
    empty = [pat for pat, n in pattern_counts if n == 0]
    if empty and not a.allow_empty_pattern:
        # One input matching nothing must not hide behind another that matched: a typo'd
        # glob or a broken build directory silently passing is the exact CI footgun the
        # whole-set zero check missed (0.6.1, external review P0)
        print(M("err_scan_pattern_empty") % "; ".join(empty), file=sys.stderr)
        return 2
    if not files:
        print(M("err_scan_none") % " ".join(a.paths), file=sys.stderr)
        return 2

    # A single CLI baseline applied across many decks is unsafe: fingerprints carry no
    # file identity, so a finding accepted in deck A would suppress the same finding in
    # deck B (0.6.0, external review P0). Per-deck baselines still work via each deck
    # folder's config file.
    if a.baseline and len(files) > 1:
        print(M("err_scan_baseline"), file=sys.stderr)
        return 2

    # A deck folder's config may set "lang"; without restoration it would leak into every
    # later file and, because messages render lazily, into the whole aggregate report
    # (0.6.0, external verification finding). The scan report renders in the language the
    # scan was invoked with; per-deck config lang is a single-file-mode feature.
    lang0 = get_lang()
    results = []   # (path, res_dict_or_None, usage_error_message_or_None)
    for path in files:
        try:
            res = _lint_one(path, a)
            results.append((path, res, None))
        except UsageError as e:
            print(str(e), file=sys.stderr)
            results.append((path, None, str(e)))
        finally:
            set_lang(lang0)

    ok = [(p, r) for (p, r, _e) in results if r is not None]
    errored = [(p, e) for (p, r, e) in results if r is None]
    failed = sum(1 for (_p, r) in ok if r["fail"]) + len(errored)

    if a.sarif:
        import json
        sarif_doc = _reporters.build_sarif_multi(
            [(p, r["errors"], r["warns"]) for (p, r) in ok])
        with open(a.sarif, "w", encoding="utf-8", newline="\n") as f:
            json.dump(sarif_doc, f, ensure_ascii=False, indent=2)

    if a.junit:
        junit_items = []
        for p, r, e in results:
            if r is None:
                junit_items.append((p, [], [], set(),
                                    {"fail_on_warning": False, "fail_incomplete": False}, e))
            else:
                junit_items.append((p, r["errors"], r["warns"],
                                    set(r["summary"]["skipped_codes"]),
                                    r["summary"]["policy"], None))
        with open(a.junit, "w", encoding="utf-8", newline="\n") as f:
            f.write(_reporters.build_junit_multi(junit_items))

    if a.html:
        _write_html_report(a.html, [(p, r, e, None) for p, r, e in results])

    if a.json:
        import json
        # The scan report is its own document type. Its root schema_version tracks the
        # per-file schema so a consumer that keys on the root does not misparse a v2 file
        # object under a v1 root (0.7.1, external review P0). file_schema_version names the
        # per-file shape explicitly.
        file_schema = "2.0" if str(getattr(a, "schema", "1.0")) in ("2", "2.0") else "1.0"
        root_schema = "scan-2.0" if file_schema == "2.0" else "scan-1.0"
        docs = []
        for p, r, e in results:
            if r is None:
                docs.append({"file": p, "status": "error", "error": e})
            else:
                doc = _reporters.build_json_doc(p, r["errors"], r["warns"], r["ghost"],
                                                r["summary"], schema=r["schema"],
                                                capabilities=r["capabilities"],
                                                abstentions=r["abstentions"],
                                                invocation=r["invocation"],
                                                rules_split=r["rules_split"])
                doc["status"] = "fail" if r["fail"] else "pass"
                docs.append(doc)
        agg = {"schema_version": root_schema,
               "kind": "scan-report",
               "file_schema_version": file_schema,
               "tool": {"name": "archforge", "version": _reporters._tool_version()},
               "lang": get_lang(),
               "scan": {"inputs": [{"pattern": pat, "matches": n}
                                   for pat, n in pattern_counts]},
               "files": docs,
               "summary": {"file_count": len(results), "failed_files": failed,
                           "error_files": len(errored),
                           "error_count": sum(len(r["errors"]) for (_p, r) in ok),
                           "warn_count": sum(len(r["warns"]) for (_p, r) in ok),
                           "pass": failed == 0,
                           "incomplete": any(r["summary"]["incomplete"] for (_p, r) in ok)}}
        print(json.dumps(agg, ensure_ascii=False, indent=2))
        return 1 if failed else 0

    for p, r, e in results:
        if r is None:
            print(M("scan_file_error") % (p, e))
            print()
            continue
        for line in _reporters.render_text(p, r["errors"], r["warns"], r["ghost"],
                                           r["profile"], r["profile_excl"], r["skip"],
                                           config_path=r["cfg_path"],
                                           baseline_suppressed=r["baseline_suppressed"],
                                           baseline_path=r["baseline_path"]):
            print(line)
        print()
    print(M("scan_summary") % (len(results), failed))
    return 1 if failed else 0


def rules_main(argv=None):
    """`archforge rules`: one line per rule so users can discover the gate set without
    opening the README (0.6.1, external review)."""
    ap = argparse.ArgumentParser(prog="archforge rules", description=M("rules_desc"))
    ap.add_argument("--json", action="store_true", help=M("help_json"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    from .rules import TITLES, category
    rows = []
    for code in sorted(ALL_CODES, key=lambda c: (c[0] != "E", int(c[1:]))):
        profiles = sorted(p for p, excl in PROFILES.items() if code not in excl)
        rows.append({"code": code, "severity": severity(code), "category": category(code),
                     "title": TITLES.get(code, code), "profiles": profiles})
    if a.json:
        import json
        print(json.dumps(rows, ensure_ascii=False, indent=2))
        return 0
    for r in rows:
        print("%-4s %-8s %-11s %s" % (r["code"], r["severity"], r["category"], r["title"]))
    return 0


def fix_main(argv=None):
    """`archforge fix deck.pptx -o fixed.pptx`: deterministic auto-fixes for the three
    mechanically safe rules (E1 a:ea font, E2 dash punctuation, E4 tracking). Everything
    that needs layout judgment stays find-only; re-lint after fixing."""
    ap = argparse.ArgumentParser(prog="archforge fix", description=M("fix_desc"))
    ap.add_argument("pptx")
    ap.add_argument("-o", "--output", required=True, metavar="PATH",
                    help=M("help_fix_output"))
    ap.add_argument("--rules", default="E1,E2,E4", metavar="CODES",
                    help=M("help_fix_rules"))
    ap.add_argument("--ea-font", default=None, metavar="FONT", help=M("help_fix_ea"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    from . import fixes as _fixes
    rules = {c.strip().upper() for c in a.rules.split(",") if c.strip()}
    bad = sorted(rules - set(_fixes.FIXABLE))
    if bad:
        print(M("err_fix_rules") % (",".join(bad), ",".join(_fixes.FIXABLE)),
              file=sys.stderr)
        return 2
    if not os.path.exists(a.pptx):
        print(M("err_notfound") % a.pptx, file=sys.stderr)
        return 2
    err = _check_out_dir(a.output)
    if err:
        print(err, file=sys.stderr)
        return 2
    try:
        changes = _fixes.apply_fixes(a.pptx, a.output, rules=rules,
                                     ea_font=a.ea_font or _fixes.DEFAULT_EA)
    except Exception as e:
        print(M("err_open") % (a.pptx, e), file=sys.stderr)
        return 2
    for ch in changes:
        print("  FIX p%02d [%s] %s" % (ch["page"], ch["code"], ch["detail"]))
    print(M("fix_summary") % (len(changes), a.output))
    return 0


def baseline_main(argv=None):
    """`archforge baseline inspect PATH`: what a baseline actually suppresses, and under
    which recorded conditions (0.8, external review: baselines were opaque blobs)."""
    ap = argparse.ArgumentParser(prog="archforge baseline",
                                 description=M("baseline_desc"))
    ap.add_argument("action", choices=("inspect",))
    ap.add_argument("path")
    ap.add_argument("--json", action="store_true", help=M("help_json"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    import json
    from collections import Counter
    try:
        with open(a.path, encoding="utf-8") as f:
            doc = json.load(f)
    except Exception as e:
        print(M("err_config") % ("baseline %s (%s)" % (a.path, e)), file=sys.stderr)
        return 2
    by_code = Counter()
    total = 0
    for e in doc.get("findings", []):
        by_code[e.get("code", "?")] += int(e.get("count", 1))
        total += int(e.get("count", 1))
    info = {"schema_version": doc.get("schema_version"),
            "tool_version": doc.get("tool_version"),
            "profile": doc.get("profile"),
            "lang": doc.get("lang"),
            "threshold_hash": doc.get("threshold_hash"),
            "artifact": doc.get("artifact"),
            "suppressed_total": total,
            "suppressed_by_code": dict(sorted(by_code.items()))}
    if a.json:
        print(json.dumps(info, ensure_ascii=False, indent=2))
        return 0
    print("baseline %s (schema %s, tool %s, profile %s)"
          % (a.path, info["schema_version"], info["tool_version"], info["profile"]))
    art = info["artifact"] or {}
    if art:
        print("  artifact: %s (sha256 %s)"
              % (art.get("file_name"), art.get("sha256_12")))
    print("  suppresses %d finding(s):" % total)
    for code, n in info["suppressed_by_code"].items():
        print("    %-4s x%d" % (code, n))
    if str(info["schema_version"]) != "3":
        print(M("err_config") % ("outdated baseline schema; regenerate"), file=sys.stderr)
        return 1
    return 0


def explain_main(argv=None):
    """`archforge explain CODE`: what a rule means, when it fires, and how to fix it,
    from the same fix guidance the agent skill pack teaches (0.6.1)."""
    ap = argparse.ArgumentParser(prog="archforge explain", description=M("explain_desc"))
    ap.add_argument("code")
    ap.add_argument("--json", action="store_true", help=M("help_json"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    from .rules import TITLES, category
    code = a.code.strip().upper()
    if code not in ALL_CODES:
        print(M("err_skip_unknown") % code, file=sys.stderr)
        return 2
    profiles = sorted(p for p, excl in PROFILES.items() if code not in excl)
    doc = {"code": code, "severity": severity(code), "category": category(code),
           "title": TITLES.get(code, code), "profiles": profiles,
           "fix": M("fix_" + code.lower()),
           "help_uri": "https://github.com/Love-Ash/archforge/blob/main/docs/rules/%s.md" % code}
    if a.json:
        import json
        print(json.dumps(doc, ensure_ascii=False, indent=2))
        return 0
    print("%s  %s" % (code, doc["title"]))
    print("  severity: %s | category: %s | profiles: %s"
          % (doc["severity"], doc["category"], ",".join(profiles)))
    print("  fix: %s" % doc["fix"])
    print("  docs: %s" % doc["help_uri"])
    return 0


def demo_main(argv=None):
    """`archforge demo`: generates a deck seeded with defects and its corrected version, then
    lints them on the spot (0.5.0 onboarding).
    Serves as a first-run experience that shows, within 30 seconds of installing, exactly
    what the tool catches."""
    ap = argparse.ArgumentParser(prog="archforge demo", description=M("demo_desc"))
    ap.add_argument("--dir", default="archforge-demo", help=M("help_demo_dir"))
    ap.add_argument("--lang", default=None, choices=("ko", "en"), help=M("help_lang"))
    a = ap.parse_args(argv)
    try:
        from . import demo as _demo
    except ImportError:
        import demo as _demo
    os.makedirs(a.dir, exist_ok=True)
    broken = os.path.join(a.dir, "broken.pptx")
    fixed = os.path.join(a.dir, "fixed.pptx")
    # Deck text also follows the report language (an English user gets an English demo deck,
    # 0.5.0)
    deck_lang = get_lang() if get_lang() in ("ko", "en") else "en"
    _demo.build_broken(broken, lang=deck_lang)
    _demo.build_fixed(fixed, lang=deck_lang)
    print(M("demo_built") % a.dir)
    print()
    rc = 0
    for path in (broken, fixed):
        errors, warns = lint(path, profile="full")
        for line in _reporters.render_text(path, errors, warns, None,
                                           "full", PROFILES["full"], set()):
            print(line)
        print()
        if path == fixed and (errors or warns):
            rc = 1   # the corrected version must always be clean (a contract pinned by tests)
    print(M("demo_next") % broken)
    return rc


if __name__ == "__main__":
    main()
