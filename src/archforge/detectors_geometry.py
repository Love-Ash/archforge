# -*- coding: utf-8 -*-
"""Geometry detectors: the rules that judge where things sit on the page.

  W15  two text frames whose effective glyph areas overlap
  W16  text glyphs or picture ink past the canvas edge
  W17  text straddling a picture's ink edge

They share one model of where a glyph actually is, built here rather than per rule: the
estimated ink box of a run after the effective size, the autofit scale and any group
transform. Three rules reading three slightly different boxes is how a geometry gate
starts contradicting itself, and the abstention path matters as much as the finding, which
is why a span the model cannot measure is pushed into W18 instead of guessed at.

Extracted from lint.py for the 0.7 decomposition (#5). Re-exported from lint for backward
compatibility.
"""
import math
from collections import namedtuple
from typing import List, Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from .findings import Finding, shape_loc
    from .messages import M
    from .ooxml import EMU_PER_IN, NS
    from .geometry import _geo_rect, frame_autofit, iter_shapes_geo
    from .scripts import is_cjk, geometry_unsupported as _geometry_unsupported
except ImportError:   # standalone execution
    from findings import Finding, shape_loc
    from messages import M
    from ooxml import EMU_PER_IN, NS
    from geometry import _geo_rect, frame_autofit, iter_shapes_geo
    from scripts import is_cjk, geometry_unsupported as _geometry_unsupported

try:
    from PIL import Image
except ImportError:   # pragma: no cover
    Image = None
import io as _io


class _FldRun:
    """Adapter for a:fld (auto fields such as slide number, date). Since CT_TextField also has
    an rPr+t structure per the schema, PowerPoint renders it with the same rules as a normal
    run, but python-pptx's para.runs only returns a:r, so field text was a blind spot for
    E1/E3/E4 (carried over from the fourth review, 0.5.0). Exposes only the minimal interface
    the checking code uses: ._r for run_fonts/run_track, .font.size for size."""
    __slots__ = ("_r", "text")

    class _Pt:
        __slots__ = ("pt",)

        def __init__(self, pt):
            self.pt = pt

    def __init__(self, fld):
        self._r = fld
        t = fld.find(NS + "t")
        self.text = (t.text or "") if t is not None else ""

    @property
    def font(self):
        return self   # only the .size access is used

    @property
    def size(self):
        try:
            v = self._r.find(NS + "rPr").get("sz")
        except Exception:
            return None
        if not v:
            return None
        try:
            return self._Pt(int(v) / 100.0)
        except (TypeError, ValueError):
            return None


# Effective glyph bbox (in) per paragraph. Turned a magic index tuple into named fields
# (external review, 2026-07-10).
# sp (owning shape) is for the loc payload of W15-W17 findings (0.5.0); the coordinates are
# already the group's absolute coordinates.
GlyphBox = namedtuple("GlyphBox", "x0 y0 x1 y1 rep max_pt frame_id sp cell para field")
GlyphBox.__new__.__defaults__ = (None, None, None, False)

# W15 text overlap: the most common defect axis in generated decks (elements pile up with
# every revision round), but the frame bbox is drawn generously by convention and can't be
# used, so this approximates the effective glyph width instead.
_W_CJK, _W_LAT, _W_SP = 0.96, 0.52, 0.28   # character-width/font-size ratio approximation
                                           # (conservative: suppresses false positives)


def _glyph_w(s, size_pt):
    w = 0.0
    for ch in s:
        if ch == " ":
            w += _W_SP
        elif is_cjk(ch) or ord(ch) > 0x2E80:
            w += _W_CJK
        else:
            w += _W_LAT
    return w * size_pt / 72.0


def _empty_para_pt(para, default_pt):
    """The effective size of an empty paragraph (a spacer): prioritizes endParaRPr/defRPr sz
    (fixes a phantom-height issue where a 4pt spacer was counted as 12pt, measured in
    adversarial verification)."""
    try:
        if para.font.size is not None:
            return para.font.size.pt
    except Exception:
        pass
    try:
        epr = para._p.find(NS + "endParaRPr")
        if epr is not None and epr.get("sz"):
            return int(epr.get("sz")) / 100.0
    except Exception:
        pass
    return default_pt


def _text_glyph_boxes(slide, default_pt=12.0, skipped=None, styler=None):
    """Approximates the effective glyph bbox (in) per paragraph. Returns
    [(x0,y0,x1,y1,representative text,max_pt,frame_id)].
    Width is summed from each run's actual size; line height reflects the actual
    line_spacing value (1.2 if absent) combined with autofit lnSpcReduction. x is placed
    using per-paragraph alignment (the frame's first explicit value if unset, otherwise
    left), which reduces misplacement in frames with mixed alignment. wrap=none
    (word_wrap=False, the python-pptx add_textbox default) extends past the frame in a
    single line, so no wrap folding is applied and the actual width is used as-is. A rotated
    frame is skipped since estimation would be invalid. Groups are converted to absolute
    coordinates via iter_shapes_geo.
    Calibrated against real-deck render comparisons plus reproduced adversarial-verification
    measurements.
    Script layer (0.2.1): vertical writing (bodyPr@vert) and frames containing RTL or
    complex-shaping scripts are skipped, since glyph-width approximation is meaningless for
    them, and if a skipped Counter is passed it tallies these and surfaces them via W18.
    0.3.1 (third external review, P0): if given a styler (StyleResolver), a run with no
    explicit size is resolved through the same inheritance chain as E3 (fixes an
    inconsistency where two different effective-style models existed in a single document).
    Native tables compute each cell's rectangle by accumulating column widths and row
    heights, so cell text is included too.
    Known limitations: a placeholder that inherits alignment from layout lstStyle falls back
    to left alignment, and skipping rotated frames is excluded from the W18 tally since it is
    standard decorative practice."""
    import math
    out = []

    def emit_frame(tframe, fx, fy, fw, fh, fid, owner_sp, cell=None, sx=1.0, sy=1.0):
        try:
            bodyPr = tframe._txBody.find(NS + "bodyPr")
            vert = bodyPr.get("vert") if bodyPr is not None else None
            if vert not in (None, "horz"):
                if skipped is not None:
                    skipped["vertical_text"] += 1
                return
        except Exception:
            bodyPr = None
        # Text-frame insets (0.6.0, external review): glyphs start inside the frame, not
        # at its edge. OOXML defaults are lIns/rIns 91440 EMU (0.1in) and tIns/bIns 45720
        # (0.05in), the same order of magnitude as W16's 0.15in tolerance, so ignoring
        # them shifted every glyph box left/up and overstated usable width. Insets live
        # in shape-local units, so group scale (sx/sy) applies.
        li, ri_, ti, bi = 91440, 91440, 45720, 45720
        try:
            if bodyPr is not None:
                li = int(bodyPr.get("lIns", li))
                ri_ = int(bodyPr.get("rIns", ri_))
                ti = int(bodyPr.get("tIns", ti))
                bi = int(bodyPr.get("bIns", bi))
        except Exception:
            pass
        fx += sx * li / EMU_PER_IN
        fy += sy * ti / EMU_PER_IN
        fw = max(fw - sx * (li + ri_) / EMU_PER_IN, 0.0)
        fh = max(fh - sy * (ti + bi) / EMU_PER_IN, 0.0)
        try:
            # All a:t descendants, not para.runs: field text (a:fld) participates in
            # geometry, so it must participate in the complex-script screen too. A
            # field-only Arabic frame used to bypass this check and get measured with
            # the Latin/CJK width model without any W18 (0.6.1, external review).
            frame_text = "".join(t.text or "" for t in tframe._txBody.iter(NS + "t"))
            if _geometry_unsupported(frame_text):
                if skipped is not None:
                    skipped["complex_script"] += 1
                return
        except Exception:
            pass
        fw2 = max(fw, 0.05)
        scale, lnred = frame_autofit(tframe)
        wrap = tframe.word_wrap is not False   # None (no attribute) = OOXML default
                                               # square = wrap
        frame_align = None
        for para in tframe.paragraphs:
            if para.alignment is not None:
                frame_align = para.alignment
                break
        paras = []   # (line_w, pmx, ptxt, factor, n, align, p_idx, field_only)
        for p_idx, para in enumerate(tframe.paragraphs):
            pmx, ptxt = 0.0, ""
            saw_field, saw_run_text = False, False
            segs = [0.0]   # widths of a:br-separated visual lines
            # Document-order walk over a:r / a:fld / a:br (0.6.0, external review): fld
            # text occupies real width and an explicit line break starts a new visual
            # line. Before this, a br-split sentence was measured as one overlong line
            # (width overstated, height understated). Falls back to para.runs on any
            # structural surprise.
            items = []
            try:
                runs_l = list(para.runs)
                r_seen = 0
                for child in para._p:
                    tag = child.tag
                    if tag == NS + "r":
                        if r_seen < len(runs_l):
                            items.append(runs_l[r_seen])
                            r_seen += 1
                    elif tag == NS + "br":
                        items.append(None)
                    elif tag == NS + "fld":
                        items.append(_FldRun(child))
                if r_seen != len(runs_l):
                    items = list(runs_l)
            except Exception:
                items = list(para.runs)
            for r in items:
                if r is None:   # a:br: next visual line
                    segs.append(0.0)
                    ptxt += " "
                    continue
                t = r.text
                if not t:
                    continue
                if isinstance(r, _FldRun):
                    saw_field = True
                else:
                    saw_run_text = True
                if r.font.size is not None:
                    sz = r.font.size.pt
                elif para.font.size is not None:
                    sz = para.font.size.pt
                else:
                    sz = None
                    if styler is not None:
                        try:
                            sz = styler.resolve(tframe, owner_sp, slide, getattr(para, "level", 0))
                        except Exception:
                            sz = None
                    if sz is None:
                        sz = default_pt
                sz *= scale
                segs[-1] += _glyph_w(t, sz)
                ptxt += t
                if sz > pmx:
                    pmx = sz
            if not ptxt.strip():
                paras.append((0.0, _empty_para_pt(para, default_pt) * scale, "", 1.2, 1,
                              None, p_idx, False))
                continue
            ls = para.line_spacing
            if ls is None:
                factor = 1.2
            elif isinstance(ls, float):
                factor = ls
            else:
                try:
                    factor = ls.pt / pmx if pmx else 1.2
                except Exception:
                    factor = 1.2
            if wrap:
                n = sum(max(1, math.ceil(w / (fw2 * 1.04))) for w in segs)
            else:
                n = len(segs)
            line_w = max(segs)
            al = para.alignment if para.alignment is not None else frame_align
            paras.append((line_w, pmx, ptxt, factor, n, al, p_idx,
                          saw_field and not saw_run_text))
        gh_total = sum(n * pmx * max(f, 0.95) * (1.0 - lnred) / 72.0
                       for (pw, pmx, ptxt, f, n, al, _pi, _fo) in paras)
        if gh_total <= 0:
            return
        va = str(tframe.vertical_anchor) if tframe.vertical_anchor is not None else ""
        if "MIDDLE" in va:
            cy = fy + max(0.0, (fh - gh_total) / 2)
        elif "BOTTOM" in va:
            cy = fy + max(0.0, fh - gh_total)
        else:
            cy = fy
        for (pw, pmx, ptxt, factor, n, al, p_idx, field_only) in paras:
            ph = n * pmx * max(factor, 0.95) * (1.0 - lnred) / 72.0
            if ptxt.strip():
                gw = pw if not wrap else min(pw, fw2)
                a = str(al) if al is not None else ""
                if "CENTER" in a:
                    x0 = fx + (fw2 - gw) / 2
                elif "RIGHT" in a:
                    x0 = fx + fw2 - gw
                else:
                    x0 = fx
                out.append(GlyphBox(x0, cy, x0 + gw, cy + ph, ptxt[:24], pmx, fid,
                                    owner_sp, cell, p_idx, field_only))
            cy += ph

    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "has_table", False):
            # Native table cells (third external review, P0: tables in auto-generated decks
            # were a geometry blind spot).
            # Cell rectangle = table origin + accumulated column widths/row heights (xf
            # scaling applied to EMU).
            geo = _geo_rect(sp, xf)
            if geo is None or geo[4]:
                continue
            tx, ty = geo[0], geo[1]
            ax = xf[0]
            ay = xf[2]
            try:
                tbl = sp.table
                col_w = [(c.width or 0) for c in tbl.columns]
                row_h = [(r.height or 0) for r in tbl.rows]
            except Exception:
                continue
            cy_off = 0
            for ri, row in enumerate(tbl.rows):
                cx_off = 0
                for ci, cell in enumerate(row.cells):
                    cw = col_w[ci] if ci < len(col_w) else 0
                    rh = row_h[ri] if ri < len(row_h) else 0
                    try:
                        # Merged regions (0.6.0, external review): continuation cells are
                        # covered by their origin; the origin's rectangle spans the merged
                        # column widths and row heights instead of a single grid cell.
                        if cell.is_spanned:
                            cx_off += cw
                            continue
                        span_w = sum(col_w[ci:ci + max(1, cell.span_width)]) or cw
                        span_h = sum(row_h[ri:ri + max(1, cell.span_height)]) or rh
                    except Exception:
                        span_w, span_h = cw, rh
                    try:
                        ctf = cell.text_frame
                    except Exception:
                        cx_off += cw
                        continue
                    emit_frame(ctf,
                               tx + ax * cx_off / EMU_PER_IN,
                               ty + ay * cy_off / EMU_PER_IN,
                               ax * span_w / EMU_PER_IN,
                               ay * span_h / EMU_PER_IN,
                               id(cell._tc), sp, cell=(ri, ci), sx=ax, sy=ay)
                    cx_off += cw
                cy_off += row_h[ri] if ri < len(row_h) else 0
            continue
        if not getattr(sp, "has_text_frame", False):
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        fx, fy, fw, fh, rotated = geo
        if rotated:
            continue
        emit_frame(sp.text_frame, fx, fy, fw, fh, id(sp), sp, sx=xf[0], sy=xf[2])
    return out


def text_overlap_check(slide, si, warns, boxes: Optional[List[GlyphBox]] = None):
    """W15: the effective glyph regions of two different text frames overlap meaningfully
    (occlusion/collision). This is approximation-based, hence WARN. Fires only when the
    intersection area exceeds 45% of the smaller box, at most 2 findings per page.
    The 45% threshold is measured against render comparisons: cases in the 30-35% range are
    all false positives (estimated as a title under a big number, or bleed-over between two
    columns), while 60%+ are all real overlaps (e.g. a chart overrunning an exhibit label, a
    caption chip sitting on a legend).
    Intentional layering is excluded: an echo of identical text (an afterimage typography
    effect), and 1-2 character oversized glyphs (drop caps, chapter numerals like I/II).
    If boxes is given, it is used as-is without recomputation (a once-per-slide computation
    cache)."""
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    hits = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            if a.frame_id == b.frame_id:      # paragraphs of the same frame are stacked,
                                               # so excluded
                continue
            if a.rep.strip() == b.rep.strip():
                continue
            if any(len(x.rep.strip()) <= 2 and x.max_pt >= 28 for x in (a, b)):
                continue
            ix = min(a.x1, b.x1) - max(a.x0, b.x0)
            iy = min(a.y1, b.y1) - max(a.y0, b.y0)
            if ix <= 0.02 or iy <= 0.02:
                continue
            area = ix * iy
            amin = min((a.x1 - a.x0) * (a.y1 - a.y0), (b.x1 - b.x0) * (b.y1 - b.y0))
            if amin > 0 and area > 0.45 * amin:
                hits.append((area / amin, a, b))
    # Sort key is frac only: GlyphBox's sp field isn't comparable, so a full tuple comparison
    # would fail (0.5.0)
    for frac, a, b in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        # loc: not the frame's raw bbox, but the effective glyph bbox (absolute, in) actually
        # used to judge the overlap, as-is.
        # related carries the counterpart frame so an agent can pinpoint which pair to move
        # (0.5.0).
        loc = shape_loc(a.sp, bbox=[a.x0, a.y0, a.x1 - a.x0, a.y1 - a.y0], cell=a.cell,
                        paragraph=a.para, field=a.field) or {}
        rel = shape_loc(b.sp, bbox=[b.x0, b.y0, b.x1 - b.x0, b.y1 - b.y0], cell=b.cell,
                        paragraph=b.para, field=b.field)
        if rel:
            loc["related"] = rel
        warns.append(Finding(si, "W15", "w15", (frac * 100,), "%r ~ %r" % (a.rep, b.rep),
                             data={"confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False},
                             loc=loc or None))


def _pic_boxes(slide, sw_in, sh_in, skipped=None):
    """The effective ink bbox (in) and z-order of non-background pictures. Returns
    [(x0,y0,x1,y1,z)].
    Full-bleed or mesh backgrounds covering 70%+ of the slide are excluded. A transparent PNG
    (e.g. a matplotlib chart) has a frame bbox much larger than its ink, which is a source of
    false positives, so it is trimmed to the alpha-opaque bbox.
    Performance budget (0.4.0, third review): for images over 25MP, alpha trimming is skipped,
    the frame bbox is used as-is, and this is disclosed via the skipped counter (prevents
    decode blowup on large batches).
    Reflects adversarial-verification measurements (2026-07-03): P mode + tRNS is converted
    to RGBA before trimming, srcRect crop is mapped by narrowing to the visible source
    window, flipH/flipV mirrors within that window, a rotated picture uses the axis-aligned
    expanded bbox but skips ink trimming since it would be invalid, and group children are
    converted to absolute coordinates."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, rotated = geo
        if w * h >= 0.7 * sw_in * sh_in:
            continue
        if not rotated:
            try:
                from PIL import Image
                import io as _io
                im = Image.open(_io.BytesIO(sp.image.blob))
                if im.width * im.height > 25_000_000:
                    if skipped is not None:
                        skipped["image_decode_budget"] += 1
                    raise RuntimeError("image decode budget")   # the except below keeps the
                                                                 # frame bbox
                if im.mode == "P" and "transparency" in im.info:
                    im = im.convert("RGBA")
                if "A" in im.getbands():
                    bb = im.getchannel("A").point(lambda a: 255 if a > 16 else 0).getbbox()
                    if bb is None:
                        continue   # fully transparent = no ink
                    iw, ih = im.size
                    wl, wt = iw * float(sp.crop_left or 0), ih * float(sp.crop_top or 0)
                    wr = iw * (1.0 - float(sp.crop_right or 0))
                    wb = ih * (1.0 - float(sp.crop_bottom or 0))
                    l, t, r, b = bb
                    l, r = max(l, wl), min(r, wr)
                    t, b = max(t, wt), min(b, wb)
                    if l >= r or t >= b:
                        continue   # no ink within the crop window
                    try:
                        x2 = sp._element.spPr.find(NS + "xfrm")
                        if x2 is not None and x2.get("flipH") == "1":
                            l, r = wl + (wr - r), wl + (wr - l)
                        if x2 is not None and x2.get("flipV") == "1":
                            t, b = wt + (wb - b), wt + (wb - t)
                    except Exception:
                        pass
                    ww, wh = (wr - wl) or 1.0, (wb - wt) or 1.0
                    x, y, w, h = (x + w * (l - wl) / ww, y + h * (t - wt) / wh,
                                  w * (r - l) / ww, h * (b - t) / wh)
            except Exception as ex:
                # The frame bbox is kept as the fallback either way, but a real decode
                # failure is no longer silent (0.6.0, external review: "nothing dies
                # silently" only held for the budget path). The budget path already
                # tallied itself above.
                if skipped is not None and str(ex) != "image decode budget":
                    skipped["image_decode"] += 1
        out.append((x, y, x + w, y + h, z, sp))
    return out


def overflow_check(slide, si, sw_in, sh_in, warns,
                   boxes: Optional[List[GlyphBox]] = None, pics: Optional[list] = None):
    """W16: off-canvas overflow. Using the frame bbox as the criterion was previously
    rejected because of the large false-positive rate from the generous-frame convention, but
    W15's effective glyph bbox resolved that objection (2026-07-03): text fires only when the
    actual character area breaches the boundary. For non-text, only pictures (ink bbox,
    trimmed) are checked: bleed where a decorative shape (e.g. a glow circle) spills off a
    corner is standard technique and not a defect (checking shapes was rejected after
    corpus render measurements)."""
    TOL_T, TOL_S = 0.15, 0.12
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    if pics is None:
        pics = _pic_boxes(slide, sw_in, sh_in)
    hits = []
    for gb in boxes:
        over = max(-gb.x0, -gb.y0, gb.x1 - sw_in, gb.y1 - sh_in)
        if over > TOL_T:
            hits.append((over, M("w16_text") % gb.rep, "t|%r" % gb.rep,
                         shape_loc(gb.sp, bbox=[gb.x0, gb.y0, gb.x1 - gb.x0, gb.y1 - gb.y0],
                                   cell=gb.cell, paragraph=gb.para, field=gb.field)))
    for (px0, py0, px1, py1, _z, psp) in pics:
        over = max(-px0, -py0, px1 - sw_in, py1 - sh_in)
        if over > TOL_S:
            hits.append((over, M("w16_pic") % (px1 - px0, py1 - py0),
                         "p|%.1fx%.1f" % (px1 - px0, py1 - py0),
                         shape_loc(psp, bbox=[px0, py0, px1 - px0, py1 - py0])))
    # Sort key is over only: a loc dict isn't comparable, so a full tuple comparison would
    # fail (0.5.0)
    for over, what, fpk, loc in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        # fp_key: detail (what) is a locale-dependent string, so it's excluded from the
        # baseline fingerprint (fourth review)
        warns.append(Finding(si, "W16", "w16", (over,), what, fp_key=fpk, loc=loc,
                             data={"kind": "text" if fpk.startswith("t|") else "picture",
                                   "confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False}))


def _occluder_boxes(slide, sw_in, sh_in):
    """The bbox and z of solid-fill shapes (cards, panels) sitting on top of a picture. Used
    to suppress a legitimate layout, a caption card over a photo, that was falsely caught by
    W17 (measured in adversarial verification)."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            continue
        try:
            if sp.fill.type is None or "SOLID" not in str(sp.fill.type):
                continue
        except Exception:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, _rot = geo
        if w * h >= 0.9 * sw_in * sh_in:
            continue
        out.append((x, y, x + w, y + h, z))
    return out


def text_image_straddle_check(slide, si, sw_in, sh_in, warns,
                              boxes: Optional[List[GlyphBox]] = None, pics: Optional[list] = None):
    """W17: text straddles the ink boundary of a non-background picture (only 25-75% of the
    glyph is inside the image) = half on, half off, so it looks cropped or the background
    looks split. Fully on top (an overlay caption) is the contrast-gate's (W7's) jurisdiction,
    and fully off is irrelevant. Pictures under 1 square inch (icon/logo scale) are ignored.
    If a solid card sits in z-order between the photo and the text and backs 90%+ of the text
    area, this is excluded as a caption on a card rather than a straddle.
    At most 2 findings per page."""
    if pics is None:
        pics = _pic_boxes(slide, sw_in, sh_in)
    if not pics:
        return
    occl = _occluder_boxes(slide, sw_in, sh_in)
    if boxes is None:
        boxes = _text_glyph_boxes(slide)
    hits = []
    for gb in boxes:
        rep = gb.rep
        if len(rep.strip()) < 3:
            continue
        ta = (gb.x1 - gb.x0) * (gb.y1 - gb.y0)
        if ta <= 0:
            continue
        for (px0, py0, px1, py1, pz, psp) in pics:
            if (px1 - px0) * (py1 - py0) < 1.0:
                continue
            ix = min(gb.x1, px1) - max(gb.x0, px0)
            iy = min(gb.y1, py1) - max(gb.y0, py0)
            if ix <= 0 or iy <= 0:
                continue
            frac = (ix * iy) / ta
            if not (0.25 <= frac <= 0.75):
                continue
            carded = False
            for (ox0, oy0, ox1, oy1, oz) in occl:
                if oz <= pz:
                    continue
                cx = min(gb.x1, ox1) - max(gb.x0, ox0)
                cy2 = min(gb.y1, oy1) - max(gb.y0, oy0)
                if cx > 0 and cy2 > 0 and cx * cy2 >= 0.9 * ta:
                    carded = True
                    break
            if not carded:
                hits.append((frac, gb, (px0, py0, px1, py1, psp)))
    # Sort key is frac only (sp field isn't comparable, 0.5.0)
    for frac, gb, pic in sorted(hits, key=lambda h: h[0], reverse=True)[:2]:
        loc = shape_loc(gb.sp, bbox=[gb.x0, gb.y0, gb.x1 - gb.x0, gb.y1 - gb.y0],
                        cell=gb.cell, paragraph=gb.para, field=gb.field) or {}
        rel = shape_loc(pic[4], bbox=[pic[0], pic[1], pic[2] - pic[0], pic[3] - pic[1]])
        if rel:
            loc["related"] = rel
        warns.append(Finding(si, "W17", "w17", (frac * 100,), "%r" % gb.rep, loc=loc or None,
                             data={"confidence": "estimate",
                                   "evidence_source": "xml_geometry",
                                   "render_confirmed": False}))
