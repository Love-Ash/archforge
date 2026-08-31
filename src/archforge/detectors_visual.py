# -*- coding: utf-8 -*-
"""Appearance detectors: the rules that judge how a page is drawn.

  W7   low contrast between text and the image under it (needs rendered pages)
  W9   accent bars repeated as list markers
  W12  footer baselines drifting from the dominant baseline
  W13  native PowerPoint shadow, glow and 3D effects
  W21  a rare color sitting next to a dominant near-identical one

Plus the two signature helpers the repetition rules read, the layout skeleton and the
fill-token set, which W6 and W10 compare across pages.

Extracted from lint.py for the 0.7 decomposition (#5). W7 is the one detector here that
touches the filesystem, because it reads the PNGs the caller exported; everything else
works off the shape tree. Re-exported from lint for backward compatibility.
"""
import glob
import math
import os
from collections import Counter

from pptx import Presentation

try:
    from .findings import Finding, shape_loc
    from .ooxml import EMU_PER_IN, NS, NS_P
    from .colors import (_hex_rgb, _is_accent, _luma, _resolve_run_rgb,
                         _shape_fill_hex, _shape_line_hex, _COLOR_UNKNOWN)
    from .geometry import (_geo_rect, _is_pic, iter_shapes, iter_shapes_geo,
                           collect_frames as _collect_frames)
    from .inline import iter_inline_items
except ImportError:   # standalone execution
    from findings import Finding, shape_loc
    from ooxml import EMU_PER_IN, NS, NS_P
    from colors import (_hex_rgb, _is_accent, _luma, _resolve_run_rgb,
                        _shape_fill_hex, _shape_line_hex, _COLOR_UNKNOWN)
    from geometry import (_geo_rect, _is_pic, iter_shapes, iter_shapes_geo,
                          collect_frames as _collect_frames)
    from inline import iter_inline_items

try:
    from PIL import Image
except ImportError:   # pragma: no cover - W7 abstains without Pillow
    Image = None

_EFFECT_TAGS = tuple(NS + t for t in ("outerShdw", "innerShdw", "glow", "reflection"))

_3D_TAGS = tuple(NS + t for t in ("sp3d", "scene3d"))


def accent_vbars_check(slide, si, sw, sh, warns):
    """W9: an AI-generated-deck tell where accent-colored vertical bars are repeated as list
    markers, using color to build item structure (measured in real decks). Reflects the
    2026-07-02 adversarial audit: colors are read via the sp.line/sp.fill accessors to cover
    namespaces and connectors, vertical bars explicitly include zero-width connectors (w~0),
    and adjacent text to the right confirms it is really a list marker."""
    bars, texts = [], []
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht):
            continue
        x, y, w, h = L / EMU_PER_IN, T / EMU_PER_IN, Wd / EMU_PER_IN, Ht / EMU_PER_IN
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            texts.append((x, y, w, h))
        hexc = _shape_line_hex(sp) or _shape_fill_hex(sp)
        if not _is_accent(hexc):
            continue
        if 0.2 <= h <= 1.0 and (w < 0.05 or h > 3 * w):   # vertical bar (including
                                                           # zero-width connectors)
            bars.append((x, y, w, h, hexc))
    if len(bars) < 3:
        return
    hues = {b[4] for b in bars}
    if len(hues) != 1:                       # multiple colors is legitimate data encoding
                                              # (a legend)
        return
    xs = [b[0] for b in bars]
    if max(xs) - min(xs) > 0.15:             # vertically aligned stack only (horizontal
                                              # spread = chart/divider, excluded)
        return
    rt = 0
    for bx, by, bw, bh, _hx in bars:
        for tx, ty, tw, th in texts:
            if tx > bx and (tx - (bx + bw)) < 0.6 and not (ty > by + bh or ty + th < by):
                rt += 1
                break
    if rt >= len(bars) - 1:
        warns.append(Finding(si, "W9", "w9", (len(bars),),
                         "x=%.2fin hue=%s" % (min(xs), next(iter(hues)))))


def _fill_tokens(slide, sw, sh):
    """Turns the slide's solid-fill shapes into a multiset of (color, 24-grid position/size)
    tokens. Full-bleed backgrounds are excluded."""
    t = Counter()
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:
            continue
        fh = _shape_fill_hex(sp)
        if fh is None:
            continue
        t[(fh, round(L / sw * 24), round(T / sh * 24), round(Wd / sw * 24), round(Ht / sh * 24))] += 1
    return t


def footer_top(slide, sw, sh):
    """The top (in) of the bottommost text in the slide's bottom band (y>0.88H). None if there
    is no footer."""
    best = None
    for sp in iter_shapes(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        try:
            if not sp.text_frame.text.strip():
                continue
            t = sp.top
        except Exception:
            continue
        if t is None or t <= 0.88 * sh:
            continue
        ti = t / EMU_PER_IN
        if best is None or ti > best:
            best = ti
    return best


def footer_check(foot_tops, warns):
    """W12: footer baseline misalignment across pages. In a 50-deck measurement (2026-07-02),
    the absolute-deviation approach mistook cover-page credits and bottom captions for footers,
    producing false positives in 17 decks. Instead this excludes the cover page (p1), treats
    the dominant baseline (mode of the 0.05in quantization buckets, 3+ pages) as the house
    footer, and flags only pages that deviate "slightly" (0.03-0.25in) from it. Anything off
    by more than 0.25in is assumed to be a different element such as a caption or divider and
    is ignored (its existence is not even checked)."""
    tops = [(si, t) for si, t in foot_tops.items() if t is not None and si > 1]
    if len(tops) < 4:
        return
    q = Counter(round(t / 0.05) for _si, t in tops)
    qv, cnt = q.most_common(1)[0]
    if cnt < 3:
        return
    # House baseline = the median of the actual values in the dominant bucket (using the
    # bucket center instead would make a majority of 7.08 values sit exactly 0.02 away from
    # 7.10, a floating-point boundary false positive confirmed by rescanning the 50 decks)
    bvals = sorted(t for _si, t in tops if round(t / 0.05) == qv)
    base = bvals[len(bvals) // 2]
    off = [(si, t) for si, t in tops if 0.03 < abs(t - base) <= 0.25]
    if off:
        ex = " ".join("p%d=%.2f" % (si, t) for si, t in off[:4])
        warns.append(Finding(0, "W12", "w12", (base, len(off)), ex))


def effects_count(slide):
    """The count and kinds of the slide's effective PPT effects (shadow, glow, 3D). Some
    generators leave a childless empty effectLst purely to block inheritance, so this only
    counts elements that actually have effect children (prevents an empty-element false
    positive)."""
    n = 0
    kinds = set()
    for sp in iter_shapes(slide.shapes):
        spPr = getattr(sp._element, "spPr", None)
        if spPr is None:
            continue
        eff = spPr.find(NS + "effectLst")
        if eff is not None:
            for ch in eff:
                if ch.tag in _EFFECT_TAGS:
                    n += 1
                    kinds.add(ch.tag.split("}")[1])
        for tag in _3D_TAGS:
            if spPr.find(tag) is not None:
                n += 1
                kinds.add(tag.split("}")[1])
    return n, kinds


def effects_check_deck(per_page, warns):
    """W13: aggregated once per deck (firing repeatedly per page is noise: measured on the
    50-deck corpus). Could be an intentional neon/glow style, so this is a WARN and the
    judgment is left to human eyes."""
    hits = [(si, n, kinds) for si, (n, kinds) in per_page.items() if n >= 2]
    if not hits:
        return
    total = sum(n for _si, n, _k in hits)
    kinds = sorted(set().union(*[k for _si, _n, k in hits]))
    pages = ",".join("p%d" % si for si, _n, _k in hits[:6])
    warns.append(Finding(0, "W13", "w13", (total, len(hits)),
                         "%s | %s" % (pages, ",".join(kinds))))


def _diagram_clone_marks(inter):
    """Counts "decorative texture clones" (small dots, joints, etc., 1x1 or smaller on the
    24-grid) in the multiset of fill shapes shared between two pages.
    Card-shaped (mid-size block) and table (full-width band) elements are W6's jurisdiction
    and are not counted here, avoiding a false positive on three-column comparison cards
    (reflecting the adversarial audit)."""
    marks = area = 0
    for (fh, gx, gy, gw, gh), c in inter.items():
        if gw <= 1 and gh <= 1:
            marks += c
        if gw >= 1 and gh >= 1:
            area += gw * gh * c
    if marks >= 8 and area / (24.0 * 24.0) >= 0.06:
        return marks
    return 0


def slide_layout_sig(slide, sw, sh, gw=6, gh=4):
    """Turns slide element placement into a gw x gh grid occupancy vector. Full-bleed
    backgrounds are excluded.
    Weights: text=1, shape=0.5, image=2. Used to compare "what is where" (the skeleton)."""
    sig = [0.0] * (gw * gh)
    n = 0
    for sp in iter_shapes(slide.shapes):
        try:
            L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
        except Exception:
            continue
        if None in (L, T, Wd, Ht) or not Wd or not Ht:
            continue
        if Wd > 0.9 * sw and Ht > 0.9 * sh:      # excludes full-bleed backgrounds/images
            continue
        cx = (L + Wd / 2) / sw; cy = (T + Ht / 2) / sh
        if not (0 <= cx <= 1.0 and 0 <= cy <= 1.0):
            continue
        gc = min(gw - 1, max(0, int(cx * gw))); gr = min(gh - 1, max(0, int(cy * gh)))
        wgt = 2.0 if _is_pic(sp) else (1.0 if getattr(sp, "has_text_frame", False) else 0.5)
        sig[gr * gw + gc] += wgt
        n += 1
    return sig, n


def contrast_check(slide, si, sw, sh, render_dir, warns, styler=None, thm_colors=None,
                   skipped=None):
    """Detects low-contrast text over an image (approximated from the rendered PNG). Only
    text frames overlapping a picture, once per slide.
    Returns: "no_pics" (nothing to check) / "no_png" (there is a picture but no conventional
    render = incomplete) / "ok" (checked).
    Coordinates are absolute, including the group transform (third review P1: fixes pictures
    and text inside a group being misaligned when using raw coordinates)."""
    from PIL import Image
    pics = []
    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if _is_pic(sp):
            geo = _geo_rect(sp, xf)
            if geo is not None:
                x, y, w, h, _rot = geo
                pics.append((x * EMU_PER_IN, y * EMU_PER_IN, w * EMU_PER_IN, h * EMU_PER_IN))
    if not pics:
        return "no_pics"
    cand = glob.glob(os.path.join(render_dir, "p%02d.png" % si))
    if not cand:
        return "no_png"
    try:
        im = Image.open(cand[0]).convert("RGB"); px = im.load(); PW, PH = im.size
    except Exception:
        return "no_png"
    for sp, _z, xf in iter_shapes_geo(slide.shapes):
        if not getattr(sp, "has_text_frame", False):
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        gx, gy, gw_, gh_, _rot = geo
        L, T, Wd, Ht = gx * EMU_PER_IN, gy * EMU_PER_IN, gw_ * EMU_PER_IN, gh_ * EMU_PER_IN
        over = any(not (L + Wd <= p0 or L >= p0 + pw or T + Ht <= q0 or T >= q0 + ph)
                   for p0, q0, pw, ph in pics)
        if not over:
            continue
        rgbs = [_resolve_run_rgb(r, para, sp.text_frame, sp, slide, styler, thm_colors)
                for para in sp.text_frame.paragraphs for r in para.runs if r.text.strip()]
        if any(c is _COLOR_UNKNOWN for c in rgbs):
            # An explicit color we cannot decode: judging with an inherited color instead
            # produced false positives (0.6.1). Abstain on this frame and surface it.
            if skipped is not None:
                skipped["w7_color_unknown"] += 1
            continue
        rgbs = [c for c in rgbs if c]
        if not rgbs:
            continue
        txt_rgb = rgbs[0]
        x0 = max(0, int(L / sw * PW)); y0 = max(0, int(T / sh * PH))
        x1 = min(PW, int((L + Wd) / sw * PW)); y1 = min(PH, int((T + Ht) / sh * PH))
        if x1 <= x0 or y1 <= y0:
            continue
        sx = max(1, (x1 - x0) // 24); sy = max(1, (y1 - y0) // 24)
        lumas = sorted(_luma(px[x, y]) for x in range(x0, x1, sx) for y in range(y0, y1, sy))
        if not lumas:
            continue
        L_txt = _luma(txt_rgb)
        # Background luma is taken from the quantile at the opposite extreme of the text
        # color: this measures worst-case local contrast (WCAG is based on the worst case)
        # and avoids mean contamination from mixed-in text ink (dark 15th percentile for light
        # text, light 85th percentile for dark text).
        if L_txt >= 0.5:
            L_bg = lumas[int(len(lumas) * 0.15)]
        else:
            L_bg = lumas[min(len(lumas) - 1, int(len(lumas) * 0.85))]
        hi = max(L_bg, L_txt); lo = min(L_bg, L_txt)
        ratio = (hi + 0.05) / (lo + 0.05)
        if ratio < 2.5:
            warns.append(Finding(si, "W7", "w7", (ratio,),
                                 "text=%r" % sp.text_frame.text[:20],
                                 data={"confidence": "measured",
                                       "evidence_source": "render",
                                       "render_confirmed": True},
                                 loc=shape_loc(sp, bbox=[gx, gy, gw_, gh_])))
            return "ok"
    return "ok"   # found and checked the render PNG for this page (regardless of whether
                  # W7 actually fired)


def solid_contrast_check(slide, si, warns, styler=None, thm_colors=None, skipped=None,
                         skipped_locs=None):
    """W19: text whose color is nearly indistinguishable from its own shape's solid fill.

    The renderless sibling of W7. W7 needs a rendered PNG because an image background is
    unknowable from XML; a solid fill is exactly knowable, and the most common AI-deck
    contrast defect -- light gray text in a white box, or ghost placeholder text left the
    same color as its fill -- never touches an image at all.

    Scope is deliberately narrow so silence stays honest: only the run's OWN shape fill
    (slide backgrounds and shapes underneath need occlusion logic and belong to the
    rendered path), and only when both the fill and the effective run color resolve to
    definite RGB values. An explicit-but-undecodable color (hslClr and friends,
    _COLOR_UNKNOWN) is surfaced as a w19_color_unknown abstention; a run with no explicit
    color anywhere is skipped without noise, because no claim about its contrast is
    checkable from the XML.

    The 2.0:1 threshold is calibrated, not chosen: across the 29-deck private set,
    everything under 2.0 was a same-color ghost placeholder (1.0:1) or a near-invisible
    watermark (1.85:1), and everything intentional -- white display type on brand colors --
    sat at 2.3:1 or higher. Findings cap at 2 per page with a w19_capped disclosure, the
    same pattern the geometry gates use."""
    hits = []
    for tf, _w, sp, cell_rc, sp_xf in _collect_frames(slide.shapes):
        try:
            fill = _shape_fill_hex(sp)
        except Exception:
            fill = None
        if not fill:
            continue
        try:
            bg = (int(fill[0:2], 16), int(fill[2:4], 16), int(fill[4:6], 16))
        except Exception:
            continue
        for pi, para in enumerate(tf.paragraphs):
            for ri, (run_like, _rix, is_fld) in enumerate(iter_inline_items(para)):
                if run_like is None or not (run_like.text or "").strip():
                    continue
                try:
                    rgb = _resolve_run_rgb(run_like, para, tf, sp, slide, styler, thm_colors)
                except Exception:
                    continue
                if rgb is None:
                    continue
                if rgb is _COLOR_UNKNOWN:
                    if skipped is not None:
                        skipped["w19_color_unknown"] += 1
                    continue
                lt, lb = _luma(rgb), _luma(bg)
                ratio = (max(lt, lb) + 0.05) / (min(lt, lb) + 0.05)
                if ratio < 2.0:
                    hits.append((ratio, si, fill, rgb, run_like.text, sp, pi, ri,
                                 cell_rc, sp_xf, is_fld))
    if skipped is not None and len(hits) > 2:
        skipped["w19_capped"] += len(hits) - 2
    for ratio, si_, fill, rgb, text, sp, pi, ri, cell_rc, sp_xf, is_fld in \
            sorted(hits, key=lambda h: h[0])[:2]:
        warns.append(Finding(si_, "W19", "w19", (ratio,),
                             "bg=#%s fg=#%02X%02X%02X text=%r"
                             % (fill, rgb[0], rgb[1], rgb[2], text[:24]),
                             data={"contrast_ratio": round(ratio, 2),
                                   "bg_hex": fill,
                                   "fg_hex": "%02X%02X%02X" % rgb,
                                   "confidence": "estimate",
                                   "evidence_source": "xml_colors"},
                             loc=shape_loc(sp, paragraph=pi, run=ri, cell=cell_rc,
                                           xf=sp_xf, field=is_fld)))


_BG_UNKNOWN = object()   # explicit background present but not decodable from XML


def _filled_backdrops(slide, sw_in, sh_in):
    """Solid-fill shapes that can act as the background of text drawn over them.

    Close cousin of the occluder list W17 builds, but kept separate because the question is
    different: W17 asks whether a card hides a photo, this asks what color is behind a run.
    So the fill hex travels with the box, and the same exclusions apply for the same reasons.
    Pictures are out (their color is not knowable from XML, that is W7's job), shapes that
    carry their own text are out (a run inside its own fill is W19, and two texts colliding
    is W15), and a shape covering almost the whole page is the background rather than an
    element."""
    out = []
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        if getattr(sp, "shape_type", None) is not None and _is_pic(sp):
            continue
        if getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip():
            continue
        try:
            fill = _shape_fill_hex(sp)
        except Exception:
            continue
        if not fill:
            continue
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, _rot = geo
        if w <= 0 or h <= 0 or w * h >= 0.9 * sw_in * sh_in:
            continue
        out.append((x, y, x + w, y + h, z, fill, sp))
    return out


def _shape_z(slide):
    """Map each shape to its z index, so a run can be compared against what is under it.

    Keyed on the underlying XML element, never on the shape object. python-pptx builds a
    fresh proxy on every traversal, so two walks over the same slide hand back different
    objects for the same shape and an id(sp) lookup silently misses every time -- the gate
    stays quiet and looks like it found nothing. lxml caches its element proxies, so the
    element is stable across walks (verified on this slide before relying on it)."""
    return {id(sp._element): z for sp, z, _xf in iter_shapes_geo(slide.shapes)}


def _slide_bg_hex(slide, thm_colors):
    """The slide's own background color, resolved the way PowerPoint renders it.

    Walks slide -> layout -> master and takes the first p:bg found, because a lower
    level's background only shows through when nothing above declares one. Two forms
    cover what actually ships (measured on this corpus and the decks this gate was
    built against): an explicit p:bgPr solidFill, srgbClr or schemeClr through the
    theme, and the stock `p:bgRef idx="1001"` whose color child names the scheme
    slot. That second form leans on the built-in theme's first background fill style
    being a solid phClr; a theme that redefines style 1001 gets _BG_UNKNOWN, never a
    guess, and so does any explicit gradient or picture background or unresolvable
    scheme name. No p:bg at any level falls through to the theme's bg1/lt1, which is
    what an untouched deck renders."""
    parts = [slide._element]
    try:
        parts.append(slide.slide_layout._element)
        parts.append(slide.slide_layout.slide_master._element)
    except Exception:
        pass
    for el in parts:
        csld = el.find(NS_P + "cSld")
        bg = csld.find(NS_P + "bg") if csld is not None else None
        if bg is None:
            continue
        bgpr = bg.find(NS_P + "bgPr")
        if bgpr is not None:
            solid = bgpr.find(NS + "solidFill")
            if solid is None:
                return _BG_UNKNOWN
            return _bg_color_hex(solid, thm_colors)
        ref = bg.find(NS_P + "bgRef")
        if ref is not None:
            if ref.get("idx") != "1001":
                return _BG_UNKNOWN
            return _bg_color_hex(ref, thm_colors)
        return _BG_UNKNOWN
    if thm_colors:
        return thm_colors.get("bg1") or thm_colors.get("lt1")
    return None


def _bg_color_hex(parent, thm_colors):
    srgb = parent.find(NS + "srgbClr")
    if srgb is not None and srgb.get("val"):
        return srgb.get("val").upper()
    sch = parent.find(NS + "schemeClr")
    if sch is not None and sch.get("val"):
        hexv = (thm_colors or {}).get(sch.get("val"))
        return hexv.upper() if hexv else _BG_UNKNOWN
    return _BG_UNKNOWN


def _bg_exposure_walk(slide, sw_in, sh_in):
    """Everything standing between a text frame and the page background.

    Returns (blockers, page). blockers are (x0, y0, x1, y1, z, solid-hex-or-None, element)
    for every shape that renders ink of its own: pictures, and any shape whose fill
    python-pptx does not report as the explicit no-fill. Type 5 covers both an
    a:noFill and an unfilled text box (measured: a bare add_textbox reports 5), and
    None means the fill is inherited from the shape style, which renders filled, so
    None blocks too. page is the topmost full-bleed solid textless shape promoted to
    page background, (hex-or-unknown, z, element): the case where a deck paints its
    ground with a rectangle instead of p:bg."""
    blockers = []
    page = None
    for sp, z, xf in iter_shapes_geo(slide.shapes):
        geo = _geo_rect(sp, xf)
        if geo is None:
            continue
        x, y, w, h, _rot = geo
        if w <= 0 or h <= 0:
            continue
        if _is_pic(sp):
            blockers.append((x, y, x + w, y + h, z, True, sp._element))
            continue
        try:
            ftype = sp.fill.type
        except Exception:
            continue
        if ftype is not None and int(ftype) == 5:
            continue
        try:
            solid = _shape_fill_hex(sp)
        except Exception:
            solid = None
        if w * h >= 0.9 * sw_in * sh_in and not (
                getattr(sp, "has_text_frame", False) and sp.text_frame.text.strip()):
            hexv = solid if solid else _BG_UNKNOWN
            if page is None or z > page[1]:
                page = (hexv, z, sp._element)
            continue
        blockers.append((x, y, x + w, y + h, z, solid, sp._element))
    return blockers, page


def underlying_contrast_check(slide, si, sw_in, sh_in, warns, boxes=None, styler=None,
                              thm_colors=None, skipped=None, min_cover=0.5):
    """W20: text buried on what is actually visible behind it.

    W19 answers the same question one shape earlier: a run inside its own fill. It stops
    there on purpose, and its docstring says why -- anything underneath "needs occlusion
    logic and belongs to the rendered path". That was true when W19 shipped. It stopped
    being true once W17 built the z-ordered box list, so the case is now decidable from
    XML with no render at all, and this gate closes it.

    The background of a glyph is whatever is painted directly beneath it, so the layers
    under the text are consumed in descending z: an upper card claims its overlap first,
    and a shape below it only counts for whatever area is still exposed. The first cut
    of this gate compared the text against every lower layer independently, and a
    readable label on a bright bar was flagged against the page behind the bar -- the
    fixture that caught it is in the corpus. Whatever area no shape claims falls through
    to the resolved slide background (an explicit solid p:bg at slide, layout, or master
    level, the stock bgRef-1001 style through the theme, or a full-bleed solid shape
    painted as the ground), which is how ghost text left on an empty slide gets caught.

    Low-contrast exposure is summed across layers, not judged per shape: a caption
    spanning two bars is buried by both (measured on a real deck, 27.5% + 27.5%, zero
    on a per-shape test). Undecodable paint under the text -- gradients, pictures,
    style-inherited fills, an undecodable background -- abstains as w20_fill_unknown
    once it claims enough of the glyph area, because no honest verdict exists there.
    A text frame with its own fill is out entirely: a solid one is W19's verdict, and
    a style-inherited one renders as a color this gate cannot know.

    Contrast uses the same 2.0:1 line as W19; the 50% floor is the parameter still
    soaking. Findings cap at 2 per page with a w20_capped disclosure."""
    if boxes is None:
        return
    blockers, page = _bg_exposure_walk(slide, sw_in, sh_in)
    slide_bg = _slide_bg_hex(slide, thm_colors)
    zmap = _shape_z(slide)
    hits = []
    for gb in boxes:
        if not (gb.rep or "").strip():
            continue
        area = (gb.x1 - gb.x0) * (gb.y1 - gb.y0)
        if area <= 0:
            continue
        tz = zmap.get(id(gb.sp._element))
        if tz is None:
            continue
        try:
            own = gb.sp.fill.type
        except Exception:
            continue
        if own is not None and int(own) != 5:
            continue
        rgb = _paragraph_rgb(gb, slide, styler, thm_colors, skipped)
        if rgb is None:
            continue
        # Consume the area top-down: each lower layer claims what is left of the
        # glyph box after every layer above it took its share.
        layers = []
        for (bx0, by0, bx1, by1, bz, fill, bel) in blockers:
            if bz >= tz or bel is gb.sp._element:
                continue
            if page is not None and bel is page[2]:
                continue
            ix = min(gb.x1, bx1) - max(gb.x0, bx0)
            iy = min(gb.y1, by1) - max(gb.y0, by0)
            if ix <= 0 or iy <= 0:
                continue
            layers.append((bz, ix * iy, fill, bel))
        layers.sort(key=lambda l: l[0], reverse=True)
        remaining = area
        buried = 0.0
        unknown = 0.0
        worst = None
        for (bz, ov, fill, bel) in layers:
            take = min(ov, remaining)
            if take <= 0:
                continue
            if fill is None:
                unknown += take
            else:
                try:
                    bg = (int(fill[0:2], 16), int(fill[2:4], 16), int(fill[4:6], 16))
                except Exception:
                    unknown += take
                    remaining -= take
                    continue
                lt, lb = _luma(rgb), _luma(bg)
                ratio = (max(lt, lb) + 0.05) / (min(lt, lb) + 0.05)
                if ratio < 2.0:
                    buried += take
                    if worst is None or ratio < worst[0]:
                        worst = (ratio, fill, bel, "w20")
            remaining -= take
            if remaining <= 0:
                break
        if remaining > 0:
            bg_hex = slide_bg
            if page is not None:
                if tz <= page[1]:
                    remaining = 0       # under a full-bleed overlay: not visible at all
                    bg_hex = None
                else:
                    bg_hex = page[0]
            if remaining > 0 and bg_hex is _BG_UNKNOWN:
                unknown += remaining
            elif remaining > 0 and bg_hex:
                try:
                    bg = (int(bg_hex[0:2], 16), int(bg_hex[2:4], 16),
                          int(bg_hex[4:6], 16))
                    lt, lb = _luma(rgb), _luma(bg)
                    ratio = (max(lt, lb) + 0.05) / (min(lt, lb) + 0.05)
                    if ratio < 2.0:
                        buried += remaining
                        if worst is None or ratio < worst[0]:
                            worst = (ratio, bg_hex, None, "w20_bg")
                except Exception:
                    pass
        if unknown / area >= min_cover:
            if skipped is not None:
                skipped["w20_fill_unknown"] += 1
            continue
        if worst is None:
            continue
        cover = min(buried / area, 1.0)
        if cover >= min_cover:
            hits.append((worst[0], cover, gb, rgb, worst[1], worst[2], worst[3]))
    if skipped is not None and len(hits) > 2:
        skipped["w20_capped"] += len(hits) - 2
    for ratio, cover, gb, rgb, fill, bel, reason in sorted(hits, key=lambda h: h[0])[:2]:
        loc = shape_loc(gb.sp, bbox=[gb.x0, gb.y0, gb.x1 - gb.x0, gb.y1 - gb.y0],
                        cell=gb.cell, paragraph=gb.para, field=gb.field) or {}
        data = {"contrast_ratio": round(ratio, 2),
                "bg_hex": fill,
                "fg_hex": "%02X%02X%02X" % rgb,
                "confidence": "estimate",
                "evidence_source": "xml_colors_geometry"}
        if reason == "w20":
            args = (cover * 100, ratio)
            data["covered_pct"] = round(cover * 100, 1)
        else:
            args = (ratio,)
            data["exposed_pct"] = round(cover * 100, 1)
        warns.append(Finding(si, "W20", reason, args,
                             "bg=#%s fg=%s text=%r"
                             % (fill, data["fg_hex"], (gb.rep or "")[:24]),
                             data=data, loc=loc or None))


_SVG_NS = "{http://www.w3.org/2000/svg}"
_SVGBLIP_NS = "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}"
_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
# Mean advance width as a fraction of font size, for the em-box estimate of an SVG
# text run. The deck engine's glyph model uses per-script tables; inside an SVG we
# know only the string and the size, so one factor covers the mixed case.
_SVG_CHAR_W = 0.6


def _svg_style(el):
    out = {}
    for part in (el.get("style") or "").split(";"):
        if ":" in part:
            k, v = part.split(":", 1)
            out[k.strip()] = v.strip()
    return out


def _svg_hex(value):
    """'#rrggbb' -> RGB tuple; None for no paint; _BG_UNKNOWN for url()/named paint."""
    if not value or value == "none":
        return None
    v = value.strip()
    if v.startswith("#") and len(v) == 7:
        try:
            return (int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16))
        except ValueError:
            return _BG_UNKNOWN
    if v.startswith("#") and len(v) == 4:
        try:
            return tuple(int(c * 2, 16) for c in v[1:])
        except ValueError:
            return _BG_UNKNOWN
    return _BG_UNKNOWN


def _svg_path_bbox(d):
    """Bbox of an M/L/z-only path (matplotlib emits rectangles this way). A path with
    curves (C/Q/A/S/T) returns None: it is a line or a marker, not the kind of flat
    panel that buries a caption, and pretending to know its footprint would be a
    guess."""
    tokens = (d or "").replace(",", " ").split()
    xs, ys = [], []
    i = 0
    cur = None
    while i < len(tokens):
        t = tokens[i]
        if t in ("M", "L"):
            cur = t
            i += 1
            continue
        if t in ("z", "Z"):
            i += 1
            continue
        if t in ("C", "Q", "A", "S", "T", "H", "V", "m", "l", "c", "q", "a",
                 "s", "t", "h", "v"):
            return None
        if cur in ("M", "L"):
            try:
                xs.append(float(t))
                ys.append(float(tokens[i + 1]))
            except (ValueError, IndexError):
                return None
            i += 2
            continue
        return None
    if not xs:
        return None
    return (min(xs), min(ys), max(xs), max(ys))


def _svg_parts_of(slide):
    """(shape, svg-bytes) for every picture whose blip carries an svgBlip extension."""
    out = []
    for sp in slide.shapes:
        if not _is_pic(sp):
            continue
        try:
            blip = sp._element.blipFill.find(NS + "blip")
            if blip is None:
                continue
            svg = blip.find(NS + "extLst/" + NS + "ext/" + _SVGBLIP_NS + "svgBlip")
            if svg is None:
                continue
            rid = svg.get(_R_NS + "embed")
            if not rid:
                continue
            out.append((sp, slide.part.rels[rid].target_part.blob))
        except Exception:
            continue
    return out


def svg_buried_text_check(slide, si, warns, skipped=None, min_cover=0.5):
    """W20 inside a vector picture (reason w20_svg): text buried on a shape in an SVG.

    A chart exported as a PNG erases its own text; the same chart carried as an
    svgBlip keeps every string and fill as XML, so the buried-caption judgment that
    W20 makes on the slide can be made inside the picture too. Scope is the SVG that
    tools actually emit (measured against matplotlib output before writing this):
    solid style/attribute fills in hex, axis-aligned M/L/z rectangle paths and rect
    elements, unrotated text with x/y/font-size. Paint order is document order, so
    only shapes before the text are under it. A shape whose paint is not a solid hex
    (gradients, url() references) abstains as w20_fill_unknown when it covers the
    text, and a path with curves is ignored as a line rather than a panel. Text
    outlined to paths (svg.fonttype="path", the matplotlib default) contains no text
    elements at all, and this gate stays honestly silent -- preserving text is the
    exporter's side of the contract."""
    try:
        from lxml import etree
    except ImportError:
        return
    for sp, blob in _svg_parts_of(slide):
        try:
            root = etree.fromstring(blob)
        except Exception:
            continue
        shapes = []      # (x0, y0, x1, y1, rgb-or-unknown) in document order
        texts = []       # (x0, y0, x1, y1, rgb, rep, order)
        order = 0
        for el in root.iter():
            tag = el.tag if isinstance(el.tag, str) else ""
            order += 1
            if tag == _SVG_NS + "path" or tag == _SVG_NS + "rect":
                style = _svg_style(el)
                fill = _svg_hex(style.get("fill") or el.get("fill"))
                if fill is None:
                    continue
                if tag == _SVG_NS + "rect":
                    try:
                        x, y = float(el.get("x", 0)), float(el.get("y", 0))
                        w, h = float(el.get("width")), float(el.get("height"))
                    except (TypeError, ValueError):
                        continue
                    box = (x, y, x + w, y + h)
                else:
                    box = _svg_path_bbox(el.get("d"))
                if box is None or box[2] <= box[0] or box[3] <= box[1]:
                    continue
                shapes.append((box[0], box[1], box[2], box[3], fill, order))
            elif tag == _SVG_NS + "text":
                style = _svg_style(el)
                fill = _svg_hex(style.get("fill") or el.get("fill"))
                if fill is None or fill is _BG_UNKNOWN:
                    continue
                rep = "".join(el.itertext()).strip()
                if not rep:
                    continue
                tr = el.get("transform") or ""
                if "rotate" in tr:
                    ang = tr.split("rotate(", 1)[1].split()[0].lstrip("(")
                    try:
                        if abs(float(ang)) > 0.5:
                            continue        # rotated text: out of scope, same as W15-17
                    except ValueError:
                        continue
                try:
                    x, y = float(el.get("x", 0)), float(el.get("y", 0))
                except (TypeError, ValueError):
                    continue
                size = 12.0
                fs = (style.get("font-size") or el.get("font-size") or "").rstrip("px")
                try:
                    size = float(fs)
                except ValueError:
                    pass
                w = len(rep) * size * _SVG_CHAR_W
                anchor_ = style.get("text-anchor") or el.get("text-anchor") or "start"
                x0 = x - w / 2 if anchor_ == "middle" else (x - w if anchor_ == "end" else x)
                texts.append((x0, y - 0.75 * size, x0 + w, y + 0.25 * size,
                              fill, rep, order))
        hits = []
        for (tx0, ty0, tx1, ty1, rgb, rep, torder) in texts:
            area = (tx1 - tx0) * (ty1 - ty0)
            if area <= 0:
                continue
            # SVG paint order is document order, and the visible background of a
            # glyph is the LAST shape painted beneath it, so layers are consumed
            # from the top down -- the same greedy the slide-side check uses, and
            # for the same reason: judging against every lower layer independently
            # flagged a readable label on a bright bar for the page behind the bar.
            layers = []
            for (sx0, sy0, sx1, sy1, fill, sorder) in shapes:
                if sorder >= torder:
                    continue
                ix = min(tx1, sx1) - max(tx0, sx0)
                iy = min(ty1, sy1) - max(ty0, sy0)
                if ix <= 0 or iy <= 0:
                    continue
                layers.append((sorder, ix * iy, fill))
            layers.sort(key=lambda l: l[0], reverse=True)
            remaining = area
            buried = 0.0
            unknown = 0.0
            worst = None
            for (_o, ov, fill) in layers:
                take = min(ov, remaining)
                if take <= 0:
                    continue
                if fill is _BG_UNKNOWN:
                    unknown += take
                else:
                    lt, lb = _luma(rgb), _luma(fill)
                    ratio = (max(lt, lb) + 0.05) / (min(lt, lb) + 0.05)
                    if ratio < 2.0:
                        buried += take
                        if worst is None or ratio < worst[0]:
                            worst = (ratio, fill)
                remaining -= take
                if remaining <= 0:
                    break
            if unknown / area >= min_cover:
                if skipped is not None:
                    skipped["w20_fill_unknown"] += 1
                continue
            if worst is None:
                continue
            cover = min(buried / area, 1.0)
            if cover >= min_cover:
                hits.append((worst[0], cover, rep, rgb, worst[1], sp))
        if skipped is not None and len(hits) > 2:
            skipped["w20_capped"] += len(hits) - 2
        for ratio, cover, rep, rgb, fill, sp_ in sorted(hits, key=lambda h: h[0])[:2]:
            warns.append(Finding(
                si, "W20", "w20_svg", (cover * 100, ratio),
                "bg=#%02X%02X%02X fg=%02X%02X%02X text=%r"
                % (fill[0], fill[1], fill[2], rgb[0], rgb[1], rgb[2], rep[:24]),
                data={"contrast_ratio": round(ratio, 2),
                      "covered_pct": round(cover * 100, 1),
                      "bg_hex": "%02X%02X%02X" % fill,
                      "fg_hex": "%02X%02X%02X" % rgb,
                      "confidence": "estimate",
                      "evidence_source": "svg_vector"},
                loc=shape_loc(sp_) or None))


def _paragraph_rgb(gb, slide, styler, thm_colors, skipped):
    """The color of the paragraph a glyph box stands for, or None when it is not knowable.

    A box is a paragraph, so it can hold runs of different colors. The darkest-against-its-
    backdrop case is what a reader notices, but mixing colors inside one verdict would make
    the report lie about which run is buried, so a paragraph whose runs disagree is skipped
    rather than guessed. Undecodable colors surface as an abstention, matching W19."""
    tf = getattr(gb.sp, "text_frame", None)
    if tf is None:
        return None
    try:
        paras = list(tf.paragraphs)
    except Exception:
        return None
    if gb.para is None or gb.para >= len(paras):
        return None
    para = paras[gb.para]
    seen = set()
    for run_like, _rix, _is_fld in iter_inline_items(para):
        if run_like is None or not (run_like.text or "").strip():
            continue
        try:
            rgb = _resolve_run_rgb(run_like, para, tf, gb.sp, slide, styler, thm_colors)
        except Exception:
            return None
        if rgb is None:
            continue
        if rgb is _COLOR_UNKNOWN:
            if skipped is not None:
                skipped["w20_color_unknown"] += 1
            return None
        seen.add(rgb)
    if len(seen) != 1:
        return None
    return seen.pop()


# W21 thresholds, from the 148-deck corpus sweep (archforge corpus + a private deck set).
# At these values 95% of decks report nothing and the worst reports two.
_W21_DIST = 16.0        # redmean distance below which two colors read as the same one
_W21_RATIO = 6.0        # dominant uses / stray uses
_W21_NEIGHBORS = 2      # more near neighbours than this reads as a deliberate ramp
_W21_MIN_USES = 2       # a color painted once is not yet a pattern


def _redmean(a, b):
    """Cheap perceptual distance between two RGB triples (0 = identical). Closer to the eye
    than plain sRGB Euclid and needs no LAB conversion, which matters because this runs over
    every pair of colors in the deck."""
    rm = (a[0] + b[0]) / 2.0
    dr, dg, db = a[0] - b[0], a[1] - b[1], a[2] - b[2]
    return math.sqrt((2 + rm / 256.0) * dr * dr + 4.0 * dg * dg
                     + (2 + (255 - rm) / 256.0) * db * db)


def collect_palette(slide, palette, styler=None, thm_colors=None):
    """Tallies the colors this slide actually paints into a deck-wide Counter.

    Fill and line come from the XML and run color from the same resolution chain W19 uses,
    so W21 judges what gets painted rather than what the theme declares. Reading
    run.font.color instead would materialize an empty solidFill on the run and change the
    file being linted."""
    for sp in iter_shapes(slide.shapes):
        for hexc in (_shape_fill_hex(sp), _shape_line_hex(sp)):
            rgb = _hex_rgb(hexc)
            if rgb:
                palette[rgb] += 1
    for tf, _w, sp, _rc, _xf in _collect_frames(slide.shapes):
        for para in tf.paragraphs:
            for run in para.runs:
                if not (run.text or "").strip():
                    continue
                rgb = _resolve_run_rgb(run, para, tf, sp, slide, styler=styler,
                                       thm_colors=thm_colors)
                if isinstance(rgb, tuple):
                    palette[rgb] += 1


def palette_drift_check(palette, warns, dist=_W21_DIST, ratio=_W21_RATIO,
                        neighbors=_W21_NEIGHBORS, min_uses=_W21_MIN_USES):
    """W21: a rare color sitting right beside a dominant near-identical one, which is what a
    mistyped hex looks like from the outside.

    Counting distinct colors would only measure genre, since an infographic legitimately
    carries many and a minimal deck few. So this measures the shape of the mistake instead:
    #8C8C8C painted 116 times next to #888888 painted 3 times is not a decision.

    The first corpus sweep flagged deliberate gradients, where a ramp of stacked rectangles
    yields many near-identical colors with equal use counts (one deck: 24 pairs, every pair
    26-vs-26 or 17-vs-17). Two guards drop that class. The pair must be lopsided, and
    neither member may carry more than `neighbors` near neighbours. Deck-level like W13,
    because a palette is a property of the deck rather than of a page."""
    keys = [c for c, n in palette.items() if n >= min_uses]
    near = {c: [] for c in keys}
    for i, a in enumerate(keys):
        for b in keys[i + 1:]:
            d = _redmean(a, b)
            if 0 < d <= dist:
                near[a].append((d, b))
                near[b].append((d, a))

    for stray, adj in near.items():
        if not adj or len(adj) > neighbors:
            continue
        d, dom = min(adj)
        if len(near[dom]) > neighbors:
            continue
        n_stray, n_dom = palette[stray], palette[dom]
        if n_stray >= n_dom or n_dom < n_stray * ratio:
            continue
        warns.append(Finding(0, "W21", "w21",
                             ("%02X%02X%02X" % stray, n_stray,
                              "%02X%02X%02X" % dom, n_dom),
                             "dC %.1f | %.0fx" % (d, n_dom / float(n_stray))))
