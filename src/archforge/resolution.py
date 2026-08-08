# -*- coding: utf-8 -*-
"""Style resolution: the theme font and colour maps per master, the lstStyle inheritance
chain, and the effective size of a run once every default has been applied.

Extracted from lint.py for the 0.7 decomposition (#5). This is the layer the CALIBRATION
log describes as measured rather than specified: the priority it implements (run rPr >
paragraph pPr/defRPr > lstStyle chain > non-empty theme a:ea > run a:latin > OS fallback)
came out of probe decks rendered through PowerPoint COM, not out of the OOXML spec.

Pure functions and one resolver object over lxml elements; no Finding, no I/O, no CLI.
Re-exported from lint for backward compatibility.
"""
from typing import Dict, Optional, Tuple

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

try:
    from .ooxml import NS, NS_P
except ImportError:   # standalone execution
    from ooxml import NS, NS_P

def _theme_fonts_from_blob(blob: bytes) -> Optional[Dict[str, str]]:
    """The 4 major/minor font slots from the theme XML: {"mn-ea","mn-lt","mj-ea","mj-lt"}.
    Because this uses XML parsing, it is unaffected by quote serialization or attribute order
    (fixes a vulnerability from the earlier byte-regex era, external review 2026-07-10). Also
    used to resolve theme tokens like "+mn-lt" in run rPr (confirmed in the adversarial panel:
    treating the token as a literal font name causes an E1 false negative). None on parse
    failure."""
    try:
        from lxml import etree
        root = etree.fromstring(blob)
        out = {}
        for prefix, tag in (("mn", "minorFont"), ("mj", "majorFont")):
            base = root.find(".//" + NS + "fontScheme/" + NS + tag)
            if base is None:
                continue
            for suffix, slot in (("ea", "ea"), ("lt", "latin")):
                el = base.find(NS + slot)
                if el is not None:
                    out["%s-%s" % (prefix, suffix)] = el.get("typeface") or ""
        return out
    except Exception:
        return None

def theme_fonts_by_master(prs) -> Dict[str, Optional[Dict[str, str]]]:
    """A map of theme font slots per slide master. Key = master partname string, value = None
    on parse failure. Fixes an issue where, in a multi-master deck, grabbing the first theme
    from iter_parts falsely fired E1 based on an empty slot from an unrelated master: this is
    resolved via the master-to-theme relationship (rels) instead (external review,
    2026-07-10)."""
    out = {}
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        for master in prs.slide_masters:
            try:
                theme_part = master.part.part_related_by(RT.THEME)
                out[str(master.part.partname)] = _theme_fonts_from_blob(theme_part.blob)
            except Exception:
                out[str(master.part.partname)] = None
    except Exception:
        pass
    return out

def theme_ea_by_master(prs) -> Dict[str, Optional[str]]:
    """Backward compatibility: a map with only the minorFont a:ea per master (value None =
    parse failure, "" = empty slot)."""
    return {k: (v.get("mn-ea", "") if v is not None else None)
            for k, v in theme_fonts_by_master(prs).items()}

def theme_ea_font(prs) -> Optional[str]:
    """Backward-compatible entry point: the theme a:ea of the first master (relationship-based
    resolution).
    Returns: font name / "" (empty slot = Malgun fallback on Windows) / None (theme parse
    failure)."""
    ea_map = theme_ea_by_master(prs)
    for v in ea_map.values():
        if v is not None:
            return v
    return None

def _sz_from_defrpr(d) -> Optional[float]:
    """Converts defRPr@sz (1/100pt integer) to pt. None if absent or garbage."""
    if d is None:
        return None
    sz = d.get("sz")
    if sz is None:
        return None
    try:
        return int(sz) / 100.0
    except ValueError:
        return None

def _lst_defrpr(lst_el, lvl: int):
    """The defRPr element for the given level from an lstStyle-type container (a:lvlXpPr
    children)."""
    if lst_el is None:
        return None
    lvl = min(max(int(lvl or 0), 0), 8)
    p = lst_el.find(NS + "lvl%dpPr" % (lvl + 1))
    if p is None:
        return None
    return p.find(NS + "defRPr")

def _lst_sz_pt(lst_el, lvl: int) -> Optional[float]:
    """Backward-compatible shim: the defRPr sz (pt) at the given level in lstStyle."""
    return _sz_from_defrpr(_lst_defrpr(lst_el, lvl))

class StyleResolver:
    """When a run or paragraph has no explicit attribute, resolves the effective style via the
    OOXML inheritance chain. Size and font (a:ea/a:latin) are resolved through the same
    chain, but each attribute is searched independently.

    Chain (ECMA-376 text style hierarchy): shape txBody lstStyle -> (if a placeholder) the
    layout's same-idx placeholder lstStyle -> master placeholder lstStyle -> master
    txStyles (title/body/other) -> presentation defaultTextStyle.

    Size: None (= W5) if absent everywhere. 0.1.0 only looked at run/paragraph and let
    everything fall through to W5, which effectively killed E3/W1/W8 in placeholder-based
    decks (external review, 2026-07-10).
    Font: confirmed by measured-by-render COM probing that the master lstStyle's a:ea is
    actually inherited into rendering (2026-07-10, probe 6, docs/CALIBRATION.md). 0.2.0 only
    looked at run rPr and produced a "confirmed Malgun fallback" false positive on standard
    corporate templates (confirmed in the second external re-check)."""

    def __init__(self, prs):
        self._prs = prs
        self._default_el = None
        self._default_loaded = False
        # An inheritance-dependent deck has thousands of runs repeatedly querying the same
        # layout/master: memoize the defRPr element keyed by partname (fixes an asymmetry the
        # adversarial panel flagged, where every run re-traversed the chain)
        self._default_cache: Dict[int, object] = {}
        self._layout_cache: Dict[Tuple[str, int, int], object] = {}
        self._master_ph_cache: Dict[Tuple[str, str, int], object] = {}
        self._master_tx_cache: Dict[Tuple[str, str, int], object] = {}

    def _default_defrpr(self, lvl: int):
        if not self._default_loaded:
            self._default_loaded = True
            try:
                from lxml import etree
                root = etree.fromstring(self._prs.part.blob)
                self._default_el = root.find(NS_P + "defaultTextStyle")
            except Exception:
                self._default_el = None
        if lvl not in self._default_cache:
            self._default_cache[lvl] = _lst_defrpr(self._default_el, lvl)
        return self._default_cache[lvl]

    @staticmethod
    def _ph_family(ph_type) -> str:
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER as PH
            if ph_type in (PH.TITLE, PH.CENTER_TITLE, PH.VERTICAL_TITLE):
                return "title"
            if ph_type in (PH.BODY, PH.SUBTITLE, PH.VERTICAL_BODY, PH.OBJECT):
                return "body"
        except Exception:
            pass
        return "other"

    @staticmethod
    def _ph_lst(shape_like):
        try:
            txBody = shape_like.text_frame._txBody
            return txBody.find(NS + "lstStyle")
        except Exception:
            return None

    def ph_family_of(self, sp) -> Optional[str]:
        """title/body/other if a placeholder, otherwise None. Used for E1's majorFont branch
        (the title family uses the theme majorFont ea: measured in probe 6, Q1)."""
        try:
            if getattr(sp, "is_placeholder", False):
                return self._ph_family(sp.placeholder_format.type)
        except Exception:
            pass
        return None

    def _layout_ph_defrpr(self, slide, idx: int, lvl: int):
        # The guard is per placeholder: wrapping the whole loop in one try lets a single
        # corrupted placeholder abort the entire search, missing the real style at a later
        # index (confirmed in the adversarial panel, 2026-07-10)
        try:
            layout = slide.slide_layout
            key = (str(layout.part.partname), idx, lvl)
        except Exception:
            return None
        if key in self._layout_cache:
            return self._layout_cache[key]
        found = None
        try:
            phs = list(layout.placeholders)
        except Exception:
            phs = []
        for ph in phs:
            try:
                if ph.placeholder_format.idx == idx:
                    found = _lst_defrpr(self._ph_lst(ph), lvl)
                    break
            except Exception:
                continue
        self._layout_cache[key] = found
        return found

    def _master_ph_defrpr(self, slide, family: str, lvl: int):
        try:
            master = slide.slide_layout.slide_master
            key = (str(master.part.partname), family, lvl)
        except Exception:
            return None
        if key in self._master_ph_cache:
            return self._master_ph_cache[key]
        found = None
        try:
            phs = list(master.placeholders)
        except Exception:
            phs = []
        for ph in phs:
            try:
                if self._ph_family(ph.placeholder_format.type) == family:
                    found = _lst_defrpr(self._ph_lst(ph), lvl)
                    break
            except Exception:
                continue
        self._master_ph_cache[key] = found
        return found

    def _master_tx_defrpr(self, slide, family: str, lvl: int):
        try:
            master = slide.slide_layout.slide_master
            key = (str(master.part.partname), family, lvl)
        except Exception:
            return None
        if key in self._master_tx_cache:
            return self._master_tx_cache[key]
        found = None
        try:
            tx = master.element.find(NS_P + "txStyles")
            if tx is not None:
                tag = {"title": "titleStyle", "body": "bodyStyle"}.get(family, "otherStyle")
                found = _lst_defrpr(tx.find(NS_P + tag), lvl)
        except Exception:
            found = None
        self._master_tx_cache[key] = found
        return found

    def _chain(self, tf, sp, slide, lvl: int):
        """Iterates the inheritance chain as (defRPr element, source). Nodes with no
        element are skipped."""
        try:
            own = _lst_defrpr(tf._txBody.find(NS + "lstStyle"), lvl)
            if own is not None:
                yield own, "own"
        except Exception:
            pass
        idx = ph_type = None
        try:
            if getattr(sp, "is_placeholder", False):
                pf = sp.placeholder_format
                idx, ph_type = pf.idx, pf.type
        except Exception:
            idx = ph_type = None
        if idx is not None:
            el = self._layout_ph_defrpr(slide, idx, lvl)
            if el is not None:
                yield el, "layout"
            family = self._ph_family(ph_type)
            el = self._master_ph_defrpr(slide, family, lvl)
            if el is not None:
                yield el, "master_ph"
            el = self._master_tx_defrpr(slide, family, lvl)
            if el is not None:
                yield el, "master_tx"
        el = self._default_defrpr(lvl)
        if el is not None:
            yield el, "default"

    def resolve_size(self, tf, sp, slide, lvl: int) -> Tuple[Optional[float], Optional[str]]:
        """(effective size in pt, source) or (None, None). The source distinction exists to
        prevent the title collector from flooding: an 18pt defaultTextStyle fallback is valid
        for gating purposes but is not "the title size this deck intended," so it is excluded
        from title candidacy (fixes a regression confirmed in the second external
        re-check)."""
        for el, src in self._chain(tf, sp, slide, lvl):
            v = _sz_from_defrpr(el)
            if v is not None:
                return v, src
        return None, None

    def resolve_font(self, tf, sp, slide, lvl: int, slot: str) -> Optional[str]:
        """The a:{slot} (ea/latin) typeface of defRPr from the inheritance chain. May be a
        theme token."""
        for el, _src in self._chain(tf, sp, slide, lvl):
            f = el.find(NS + slot)
            if f is not None:
                name = f.get("typeface")
                if name:
                    return name
        return None

    def resolve(self, tf, sp, slide, lvl: int) -> Optional[float]:
        """Backward-compatible entry point: effective size (pt) only. tf = the text frame
        being checked (including table cells), sp = the owning shape."""
        return self.resolve_size(tf, sp, slide, lvl)[0]


SizeResolver = StyleResolver   # backward-compatible alias (0.2.0 public name)
