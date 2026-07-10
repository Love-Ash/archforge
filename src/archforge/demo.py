# -*- coding: utf-8 -*-
"""Demo deck generator (0.5.0): an onboarding asset that shows what the linter catches
within 30 seconds of install.

`archforge demo` uses this module to build broken.pptx (representative defects
deliberately seeded) and fixed.pptx (the same content, corrected; clean under the full
profile), then lints both on the spot. The committed copies under the repo's examples/
are also generated from this module by scripts/make_examples.py.

Deck text follows the same locale as the report language (lang="ko"/"en"), and each
variant is monolingual (user decision): the Korean deck seeds all 6 defect axes, while
the English deck seeds the 4 that exist without Hangul (E1 font fallback and E4 Hangul
tracking physically require Hangul text, so they are demonstrated by the Korean deck
and by examples/broken.pptx only).

The seeded defects cover the axes most common in real generated decks:
- E1 silent Hangul font fallback (Hangul carried in Arial's a:latin, empty theme ea slot; ko only)
- E2 an em dash used as punctuation (the #1 tell of an AI-generated deck; full profile)
- E3 unreadable size (4pt source attribution)
- E4 positive tracking on consecutive Hangul (ko only)
- W15 text frame overlap / W16 off-canvas overflow
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Do not keep the banned character as a literal in the source (same convention as
# tests): it exists only inside the pptx under test
EM_DASH = chr(0x2014)

# Per-language deck text. Each variant is monolingual (user decision, 2026-07-10):
# the English README/demo must contain English only, so the en deck drops the
# Hangul-dependent defects (E1/E4) instead of mixing scripts.
_T = {
    "ko": {
        "title1": "3분기 실적 요약",
        "e1": "3분기 실적 요약",   # the title itself is E1 (Hangul in Arial)
        "e2": "핵심 지표는 전 분기 대비 개선" + EM_DASH + "특히 구독 매출이 견인했습니다",
        "e2_fixed": "핵심 지표는 전 분기 대비 개선: 특히 구독 매출이 견인했습니다",
        "e4": "자간이 벌어진 한글 강조",
        "e4_fixed": "자간을 되돌린 한글 강조",
        "e3": "출처: 사내 관리회계, 2026-06",
        "title2": "분기 핵심 지표",
        "kpi1": "매출 성장률 +18%",
        "kpi2": "영업이익률 12.4%",
        "w16": "다음 분기 가이던스",
    },
    "en": {
        "title1": "Q3 Results Summary",
        "e2": "Key metrics improved quarter over quarter" + EM_DASH + "subscriptions led the growth",
        "e2_fixed": "Key metrics improved quarter over quarter: subscriptions led the growth",
        "body": "Subscriptions now account for 41% of total revenue",
        "e3": "Source: internal management accounts, 2026-06",
        "title2": "Key metrics this quarter",
        "kpi1": "Revenue growth +18%",
        "kpi2": "Operating margin 12.4%",
        "w16": "Next quarter guidance",
    },
}


def _prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def _slide(p):
    return p.slides.add_slide(p.slide_layouts[6])


def _tb(slide, x, y, w, h, text, size=14, font=None, ea=None, spc=None,
        color="1A1A1A", bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    if font:
        r.font.name = font
    if ea:
        rPr = r._r.get_or_add_rPr()
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": ea}))
    if spc is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(spc))
    return box


def build_broken(path, lang="ko"):
    """Deck seeded with representative defects. Under the full profile the ko variant
    yields ERROR 4 (E1/E2/E3/E4) + WARN 2 (W15/W16); the en variant, which is
    monolingual English and therefore cannot carry the Hangul-only defects, yields
    ERROR 2 (E2/E3) + WARN 2."""
    t = _T[lang]
    p = _prs()
    s = _slide(p)   # p1: typography defects
    if lang == "ko":
        # The title itself is E1 (Hangul carried in Arial). The remaining runs specify ea
        # to keep one code per defect (the default Office theme has an empty ea slot:
        # measured confirmation that every Hangul run without ea fires E1 correctly).
        _tb(s, 0.8, 0.6, 8.0, 0.9, t["title1"], size=26, font="Arial", bold=True)   # E1
        _tb(s, 0.8, 2.8, 6.0, 0.5, t["e4"], size=16, spc=300, ea="맑은 고딕")        # E4
    else:
        _tb(s, 0.8, 0.6, 8.0, 0.9, t["title1"], size=26, font="Arial", bold=True)
        _tb(s, 0.8, 2.8, 8.0, 0.5, t["body"], size=16)   # neutral filler, no defect
    _tb(s, 0.8, 1.8, 10.0, 0.6, t["e2"], size=14, ea="맑은 고딕")                    # E2
    _tb(s, 0.8, 6.9, 5.0, 0.3, t["e3"], size=4, ea="맑은 고딕")                      # E3
    s = _slide(p)   # p2: geometry defects
    _tb(s, 0.8, 0.6, 8.0, 0.9, t["title2"], size=26, ea="맑은 고딕", bold=True)
    _tb(s, 1.0, 2.4, 5.0, 1.0, t["kpi1"], size=24, ea="맑은 고딕")                   # W15 pair
    _tb(s, 1.2, 2.5, 5.0, 1.0, t["kpi2"], size=24, ea="맑은 고딕")
    _tb(s, 12.0, 4.5, 3.0, 0.6, t["w16"], size=18, ea="맑은 고딕")                   # W16
    p.save(path)
    return path


def build_fixed(path, lang="ko"):
    """The corrected version with the same content as broken. ERROR 0, WARN 0 under the
    full profile."""
    t = _T[lang]
    p = _prs()
    s = _slide(p)
    # E1 fix: leave a:latin as-is and specify a Hangul font on a:ea (the most robust fix)
    _tb(s, 0.8, 0.6, 8.0, 0.9, t["title1"], size=26, font="Arial",
        ea="맑은 고딕", bold=True)
    if lang == "ko":
        # E4 fix: zero tracking on the Hangul run
        _tb(s, 0.8, 2.8, 6.0, 0.5, t["e4_fixed"], size=16, ea="맑은 고딕")
    else:
        _tb(s, 0.8, 2.8, 8.0, 0.5, t["body"], size=16)
    # E2 fix: replace prose dashes with a colon, comma, or parentheses
    _tb(s, 0.8, 1.8, 10.0, 0.6, t["e2_fixed"], size=14, ea="맑은 고딕")
    # E3 fix: sources/captions at least 9pt too
    _tb(s, 0.8, 6.9, 5.0, 0.3, t["e3"], size=9, ea="맑은 고딕")
    s = _slide(p)
    _tb(s, 0.8, 0.6, 8.0, 0.9, t["title2"], size=26, ea="맑은 고딕", bold=True)
    _tb(s, 1.0, 2.4, 5.0, 1.0, t["kpi1"], size=24, ea="맑은 고딕")
    _tb(s, 6.8, 2.4, 5.0, 1.0, t["kpi2"], size=24, ea="맑은 고딕")
    _tb(s, 9.2, 4.5, 3.4, 0.6, t["w16"], size=18, ea="맑은 고딕")
    p.save(path)
    return path


def build_warnings(path):
    """A deck (for examples/) that is clean under the core profile and only produces
    style WARNs under the full profile: W13 native shadows, W14 nominal-phrase titles
    listed in a row. A teaching aid for the profile split. Since W14 is a Hangul-title
    heuristic, this deck has only a Korean-language version."""
    p = _prs()
    for i, title in enumerate(("시장 현황", "경쟁 구도 분석", "성장 전략 로드맵")):
        s = _slide(p)
        _tb(s, 0.8, 0.6, 8.0, 0.9, title, size=24, ea="맑은 고딕", bold=True)
        _tb(s, 0.8, 2.0, 10.0, 0.6,
            "본문 요지는 페이지마다 다르게 구성했습니다 (%d)" % (i + 1),
            size=14, ea="맑은 고딕")
        from pptx.enum.shapes import MSO_SHAPE
        # W13 is counted when a page has 2+ effective effects (designed to avoid false
        # positives from an empty effectLst)
        for j in range(2):
            sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.8 + i * 0.7 + j * 3.2), Inches(4.0 + i * 0.4),
                                    Inches(2.5), Inches(1.0))
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor.from_string("DDE3EA")
            sp.line.fill.background()
            spPr = sp._element.spPr
            eff = spPr.makeelement(qn("a:effectLst"), {})
            sh = spPr.makeelement(qn("a:outerShdw"),
                                  {"blurRad": "40000", "dist": "20000"})
            eff.append(sh)
            spPr.append(eff)
    p.save(path)
    return path
