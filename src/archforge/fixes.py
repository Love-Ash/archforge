# -*- coding: utf-8 -*-
"""Deterministic auto-fixes for the three mechanically safe rules (0.8.x, external
strategy review: find-and-fix demos better than find-only, and these three need no
layout judgment).

  E4  remove positive tracking from Hangul runs (pure formatting)
  E1  set a Hangul-capable font on the run's a:ea slot (pure formatting; the robust
      fix the E1 guidance already teaches)
  E2  replace a dash used as sentence punctuation with ", " (text edit, but exactly
      the substitution the rule's own fix guidance prescribes; range dashes that the
      rule exempts are never touched)

Everything else (E3 sizes, W15 collisions, W16 overflow...) needs design judgment and
stays agent territory: the JSON location payload is the contract for that. This module
never touches a run the corresponding detector would not flag."""
import os
from typing import Dict, List, Optional, Set

try:
    from .fonts import NS, e1_violation, run_fonts, para_fonts, run_track, resolve_font_tokens
    from .dashes import LONG_DASHES, dash_violations
    from .scripts import is_hangul, is_kana, is_hanja
except ImportError:   # pragma: no cover
    from fonts import NS, e1_violation, run_fonts, para_fonts, run_track, resolve_font_tokens
    from dashes import LONG_DASHES, dash_violations
    from scripts import is_hangul, is_kana, is_hanja

FIXABLE = ("E1", "E2", "E4")
DEFAULT_EA = "맑은 고딕"


def _fix_e4(run) -> bool:
    rPr = run._r.find(NS + "rPr")
    if rPr is not None and rPr.get("spc") is not None:
        del rPr.attrib["spc"]
        return True
    return False


def _fix_e1(run, ea_font: str) -> bool:
    rPr = run._r.find(NS + "rPr")
    if rPr is None:
        rPr = run._r.makeelement(NS + "rPr", {})
        run._r.insert(0, rPr)
    ea = rPr.find(NS + "ea")
    if ea is None:
        ea = rPr.makeelement(NS + "ea", {})
        rPr.append(ea)
    ea.set("typeface", ea_font)
    return True


def _fix_e2_text(ptext: str, lo: int, hi: int) -> Optional[str]:
    """The run's replacement text: each dash the detector flags (paragraph context, run
    span) becomes ", " with adjacent spaces collapsed. None when nothing changes."""
    flagged = set()
    for i in range(lo, min(hi, len(ptext))):
        if ptext[i] in LONG_DASHES and dash_violations(ptext, span=(i, i + 1)):
            flagged.add(i)
    if not flagged:
        return None
    out = []
    i = lo
    while i < hi:
        if i in flagged:
            while out and out[-1] == " ":
                out.pop()
            out.append(", ")
            i += 1
            while i < hi and ptext[i] == " ":
                i += 1
            continue
        out.append(ptext[i])
        i += 1
    return "".join(out)


def apply_fixes(in_path: str, out_path: str, rules: Optional[Set[str]] = None,
                ea_font: str = DEFAULT_EA) -> List[Dict]:
    """Returns the change list [{page, code, detail}]. Walks the same frames the
    detectors walk (tables and groups included) and mirrors each detector's firing
    condition before touching anything."""
    from pptx import Presentation
    try:
        from .lint import collect_frames, theme_ea_by_master, theme_fonts_by_master
    except ImportError:   # pragma: no cover
        from lint import collect_frames, theme_ea_by_master, theme_fonts_by_master

    rules = set(rules or FIXABLE)
    prs = Presentation(in_path)
    thm_ea_map = theme_ea_by_master(prs)
    thm_fonts_map = theme_fonts_by_master(prs)
    thm_ea_default = next((v for v in thm_ea_map.values() if v is not None), None)
    thm_fonts_default = next((v for v in thm_fonts_map.values() if v is not None), None)
    changes = []
    for si, slide in enumerate(prs.slides, 1):
        try:
            key = str(slide.slide_layout.slide_master.part.partname)
            thm_ea = thm_ea_map.get(key, thm_ea_default)
            thm_fonts = thm_fonts_map.get(key, thm_fonts_default)
        except Exception:
            thm_ea, thm_fonts = thm_ea_default, thm_fonts_default
        for tframe, _w, _sp, _cell, _xf in collect_frames(slide.shapes):
            for para in tframe.paragraphs:
                ptext = "".join(r.text or "" for r in para.runs)
                pfonts = para_fonts(para)
                off = 0
                for run in para.runs:
                    t = run.text or ""
                    lo, hi = off, off + len(t)
                    off = hi
                    if not t:
                        continue
                    if "E4" in rules \
                            and not any(is_kana(c) for c in t) \
                            and any(is_hangul(c) for c in t) \
                            and sum(1 for c in t if is_hangul(c) or is_hanja(c)) >= 2:
                        tr = run_track(run)
                        if tr is not None and tr > 50 and _fix_e4(run):
                            changes.append({"page": si, "code": "E4",
                                            "detail": "removed tracking %d on %r"
                                                      % (tr, t[:20])})
                    if "E1" in rules and any(is_hangul(c) or is_kana(c) or is_hanja(c)
                                             for c in t):
                        fonts = dict(pfonts)
                        fonts.update(run_fonts(run))
                        fonts = resolve_font_tokens(fonts, thm_fonts)
                        script = ("hangul" if any(is_hangul(c) for c in t)
                                  else "cjk_other")
                        if e1_violation(t, fonts, thm_ea, script) is not None \
                                and _fix_e1(run, ea_font):
                            changes.append({"page": si, "code": "E1",
                                            "detail": "set a:ea=%r on %r"
                                                      % (ea_font, t[:20])})
                    if "E2" in rules:
                        new = _fix_e2_text(ptext, lo, hi)
                        if new is not None:
                            run.text = new
                            changes.append({"page": si, "code": "E2",
                                            "detail": "replaced dash punctuation in %r"
                                                      % t[:24]})
    prs.save(out_path)
    return changes
