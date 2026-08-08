# -*- coding: utf-8 -*-
"""Shape traversal and geometry: walking a slide's shape tree, resolving group transforms
into absolute coordinates, and reading the autofit state of a text frame.

Extracted from lint.py for the 0.7 decomposition (#5). Pure functions over python-pptx
objects and lxml elements; no Finding, no I/O, no CLI. Every geometry gate (W15, W16, W17)
and the size gates read the world through here, which is why the group-transform handling
lives in one place: a rule that walks shapes itself is a rule that will disagree with the
others about where a shape is.

Re-exported from lint for backward compatibility.
"""
import math

from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from .ooxml import EMU_PER_IN, NS
except ImportError:   # standalone execution
    from ooxml import EMU_PER_IN, NS

def _pct_attr(v, default):
    """OOXML percentage union type: both '62500' (1/1000 %) and '62.5%' (string form) are
    valid (ST_TextFontScalePercentOrPercentString). int() alone dies with ValueError on the
    latter, and a blanket except swallowed it to 1.0 (measured in adversarial verification,
    2026-07-03)."""
    if not v:
        return default
    v = v.strip()
    if v.endswith("%"):
        return float(v[:-1]) / 100.0
    return int(v) / 100000.0

def frame_autofit(tf):
    """The (fontScale, lnSpcReduction) ratio pair. (1.0, 0.0) if there is no normAutofit.

    Absence is handled explicitly; a malformed value now PROPAGATES instead of silently
    reading as scale 1.0 (0.8.x exception audit: a garbled fontScale used to hide a real
    autofit shrink from E3, a silent false negative). Every in-engine caller sits under
    a page/run guard that converts the raise into a W18 abstention."""
    bodyPr = tf._txBody.find(NS + "bodyPr")
    if bodyPr is not None:
        na = bodyPr.find(NS + "normAutofit")
        if na is not None:
            return (_pct_attr(na.get("fontScale"), 1.0),
                    _pct_attr(na.get("lnSpcReduction"), 0.0))
    return 1.0, 0.0

def frame_font_scale(tf):
    """The text frame's autofit fontScale (ratio). Kept for backward compatibility with
    existing consumers such as E3."""
    return frame_autofit(tf)[0]

def _group_xf(sp, xf):
    """Composes the group shape's off/ext vs chOff/chExt affine with the parent xf.
    Coefficients (ax,bx,ay,by): abs = a*raw + b (EMU). Falls back to the parent xf unchanged
    (identity fallback) if parsing fails.

    grpSpPr is in the p: namespace at the slide level (only the inner xfrm/off elements are
    a:). Fixes a latent bug where the previous code did find(a:grpSpPr), which always
    returned None, causing a silent identity fallback (measured in 0.5.0: found by
    reproducing a case where a moved, desynced group's loc bbox came out in raw
    coordinates)."""
    try:
        gsp = None
        for ch in sp._element:
            if isinstance(ch.tag, str) and ch.tag.endswith("}grpSpPr"):
                gsp = ch
                break
        x = gsp.find(NS + "xfrm")
        off, ext = x.find(NS + "off"), x.find(NS + "ext")
        cho, che = x.find(NS + "chOff"), x.find(NS + "chExt")
        ox, oy = int(off.get("x")), int(off.get("y"))
        ew, eh = int(ext.get("cx")), int(ext.get("cy"))
        cx, cy = int(cho.get("x")), int(cho.get("y"))
        cw_, ch_ = int(che.get("cx")) or ew, int(che.get("cy")) or eh
        sx, sy = ew / float(cw_), eh / float(ch_)
        ax, bx, ay, by = xf
        return (ax * sx, ax * (ox - cx * sx) + bx,
                ay * sy, ay * (oy - cy * sy) + by)
    except Exception:
        # Accepted-soft fallback (0.8.x exception audit): a malformed group transform
        # falls back to the parent transform, so child coordinates are approximately
        # placed instead of the whole page's geometry aborting into W18 over one odd
        # group. The trade-off is documented in docs/EXCEPTION_AUDIT.md.
        return xf

def collect_frames(shapes, xf=(1.0, 0.0, 1.0, 0.0)):
    """A list of (text_frame, width_emu, owner_shape, cell_rc, xf). Recurses into groups and
    includes native table cells. cell_rc is (row, col), 0-based, for table cells, otherwise
    None. xf is the group absolute-coordinate affine (same coefficients as iter_shapes_geo),
    used to turn the loc bbox of a run-level finding into real slide coordinates instead of
    the group's chOff coordinate space (carried over from the fourth review, 0.5.0)."""
    out = []
    for sp in shapes:
        try:
            st = sp.shape_type
        except Exception:
            st = None
        if st == MSO_SHAPE_TYPE.GROUP:
            out += collect_frames(sp.shapes, _group_xf(sp, xf))
            continue
        if getattr(sp, "has_table", False):
            tbl = sp.table
            ncol = len(tbl.columns) or 1
            try:
                col_w = [(c.width or 0) for c in tbl.columns]
            except Exception:
                col_w = []
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    # Merged regions (0.6.0): continuation cells mirror their origin's
                    # text frame, so walking them double-counted the same runs. The
                    # origin's usable width spans the merged columns (real column widths
                    # when available; the old even-split approximation as fallback).
                    try:
                        if cell.is_spanned:
                            continue
                        w_emu = sum(col_w[ci:ci + max(1, cell.span_width)]) \
                            or (sp.width or 0) // ncol
                    except Exception:
                        w_emu = (sp.width or 0) // ncol
                    out.append((cell.text_frame, w_emu, sp, (ri, ci), xf))
            continue
        if sp.has_text_frame:
            out.append((sp.text_frame, sp.width or 0, sp, None, xf))
    return out

def iter_shapes(shapes):
    """Flattens all shapes into a single traversal, recursing into groups."""
    for sp in shapes:
        try:
            if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
                for inner in iter_shapes(sp.shapes):
                    yield inner
                continue
        except Exception:
            pass
        yield sp

def _is_pic(sp):
    try:
        return sp.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False

# Absolute-coordinate traversal for W15-W17 geometry consumers. A group child's raw left/top
# is in the group's chOff coordinate space, which drifts from slide coordinates in a pptx
# where the group has been moved or resized (off!=chOff desync, standard behavior when
# dragging in PowerPoint) (measured in adversarial verification, 2026-07-03). Composes the
# off/ext vs chOff/chExt affine and yields (shape, z-order, absolute xfrm function
# coefficients). xf=(ax,bx,ay,by): abs = a*raw + b (EMU).
def iter_shapes_geo(shapes, xf=(1.0, 0.0, 1.0, 0.0), _z=None):
    if _z is None:
        _z = [0]
    for sp in shapes:
        try:
            is_grp = sp.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_grp = False
        if is_grp:
            # Group-internal coordinate g: composes abs_g = ox + (g - cx)*sx with the parent
            # xf (_group_xf)
            for inner in iter_shapes_geo(sp.shapes, _group_xf(sp, xf), _z):
                yield inner
            continue
        z = _z[0]
        _z[0] += 1
        yield sp, z, xf

def _geo_rect(sp, xf):
    """The absolute bbox (in) with xf applied. A rotated shape is expanded to an axis-aligned
    bbox around a fixed center. None on failure."""
    try:
        L, T, Wd, Ht = sp.left, sp.top, sp.width, sp.height
    except Exception:
        return None
    if None in (L, T, Wd, Ht):
        return None
    ax, bx, ay, by = xf
    x = (ax * L + bx) / EMU_PER_IN
    y = (ay * T + by) / EMU_PER_IN
    w = ax * Wd / EMU_PER_IN
    h = ay * Ht / EMU_PER_IN
    rot = 0.0
    try:
        rot = float(sp.rotation or 0.0)
    except Exception:
        pass
    if rot % 360.0:
        import math
        r = math.radians(rot)
        w2 = abs(w * math.cos(r)) + abs(h * math.sin(r))
        h2 = abs(w * math.sin(r)) + abs(h * math.cos(r))
        x, y, w, h = x + (w - w2) / 2, y + (h - h2) / 2, w2, h2
    return x, y, w, h, (rot % 360.0 != 0.0)
