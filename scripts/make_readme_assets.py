# -*- coding: utf-8 -*-
"""Build the README showcase assets: decks, renders, comparison PNGs, and the GIF.

The README's hero assets are purpose-built showcase decks, deliberately separate
from the ``archforge demo`` decks (which stay a minimal 30-second introduction).
This script is the single reproduction path the README promises: it builds the
before/after decks, renders them through PowerPoint COM, composes the light and
dark comparison images, and assembles the animation.

Design decisions (2026-09-01, external design review):
- Each showcase deck carries exactly four defects, two visible per slide:
  page 1 seeds E3 (a 4pt source line) and W20 (a caption buried on a KPI card),
  page 2 seeds W15 (two frames collided) and W22 (a rule impaling a label).
  That yields ERROR 1 / WARN 3 and exit 1 under the full profile, which is the
  narrative the animation tells; the ko deck adds E4 on the title. A self-check
  below asserts the exact finding set so an engine change cannot silently turn
  the story false.
- The animation uses hard match cuts on identical coordinates instead of
  crossfades (double exposure was blurring the very contrast it should show),
  plus magnified insets, cropped from the full-resolution render, for the two
  defects too small to read at README width. The stage shares the banner's
  palette (near-black, bronze, a terminal bar), defects get a spotlight dim
  and a filled code badge, and each match cut is confirmed by green
  code-fixed badges on the patched regions.

Requires PowerPoint for Windows (COM render). Output goes to docs/assets/.
"""
import json
import os
import subprocess
import sys

from PIL import Image, ImageDraw, ImageFont

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(ROOT, "docs", "assets")
WORK = os.path.join(ROOT, "build", "readme_assets")
sys.path.insert(0, os.path.join(ROOT, "src"))

from pptx import Presentation                     # noqa: E402
from pptx.dml.color import RGBColor               # noqa: E402
from pptx.enum.shapes import MSO_SHAPE            # noqa: E402
from pptx.util import Inches, Pt                  # noqa: E402

# ---------------------------------------------------------------- palette
INK = "1F2430"          # body text
SUB = "5A6472"          # secondary text
CARD = "3E4A5A"         # KPI card fill (working-deck slate)
CARD_BURIED = "46525F"  # W20: caption nearly the card's own color
ACCENT = "2F6B4F"       # one accent figure

T = {
    "en": {
        "title1": "Q3 Results Summary",
        "sub1": "Subscriptions carried the quarter and churn stayed flat",
        "bullets1": [
            "Revenue grew 18% QoQ to $4.2M, a third straight quarter of double digit growth",
            "Subscriptions now account for 41% of total revenue, up from 33% in Q2",
            "Gross margin improved to 71% on infrastructure consolidation",
            "Enterprise pipeline doubled, with 14 deals above $50k in stage 3 or later",
            "Churn held at 2.1% monthly through the March price change",
        ],
        "cards1": [("+18%", "Revenue QoQ"), ("41%", "Subscription mix"), ("71%", "Gross margin")],
        "buried": "was 68% in Q2",
        "e3": "Source: internal management accounts, June 2026. Unaudited, excludes FX effects.",
        "title2": "Key metrics this quarter",
        "sub2": "Every figure below is management-accounts basis, not audited",
        "kpi_a": "Revenue growth +18%",
        "kpi_b": "Operating margin 12.4%",
        "bullets2": [
            "Cash conversion at 94%, receivables down 11 days",
            "Headcount flat at 63, revenue per head up 19%",
            "Support load down 8% after the onboarding rework",
        ],
        "impaled": "Next quarter guidance",
        "gif_lines": {
            "open": "$ archforge showcase.pptx --profile full",
            "p1": "page 1: E3 4pt source line, W20 caption buried on its card",
            "inset1": "the 4pt source line, magnified",
            "cut1": "same coordinates, patched",
            "p2": "page 2: W15 frames collided, W22 rule through the label",
            "inset2": "the impaled label, magnified",
            "end": "4 found  >  4 patched  >  CLEAN",
        },
        "ba_caption": "PowerPoint renders both without a single error message. "
                      "archforge reads the .pptx and blocks the left one.",
    },
    "ko": {
        "title1": "3분기 실적 요약",
        "sub1": "구독이 분기를 끌었고 이탈률은 유지됐습니다",
        "bullets1": [
            "매출은 전분기 대비 18% 성장한 42억 원, 세 분기 연속 두 자릿수 성장",
            "구독 매출 비중 41%, 2분기 33%에서 상승",
            "인프라 통합으로 매출총이익률 71%로 개선",
            "엔터프라이즈 파이프라인 2배, 5천만 원 이상 딜 14건이 3단계 이후",
            "3월 가격 개편에도 월 이탈률 2.1% 유지",
        ],
        "cards1": [("+18%", "매출 성장"), ("41%", "구독 비중"), ("71%", "매출총이익률")],
        "buried": "2분기는 68%",
        "e3": "출처: 내부 관리회계, 2026년 6월. 비감사 수치이며 환효과 제외.",
        "title2": "이번 분기 핵심 지표",
        "sub2": "아래 수치는 전부 관리회계 기준이며 감사 전입니다",
        "kpi_a": "매출 성장 +18%",
        "kpi_b": "영업이익률 12.4%",
        "bullets2": [
            "현금전환율 94%, 매출채권 회수 11일 단축",
            "인원 63명 유지, 인당 매출 19% 상승",
            "온보딩 개편 후 지원 문의 8% 감소",
        ],
        "impaled": "다음 분기 가이던스",
        "gif_lines": {
            "open": "$ archforge showcase.pptx --profile full",
            "p1": "1쪽: E3 4pt 출처, W20 카드색에 묻힌 캡션",
            "inset1": "4pt 출처 표기, 확대",
            "cut1": "같은 좌표, 교정본",
            "p2": "2쪽: W15 프레임 충돌, W22 괘선이 라벨 관통",
            "inset2": "관통당한 라벨, 확대",
            "end": "4 found  >  4 patched  >  CLEAN",
        },
        "ba_caption": "파워포인트는 둘 다 에러 없이 엽니다. archforge는 pptx를 읽어 왼쪽을 차단합니다.",
    },
}

# defect geometry shared by the deck builder and the marker/inset composer (inches)
GEO = {
    "e3": (0.8, 6.95, 8.6, 0.28),
    "e3_inset": (0.85, 6.985, 2.4, 0.1),
    "buried_card": (9.05, 4.7, 3.4, 1.9),
    "buried_label": (9.25, 6.05, 3.0, 0.35),
    "kpi_a": (1.0, 2.6, 5.0, 0.9),
    "kpi_b_broken": (1.25, 2.75, 5.0, 0.9),
    "kpi_b_fixed": (1.0, 3.45, 5.0, 0.7),
    "impaled": (9.6, 5.4, 2.9, 0.4),
    "vrule_broken": (9.98, 1.9, 0.012, 4.1),
    "vrule_fixed": (9.98, 1.9, 0.012, 3.3),
}


def _tb(slide, x, y, w, h, text, size=16, color=INK, bold=False, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = RGBColor.from_string(color)
    return box


def _run_ea(box, ea):
    r = box.text_frame.paragraphs[0].runs[0]
    rPr = r._r.get_or_add_rPr()
    e = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea", {"typeface": ea})
    rPr.append(e)


def _run_spc(box, spc):
    r = box.text_frame.paragraphs[0].runs[0]
    r._r.get_or_add_rPr().set("spc", str(spc))


def _rect(slide, x, y, w, h, fill, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if not line:
        sp.line.fill.background()
    return sp


def build_showcase(path, lang, fixed):
    t = T[lang]
    ko = lang == "ko"
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)

    # ---- page 1: E3 + W20 (+E4 title on ko)
    s = p.slides.add_slide(p.slide_layouts[6])
    title = _tb(s, 0.8, 0.5, 10.5, 0.8, t["title1"], size=30, bold=True)
    if ko:
        _run_ea(title, "맑은 고딕")
        if not fixed:
            _run_spc(title, 300)                          # E4
    sub1 = _tb(s, 0.8, 1.25, 10.5, 0.4, t["sub1"], size=15, color=SUB)
    if ko:
        _run_ea(sub1, "맑은 고딕")
    for i, b in enumerate(t["bullets1"]):
        bb = _tb(s, 0.95, 1.95 + i * 0.52, 7.6, 0.45, b, size=14)
        if ko:
            _run_ea(bb, "맑은 고딕")
    for i, (num, lab) in enumerate(t["cards1"]):
        cx = 9.05
        cy = 1.95 + i * 1.42
        if i < 2:
            _rect(s, cx, cy, 3.4, 1.15, CARD)
            _tb(s, cx + 0.2, cy + 0.12, 3.0, 0.55, num, size=26, bold=True, color="FFFFFF")
            cl = _tb(s, cx + 0.2, cy + 0.68, 3.0, 0.35, lab, size=12, color="C9D2DC")
            if ko:
                _run_ea(cl, "맑은 고딕")
    # third card carries the buried caption (W20)
    bx, by, bw, bh = GEO["buried_card"]
    _rect(s, bx, by, bw, bh, CARD)
    _tb(s, bx + 0.2, by + 0.12, 3.0, 0.55, t["cards1"][2][0], size=26, bold=True, color="FFFFFF")
    cl3 = _tb(s, bx + 0.2, by + 0.68, 3.0, 0.35, t["cards1"][2][1], size=12, color="C9D2DC")
    if ko:
        _run_ea(cl3, "맑은 고딕")
    lx, ly, lw, lh = GEO["buried_label"]
    cap = _tb(s, lx, ly, lw, lh, t["buried"], size=13,
              color="FFFFFF" if fixed else CARD_BURIED)          # W20 when broken
    if ko:
        _run_ea(cap, "맑은 고딕")
    ex, ey, ew, eh = GEO["e3"]
    src = _tb(s, ex, ey, ew, eh, t["e3"], size=9 if fixed else 4, color=SUB)  # E3
    if ko:
        _run_ea(src, "맑은 고딕")

    # ---- page 2: W15 + W22
    s = p.slides.add_slide(p.slide_layouts[6])
    t2 = _tb(s, 0.8, 0.5, 10.5, 0.8, t["title2"], size=30, bold=True)
    if ko:
        _run_ea(t2, "맑은 고딕")
    sub2 = _tb(s, 0.8, 1.25, 10.5, 0.4, t["sub2"], size=15, color=SUB)
    if ko:
        _run_ea(sub2, "맑은 고딕")
    ax, ay, aw, ah = GEO["kpi_a"]
    ka = _tb(s, ax, ay, aw, ah, t["kpi_a"], size=24, color=ACCENT, bold=True)
    bxy = GEO["kpi_b_fixed"] if fixed else GEO["kpi_b_broken"]
    kb = _tb(s, bxy[0], bxy[1], bxy[2], bxy[3], t["kpi_b"], size=24, bold=True)  # W15 pair
    for i, b in enumerate(t["bullets2"]):
        bb = _tb(s, 0.95, 4.3 + i * 0.52, 7.4, 0.45, b, size=14)
        if ko:
            _run_ea(bb, "맑은 고딕")
    if ko:
        for box in (ka, kb):
            _run_ea(box, "맑은 고딕")
    vx, vy, vw, vh = GEO["vrule_fixed"] if fixed else GEO["vrule_broken"]
    _rect(s, vx, vy, vw, vh, "9AA4B0")                            # W22 rule when broken
    ix, iy, iw, ih = GEO["impaled"]
    imp = _tb(s, ix, iy, iw, ih, t["impaled"], size=14, color=SUB)
    if ko:
        _run_ea(imp, "맑은 고딕")

    p.save(path)
    return path


# ---------------------------------------------------------------- lint self-check
EXPECT = {"en": {"E3": 1, "W20": 1, "W15": 1, "W22": 1},
          "ko": {"E3": 1, "W20": 1, "W15": 1, "W22": 1, "E4": 1}}


def selfcheck(path, lang, fixed):
    out = subprocess.run(
        [sys.executable, "-m", "archforge", path, "--profile", "full",
         "--json", "--schema", "2"],
        capture_output=True, text=True, encoding="utf-8", cwd=ROOT)
    doc = json.loads(out.stdout)
    got = {}
    for f in doc["findings"]:
        got[f["code"]] = got.get(f["code"], 0) + 1
    if fixed:
        assert not got, "fixed showcase (%s) must be clean, got %r" % (lang, got)
        assert doc["summary"]["pass"] is True
    else:
        assert got == EXPECT[lang], \
            "broken showcase (%s) finding drift: want %r got %r" % (lang, EXPECT[lang], got)
        assert doc["summary"]["pass"] is False
    return doc


# ---------------------------------------------------------------- rendering
def render(pptx_path, out_dir, width=4200):
    import win32com.client
    os.makedirs(out_dir, exist_ok=True)
    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True,
                                  Untitled=False, WithWindow=False)
    try:
        for i in range(1, pres.Slides.Count + 1):
            pres.Slides(i).Export(os.path.join(out_dir, "s%d.png" % i), "PNG",
                                  width, int(width * 7.5 / 13.333))
    finally:
        pres.Close()
    return [os.path.join(out_dir, "s%d.png" % i) for i in (1, 2)]


# ---------------------------------------------------------------- composition
def _font(size, mono=True):
    name = "consola.ttf" if mono else "arial.ttf"
    try:
        return ImageFont.truetype(os.path.join(os.environ["WINDIR"], "Fonts", name), size)
    except Exception:
        return ImageFont.load_default()


def _kfont(size):
    for name in ("malgun.ttf", "malgunbd.ttf"):
        try:
            return ImageFont.truetype(os.path.join(os.environ["WINDIR"], "Fonts", name), size)
        except Exception:
            continue
    return ImageFont.load_default()


DARK_BG = (13, 15, 19)
LIGHT_BG = (255, 255, 255)


def _slide_xy(inch_geo, slide_box):
    x, y, w, h = inch_geo
    sx, sy, sw, sh = slide_box
    fx = sw / 13.333
    fy = sh / 7.5
    return (int(sx + x * fx), int(sy + y * fy), int(sx + (x + w) * fx), int(sy + (y + h) * fy))


def compose_ba(lang, renders_broken, renders_fixed, dark, out_path):
    t = T[lang]
    bg = DARK_BG if dark else LIGHT_BG
    ink = (235, 238, 242) if dark else (31, 36, 48)
    sub = (160, 168, 178) if dark else (90, 100, 114)
    red = (196, 60, 46)
    green = (47, 122, 84) if not dark else (98, 190, 138)
    W, H = 2463, 1496
    card_w, card_h = 1170, 658
    xs, ys = (39, 1254), (69, 760)
    im = Image.new("RGB", (W, H), bg)
    d = ImageDraw.Draw(im)
    mono28 = _font(30)
    mono22 = _font(22)
    kf = _kfont(22) if lang == "ko" else mono22

    slide_boxes = {}
    for col, pages in ((0, renders_broken), (1, renders_fixed)):
        for row, png in enumerate(pages):
            s = Image.open(png).convert("RGB").resize((card_w, card_h), Image.LANCZOS)
            x, y = xs[col], ys[row]
            im.paste(s, (x, y))
            d.rectangle([x - 1, y - 1, x + card_w, y + card_h],
                        outline=(70, 76, 86) if dark else (208, 212, 218), width=1)
            slide_boxes[(col, row)] = (x, y, card_w, card_h)

    d.text((xs[0], 22), "before.pptx", font=mono28, fill=ink)
    d.text((xs[0] + 300, 26), "ERROR 1 / WARN 3, exit 1", font=mono22, fill=red)
    d.text((xs[1], 22), "after.pptx", font=mono28, fill=ink)
    d.text((xs[1] + 260, 26), "CLEAN, exit 0", font=mono22, fill=green)

    markers = [("E3", GEO["e3"], (0, 0)), ("W20", GEO["buried_label"], (0, 0)),
               ("W15", GEO["kpi_b_broken"], (0, 1)), ("W22", GEO["impaled"], (0, 1))]
    for code, geo, cell in markers:
        bx0, by0, bx1, by1 = _slide_xy(geo, slide_boxes[cell])
        cy = (by0 + by1) // 2
        lx = slide_boxes[cell][0] + card_w
        d.line([bx1 + 6, cy, lx + 34, cy], fill=red, width=2)
        d.text((lx + 42, cy - 12), code, font=mono22, fill=red)

    d.text((xs[0], H - 46), t["ba_caption"], font=kf, fill=sub)
    im.save(out_path)
    return out_path


# The animation stage borrows the banner's look: a near-black canvas with faint
# bronze blueprint lines, the slide floated as a lit card, and a terminal bar that
# carries the narration. Defects get a spotlight (everything else dims), a filled
# code badge, and after each match cut the patched region flashes a green check so
# the change itself is pointed at, not merely present.
STAGE_BG = (11, 13, 17)
STAGE_LINE = (176, 141, 87)
TERM_BG = (17, 20, 26)
TERM_INK = (214, 220, 228)
RED = (206, 74, 58)
GREEN = (86, 176, 128)


def _stage(base_size):
    W, H = base_size
    im = Image.new("RGB", (W, H), STAGE_BG)
    d = ImageDraw.Draw(im)
    faint = (34, 30, 24)
    d.ellipse([W - 260, -140, W + 120, 240], outline=faint, width=2)
    d.rectangle([-120, H - 300, 140, H - 40], outline=faint, width=2)
    return im


def _card_box(base_size):
    W, H = base_size
    card_h = H - 66 - 52 - 10          # terminal bar, top margin, breathing room
    card_w = int(card_h * 13.333 / 7.5)
    return (W - card_w) // 2, 52, card_w, card_h


def _badge(d, x, y, text, fill, fnt):
    tw = int(d.textlength(text, font=fnt))
    d.rectangle([x, y, x + tw + 22, y + 38], fill=fill)
    d.text((x + 11, y + 6), text, font=fnt, fill=(255, 255, 255))


def _frame(base_size, slide_png, term_line, kf, marks=None, inset=None,
           fixed_marks=None, verdict="exit 1"):
    W, H = base_size
    im = _stage(base_size)
    s = Image.open(slide_png).convert("RGB")
    sx, sy, card_w, card_h = _card_box(base_size)
    s = s.resize((card_w, card_h), Image.LANCZOS)

    if marks:
        dim = Image.new("L", s.size, 84)          # spotlight: dim all but the defects
        dd = ImageDraw.Draw(dim)
        for _, geo in marks:
            x0, y0, x1, y1 = _slide_xy(geo, (0, 0, card_w, card_h))
            dd.rectangle([x0 - 14, y0 - 10, x1 + 14, y1 + 10], fill=0)
        s = Image.composite(Image.new("RGB", s.size, (24, 26, 31)), s,
                            dim.point(lambda v: v))
        s = Image.blend(Image.open(slide_png).convert("RGB").resize(s.size, Image.LANCZOS),
                        s, 0.72)

    im.paste(s, (sx, sy))
    d = ImageDraw.Draw(im)
    d.rectangle([sx - 2, sy - 2, sx + card_w + 1, sy + card_h + 1],
                outline=(107, 87, 63), width=2)
    d.text((sx, sy - 34), "showcase.pptx", font=_font(21), fill=(150, 128, 96))

    code_f = _font(24)
    if marks:
        for code, geo in marks:
            x0, y0, x1, y1 = _slide_xy(geo, (sx, sy, card_w, card_h))
            d.rectangle([x0 - 14, y0 - 10, x1 + 14, y1 + 10], outline=RED, width=4)
            _badge(d, x0 - 14, max(sy + 4, y0 - 54), code, RED, code_f)
    if fixed_marks:
        for code, geo in fixed_marks:
            x0, y0, x1, y1 = _slide_xy(geo, (sx, sy, card_w, card_h))
            d.rectangle([x0 - 14, y0 - 10, x1 + 14, y1 + 10], outline=GREEN, width=4)
            _badge(d, x0 - 14, max(sy + 4, y0 - 54), code + " fixed", (47, 122, 84), code_f)
    if inset:
        geo, scale, code, pad = inset
        x0, y0, x1, y1 = _slide_xy(geo, (sx, sy, card_w, card_h))
        # crop from the full-resolution render, not the resized card, so the
        # magnified region enlarges real pixels instead of card-size ones
        src = Image.open(slide_png).convert("RGB")
        fx = src.size[0] / card_w
        crop = src.crop((int((x0 - sx - pad) * fx), int((y0 - sy - pad) * fx),
                         int((x1 - sx + pad) * fx), int((y1 - sy + pad) * fx)))
        cw = int((x1 - x0 + 2 * pad) * scale)
        ch = int((y1 - y0 + 2 * pad) * scale)
        crop = crop.resize((cw, ch), Image.LANCZOS)
        px, py = (W - cw) // 2, sy + card_h - ch - 44
        d.rectangle([px - 8, py - 8, px + cw + 7, py + ch + 7], fill=STAGE_BG)
        im.paste(crop, (px, py))
        d.rectangle([px - 3, py - 3, px + cw + 2, py + ch + 2], outline=RED, width=4)
        _badge(d, px - 3, py - 54, code, RED, code_f)

    # terminal bar
    ty = H - 66
    d.rectangle([0, ty, W, H], fill=TERM_BG)
    d.line([0, ty, W, ty], fill=(45, 50, 58), width=1)
    d.text((40, ty + 18), "$", font=_font(26), fill=STAGE_LINE)
    d.text((72, ty + 18), term_line, font=kf, fill=TERM_INK)
    vcol = GREEN if "0" in verdict else RED
    vw = int(d.textlength(verdict, font=_font(24)))
    d.text((W - vw - 44, ty + 19), verdict, font=_font(24), fill=vcol)
    return im


def _end_card(base_size, g, lang):
    W, H = base_size
    im = _stage(base_size)
    d = ImageDraw.Draw(im)
    big = _font(52)
    txt = g["end"]
    tw = int(d.textlength(txt, font=big))
    d.text(((W - tw) // 2, H // 2 - 88), txt, font=big, fill=(232, 226, 210))
    ok = "exit 0"
    of = _font(34)
    ow = int(d.textlength(ok, font=of))
    d.text(((W - ow) // 2, H // 2), ok, font=of, fill=GREEN)
    try:
        serif = ImageFont.truetype(
            os.path.join(os.environ["WINDIR"], "Fonts", "georgia.ttf"), 40)
    except Exception:
        serif = _font(40)
    name = "Archforge"
    nw = int(d.textlength(name, font=serif))
    d.text(((W - nw) // 2, H // 2 + 110), name, font=serif, fill=(150, 128, 96))
    return im


def compose_gif(lang, rb, rf, out_path):
    g = T[lang]["gif_lines"]
    kf = _kfont(24) if lang == "ko" else _font(24)
    size = (1600, 900)
    frames = []

    def add(n, img):
        frames.extend([img] * n)

    p1_marks = [("E3", GEO["e3"]), ("W20", GEO["buried_label"])]
    p1_fixed = [("E3", GEO["e3"]), ("W20", GEO["buried_label"])]
    p2_marks = [("W15", GEO["kpi_b_broken"]), ("W22", GEO["impaled"])]
    p2_fixed = [("W15", GEO["kpi_b_fixed"]), ("W22", GEO["impaled"])]

    add(3, _frame(size, rb[0], g["open"], kf))
    add(3, _frame(size, rb[0], g["p1"], kf, marks=p1_marks))
    add(3, _frame(size, rb[0], g["inset1"], kf, inset=(GEO["e3_inset"], 3.8, "E3", 12)))
    add(1, _frame(size, rb[0], g["p1"], kf, marks=p1_marks))
    add(3, _frame(size, rf[0], g["cut1"], kf, fixed_marks=p1_fixed, verdict="exit 0"))
    add(1, _frame(size, rf[0], g["cut1"], kf, verdict="exit 0"))
    add(3, _frame(size, rb[1], g["p2"], kf, marks=p2_marks))
    add(3, _frame(size, rb[1], g["inset2"], kf, inset=(GEO["impaled"], 2.2, "W22", 26)))
    add(1, _frame(size, rb[1], g["p2"], kf, marks=p2_marks))
    add(3, _frame(size, rf[1], g["cut1"], kf, fixed_marks=p2_fixed, verdict="exit 0"))
    add(1, _frame(size, rf[1], g["cut1"], kf, verdict="exit 0"))
    add(3, _end_card(size, g, lang))

    q = [f.quantize(colors=128, dither=Image.Dither.NONE) for f in frames]
    q[0].save(out_path, save_all=True, append_images=q[1:], duration=340, loop=0,
              optimize=True)
    return out_path


# ---------------------------------------------------------------- terminal card
def compose_terminal(lang, deck_path, out_path):
    """Render the actual CLI output as a dark terminal card. The text is captured
    from a real run, not typed in, so the card cannot drift from the product."""
    out = subprocess.run(
        [sys.executable, "-m", "archforge", deck_path, "--profile", "full",
         "--lang", lang],
        capture_output=True, text=True, encoding="utf-8", cwd=ROOT)
    cmd = "archforge showcase.pptx --profile full"
    lines = [ln for ln in out.stdout.rstrip().split("\n") if ln.strip()]
    lines[0] = "=== ARCHFORGE LINT: showcase.pptx ==="   # hide the build path

    body_f = _kfont(24) if lang == "ko" else _font(24)
    mono = _font(24)
    probe = ImageDraw.Draw(Image.new("RGB", (8, 8)))

    # fold long finding lines at their " | " payload separator so the card stays
    # readable at README width instead of stretching into a ribbon
    wrapped = []
    for ln in lines:
        if int(probe.textlength(ln, font=body_f)) > 1440 and " | " in ln:
            head, tail = ln.split(" | ", 1)
            wrapped.append((head, ln))
            wrapped.append(("        | " + tail, ln))
        else:
            wrapped.append((ln, ln))
    lines = wrapped
    maxw = max(int(probe.textlength(txt, font=body_f)) for txt, _ in lines)
    W = max(1360, maxw + 120)
    lh = 42
    H = 96 + 46 + lh * len(lines) + 40

    im = Image.new("RGB", (W, H), TERM_BG)
    d = ImageDraw.Draw(im)
    d.rectangle([0, 0, W, 64], fill=(24, 28, 35))
    for i, c in enumerate([(228, 96, 84), (222, 176, 74), (98, 186, 116)]):
        d.ellipse([28 + i * 34, 24, 44 + i * 34, 40], fill=c)
    tw = int(d.textlength(cmd, font=mono))
    d.text(((W - tw) // 2, 18), cmd, font=mono, fill=(150, 158, 168))

    y = 96
    d.text((48, y), "$ " + cmd, font=mono, fill=TERM_INK)
    y += 46
    for txt, src in lines:
        if src.startswith("---"):
            col = RED
        elif "ERROR" in src.split("[")[0]:
            col = (226, 118, 104)
        elif src.lstrip().startswith("WARN"):
            col = (217, 172, 96)
        else:
            col = (168, 176, 186)
        d.text((48, y), txt, font=body_f, fill=col)
        y += lh
    d.rectangle([0, 0, W - 1, H - 1], outline=(52, 58, 66), width=1)
    im.save(out_path)
    return out_path


# ---------------------------------------------------------------- main
def main():
    os.makedirs(WORK, exist_ok=True)
    for lang in ("en", "ko"):
        broken = os.path.join(WORK, "showcase-before-%s.pptx" % lang)
        fixedp = os.path.join(WORK, "showcase-after-%s.pptx" % lang)
        build_showcase(broken, lang, fixed=False)
        build_showcase(fixedp, lang, fixed=True)
        selfcheck(broken, lang, fixed=False)
        selfcheck(fixedp, lang, fixed=True)
        rb = render(broken, os.path.join(WORK, "rb_%s" % lang))
        rf = render(fixedp, os.path.join(WORK, "rf_%s" % lang))
        compose_ba(lang, rb, rf, dark=False,
                   out_path=os.path.join(ASSETS, "showcase-%s.png" % lang))
        compose_ba(lang, rb, rf, dark=True,
                   out_path=os.path.join(ASSETS, "showcase-%s-dark.png" % lang))
        compose_gif(lang, rb, rf, os.path.join(ASSETS, "showcase-%s.gif" % lang))
        compose_terminal(lang, broken,
                         os.path.join(ASSETS, "terminal-%s.png" % lang))
        print(lang, "assets done")
    print("all showcase assets rebuilt")


if __name__ == "__main__":
    main()
