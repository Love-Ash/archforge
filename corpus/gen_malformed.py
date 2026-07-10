# -*- coding: utf-8 -*-
"""Seed corpus, malformed/atypical set: inputs that must produce a controlled outcome
(usage error or an incomplete-marked report), never a traceback or a silent pass.

Usage: python corpus/gen_malformed.py   (run after gen_python_pptx.py)
"""
import json
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DST = os.path.join(HERE, "malformed")


def main():
    os.makedirs(DST, exist_ok=True)

    # 1. Truncated package: the first 800 bytes of a valid pptx. Must be a clean usage
    #    error (zip preflight / open failure), exit 2 at the CLI.
    src = os.path.join(HERE, "python-pptx", "clean_bilingual.pptx")
    with open(src, "rb") as f:
        head = f.read(800)
    with open(os.path.join(DST, "truncated.pptx"), "wb") as f:
        f.write(head)
    with open(os.path.join(DST, "truncated.json"), "w", encoding="utf-8") as f:
        json.dump({"generator": "byte-truncated python-pptx output",
                   "expect_exit2": True,
                   "notes": "must be a controlled usage error, never a traceback"},
                  f, ensure_ascii=False, indent=2)

    # 2. Vertical text: geometry cannot estimate it, so the report must be marked
    #    incomplete (W18) instead of guessing widths.
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    s = p.slides.add_slide(p.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "세로쓰기 텍스트"
    r.font.size = Pt(18)
    rPr = r._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": "맑은 고딕"}))
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.set("vert", "eaVert")
    p.save(os.path.join(DST, "vertical_text.pptx"))
    with open(os.path.join(DST, "vertical_text.json"), "w", encoding="utf-8") as f:
        json.dump({"generator": "python-pptx (bodyPr@vert=eaVert)",
                   "expected": {}, "expect_incomplete": True,
                   "notes": "geometry abstains on vertical text and must say so"},
                  f, ensure_ascii=False, indent=2)
    print("wrote", DST)


if __name__ == "__main__":
    main()
