# -*- coding: utf-8 -*-
"""Seed corpus, python-pptx generator: one deck per flagship defect, constructed so the
ground truth is known by construction (the defect is deliberately seeded) rather than
by trusting the linter's own output. Expected findings live in the sibling .json
manifests consumed by corpus/run_corpus.py.

Usage: python corpus/gen_python_pptx.py
"""
import json
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

EM_DASH = chr(0x2014)
HERE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "python-pptx")


def _prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def _tb(s, x, y, w, h, text, size=14, font=None, ea=None, spc=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    if font:
        r.font.name = font
    if ea:
        rPr = r._r.get_or_add_rPr()
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": ea}))
    if spc is not None:
        r._r.get_or_add_rPr().set("spc", str(spc))
    return box


def emit(name, build, manifest):
    os.makedirs(HERE, exist_ok=True)
    path = os.path.join(HERE, name + ".pptx")
    p = _prs()
    build(p)
    # docProps carries python-pptx's bundled template properties, which name a
    # stranger and claim PowerPoint on a Mac wrote the file. Neither is true and both
    # ship in the sdist, so the writer is stated honestly instead. No rule reads
    # docProps, and the manifest's "generator" field is what records provenance.
    from archforge.demo import _save
    # docProps says "archforge corpus", not the library name. Provenance lives in the
    # manifest's "generator" field, which is where the corpus README points and where a
    # reader looks; putting the library name in docProps as well only ships a
    # generated-by tell in every sdist for information already recorded.
    _save(p, path, application="archforge corpus")
    manifest.setdefault("generator", "python-pptx")
    manifest.setdefault("profile", "full")
    manifest.setdefault("ground_truth", "by construction (defect deliberately seeded)")
    with open(os.path.join(HERE, name + ".json"), "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print("wrote", path)


def main():
    def e1(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 6, 1, "아리알에 실린 한글", size=20, font="Arial")
    emit("e1_arial_hangul", e1, {"expected": {"E1": 1},
         "notes": "Hangul on a Latin-only a:latin with the default theme's empty ea slot"})

    def e2(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 8, 1, "growth was structural" + EM_DASH + "not cyclical", size=16,
            ea="맑은 고딕")
    emit("e2_em_dash", e2, {"expected": {"E2": 1},
         "notes": "word-to-word em dash; numeric ranges would pass"})

    def e3(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 6.9, 5, 0.3, "source: internal accounts", size=4, ea="맑은 고딕")
    emit("e3_4pt", e3, {"expected": {"E3": 1}})

    def e4(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 6, 1, "자간이 벌어진 한글", size=16, spc=300, ea="맑은 고딕")
    emit("e4_tracked_hangul", e4, {"expected": {"E4": 1}})

    def w15(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1.0, 2.4, 5.0, 1.0, "Revenue growth +18%", size=24, ea="맑은 고딕")
        _tb(s, 1.2, 2.5, 5.0, 1.0, "Operating margin 12.4%", size=24, ea="맑은 고딕")
    emit("w15_overlap", w15, {"expected": {"W15": 1}})

    def w16(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 12.0, 4.5, 3.0, 0.6, "Next quarter guidance", size=18, ea="맑은 고딕")
    emit("w16_offcanvas", w16, {"expected": {"W16": 1}})

    def w19(p):
        from pptx.dml.color import RGBColor
        s = p.slides.add_slide(p.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Nearly invisible label"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor.from_string("DDDDDD")
    emit("w19_low_contrast", w19, {"expected": {"W19": 1},
         "notes": "Explicit #DDDDDD text on the shape's own explicit #FFFFFF solid fill. "
                  "Ground truth by construction: relative luminance of #DDDDDD is 0.723 and "
                  "of #FFFFFF is 1.0, so the WCAG ratio is 1.05/0.773 = 1.36, under the "
                  "calibrated 2.0 threshold. Both colors are explicit RGB, so no resolution "
                  "chain is involved."})

    def w19_neg(p):
        from pptx.dml.color import RGBColor
        s = p.slides.add_slide(p.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Subdued but readable caption"
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor.from_string("6E6E73")
    emit("w19_contrast_negative", w19_neg, {"expected": {},
         "notes": "Subdued gray #6E6E73 on white, the standard secondary-text idiom. Ratio "
                  "is roughly 4.9:1, far above the 2.0 threshold; a rule that fires here "
                  "would flag every well-set caption. Must stay silent."})

    def w20(p):
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        s = p.slides.add_slide(p.slide_layouts[6])
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5),
                                 Inches(3), Inches(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string("2ECC9B")
        bar.line.fill.background()
        box = s.shapes.add_textbox(Inches(1.4), Inches(4.6), Inches(4.4), Inches(0.5))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Footnote dropped onto the bar"
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor.from_string("8A8A8A")
    emit("w20_text_over_shape", w20, {"expected": {"W20": 1},
         "notes": "A separate caption box laid across a filled bar below it in z-order. "
                  "Ground truth by construction: #8A8A8A on #2ECC9B is a 1.53:1 WCAG ratio, "
                  "under the 2.0 line W19 already uses, and the glyph box sits far enough "
                  "onto the bar to pass the coverage floor. W19 cannot see this because the "
                  "text is not inside the filled shape, and W15 cannot because the bar holds "
                  "no text of its own: the defect falls between them, which is why W20 "
                  "exists."})

    def w20_neg(p):
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        s = p.slides.add_slide(p.slide_layouts[6])
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5),
                                 Inches(3), Inches(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string("2ECC9B")
        bar.line.fill.background()
        box = s.shapes.add_textbox(Inches(1.4), Inches(4.6), Inches(4.4), Inches(0.5))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Label deliberately set on the bar"
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor.from_string("0B0B0B")
    emit("w20_over_shape_negative", w20_neg, {"expected": {},
         "notes": "Same geometry as the positive fixture, near-black text instead of gray. "
                  "#0B0B0B on #2ECC9B is about 8.4:1, so the label reads cleanly. A caption "
                  "deliberately placed on a colored panel is a normal layout and the gate "
                  "must stay silent on it; this fixture is what stops W20 from becoming a "
                  "rule against overlapping text as such."})

    def w20_bg_ghost(p):
        from pptx.dml.color import RGBColor
        s = p.slides.add_slide(p.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(0.8))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Ghost note left on the empty slide"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor.from_string("DDDDDD")
    emit("w20_bg_ghost", w20_bg_ghost, {"expected": {"W20": 1},
         "notes": "A transparent text box on the bare default background: no shape under "
                  "it at all. The template carries no slide-level p:bg, so the color "
                  "resolves through the stock bgRef-1001 chain to the theme's bg1/lt1 "
                  "white, and #DDDDDD on #FFFFFF is 1.36:1 -- the same ghost-text ratio "
                  "the W19 fixture pins, one layer further down. This is the case W19 "
                  "structurally cannot see (the run's own shape has no fill) and W15 "
                  "cannot either (nothing overlaps)."})

    def w20_bg_readable(p):
        from pptx.dml.color import RGBColor
        s = p.slides.add_slide(p.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(0.8))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Subdued but readable caption on the page"
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor.from_string("6E6E73")
    emit("w20_bg_readable", w20_bg_readable, {"expected": {},
         "notes": "Same geometry as the ghost fixture with the standard secondary-text "
                  "gray: #6E6E73 on white is about 4.9:1, far above the 2.0 line. This "
                  "fixture is what stops the background variant from flagging every "
                  "well-set caption on a plain slide."})

    def w20_bg_dark(p):
        from pptx.dml.color import RGBColor
        s = p.slides.add_slide(p.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor.from_string("0A0B0C")
        box = s.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(0.8))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "Note lost against the dark page"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor.from_string("2A2E33")
    emit("w20_bg_dark", w20_bg_dark, {"expected": {"W20": 1},
         "notes": "An explicit slide-level solid background (p:bgPr srgbClr, the form "
                  "the decks this gate was built against actually use) with dark-gray "
                  "text on the near-black page: roughly 1.4:1. Pins the slide-level "
                  "branch of the background resolution, where the ghost fixture pins "
                  "the theme-fallback branch."})

    def _svg_deck(p, text_fill):
        # An svgBlip picture built by hand, no matplotlib: a green bar as an M/L/z
        # path and one caption string laid across it, the shape SVG exporters emit.
        from pptx.oxml.ns import qn
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" width="400" height="200">'
            '<path d="M 0 200 L 400 200 L 400 0 L 0 0 z" style="fill: #0A0B0C"/>'
            '<path d="M 120 180 L 280 180 L 280 30 L 120 30 z" style="fill: #20C997"/>'
            '<text x="200" y="120" style="font-size: 12px; text-anchor: middle; '
            'fill: %s">caption resting on the bar</text>'
            "</svg>" % text_fill
        ).encode("utf-8")
        s = p.slides.add_slide(p.slide_layouts[6])
        png = os.path.join(HERE, "_svg_fallback.png")
        if not os.path.exists(png):
            from pptx.util import Emu as _Emu  # noqa: local import keeps top clean
            import zlib, struct

            def _chunk(tag, data):
                c = struct.pack(">I", len(data)) + tag + data
                return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

            raw = b"".join(b"\x00" + b"\x10\x10\x10" * 4 for _ in range(4))
            body = (_chunk(b"IHDR", struct.pack(">IIBBBBB", 4, 4, 8, 2, 0, 0, 0))
                    + _chunk(b"IDAT", zlib.compress(raw))
                    + _chunk(b"IEND", b""))
            with open(png, "wb") as f:
                f.write(b"\x89PNG\r\n\x1a\n" + body)
        pic = s.shapes.add_picture(png, Inches(1), Inches(1), width=Inches(8),
                                   height=Inches(4))
        n = 1
        while any(str(pt.partname) == "/ppt/media/svgimage%d.svg" % n
                  for pt in p.part.package.iter_parts()):
            n += 1
        uri = PackURI("/ppt/media/svgimage%d.svg" % n)
        svg_part = Part(uri, "image/svg+xml", p.part.package, svg)
        rid = s.part.relate_to(
            svg_part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
        blip = pic._element.blipFill.find(qn("a:blip"))
        ext_lst = blip.makeelement(qn("a:extLst"), {})
        ext = blip.makeelement(qn("a:ext"),
                               {"uri": "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"})
        svg_el = blip.makeelement(
            "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip",
            {qn("r:embed"): rid})
        ext.append(svg_el)
        ext_lst.append(ext)
        blip.append(ext_lst)

    def w20_svg(p):
        _svg_deck(p, "#8A8F98")   # gray on the green bar: 1.5:1, buried
    emit("w20_svg_buried", w20_svg, {"expected": {"W20": 1},
         "notes": "A caption inside an svgBlip vector picture, laid across a solid "
                  "green bar at roughly 1.5:1. The same defect as the shape variant "
                  "one container deeper: a PNG chart erases its text, an SVG chart "
                  "keeps it as XML, and this fixture pins that the gate actually "
                  "reads the picture. Ground truth by construction; the PNG fallback "
                  "is a 4x4 stub so only the vector side carries content."})

    def w20_svg_neg(p):
        _svg_deck(p, "#0B0B0B")   # near-black on the green bar: ~8.4:1, readable
    emit("w20_svg_readable", w20_svg_neg, {"expected": {},
         "notes": "Same SVG geometry with near-black text on the bar, about 8.4:1. "
                  "A label deliberately set on a colored bar is normal chart design "
                  "and the gate must stay silent on it."})

    def w22_rule_cross(p):
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        s = p.slides.add_slide(p.slide_layouts[6])
        # The evidence geometry: a vertical hairline whose lower end impales a lone
        # arrow glyph sitting flush against it. Contrast is fine and nothing else
        # overlaps, so every earlier gate stays silent by design.
        vline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.9),
                                   Inches(0.012), Inches(4.4))
        vline.fill.solid()
        vline.fill.fore_color.rgb = RGBColor.from_string("444444")
        vline.line.fill.background()
        box = s.shapes.add_textbox(Inches(6.6), Inches(6.0), Inches(0.5), Inches(0.3))
        box.text_frame.margin_left = 0
        box.text_frame.margin_right = 0
        r = box.text_frame.paragraphs[0].add_run()
        r.text = chr(0x2192)
        r.font.size = Pt(14)
    emit("w22_rule_cross", w22_rule_cross, {"expected": {"W22": 1},
         "notes": "A vertical hairline (0.012in wide, 4.4in tall) whose lower end "
                  "passes through a lone arrow glyph placed flush against it -- the "
                  "geometry lifted from the deck this rule was written against. The "
                  "rule's centerline runs through the glyph box interior and covers "
                  "the full run height, so the gate must fire; W15/W17/W20 all stay "
                  "silent here by scope, which is why W22 exists."})

    def w22_rule_negative(p):
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        s = p.slides.add_slide(p.slide_layouts[6])
        # Same hairline vocabulary used the way decks legitimately use it: a divider
        # beside the text, and an underline below the glyph box. Neither crosses.
        vline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.9),
                                   Inches(0.012), Inches(4.4))
        vline.fill.solid()
        vline.fill.fore_color.rgb = RGBColor.from_string("444444")
        vline.line.fill.background()
        box = s.shapes.add_textbox(Inches(6.8), Inches(6.0), Inches(2.5), Inches(0.3))
        box.text_frame.margin_left = 0
        r = box.text_frame.paragraphs[0].add_run()
        r.text = "label beside the divider"
        r.font.size = Pt(14)
        under = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(6.42),
                                   Inches(2.2), Inches(0.01))
        under.fill.solid()
        under.fill.fore_color.rgb = RGBColor.from_string("444444")
        under.line.fill.background()
    emit("w22_rule_negative", w22_rule_negative, {"expected": {},
         "notes": "The same hairline vocabulary used legitimately: the divider sits "
                  "beside the text and an underline rule sits below the glyph box. "
                  "Neither centerline enters a glyph box, so the gate must stay "
                  "silent -- this is what keeps W22 from becoming a rule against "
                  "dividers as such."})

    def clean(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 8, 1, "Quarterly results improved", size=20, ea="맑은 고딕")
        _tb(s, 1, 2.2, 8, 0.6, "구독 매출이 성장을 견인했습니다", size=14, ea="맑은 고딕")
    emit("clean_bilingual", clean, {"expected": {},
         "notes": "negative fixture: correct ea font, no dash, readable sizes"})

    EN_DASH = chr(0x2013)

    def e2_range_negative(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 8, 0.6, "FY2020" + EN_DASH + "2024 revenue up 18%", size=16,
            ea="맑은 고딕")
    emit("e2_range_negative", e2_range_negative, {"expected": {},
         "notes": "negative fixture: a numeric-range en dash is legitimate typography "
                  "and must pass (E2's exemption contract)"})

    def w1_small_body(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 2, 9, 1.2,
            "This body-class paragraph runs well past forty characters of text "
            "so the frame is judged as body copy.", size=8, ea="맑은 고딕")
    emit("w1_small_body", w1_small_body, {"expected": {"W1": 1},
         "notes": "wide frame, long text, 8pt: below the 9pt body floor"})

    def w8_small_cjk(p):
        s = p.slides.add_slide(p.slide_layouts[6])
        _tb(s, 1, 1, 2.0, 0.4, "목업 안 라벨", size=6, ea="맑은 고딕")
    emit("w8_small_cjk", w8_small_cjk, {"expected": {"W8": 1},
         "notes": "narrow (<=4in) frame, 6pt CJK: the mockup-label pattern"})

    def w6_repeated_skeleton(p):
        for i in range(5):
            s = p.slides.add_slide(p.slide_layouts[6])
            _tb(s, 1, 0.8, 8, 0.6, "Section title %d" % i, size=24, ea="맑은 고딕")
            _tb(s, 1, 2.0, 6, 2.5, "Body block content", size=12, ea="맑은 고딕")
            _tb(s, 8, 2.0, 4, 2.5, "Side note block", size=12, ea="맑은 고딕")
    emit("w6_repeated_skeleton", w6_repeated_skeleton, {"expected": {"W6": 1, "W14": 1},
         "notes": "the same 3-frame skeleton on 5 pages: the recycled-grid tell; English noun-phrase section titles also trip W14 under the English path"})

    def w14_en_nominal(p):
        for t in ("Market Overview", "Competitive Analysis", "Product Lineup",
                  "Expansion Strategy", "Financial Plan", "Next Steps Summary"):
            s = p.slides.add_slide(p.slide_layouts[6])
            _tb(s, 1, 0.8, 9, 0.8, t, size=26, ea="맑은 고딕")
            _tb(s, 1, 2.2, 10, 3, "Supporting body copy", size=12, ea="맑은 고딕")
    emit("w14_en_nominal", w14_en_nominal, {"expected": {"W14": 1},
         "notes": "positive W14-EN fixture: English noun-phrase titles across pages"})

    def w14_en_action(p):
        for t in ("Market grows 18% on subscriptions", "Revenue reaches $12 million",
                  "Costs fall 8pp year over year", "Retention improves 2x",
                  "Margin expands after price cuts", "Pipeline adds 40% coverage"):
            s = p.slides.add_slide(p.slide_layouts[6])
            _tb(s, 1, 0.8, 9, 0.8, t, size=26, ea="맑은 고딕")
            _tb(s, 1, 2.2, 10, 3, "Supporting body copy", size=12, ea="맑은 고딕")
    emit("w14_en_action", w14_en_action, {"expected": {},
         "notes": "negative W14-EN fixture: finite verb or number+unit counts as a claim"})

    def _ko_deck(p, titles):
        for t in titles:
            s = p.slides.add_slide(p.slide_layouts[6])
            _tb(s, 1, 0.8, 9, 0.8, t, size=26, ea="맑은 고딕")
            _tb(s, 1, 2.2, 10, 3, "본문은 읽을 수 있는 크기로 들어갑니다.",
                size=12, ea="맑은 고딕")

    def w14_ko_nominal(p):
        # Every title here ends in a syllable that also ends a Korean sentence, so before
        # the noun-collision list they all read as claims and the rule stayed silent.
        _ko_deck(p, ("제품 개요", "사업 개요", "해외 투자", "핵심 사용자",
                     "모바일 게임", "시장 현황"))
    emit("w14_ko_nominal", w14_ko_nominal, {"expected": {"W14": 1},
         "notes": "positive W14 Hangul fixture: noun phrases whose final syllable collides "
                  "with a sentence ending (개요 / 투자 / 사용자 / 게임)"})

    def w14_ko_claim(p):
        # The mirror of the above: real predicates that end in the same syllables, so the
        # collision list may not swallow them.
        _ko_deck(p, ("매출이 전년 대비 늘었다", "성장세가 꺾였어요", "비용을 줄이자",
                     "점유율이 3분기에 반등했다", "이것이 최선인가",
                     "지금 진입해야 한다"))
    emit("w14_ko_claim", w14_ko_claim, {"expected": {},
         "notes": "negative W14 Hangul fixture: declarative, polite, propositive and "
                  "interrogative endings all count as claims. This one also passes on the "
                  "pre-2026-08-07 code, so it proves nothing about the noun-collision fix; "
                  "it is the regression guard against a future rewrite that enumerates "
                  "endings instead of nouns, which would read these as noun phrases and "
                  "fire on a deck of real claims"})

    def w14_ko_structural(p):
        # Three structural slides and two content noun phrases. Counting the structural
        # ones would put five nominal titles in the pool and fire; excluding them leaves
        # two, which is under the three-title minimum.
        _ko_deck(p, ("질의응답", "참고 문헌", "회사 소개", "시장 현황", "매출 추이"))
    emit("w14_ko_structural", w14_ko_structural, {"expected": {},
         "notes": "negative W14 Hangul fixture: cover/divider/closing titles are not part "
                  "of the deck argument and must not enter the majority pool"})


if __name__ == "__main__":
    main()
