# -*- coding: utf-8 -*-
"""Seed corpus, structural axes: the bypass paths the engine advertises and the corpus did
not touch.

Measured 2026-08-08 across the 31 decks then in corpus/: zero contained a table, a group, an
a:fld field, an a:br line break, a rotated shape, or a second slide master. Those are exactly
the paths the module docstring names as covered, and they are where the defects have actually
been. The group axis is the sharpest example: a shipped bug looked for grpSpPr in the a:
namespace when it lives in p:, so the group transform silently fell back to identity and
W15-W17 judged moved groups against raw coordinates (CHANGELOG, 0.5.0). One corpus deck with
a desynced group would have caught it.

The unit suite does cover these axes. What it does not do is feed docs/ACCURACY.md, which is
computed from this corpus and is the published, reproducible record. Until now that record
was built entirely from decks that walk the simple path.

Ground truth is by construction and by XML inspection, never by running the linter and
writing down what it said. Each deck seeds exactly one defect on one axis, and the reason it
must fire is written into the manifest notes.

Usage: python corpus/gen_structural.py
"""
import json
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src"))

from pptx import Presentation                     # noqa: E402
from pptx.util import Inches, Pt, Emu             # noqa: E402
from pptx.oxml.ns import qn                       # noqa: E402

EM_DASH = chr(0x2014)
HERE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "structural")
HANGUL_FONT = "맑은 고딕"


def _prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def _run(para, text, size=14, latin=None, ea=None):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    if latin:
        r.font.name = latin
    if ea:
        r._r.get_or_add_rPr().append(r._r.makeelement(qn("a:ea"), {"typeface": ea}))
    return r


def _tb(slide, x, y, w, h, text, size=14, latin=None, ea=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    _run(box.text_frame.paragraphs[0], text, size, latin, ea)
    return box


# ---------------------------------------------------------------- the axes


def group_desync(p):
    """A group whose child coordinate space is offset from where the group sits.

    off/ext is the group's place on the slide; chOff/chExt is the coordinate space its
    children are expressed in. PowerPoint desyncs them the moment a group is dragged or
    resized, and the child's raw x/y then means nothing on its own. Here the child sits at
    raw x=0.5in inside a child space that starts at 0, while the group itself is parked at
    x=12.4in, so the run's real position is 12.9in on a 13.333in canvas and its text runs off
    the right edge.

    Fires W16 only if the affine is composed. Read raw, the shape looks comfortably inside
    the canvas and nothing is reported, which is precisely the shipped bug this axis exists
    to pin down.
    """
    s = p.slides.add_slide(p.slide_layouts[6])
    grp = s.shapes.add_group_shape()
    _tb(s, 0.5, 3.0, 3.2, 0.6, "다음 분기 가이던스", size=20, ea=HANGUL_FONT)
    inner = s.shapes[-1]
    grp.shapes._spTree.append(inner._element)
    xf = grp._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
    xf.find(qn("a:off")).set("x", str(Emu(Inches(12.4)))); xf.find(qn("a:off")).set("y", str(Emu(Inches(3.0))))
    xf.find(qn("a:ext")).set("cx", str(Emu(Inches(3.2)))); xf.find(qn("a:ext")).set("cy", str(Emu(Inches(0.6))))
    ch_off = xf.find(qn("a:chOff")); ch_ext = xf.find(qn("a:chExt"))
    ch_off.set("x", str(Emu(Inches(0.5)))); ch_off.set("y", str(Emu(Inches(3.0))))
    ch_ext.set("cx", str(Emu(Inches(3.2)))); ch_ext.set("cy", str(Emu(Inches(0.6))))


def table_cell_e1(p):
    """Hangul inside a table cell, carried in Arial with no a:ea.

    Table text lives under a:tbl inside a graphicFrame, not in the shape's own text frame, so
    a walker that only visits shapes with text frames never sees it. The default Office theme
    ships an empty minorFont a:ea, so under the measured resolution model the run falls to
    a:latin = Arial, which has no Hangul glyphs.
    """
    s = p.slides.add_slide(p.slide_layouts[6])
    _tb(s, 0.8, 0.5, 8.0, 0.8, "Quarterly summary", size=24, latin="Arial")
    shape = s.shapes.add_table(2, 2, Inches(0.8), Inches(1.8), Inches(8.0), Inches(1.4))
    tbl = shape.table
    for (rr, cc), text in {(0, 0): "Region", (0, 1): "Revenue",
                           (1, 0): "매출 성장률", (1, 1): "+18%"}.items():
        cell = tbl.cell(rr, cc)
        para = cell.text_frame.paragraphs[0]
        _run(para, text, size=14, latin="Arial")      # no a:ea: the seeded defect


def field_run_e3(p):
    """A slide-number field rendered at 4pt.

    a:fld is a run whose text the renderer supplies. It is not an a:r, so a walker that
    iterates runs alone skips it, and an unreadable field slips through. 4pt is below the
    5.0pt hard floor.
    """
    s = p.slides.add_slide(p.slide_layouts[6])
    _tb(s, 0.8, 0.5, 8.0, 0.8, "Appendix", size=24, latin="Arial")
    box = s.shapes.add_textbox(Inches(12.2), Inches(6.9), Inches(0.8), Inches(0.3))
    para = box.text_frame.paragraphs[0]._p
    fld = para.makeelement(qn("a:fld"), {"id": "{1B7B1F2A-6E2E-4E5C-9C3E-4F1B2D3C4E5F}",
                                         "type": "slidenum"})
    rPr = fld.makeelement(qn("a:rPr"), {"lang": "en-US", "sz": "400"})
    fld.append(rPr)
    t = fld.makeelement(qn("a:t"), {})
    t.text = "2"
    fld.append(t)
    para.append(fld)


def line_break_e2(p):
    """An em dash used as punctuation, with an a:br between the halves of the sentence.

    E2 reads the whole paragraph rather than the run, because run splitting would otherwise
    hide the context that decides whether a dash is a range or a parenthetical. a:br is a
    visual line break inside that same paragraph: if the walker treats it as a paragraph
    boundary the two halves stop being neighbours and the dash reads as something else.
    """
    s = p.slides.add_slide(p.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.0), Inches(1.2))
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]
    _run(para, "핵심 지표는 전 분기 대비 개선", size=16, ea=HANGUL_FONT)
    para._p.append(para._p.makeelement(qn("a:br"), {}))
    _run(para, EM_DASH + "특히 구독 매출이 견인했습니다", size=16, ea=HANGUL_FONT)


def rotated_clean(p):
    """A negative: a rotated text shape must not manufacture geometry findings.

    Rotated text is out of scope by decision, not by oversight, because the glyph box model
    is axis-aligned and rotating it would report overlaps that do not exist. This deck holds
    two rotated labels whose axis-aligned bounding boxes intersect while the drawn text does
    not. Expected findings: none.
    """
    s = p.slides.add_slide(p.slide_layouts[6])
    for i, (x, y, rot) in enumerate(((4.0, 3.0, "2700000"), (4.6, 3.0, "18900000"))):
        box = _tb(s, x, y, 4.0, 0.5, "회전된 라벨 %d" % (i + 1), size=18, ea=HANGUL_FONT)
        box._element.find(qn("p:spPr")).find(qn("a:xfrm")).set("rot", rot)


AXES = (
    ("group_desync_w16", group_desync, {"W16": 1},
     "A group whose off/ext and chOff/chExt disagree, holding a run whose absolute position "
     "is 12.9in on a 13.333in canvas. Ground truth by construction and by reading the XML: "
     "the child's raw x is 0.5in, the group's off x is 12.4in and its chOff x is 0.5in, so "
     "the composed x is 12.4 + (0.5 - 0.5) * (3.2/3.2) = 12.4in and the 3.2in-wide frame "
     "ends past the canvas edge. W16 must fire. If the group transform is ignored the shape "
     "reads as sitting at 0.5in, comfortably inside, and nothing is reported. That is the "
     "exact shape of the grpSpPr namespace bug fixed in 0.5.0."),
    ("table_cell_e1", table_cell_e1, {"E1": 1},
     "Hangul in a table cell with a:latin=Arial and no a:ea. Ground truth by construction: "
     "the default Office theme's minorFont a:ea is empty, so under the measured resolution "
     "model (docs/CALIBRATION.md) the run falls through to a:latin, and Arial has no Hangul "
     "glyphs. Table text lives under a:tbl in a graphicFrame rather than in the shape's own "
     "text frame, so a walker that visits only shapes with text frames reports nothing."),
    ("field_run_e3", field_run_e3, {"E3": 1},
     "A slidenum a:fld with sz=400, i.e. 4.0pt, below the 5.0pt hard floor. Ground truth by "
     "construction: the size is written explicitly on the field's a:rPr, so no inheritance "
     "is involved. a:fld is not an a:r, so a run walker that matches on a:r alone never sees "
     "it and an unreadable field passes."),
    ("line_break_e2", line_break_e2, {"E2": 1},
     "An em dash used as sentence punctuation, with an a:br between the two halves of the "
     "paragraph. Ground truth by construction: the dash joins two Hangul clauses with no "
     "numeric neighbour on either side, which is the blocked case rather than the exempt "
     "range case. The a:br is the axis: it is a visual break inside one paragraph, and if it "
     "is read as a paragraph boundary the clauses stop being neighbours."),
    ("rotated_clean", rotated_clean, {},
     "A clean negative. Two rotated labels whose axis-aligned bounding boxes intersect while "
     "the drawn text does not. Rotated text is out of scope by decision (the glyph box model "
     "is axis-aligned), so the correct output is silence. Ground truth by construction: with "
     "rot stripped these same two boxes overlap and W15 fires, which is what makes this a "
     "test of the exclusion rather than of an empty deck."),
)


def main():
    os.makedirs(HERE, exist_ok=True)
    for name, build, expected, notes in AXES:
        path = os.path.join(HERE, name + ".pptx")
        p = _prs()
        build(p)
        from archforge.demo import _save
    # docProps says "archforge corpus", not the library name. Provenance lives in the
    # manifest's "generator" field, which is where the corpus README points and where a
    # reader looks; putting the library name in docProps as well only ships a
    # generated-by tell in every sdist for information already recorded.
        _save(p, path, application="archforge corpus")
        manifest = {"expected": expected,
                    "notes": notes,
                    "generator": "python-pptx",
                    "profile": "full",
                    "ground_truth": "by construction and by XML inspection, not by running "
                                    "the linter"}
        with open(os.path.join(HERE, name + ".json"), "w", encoding="utf-8") as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)
        print("wrote", path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
