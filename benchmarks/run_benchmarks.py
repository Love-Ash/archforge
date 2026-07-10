# -*- coding: utf-8 -*-
"""공개 벤치마크 하네스(0.4.0 시드, 3차 외부 리뷰: 규칙별 정확도의 재현 가능한 증명).

바이너리 pptx를 커밋하는 대신 결함 매트릭스에서 픽스처 덱을 생성하고, 설치된
archforge를 CLI로 돌려 규칙별 precision/recall을 계산해 출력한다. 기대값이 어긋나면
exit 1이라 CI 게이트로도 쓴다.

한계를 명시한다: 현재 생성기는 python-pptx 한 종이다. PowerPoint·LibreOffice·
Google Slides 산출물 변형은 이 하네스의 확장 슬롯(GENERATORS)에 추가하는 것이
로드맵이며, 실덱 50코퍼스는 사적 자료라 여기 포함되지 않는다(docs/CALIBRATION.md).

사용: python benchmarks/run_benchmarks.py [--keep DIR]
"""
import argparse
import json
import os
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt

EM = chr(0x2014)
EN = chr(0x2013)


def _tb(s, x, y, w, h, text, font=None, size=12, spc=None, no_size=False):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    if font:
        r.font.name = font
    if not no_size:
        r.font.size = Pt(size)
    if spc is not None:
        r._r.get_or_add_rPr().set("spc", str(spc))
    return box


def _new():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def _slide(p):
    return p.slides.add_slide(p.slide_layouts[6])


# (이름, 빌더, 기대 코드 집합, 프로파일) : 양성과 음성을 쌍으로 둔다
def _case_e1_fallback(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono")


def _case_e1_clean(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "정상 한글 본문", font="Wanted Sans", size=14)


def _case_e2_emdash(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "삽입구" + EM + "티", font="Wanted Sans")


def _case_e2_range_clean(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "FY2020" + EN + "2024 실적", font="Wanted Sans")


def _case_e3_tiny(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "tiny", font="Wanted Sans", size=4)


def _case_e4_tracking(p):
    _tb(_slide(p), 1, 1, 5, 0.5, "자간 벌어진 한글", font="Wanted Sans", spc=100)


def _case_w16_overflow(p):
    _tb(_slide(p), 12.5, 1, 4, 0.5, "overflowing long text body", font="Wanted Sans", size=18)


def _case_w15_overlap(p):
    s = _slide(p)
    _tb(s, 1, 2.0, 4, 1.2, "12,400원", font="Wanted Sans", size=44)
    _tb(s, 1.2, 2.15, 3, 0.5, "+11.7% 전년 대비", font="Wanted Sans", size=14)


def _case_clean_page(p):
    s = _slide(p)
    _tb(s, 1, 0.8, 9, 0.8, "Clean title", font="Wanted Sans", size=28)
    _tb(s, 1, 2.2, 10, 3, "Readable body.", font="Wanted Sans", size=13)


CASES = [
    ("e1_fallback", _case_e1_fallback, {"E1"}, "core"),
    ("e1_clean", _case_e1_clean, set(), "core"),
    ("e2_emdash", _case_e2_emdash, {"E2"}, "full"),
    ("e2_range_clean", _case_e2_range_clean, set(), "full"),
    ("e3_tiny", _case_e3_tiny, {"E3"}, "core"),
    ("e4_tracking", _case_e4_tracking, {"E4"}, "core"),
    ("w16_overflow", _case_w16_overflow, {"W16"}, "core"),
    ("w15_overlap", _case_w15_overlap, {"W15"}, "core"),
    ("clean_page", _case_clean_page, set(), "core"),
]

GENERATORS = {"python-pptx": lambda builder, path: _generate_pptx(builder, path)}


def _generate_pptx(builder, path):
    p = _new()
    builder(p)
    p.save(path)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--keep", default=None, help="픽스처를 이 폴더에 남김(디버그)")
    a = ap.parse_args()
    outdir = a.keep or tempfile.mkdtemp(prefix="archforge_bench_")
    os.makedirs(outdir, exist_ok=True)

    per_rule = {}
    failures = []
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = "en"
    for gen_name, gen in GENERATORS.items():
        for name, builder, expected, profile in CASES:
            path = os.path.join(outdir, "%s__%s.pptx" % (gen_name, name))
            gen(builder, path)
            r = subprocess.run([sys.executable, "-m", "archforge.lint", path,
                                "--json", "--profile", profile],
                               capture_output=True, text=True, encoding="utf-8",
                               stdin=subprocess.DEVNULL, env=env)
            doc = json.loads(r.stdout)
            got = {x["code"] for x in doc["errors"]} | {x["code"] for x in doc["warnings"]}
            got.discard("W18")
            for code in expected:
                st = per_rule.setdefault(code, {"tp": 0, "fn": 0, "fp": 0})
                if code in got:
                    st["tp"] += 1
                else:
                    st["fn"] += 1
                    failures.append("%s/%s: expected %s, missing (got %s)" % (gen_name, name, code, sorted(got)))
            for code in got - expected:
                st = per_rule.setdefault(code, {"tp": 0, "fn": 0, "fp": 0})
                st["fp"] += 1
                failures.append("%s/%s: unexpected %s" % (gen_name, name, code))

    print("rule  tp  fn  fp   precision  recall")
    for code in sorted(per_rule):
        st = per_rule[code]
        prec = st["tp"] / (st["tp"] + st["fp"]) if (st["tp"] + st["fp"]) else 1.0
        rec = st["tp"] / (st["tp"] + st["fn"]) if (st["tp"] + st["fn"]) else 1.0
        print("%-5s %3d %3d %3d   %8.2f  %6.2f" % (code, st["tp"], st["fn"], st["fp"], prec, rec))
    if failures:
        print("\nFAILURES (%d):" % len(failures))
        for f in failures:
            print(" -", f)
        sys.exit(1)
    print("\nbenchmark expectations: all green")


if __name__ == "__main__":
    main()
