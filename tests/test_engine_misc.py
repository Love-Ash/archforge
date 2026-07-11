# -*- coding: utf-8 -*-
"""Engine tests not covered by the themed buckets (split from test_lint.py, 0.8.x)."""
from helpers import *   # noqa: F401,F403
from helpers import (_by, _add_fld, _repo_root, _require_repo_file,
                     patch_theme_fonts)   # noqa: F401


def test_line_gates_positive(tmp_path):
    p = new_prs()
    s = add_slide(p)   # p1: E1 Hangul in a latin-only font
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    s = add_slide(p)   # p2: E2 long dash (fullwidth hyphen)
    tb(s, 1, 1, 5, 0.5, "dash" + FW_HYPHEN + "test", font="Wanted Sans", size=12)
    s = add_slide(p)   # p3: E3 explicit 4pt
    tb(s, 1, 1, 5, 0.5, "tiny text 4pt", font="Wanted Sans", size=4)
    s = add_slide(p)   # p4: E3 autofit 40pt * 10% = 4pt
    box = tb(s, 1, 1, 5, 1.0, "autofit shrink", font="Wanted Sans", size=40)
    set_autofit_scale(box, 10.0)
    s = add_slide(p)   # p5: E4 positive CJK tracking
    tb(s, 1, 1, 5, 0.5, "자간 벌어진 한글", font="Wanted Sans", size=12, spc=100)
    s = add_slide(p)   # p6: W1 wide frame body 8.5pt 40+ chars
    tb(s, 1, 1, 11, 0.6, "0123456789" * 5, font="Wanted Sans", size=8.5)
    s = add_slide(p)   # p7: size unspecified -> resolved to defaultTextStyle 18pt (not W5: inheritance resolution introduced)
    tb(s, 1, 1, 5, 0.5, "no size run", font="Wanted Sans", no_size=True)
    s = add_slide(p)   # p8: W8 narrow frame, small 6pt Hangul
    tb(s, 1, 1, 2, 0.4, "목업 안 작은 한글", font="Wanted Sans", size=6)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    ec = codes(errors)
    assert "E1" in ec and by_code(errors, "E1")[0][0] == 1
    assert "E2" in ec and by_code(errors, "E2")[0][0] == 2
    e3 = [si for (si, m, d) in by_code(errors, "E3")]
    assert 3 in e3 and 4 in e3
    assert "E4" in ec and by_code(errors, "E4")[0][0] == 5
    assert any(si == 6 for (si, m, d) in by_code(warns, "W1"))
    # p7: the previous version fired W5 (inheritance unresolved). With the
    # inheritance chain introduced, defaultTextStyle 18pt is resolved, so
    # neither W5 nor a false-fire of the size gate should occur.
    assert not any(si == 7 for (si, m, d) in by_code(warns, "W5"))
    assert not any(si == 7 for (si, m, d) in by_code(errors, "E3"))
    assert any(si == 8 for (si, m, d) in by_code(warns, "W8"))

def test_text_point_attr_unit():
    tpa = jl._text_point_attr
    assert tpa("150") == 150
    assert tpa("1.5pt") == 150
    assert tpa("0.5in") == 3600
    assert tpa("abc") is None
    assert tpa(None) is None

def test_partial_salvage_on_bad_frame(tmp_path):
    """Guard granularity regression (third adversarial panel confirmed, 2026-07-10):
    even if run1 in the same frame dies on a garbage attribute, run2's genuine E1
    violation must still be caught (per-run guard), and the swallowed section must
    be surfaced to JSON as W18. In the days of frame-level guards, run2's violation
    used to vanish entirely too, producing a false clean."""
    p = new_prs()
    s = add_slide(p)   # p1: one frame with a garbage run + a genuine E1 run
    box = tb(s, 1, 1, 5, 0.5, "쓰레기 크기 한글", font="Wanted Sans", no_size=True)
    r1 = box.text_frame.paragraphs[0].runs[0]
    r1._r.get_or_add_rPr().set("sz", "notanumber")
    r2 = box.text_frame.paragraphs[0].add_run()
    r2.text = "모노 폴백 한글"
    r2.font.name = "IBM Plex Mono"
    from pptx.util import Pt as _Pt
    r2.font.size = _Pt(12)
    s = add_slide(p)   # p2: a normal E1 defect (also pins survival across pages)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    e1_pages = [si for (si, m, d) in by_code(errors, "E1")]
    assert 1 in e1_pages, errors          # the neighboring run's violation in the same frame survives
    assert 2 in e1_pages, errors
    w18 = by_code(warns, "W18")
    assert any(si == 1 for (si, m, d) in w18), warns   # the swallowed section is surfaced in the output contract

def test_w5_when_chain_fully_absent(tmp_path):
    """If even defaultTextStyle is removed (simulating output from an external
    generator), an unresolved size is surfaced as W5."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "크기 미상 한글", font="Wanted Sans", no_size=True)
    pres_el = p.slides._sldIdLst.getparent()
    dts = pres_el.find(qn("p:defaultTextStyle"))
    assert dts is not None
    pres_el.remove(dts)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert by_code(warns, "W5"), warns

def test_geo_robustness_no_false_positives(tmp_path):
    """6 false-positive scenarios reproduced measured by adversarial verification:
    all must flag 0."""
    from PIL import Image, ImageDraw
    p = new_prs()
    s = add_slide(p)  # p1: wrap=none phantom wrap
    tb(s, 1.0, 1.0, 1.0, 0.4, "Quarterly revenue grew 18% on cloud momentum",
       font="Wanted Sans", size=12)
    tb(s, 1.0, 1.4, 2.0, 0.3, "Source: 10-K", font="Wanted Sans", size=9)
    s = add_slide(p)  # p2: group-move desync
    grp = s.shapes.add_group_shape()
    box = grp.shapes.add_textbox(Inches(12.5), Inches(3.0), Inches(1.5), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "group label"
    r.font.size = Pt(12)
    grp.left, grp.top = Inches(7.5), Inches(3.0)
    s = add_slide(p)  # p3: rotated picture (bbox before rotation is outside, render is inside)
    pic = s.shapes.add_picture(png(tmp_path, "rot.png", size=(400, 50)),
                               Inches(11.0), Inches(3.5), Inches(4.0), Inches(0.5))
    pic.rotation = 90
    s = add_slide(p)  # p4: flipH (mirrored ink lands on screen)
    pic = s.shapes.add_picture(png(tmp_path, "flip.png", opaque_box=(0.0, 0.0, 0.5, 1.0)),
                               Inches(-0.8), Inches(3.0), Inches(2.0), Inches(1.0))
    pic._element.spPr.find(qn("a:xfrm")).set("flipH", "1")
    s = add_slide(p)  # p5: P-mode + tRNS transparency
    im = Image.new("RGBA", (400, 400), (0, 0, 0, 0))
    ImageDraw.Draw(im).rectangle([150, 150, 250, 250], fill=(40, 60, 90, 255))
    pmode = os.path.join(str(tmp_path), "pmode.png")
    im.convert("P").save(pmode, transparency=0)
    s.shapes.add_picture(pmode, Inches(10.33), Inches(2.0), Inches(4.0), Inches(4.0))
    s = add_slide(p)  # p6: caption over a solid card over a photo
    s.shapes.add_picture(png(tmp_path, "photo.png"), Inches(4.0), Inches(0.5), Inches(6.0), Inches(6.5))
    rect(s, 6.0, 2.5, 4.0, 1.5, "FFFFFF")
    tb(s, 6.3, 3.0, 3.4, 0.4, "카드 위 캡션 문장", font="Wanted Sans", size=14)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    geo = [x for x in warns if x[1] in ("W15", "W16", "W17")]
    assert not geo, geo

def test_clean_deck_no_flags(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 0.8, 9, 0.8, "Clean title page", font="Wanted Sans", size=30)
    tb(s, 1, 2.2, 10, 3, "Readable body text at normal size.", font="Wanted Sans", size=13)
    s = add_slide(p)
    tb(s, 1, 0.8, 6, 0.6, "Second layout", font="Wanted Sans", size=22)
    rect(s, 7.5, 1.0, 4.5, 5.0, "20242B")
    tb(s, 1, 3.5, 5.5, 2.5, "Different grid here.", font="Wanted Sans", size=13)
    s = add_slide(p)
    tb(s, 4, 3.2, 6, 0.8, "Closing page", font="Wanted Sans", size=26)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert not errors and not warns, (codes(errors), codes(warns))

def test_finding_tuple_compat_and_lazy_message():
    """Finding's 4-tuple backward compatibility and locale neutrality (rendering the
    same finding in two languages)."""
    from archforge.findings import Finding
    f = Finding(3, "E2", "e2", (), "cp=U+2014")
    page, code, msg, detail = f
    assert (page, code, detail) == (3, "E2", "cp=U+2014") and f[1] == "E2" and len(f) == 4
    jmsg.set_lang("ko")
    ko = f.message
    jmsg.set_lang("en")
    en = f.message
    jmsg.set_lang("ko")
    assert ko != en and "Dash" in en


# ---------------------------------------------------------------- skill packaging

def test_examples_contract(tmp_path):
    """Pins the documented contract for examples/: style_warnings is clean under
    core, and W13/W14 under full."""
    from archforge import demo as jdemo
    p = os.path.join(str(tmp_path), "warn.pptx")
    jdemo.build_warnings(p)
    e, w = jl.lint(p, profile="core")
    assert not e and not w, (codes(e), codes(w))
    e, w = lint_full(p)
    assert not e, codes(e)
    assert {"W13", "W14"} <= set(codes(w)), codes(w)

def test_strict_split_flags(tmp_path):
    """--fail-on-warning fails a WARN-only deck; the default does not. --strict stays
    the union alias."""
    from archforge import demo as jdemo
    p = os.path.join(str(tmp_path), "warn.pptx")
    jdemo.build_warnings(p)
    r = run_cli([p, "--profile", "full"])
    assert r.returncode == 0, r.stdout + r.stderr
    r = run_cli([p, "--profile", "full", "--fail-on-warning"])
    assert r.returncode == 1
    r = run_cli([p, "--profile", "full", "--strict"])
    assert r.returncode == 1
    doc = json.loads(run_cli([p, "--profile", "full", "--fail-on-warning", "--json"]).stdout)
    assert doc["summary"]["pass"] is False
