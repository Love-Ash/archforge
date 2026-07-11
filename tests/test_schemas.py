# -*- coding: utf-8 -*-
"""Formal schema validation: the machine-validatable JSON Schemas under schemas/ are the
published contract, and every real output shape must validate against them (0.8,
external review section 5: declaring schema_version without a validatable schema file).

Skips cleanly when the schemas/ directory is absent (wheel install)."""
import json
import os
import subprocess
import sys

import pytest

try:
    import jsonschema
    from referencing import Registry, Resource
except ImportError:   # pragma: no cover
    jsonschema = None

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCHEMAS = os.path.join(ROOT, "schemas")

pytestmark = pytest.mark.skipif(
    jsonschema is None or not os.path.isdir(SCHEMAS),
    reason="jsonschema or schemas/ not available in this install")


def _schema(name):
    with open(os.path.join(SCHEMAS, name), encoding="utf-8") as f:
        return json.load(f)


def _validator(name):
    schema = _schema(name)
    # scan-2.0 $refs report-2.0 by relative id: register both
    registry = Registry().with_resources([
        ("report-2.0.schema.json",
         Resource.from_contents(_schema("report-2.0.schema.json"))),
    ])
    return jsonschema.Draft202012Validator(schema, registry=registry)


def _cli(args):
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = "en"
    env["PYTHONPATH"] = os.path.join(ROOT, "src") + os.pathsep + env.get("PYTHONPATH", "")
    return subprocess.run([sys.executable, "-m", "archforge"] + args,
                          capture_output=True, text=True, encoding="utf-8", env=env,
                          stdin=subprocess.DEVNULL)


def _make_deck(tmp_path, name="s.pptx"):
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    s = p.slides.add_slide(p.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(6.9), Inches(5), Inches(0.3))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "tiny source line"
    r.font.size = Pt(4)   # E3
    rPr = r._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": "맑은 고딕"}))
    path = os.path.join(str(tmp_path), name)
    p.save(path)
    return path


def test_report_10_validates(tmp_path):
    deck = _make_deck(tmp_path)
    doc = json.loads(_cli([deck, "--profile", "full", "--json"]).stdout)
    _validator("report-1.0.schema.json").validate(doc)


def test_report_20_validates(tmp_path):
    deck = _make_deck(tmp_path)
    doc = json.loads(_cli([deck, "--profile", "full", "--schema", "2", "--json"]).stdout)
    _validator("report-2.0.schema.json").validate(doc)


def test_scan_20_validates(tmp_path):
    d = os.path.join(str(tmp_path), "decks")
    os.makedirs(d)
    _make_deck(d, "a.pptx")
    with open(os.path.join(d, "corrupt.pptx"), "w") as f:
        f.write("nope")   # error entry must also validate
    doc = json.loads(_cli(["scan", d, "--schema", "2", "--json"]).stdout)
    _validator("scan-2.0.schema.json").validate(doc)


def test_baseline_3_validates(tmp_path):
    deck = _make_deck(tmp_path)
    bl = os.path.join(str(tmp_path), "bl.json")
    _cli([deck, "--write-baseline", bl])
    with open(bl, encoding="utf-8") as f:
        doc = json.load(f)
    _validator("baseline-3.schema.json").validate(doc)


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-q"]))
