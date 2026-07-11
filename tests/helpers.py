# -*- coding: utf-8 -*-
"""Shared fixtures-free helpers for the split test suite (0.8.x: the test monolith
mirrored the engine monolith). Deck builders, CLI runner, and patch utilities;
the autouse language fixture lives in conftest.py."""
# -*- coding: utf-8 -*-
"""archforge gate regression: fixes, per gate, both positive fixtures (does a pptx with
a planted defect get caught) and negative fixtures (does clean, intentional staging
pass without false positives).

Characters the linter blocks, such as the em dash, are not kept as literals in this
source file but built via chr(codepoint): the source must stay clean, and the
forbidden character must exist only inside the pptx under test for the gate test to
be valid.

The expected values in the E1 fixtures follow the PowerPoint COM render measured
model (2026-07-10, docs/CALIBRATION.md): run a:ea > theme minorFont a:ea (wins over
run a:latin when non-empty) > (only when the theme ea slot is empty) run a:latin >
OS fallback (Malgun).
"""
import os
import re
import json
import shutil
import subprocess
import sys
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

import archforge.lint as jl
from archforge import messages as jmsg


EN_DASH = chr(0x2013)
EM_DASH = chr(0x2014)
MINUS = chr(0x2212)
FW_HYPHEN = chr(0xFF0D)


# ---------------------------------------------------------------- helpers
def new_prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def add_slide(p):
    return p.slides.add_slide(p.slide_layouts[6])


def tb(s, x, y, w, h, text, font=None, size=12, color="222222", spc=None, no_size=False, ea=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    if font:
        r.font.name = font
    if ea:
        rPr = r._r.get_or_add_rPr()
        el = rPr.makeelement(qn("a:ea"), {"typeface": ea})
        rPr.append(el)
    if not no_size:
        r.font.size = Pt(size)
    r.font.color.rgb = RGBColor.from_string(color)
    if spc is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(spc))
    return box


def rect(s, x, y, w, h, fill_hex):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    sp.line.fill.background()
    return sp


def vconn(s, x, y0, y1, line_hex):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y0), Inches(x), Inches(y1))
    cn.line.color.rgb = RGBColor.from_string(line_hex)
    cn.line.width = Pt(2)
    return cn


def set_autofit_scale(box, pct):
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    na = bodyPr.makeelement(qn("a:normAutofit"), {"fontScale": str(int(pct * 1000))})
    bodyPr.append(na)


def png(d, name, size=(400, 300), opaque_box=None):
    """A fully opaque RGB PNG, or an RGBA PNG opaque only in the opaque_box=(l,t,r,b
    ratio) region."""
    from PIL import Image, ImageDraw
    path = os.path.join(str(d), name)
    if opaque_box is None:
        Image.new("RGB", size, (40, 60, 90)).save(path)
    else:
        im = Image.new("RGBA", size, (0, 0, 0, 0))
        dr = ImageDraw.Draw(im)
        l, t, r, b = opaque_box
        dr.rectangle([size[0] * l, size[1] * t, size[0] * r, size[1] * b], fill=(40, 60, 90, 255))
        im.save(path)
    return path


def lint_full(*args, **kw):
    """In 0.4.0 the default profile changed to core, so existing fixtures that assume
    style rules too must specify full explicitly (the default value itself is
    verified by test_default_profile_core)."""
    kw.setdefault("profile", "full")
    return jl.lint(*args, **kw)


def codes(items):
    return [c for (_si, c, _m, _d) in items]


def by_code(items, code):
    return [(si, m, d) for (si, c, m, d) in items if c == code]


def save(p, d, name):
    path = os.path.join(str(d), name)
    p.save(path)
    return path


def patch_theme_ea(path, typeface):
    """Replace the theme a:ea empty slot in a saved pptx with typeface (rewrites the
    zip). For the default template with no multi-master: the only reliably stable
    way to build an E1 theme-branch fixture."""
    tmp = path + ".patched.pptx"
    new = ('<a:ea typeface="%s"/>' % typeface).encode("utf-8")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if re.search(r"theme/theme\d+\.xml$", item.filename):
                data = data.replace(b'<a:ea typeface=""/>', new)
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


def run_cli(args, lang="ko"):
    """CLI run helper. stdin=DEVNULL: avoids Windows handle-inheritance failure
    (WinError 6) under pytest capture. Defaults to ko (assertions use Korean
    phrasing); verify English with lang="en"."""
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = lang
    return subprocess.run([sys.executable, "-m", "archforge.lint"] + args,
                          capture_output=True, text=True, encoding="utf-8",
                          stdin=subprocess.DEVNULL, env=env)


# ---------------------------------------------------------------- line-level gates


def patch_theme_fonts(path, major_ea=None, minor_ea=None):
    """Patch the theme's major/minor ea empty slots to different values respectively
    (rewrites the zip)."""
    tmp = path + ".patched.pptx"
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if re.search(r"theme/theme\d+\.xml$", item.filename):
                if major_ea is not None:
                    m = re.search(rb"<a:majorFont>.*?</a:majorFont>", data, re.S)
                    seg = m.group(0).replace(b'<a:ea typeface=""/>',
                                             ('<a:ea typeface="%s"/>' % major_ea).encode("utf-8"))
                    data = data[:m.start()] + seg + data[m.end():]
                if minor_ea is not None:
                    m = re.search(rb"<a:minorFont>.*?</a:minorFont>", data, re.S)
                    seg = m.group(0).replace(b'<a:ea typeface=""/>',
                                             ('<a:ea typeface="%s"/>' % minor_ea).encode("utf-8"))
                    data = data[:m.start()] + seg + data[m.end():]
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path

def _repo_root():
    return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def _require_repo_file(*parts):
    """Skip (not fail) when a repo-root asset is absent: an installed wheel ships only the
    package, not repo-root skills/ or action_runner.py, so those tests apply to the repo
    and the sdist, not the wheel (0.7.1, external review P0-5)."""
    path = os.path.join(_repo_root(), *parts)
    if not os.path.exists(path):
        pytest.skip("repo-root asset not present in this install: %s" % os.path.join(*parts))
    return path

def _by(items, code):
    return [f for f in items if f.code == code]

def _add_fld(para, text, sz=None, latin=None, spc=None):
    """Plant an a:fld (auto field) directly into the paragraph. python-pptx runs
    skip this element, so it's done via raw XML."""
    fld = para._p.makeelement(qn("a:fld"), {
        "id": "{93C0AD05-B0B5-4E56-9A2A-9F2113D1B94A}", "type": "slidenum"})
    attrs = {}
    if sz is not None:
        attrs["sz"] = str(sz)
    if spc is not None:
        attrs["spc"] = str(spc)
    rPr = fld.makeelement(qn("a:rPr"), attrs)
    if latin:
        rPr.append(rPr.makeelement(qn("a:latin"), {"typeface": latin}))
    fld.append(rPr)
    t = fld.makeelement(qn("a:t"), {})
    t.text = text
    fld.append(t)
    para._p.append(fld)
    return fld
