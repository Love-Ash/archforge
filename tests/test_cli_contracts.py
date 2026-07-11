# -*- coding: utf-8 -*-
"""CLI, scan, reporters, and machine-contract tests (split from test_lint.py, 0.8.x)."""
from helpers import *   # noqa: F401,F403
from helpers import (_by, _add_fld, _repo_root, _require_repo_file,
                     patch_theme_fonts)   # noqa: F401


def test_cli_json_and_exit_codes(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    bad = save(p, tmp_path, "bad.pptx")
    r = run_cli([bad, "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 1
    assert doc["summary"]["error_count"] >= 1 and not doc["summary"]["pass"]
    assert any(e["code"] == "E1" for e in doc["errors"])
    assert "ghost" in doc and isinstance(doc["ghost"], list)   # documented JSON contract

    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.6, "clean page", font="Wanted Sans", size=20)
    good = save(p, tmp_path, "good.pptx")
    r = run_cli([good, "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 0 and doc["summary"]["pass"]

def test_cli_exit2_missing_and_corrupt(tmp_path):
    """exit 2 path regression (untested in the old version): missing file and a
    corrupt pptx."""
    r = run_cli([os.path.join(str(tmp_path), "no_such.pptx")])
    assert r.returncode == 2
    assert "찾을 수 없습니다" in r.stderr
    corrupt = os.path.join(str(tmp_path), "corrupt.pptx")
    with open(corrupt, "wb") as f:
        f.write(b"this is not a pptx file at all")
    r = run_cli([corrupt])
    assert r.returncode == 2
    assert "열 수 없습니다" in r.stderr

def test_cli_strict_gates(tmp_path):
    """--strict: a deck with only WARNs switches to exit 1, and the E2 numeric
    context exception is lifted."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 11, 0.6, "0123456789" * 5, font="Wanted Sans", size=8.5)   # W1 only
    warn_only = save(p, tmp_path, "warn.pptx")
    assert run_cli([warn_only]).returncode == 0
    assert run_cli([warn_only, "--strict"]).returncode == 1

    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 6, 0.5, "FY2020" + EN_DASH + "2024 실적", font="Wanted Sans", size=14)
    ranged = save(p, tmp_path, "range.pptx")
    assert run_cli([ranged, "--profile", "full"]).returncode == 0
    r = run_cli([ranged, "--strict", "--profile", "full", "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 1
    assert any(e["code"] == "E2" for e in doc["errors"])

def test_cli_lang_edge_cases(tmp_path):
    """--lang CLI edge cases (third panel): repeated flags let the last one win, and
    it composes safely with the skill subcommand."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json", "--lang", "en", "--lang", "ko"]).stdout)
    assert doc["lang"] == "ko"
    r = run_cli(["skill", "--lang", "ko", "--path"])
    assert r.returncode == 0 and "SKILL.md" in r.stdout
    r = run_cli(["--lang", "ko", "skill", "--path"])   # leading flag + subcommand
    assert r.returncode == 0 and "SKILL.md" in r.stdout

def test_lang_english_output(tmp_path):
    """0.3.0 i18n: in an English environment, messages are in English and lang is
    recorded in the JSON. Codes are language-independent."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json"], lang="en").stdout)
    assert doc["lang"] == "en"
    e1 = [e for e in doc["errors"] if e["code"] == "E1"][0]
    assert "Hangul" in e1["message"] and "Malgun" in e1["message"]
    doc_ko = json.loads(run_cli([path, "--json"], lang="ko").stdout)
    assert doc_ko["lang"] == "ko" and "한글" in [e for e in doc_ko["errors"] if e["code"] == "E1"][0]["message"]
    # the --lang flag beats the environment variable
    doc_flag = json.loads(run_cli([path, "--json", "--lang", "en"], lang="ko").stdout)
    assert doc_flag["lang"] == "en"

def test_lang_catalog_consistency():
    """Catalog discipline: every entry has both ko/en, and the % format-specifier
    order is identical."""
    fmt = re.compile(r"%[-#0-9.]*[sdfr%]")
    for mid, entry in jmsg.MESSAGES.items():
        assert set(entry) == {"ko", "en"}, mid
        assert fmt.findall(entry["ko"]) == fmt.findall(entry["en"]), mid

def test_json_schema_contract(tmp_path):
    """Start of the JSON version contract (third review): schema_version, tool,
    target_renderer."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    doc = json.loads(run_cli([save(p, tmp_path, "fx.pptx"), "--json"]).stdout)
    assert doc["schema_version"] == "1.0"
    assert doc["tool"]["name"] == "archforge" and doc["tool"]["version"]
    assert doc["target_renderer"] == "powerpoint-windows"

def test_sarif_output(tmp_path):
    """SARIF 2.1.0 minimal contract: version, rules, results, ruleId, level."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    deck = save(p, tmp_path, "fx.pptx")
    out = os.path.join(str(tmp_path), "out.sarif")
    run_cli([deck, "--json", "--sarif", out])
    with open(out, encoding="utf-8") as f:
        doc = json.load(f)
    assert doc["version"] == "2.1.0"
    run0 = doc["runs"][0]
    assert run0["tool"]["driver"]["name"] == "archforge"
    res = run0["results"]
    assert res and res[0]["ruleId"] == "E1" and res[0]["level"] == "error"
    assert any(r["id"] == "E1" for r in run0["tool"]["driver"]["rules"])

def test_finding_location_payload(tmp_path):
    """Structured location (third review): shape_id, bbox, part, paragraph, run are
    carried in the JSON location."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    doc = json.loads(run_cli([save(p, tmp_path, "fx.pptx"), "--json"]).stdout)
    e1 = [e for e in doc["errors"] if e["code"] == "E1"][0]
    loc = e1["location"]
    assert isinstance(loc["shape_id"], int)
    assert loc["part"].endswith("slide1.xml")
    assert len(loc["bbox"]) == 4 and abs(loc["bbox"][0] - 1.0) < 0.01
    assert loc["paragraph"] == 0 and loc["run"] == 0

def test_skill_pack_sync():
    """The repo-root skills/ (for discovery) and the package-bundled copy (the
    canonical one) must not drift apart."""
    a = _require_repo_file("skills", "archforge-pptx-lint", "SKILL.md")
    b = os.path.join(_repo_root(), "src", "archforge", "skills",
                     "archforge-pptx-lint", "SKILL.md")
    with open(a, "rb") as fa, open(b, "rb") as fb:
        assert fa.read() == fb.read(), "루트 skills/ 사본과 패키지 정본이 다릅니다"

def test_skill_frontmatter_name_matches_dir():
    """Agent Skills spec: frontmatter name == the skill directory name (regression
    for an external review finding)."""
    path = _require_repo_file("skills", "archforge-pptx-lint", "SKILL.md")
    with open(path, encoding="utf-8") as f:
        head = f.read(500)
    m = re.search(r"^name:\s*(\S+)", head, re.M)
    assert m and m.group(1) == "archforge-pptx-lint"

def test_cli_skill_subcommand(tmp_path):
    """archforge skill: both the output and install paths match the package-bundled
    copy."""
    r = run_cli(["skill"])
    assert r.returncode == 0
    assert "name: archforge-pptx-lint" in r.stdout
    r = run_cli(["skill", "--install", str(tmp_path)])
    assert r.returncode == 0
    dst = os.path.join(str(tmp_path), "archforge-pptx-lint", "SKILL.md")
    assert os.path.exists(dst)
    with open(dst, encoding="utf-8") as f:
        installed = f.read()
    assert "name: archforge-pptx-lint" in installed


# ---------------------------------------------------------------- 0.5.0: loc reinforcement + a:fld/a:br

def test_table_cell_loc(tmp_path):
    """A table-cell finding's loc.cell == [row, col], 0-based (carried over from
    the fourth review)."""
    p = new_prs()
    s = add_slide(p)
    gfx = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    cell = gfx.table.cell(1, 0)
    cell.text = "한글"
    r = cell.text_frame.paragraphs[0].runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(12)
    errors, _w = lint_full(save(p, tmp_path, "tbl_loc.pptx"))
    e1 = _by(errors, "E1")
    assert e1, codes(errors)
    assert e1[0].loc.get("cell") == [1, 0], e1[0].loc

def test_cli_demo_and_scan(tmp_path):
    """demo: generates broken (fires 6 kinds of defects) + fixed (clean) and lints
    them on the spot, exit 0.
    scan: recursive directory aggregate JSON, exit 1 if even one fails, exit 2 if
    there are 0 matches."""
    d = os.path.join(str(tmp_path), "demo")
    r = run_cli(["demo", "--dir", d])
    assert r.returncode == 0, r.stderr
    assert os.path.exists(os.path.join(d, "broken.pptx"))
    assert os.path.exists(os.path.join(d, "fixed.pptx"))
    assert "ERROR 0, WARN 0" in r.stdout   # fixed must be clean (demo contract)

    r = run_cli(["scan", d, "--profile", "full", "--json"])
    assert r.returncode == 1, r.stderr
    doc = json.loads(r.stdout)
    assert doc["summary"]["file_count"] == 2
    assert doc["summary"]["failed_files"] == 1
    assert doc["summary"]["pass"] is False
    names = {os.path.basename(f["file"]) for f in doc["files"]}
    assert names == {"broken.pptx", "fixed.pptx"}
    for fdoc in doc["files"]:
        assert fdoc["summary"]["profile"] == "full"

    # fixed only: exit 0 + a text summary line
    r = run_cli(["scan", os.path.join(d, "fixed.pptx"), "--profile", "full"])
    assert r.returncode == 0, r.stdout + r.stderr
    assert "1" in r.stdout and "0" in r.stdout   # scan summary (1 file, 0 failed)

    # 0 glob matches = exit 2, not a silent pass (prevents a CI footgun)
    r = run_cli(["scan", os.path.join(str(tmp_path), "none_*.pptx")])
    assert r.returncode == 2

def test_demo_en_variant(tmp_path):
    """English-edition demo deck contract: monolingual English (user decision), so it
    seeds only the script-independent defects. broken_en is exactly E2+E3 plus W15/W16;
    the Hangul-only defects (E1/E4) are the Korean deck's job. fixed_en is clean under
    full. The README(en) assets are this deck's actual renders."""
    from archforge import demo as jdemo
    b = os.path.join(str(tmp_path), "b.pptx")
    fx = os.path.join(str(tmp_path), "f.pptx")
    jdemo.build_broken(b, lang="en")
    jdemo.build_fixed(fx, lang="en")
    e, w = lint_full(b)
    assert set(codes(e)) == {"E2", "E3"}, codes(e)
    assert {"W15", "W16"} <= set(codes(w)), codes(w)
    e, w = lint_full(fx)
    assert not e and not w, (codes(e), codes(w))

def test_cli_scan_sarif_multi(tmp_path):
    """scan --sarif: multiple files are merged into one SARIF run, each with its
    own per-file artifactLocation."""
    d = os.path.join(str(tmp_path), "demo")
    r = run_cli(["demo", "--dir", d])
    assert r.returncode == 0, r.stderr
    sarif_path = os.path.join(str(tmp_path), "out.sarif")
    r = run_cli(["scan", d, "--profile", "full", "--sarif", sarif_path])
    assert r.returncode == 1
    with open(sarif_path, encoding="utf-8") as f:
        doc = json.load(f)
    uris = {res["locations"][0]["physicalLocation"]["artifactLocation"]["uri"]
            for res in doc["runs"][0]["results"]}
    assert any(u.endswith("broken.pptx") for u in uris), uris


# ---------------------------------------------------------------- 0.6.0: hardening batch

def test_scan_isolates_broken_files(tmp_path):
    """One corrupt deck must not abort the batch (external review P0): it becomes a
    per-file error entry, the rest is still checked, and the aggregate JSON survives."""
    d = os.path.join(str(tmp_path), "batch")
    os.makedirs(d)
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "clean deck", size=14)
    save(p, d, "ok.pptx")
    with open(os.path.join(d, "corrupt.pptx"), "w") as f:
        f.write("not a zip")
    r = run_cli(["scan", d, "--json"])
    assert r.returncode == 1, r.stderr
    doc = json.loads(r.stdout)
    assert doc["summary"]["file_count"] == 2
    assert doc["summary"]["error_files"] == 1
    statuses = {os.path.basename(f["file"]): f["status"] for f in doc["files"]}
    assert statuses["corrupt.pptx"] == "error"
    assert statuses["ok.pptx"] == "pass"

def test_version_flag():
    r = run_cli(["--version"])
    assert r.returncode == 0
    assert "archforge" in r.stdout

def test_zip_preflight_blocks_bomb(tmp_path):
    """A pptx-shaped decompression bomb is a usage error with a stated reason, not a
    hang (external review: no package-level budgets existed)."""
    import zipfile
    bomb = os.path.join(str(tmp_path), "bomb.pptx")
    with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("ppt/huge.bin", b"\x00" * 30_000_000)
    r = run_cli([bomb])
    assert r.returncode == 2
    assert "preflight" in (r.stderr or "")

def test_sarif_rule_metadata_static(tmp_path):
    """SARIF rule titles must be static (no %-placeholders) and results carry
    partialFingerprints for cross-run tracking (external review)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "tiny", size=3)
    deck = save(p, tmp_path, "s.pptx")
    out = os.path.join(str(tmp_path), "o.sarif")
    r = run_cli([deck, "--sarif", out])
    assert r.returncode == 1
    with open(out, encoding="utf-8") as f:
        doc = json.load(f)
    run0 = doc["runs"][0]
    for rule in run0["tool"]["driver"]["rules"]:
        assert "%" not in rule["shortDescription"]["text"], rule
        assert rule["helpUri"].startswith("https://")
    assert all("partialFingerprints" in res for res in run0["results"])

def test_sarif_v3_fingerprint_and_related(tmp_path):
    """0.7.1: the SARIF cross-run identity matches the baseline v3 structural fingerprint,
    and pair findings (W15/W17) use the SARIF-standard relatedLocations."""
    import archforge.lint as _jl
    # E3 finding: fingerprint parity with baseline v3
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "tiny", size=3)          # E3
    # a W15 overlap pair for relatedLocations
    tb(s, 2, 2, 5, 1, "overlap sentence one here", size=24)
    tb(s, 2.1, 2.05, 5, 1, "overlap sentence two here", size=24)
    deck = save(p, tmp_path, "s.sarif.pptx")
    out = os.path.join(str(tmp_path), "o.sarif")
    run_cli([deck, "--profile", "full", "--sarif", out])
    with open(out, encoding="utf-8") as f:
        doc = json.load(f)
    results = doc["runs"][0]["results"]
    for res in results:
        assert "archforgeFinding/v3" in res["partialFingerprints"]
    w15 = [r for r in results if r["ruleId"] == "W15"]
    assert w15 and "relatedLocations" in w15[0], "W15 must carry relatedLocations"
    # parity: recompute the E3 finding's structural_fp and match the SARIF v3 print
    errors, warns = _jl.lint(deck, profile="full")
    e3 = [f for f in errors if f.code == "E3"][0]
    e3res = [r for r in results if r["ruleId"] == "E3"][0]
    assert e3res["partialFingerprints"]["archforgeFinding/v3"] == e3.structural_fp()


# ---------------------------------------------------------------- 0.6.1: contract batch

def test_scan_empty_pattern_fails(tmp_path):
    """One input matching nothing must not hide behind another that matched: exit 2
    unless --allow-empty-pattern (external review P0)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "ok.pptx")
    ghost_glob = os.path.join(str(tmp_path), "nope", "**", "*.pptx")
    r = run_cli(["scan", deck, ghost_glob])
    assert r.returncode == 2, (r.returncode, r.stderr)
    r = run_cli(["scan", deck, ghost_glob, "--allow-empty-pattern"])
    assert r.returncode == 0, (r.returncode, r.stderr)
    r = run_cli(["scan", deck, "--json"])
    doc = json.loads(r.stdout)
    assert doc["scan"]["inputs"][0]["matches"] == 1

def test_scan_global_usage_error_exits_2(tmp_path):
    """A bad CLI flag is a global usage error (exit 2), not N per-file error entries
    (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "g.pptx")
    for extra in (["--skip", "W999"], ["--skip", "E1"],
                  ["--config", os.path.join(str(tmp_path), "missing.json")],
                  ["--hard-min", "nan"]):
        r = run_cli(["scan", deck] + extra)
        assert r.returncode == 2, (extra, r.returncode, r.stdout, r.stderr)

def test_summary_policy_recorded(tmp_path):
    """The active failure policy travels with the verdict so JSON consumers can tell
    why identical counts pass or fail (external review P0)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "pol.pptx")
    doc = json.loads(run_cli([deck, "--json", "--fail-on-warning"]).stdout)
    assert doc["summary"]["policy"] == {"fail_on_warning": True,
                                        "fail_incomplete": False,
                                        "e2_no_exemptions": False}
    doc = json.loads(run_cli([deck, "--json", "--strict"]).stdout)
    assert all(doc["summary"]["policy"].values())

def test_rules_and_explain_subcommands():
    r = run_cli(["rules"])
    assert r.returncode == 0 and "E1" in r.stdout and "W18" in r.stdout
    r = run_cli(["explain", "w15", "--lang", "en"])
    assert r.returncode == 0 and "overlap" in r.stdout.lower()
    r = run_cli(["explain", "W99"])
    assert r.returncode == 2

def test_module_entry_no_runpy_warning(tmp_path):
    """`python -m archforge` must not emit the runpy RuntimeWarning that
    `python -m archforge.lint` triggers (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "m.pptx")
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = "en"
    r = subprocess.run([sys.executable, "-m", "archforge", deck], capture_output=True,
                       text=True, encoding="utf-8", stdin=subprocess.DEVNULL, env=env)
    assert r.returncode == 0, r.stderr
    assert "RuntimeWarning" not in (r.stderr or "")

def test_lint_subcommand_alias(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "alias.pptx")
    r = run_cli(["lint", deck, "--json"])
    assert r.returncode == 0, r.stderr
    assert json.loads(r.stdout)["summary"]["pass"] is True

def test_action_runner_rejects_bool_typo(tmp_path):
    """A typo'd boolean input must fail the job, not silently disable a safety default
    (external review P0)."""
    runner = _require_repo_file("action_runner.py")
    env = dict(os.environ)
    env.update({"AF_FILES": "whatever.pptx", "AF_FAIL_INCOMPLETE": "ture"})
    r = subprocess.run([sys.executable, runner], capture_output=True, text=True,
                       encoding="utf-8", env=env, stdin=subprocess.DEVNULL)
    assert r.returncode == 2
    assert "fail-incomplete" in (r.stderr or "")

def test_timeout_flag(tmp_path):
    """--timeout runs in a child process and returns a normal result under budget; a
    zero/negative value is a usage error. (The hang path is not exercised here to keep
    the suite fast; the mechanism is subprocess.call(timeout=...).)"""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "t.pptx")
    r = run_cli([deck, "--timeout", "60", "--json"])
    assert r.returncode == 0, r.stderr
    assert json.loads(r.stdout)["summary"]["pass"] is True
    r = run_cli([deck, "--timeout", "0"])
    assert r.returncode == 2
    # 0.7.1: inf passes ">0" and would disable the timeout; nan fails comparisons. Both
    # must be controlled usage errors, not a silently-unbounded run.
    for bad in ("inf", "nan", "-inf"):
        r = run_cli([deck, "--timeout", bad])
        assert r.returncode == 2, (bad, r.returncode)

def test_output_path_missing_dir_is_controlled(tmp_path):
    """0.7.1: a --sarif/--junit/--write-baseline target under a nonexistent directory is
    a controlled exit 2, not a traceback after linting."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "o.pptx")
    ghost = os.path.join(str(tmp_path), "nope", "out")
    for flag, ext in (("--sarif", ".sarif"), ("--junit", ".xml"),
                      ("--write-baseline", ".json")):
        r = run_cli([deck, flag, ghost + ext])
        assert r.returncode == 2, (flag, r.returncode, r.stderr)
        assert "does not exist" in (r.stderr or "") or "없습니다" in (r.stderr or "")
    # scan mode too
    d = os.path.join(str(tmp_path), "decks")
    os.makedirs(d)
    p2 = new_prs()
    s2 = add_slide(p2)
    tb(s2, 1, 1, 4, 1, "x", size=14)
    save(p2, d, "a.pptx")
    r = run_cli(["scan", d, "--sarif", ghost + ".sarif"])
    assert r.returncode == 2, (r.returncode, r.stderr)

def test_junit_reporter(tmp_path):
    """JUnit mapping contract: a testcase per executed rule (skipped = excluded by
    profile), ERROR findings as <failure>, WARN findings in <system-out> unless a
    warn-failing policy is active, and a corrupt file as an <error> testcase."""
    import xml.etree.ElementTree as ET
    d = os.path.join(str(tmp_path), "batch")
    os.makedirs(d)
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "tiny", size=3)   # E3
    save(p, d, "bad.pptx")
    with open(os.path.join(d, "corrupt.pptx"), "w") as f:
        f.write("nope")
    out = os.path.join(str(tmp_path), "o.xml")
    r = run_cli(["scan", d, "--junit", out])   # core profile: E2 etc. excluded
    assert r.returncode == 1
    tree = ET.parse(out)
    suites = tree.getroot().findall("testsuite")
    assert len(suites) == 2
    by_name = {sx.get("name"): sx for sx in suites}
    bad = by_name[os.path.join(d, "bad.pptx")]
    cases = {c.get("name").split()[0]: c for c in bad.findall("testcase")}
    assert cases["E3"].find("failure") is not None
    assert cases["E2"].find("skipped") is not None   # excluded by core
    corrupt = by_name[os.path.join(d, "corrupt.pptx")]
    assert corrupt.find("testcase/error") is not None
    # single-file mode + warn-fail policy: WARNs become failures
    p2 = new_prs()
    s2 = add_slide(p2)
    tb(s2, 12.8, 3.0, 3.0, 0.5, "off the canvas edge text", size=18)   # W16
    deck2 = save(p2, tmp_path, "warn.pptx")
    out2 = os.path.join(str(tmp_path), "o2.xml")
    run_cli([deck2, "--junit", out2])
    root2 = ET.parse(out2).getroot()
    w16 = [c for c in root2.iter("testcase") if c.get("name").startswith("W16")][0]
    assert w16.find("failure") is None and w16.find("system-out") is not None
    run_cli([deck2, "--junit", out2, "--fail-on-warning"])
    root3 = ET.parse(out2).getroot()
    w16b = [c for c in root3.iter("testcase") if c.get("name").startswith("W16")][0]
    assert w16b.find("failure") is not None

def test_scan_schema2_aggregate_root(tmp_path):
    """0.7.1 P0: scan --schema 2 must not declare schema 1.0 at the root over 2.0 file
    objects. The root is a distinct scan document type carrying file_schema_version."""
    d = os.path.join(str(tmp_path), "decks")
    os.makedirs(d)
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    save(p, d, "a.pptx")
    agg = json.loads(run_cli(["scan", d, "--json", "--schema", "2"]).stdout)
    assert agg["schema_version"] == "scan-2.0"
    assert agg["kind"] == "scan-report"
    assert agg["file_schema_version"] == "2.0"
    assert agg["files"][0]["schema_version"] == "2.0"
    agg1 = json.loads(run_cli(["scan", d, "--json"]).stdout)
    assert agg1["schema_version"] == "scan-1.0" and agg1["file_schema_version"] == "1.0"


# ---------------------------------------------------------------- 0.7: schema 2.0 + baseline v3

def test_schema_2_shape(tmp_path):
    """--schema 2: single findings[] with severity + structured data, plus capabilities
    and abstentions; the verdict (which codes) matches schema 1."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 6.9, 5, 0.3, "src", size=4, ea="맑은 고딕")   # E3
    tb(s, 1, 1, 8, 0.6, "핵심 개선" + EM_DASH + "특히", size=14, ea="맑은 고딕")  # E2
    deck = save(p, tmp_path, "s2.pptx")
    v1 = json.loads(run_cli([deck, "--profile", "full", "--json"]).stdout)
    v2 = json.loads(run_cli([deck, "--profile", "full", "--schema", "2", "--json"]).stdout)
    assert v2["schema_version"] == "2.0"
    assert "findings" in v2 and "errors" not in v2
    v1codes = sorted(f["code"] for f in v1["errors"] + v1["warnings"])
    v2codes = sorted(f["code"] for f in v2["findings"])
    assert v1codes == v2codes                      # verdict-identical across schemas
    e3 = [f for f in v2["findings"] if f["code"] == "E3"][0]
    assert e3["severity"] == "error"
    assert e3["data"]["effective_pt"] == 4.0 and e3["data"]["hard_min_pt"] == 5.0
    assert set(v2["capabilities"]) == {"typography", "geometry", "structure",
                                       "render_contrast"}

def test_schema_2_invocation_and_rules(tmp_path):
    """0.7.1: schema 2.0 records the invocation (profile/policy/config/thresholds) and a
    rules split (executed / profile_excluded / user_suppressed), so a consumer can tell a
    rule that ran-and-passed from one the profile never ran (external review section 5)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "inv.pptx")
    v2 = json.loads(run_cli([deck, "--profile", "core", "--skip", "W14",
                             "--schema", "2", "--json"]).stdout)
    inv = v2["invocation"]
    assert inv["profile"] == "core"
    assert inv["policy"]["fail_incomplete"] is False
    assert inv["thresholds"]["hard_min"] == 5.0
    rules = v2["rules"]
    assert "E1" in rules["executed"] and "W14" in rules["user_suppressed"]
    assert "E2" in rules["profile_excluded"]        # core excludes the AI-tell rules
    assert "W14" not in rules["executed"]

def test_typed_finding_data(tmp_path):
    """0.8: detection sites emit typed data directly. E1 carries script/effective_font/
    font_source/fallback, E2 the offending code points, E3 nominal+autofit scale, E4
    unit-explicit tracking, W16 the target kind (external review: E1/E2 had empty data,
    E4's unit was ambiguous, W16 lacked kind)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "아리알 한글", font="Arial", size=12)                 # E1
    tb(s, 1, 2, 8, 0.5, "구조적" + EM_DASH + "개선", size=14, ea="맑은 고딕")  # E2
    tb(s, 1, 3, 5, 0.3, "tiny src", size=4, ea="맑은 고딕")                    # E3
    tb(s, 1, 4, 5, 0.5, "자간 벌어진 한글", size=14, spc=300, ea="맑은 고딕")  # E4
    tb(s, 12.8, 5, 3, 0.5, "off canvas text run", size=18, ea="맑은 고딕")     # W16 text
    deck = save(p, tmp_path, "typed.pptx")
    v2 = json.loads(run_cli([deck, "--profile", "full", "--schema", "2", "--json"]).stdout)
    by = {}
    for f in v2["findings"]:
        by.setdefault(f["code"], []).append(f)
    e1 = by["E1"][0]["data"]
    assert e1["script"] == "hangul" and e1["effective_font"].lower() == "arial"
    assert e1["font_source"] in ("run.latin", "run.ea", "theme.ea")
    assert e1["fallback_font"] == "Malgun Gothic"
    e2 = by["E2"][0]["data"]
    assert e2["characters"] == ["U+2014"] and e2["function"] == "sentence_punctuation"
    e3 = by["E3"][0]["data"]
    assert e3["effective_pt"] == 4.0 and e3["nominal_pt"] == 4.0 and e3["autofit_scale"] == 1.0
    e4 = by["E4"][0]["data"]
    assert e4["tracking_raw_hundredths_pt"] == 300 and e4["tracking_pt"] == 3.0
    w16 = by["W16"][0]["data"]
    assert w16["kind"] == "text" and w16["overflow_in"] > 0

def test_schema_2_abstentions(tmp_path):
    """A vertical-text frame abstains: schema 2 lists it in abstentions[] with the
    affected rules and marks geometry partial, alongside the W18 finding."""
    p = new_prs()
    s = add_slide(p)
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "세로쓰기"
    r.font.size = Pt(18)
    r._r.get_or_add_rPr().append(box._element.makeelement(qn("a:ea"), {"typeface": "맑은 고딕"}))
    box.text_frame._txBody.find(qn("a:bodyPr")).set("vert", "eaVert")
    deck = save(p, tmp_path, "vt.pptx")
    v2 = json.loads(run_cli([deck, "--profile", "full", "--schema", "2", "--json"]).stdout)
    assert v2["capabilities"]["geometry"] == "partial"
    reasons = {a["reason"] for a in v2["abstentions"]}
    assert "vertical_text" in reasons

def test_reason_registry_covers_all_keys():
    """0.7.1 P0: every skip-reason key a detector can write must be registered, so a
    structural abstention never lands with no affected rules while structure reads
    'complete'. Scans the source for the Counter keys the guards use."""
    import re
    with open(jl.__file__, encoding="utf-8") as _f:
        src = _f.read()
    used = set(re.findall(r'(?:skipped|deck_skipped)\[\"([a-z0-9_]+)\"\]', src))
    used |= set(re.findall(r'skipped\.get\(\"([a-z0-9_]+)\"', src))
    missing = sorted(k for k in used if k not in jl.KNOWN_REASON_KEYS)
    assert not missing, "unregistered skip reasons: %s" % missing


def test_html_report(tmp_path):
    """0.8.x (#4): the annotated visual report is one self-contained static HTML file:
    an SVG wireframe per slide, finding bboxes overlaid with their codes, error entry
    for an unreadable file in scan mode, and no external requests."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 6.9, 5, 0.3, "tiny", size=4, ea="맑은 고딕")            # E3
    tb(s, 12.8, 3.0, 3.0, 0.5, "off canvas", size=18, ea="맑은 고딕")  # W16
    s2 = add_slide(p)
    tb(s2, 1, 1, 4, 1, "clean page", size=14, ea="맑은 고딕")
    deck = save(p, tmp_path, "h.pptx")
    out = os.path.join(str(tmp_path), "r.html")
    run_cli([deck, "--profile", "full", "--html", out])
    with open(out, encoding="utf-8") as f:
        html = f.read()
    assert html.count("<svg") == 2                      # one wireframe per slide
    assert ">E3</text>" in html and ">W16</text>" in html
    assert "clean" in html                              # the clean slide says so
    assert "http://" not in html and "https://" not in html   # self-contained
    # scan mode: an unreadable file becomes an error entry, not a crash
    d = os.path.join(str(tmp_path), "batch")
    os.makedirs(d)
    save(p, d, "ok.pptx")
    with open(os.path.join(d, "corrupt.pptx"), "w") as f:
        f.write("nope")
    out2 = os.path.join(str(tmp_path), "r2.html")
    run_cli(["scan", d, "--html", out2])
    with open(out2, encoding="utf-8") as f:
        html2 = f.read()
    assert "could not check" in html2 and html2.count("<svg") == 2


def test_fix_subcommand(tmp_path):
    """0.8.x: deterministic auto-fix for E1/E2/E4 only. The fixed copy re-lints clean of
    those codes, judgment rules (E3) are untouched, exempt range dashes survive, and a
    non-fixable rule is a usage error."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.6, "아리알 한글 제목", font="Arial", size=20)              # E1
    tb(s, 1, 2, 8, 0.6, "구조적" + EM_DASH + "개선 흐름", size=14, ea="맑은 고딕")  # E2
    tb(s, 1, 3, 8, 0.6, "FY2020" + EN_DASH + "2024 실적", size=14, ea="맑은 고딕")  # exempt range
    tb(s, 1, 4, 5, 0.6, "자간 벌어진 한글", size=14, spc=300, ea="맑은 고딕")   # E4
    tb(s, 1, 6.9, 5, 0.3, "tiny src", size=4, ea="맑은 고딕")                    # E3 (untouched)
    deck = save(p, tmp_path, "fixme.pptx")
    out = os.path.join(str(tmp_path), "fixed.pptx")
    r = run_cli(["fix", deck, "-o", out])
    assert r.returncode == 0, r.stderr
    assert os.path.exists(out)
    doc = json.loads(run_cli([out, "--profile", "full", "--json"]).stdout)
    left = sorted(f["code"] for f in doc["errors"])
    assert "E1" not in left and "E2" not in left and "E4" not in left, left
    assert "E3" in left                                  # judgment rule untouched
    # the exempt range dash survived the E2 fix
    from pptx import Presentation
    texts = " ".join(sh.text_frame.text for sh in Presentation(out).slides[0].shapes
                     if sh.has_text_frame)
    assert "FY2020" + EN_DASH + "2024" in texts
    assert "구조적, 개선" in texts                        # prose dash became a comma
    # non-fixable rule -> controlled usage error
    r = run_cli(["fix", deck, "-o", out, "--rules", "W15"])
    assert r.returncode == 2
