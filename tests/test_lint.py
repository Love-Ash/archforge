# -*- coding: utf-8 -*-
"""archforge gate regression: fixes, per gate, both positive fixtures (does a pptx with
a planted defect get caught) and negative fixtures (does clean, intentional staging
pass without false positives).

Characters the linter blocks, such as the em dash, are not kept as literals in this
source file but built via chr(codepoint): the source must stay clean, and the
forbidden character must exist only inside the pptx under test for the gate test to
be valid.

The expected values in the E1 fixtures follow the PowerPoint COM render measured
model (2026-07-10, docs/CALIBRATION.md): run a:ea > theme minorFont a:ea (wins over
run a:latin when non-empty) > (only when the theme ea slot is empty) run a:latin >
OS fallback (Malgun).
"""
import os
import re
import json
import shutil
import subprocess
import sys
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

import archforge.lint as jl
from archforge import messages as jmsg


@pytest.fixture(autouse=True)
def _force_korean_messages(monkeypatch):
    """Tests pin the lang to ko to keep existing Korean-message assertions working
    (0.3.0 i18n). English output is verified explicitly in test_lang_*."""
    monkeypatch.setenv("ARCHFORGE_LANG", "ko")
    jmsg.set_lang("ko")
    yield
    jmsg.set_lang(None)


EN_DASH = chr(0x2013)
EM_DASH = chr(0x2014)
MINUS = chr(0x2212)
FW_HYPHEN = chr(0xFF0D)


# ---------------------------------------------------------------- helpers
def new_prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def add_slide(p):
    return p.slides.add_slide(p.slide_layouts[6])


def tb(s, x, y, w, h, text, font=None, size=12, color="222222", spc=None, no_size=False, ea=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    if font:
        r.font.name = font
    if ea:
        rPr = r._r.get_or_add_rPr()
        el = rPr.makeelement(qn("a:ea"), {"typeface": ea})
        rPr.append(el)
    if not no_size:
        r.font.size = Pt(size)
    r.font.color.rgb = RGBColor.from_string(color)
    if spc is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(spc))
    return box


def rect(s, x, y, w, h, fill_hex):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    sp.line.fill.background()
    return sp


def vconn(s, x, y0, y1, line_hex):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y0), Inches(x), Inches(y1))
    cn.line.color.rgb = RGBColor.from_string(line_hex)
    cn.line.width = Pt(2)
    return cn


def set_autofit_scale(box, pct):
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    na = bodyPr.makeelement(qn("a:normAutofit"), {"fontScale": str(int(pct * 1000))})
    bodyPr.append(na)


def png(d, name, size=(400, 300), opaque_box=None):
    """A fully opaque RGB PNG, or an RGBA PNG opaque only in the opaque_box=(l,t,r,b
    ratio) region."""
    from PIL import Image, ImageDraw
    path = os.path.join(str(d), name)
    if opaque_box is None:
        Image.new("RGB", size, (40, 60, 90)).save(path)
    else:
        im = Image.new("RGBA", size, (0, 0, 0, 0))
        dr = ImageDraw.Draw(im)
        l, t, r, b = opaque_box
        dr.rectangle([size[0] * l, size[1] * t, size[0] * r, size[1] * b], fill=(40, 60, 90, 255))
        im.save(path)
    return path


def lint_full(*args, **kw):
    """In 0.4.0 the default profile changed to core, so existing fixtures that assume
    style rules too must specify full explicitly (the default value itself is
    verified by test_default_profile_core)."""
    kw.setdefault("profile", "full")
    return jl.lint(*args, **kw)


def codes(items):
    return [c for (_si, c, _m, _d) in items]


def by_code(items, code):
    return [(si, m, d) for (si, c, m, d) in items if c == code]


def save(p, d, name):
    path = os.path.join(str(d), name)
    p.save(path)
    return path


def patch_theme_ea(path, typeface):
    """Replace the theme a:ea empty slot in a saved pptx with typeface (rewrites the
    zip). For the default template with no multi-master: the only reliably stable
    way to build an E1 theme-branch fixture."""
    tmp = path + ".patched.pptx"
    new = ('<a:ea typeface="%s"/>' % typeface).encode("utf-8")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if re.search(r"theme/theme\d+\.xml$", item.filename):
                data = data.replace(b'<a:ea typeface=""/>', new)
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


def run_cli(args, lang="ko"):
    """CLI run helper. stdin=DEVNULL: avoids Windows handle-inheritance failure
    (WinError 6) under pytest capture. Defaults to ko (assertions use Korean
    phrasing); verify English with lang="en"."""
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = lang
    return subprocess.run([sys.executable, "-m", "archforge.lint"] + args,
                          capture_output=True, text=True, encoding="utf-8",
                          stdin=subprocess.DEVNULL, env=env)


# ---------------------------------------------------------------- line-level gates
def test_line_gates_positive(tmp_path):
    p = new_prs()
    s = add_slide(p)   # p1: E1 Hangul in a latin-only font
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    s = add_slide(p)   # p2: E2 long dash (fullwidth hyphen)
    tb(s, 1, 1, 5, 0.5, "dash" + FW_HYPHEN + "test", font="Wanted Sans", size=12)
    s = add_slide(p)   # p3: E3 explicit 4pt
    tb(s, 1, 1, 5, 0.5, "tiny text 4pt", font="Wanted Sans", size=4)
    s = add_slide(p)   # p4: E3 autofit 40pt * 10% = 4pt
    box = tb(s, 1, 1, 5, 1.0, "autofit shrink", font="Wanted Sans", size=40)
    set_autofit_scale(box, 10.0)
    s = add_slide(p)   # p5: E4 positive CJK tracking
    tb(s, 1, 1, 5, 0.5, "자간 벌어진 한글", font="Wanted Sans", size=12, spc=100)
    s = add_slide(p)   # p6: W1 wide frame body 8.5pt 40+ chars
    tb(s, 1, 1, 11, 0.6, "0123456789" * 5, font="Wanted Sans", size=8.5)
    s = add_slide(p)   # p7: size unspecified -> resolved to defaultTextStyle 18pt (not W5: inheritance resolution introduced)
    tb(s, 1, 1, 5, 0.5, "no size run", font="Wanted Sans", no_size=True)
    s = add_slide(p)   # p8: W8 narrow frame, small 6pt Hangul
    tb(s, 1, 1, 2, 0.4, "목업 안 작은 한글", font="Wanted Sans", size=6)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    ec = codes(errors)
    assert "E1" in ec and by_code(errors, "E1")[0][0] == 1
    assert "E2" in ec and by_code(errors, "E2")[0][0] == 2
    e3 = [si for (si, m, d) in by_code(errors, "E3")]
    assert 3 in e3 and 4 in e3
    assert "E4" in ec and by_code(errors, "E4")[0][0] == 5
    assert any(si == 6 for (si, m, d) in by_code(warns, "W1"))
    # p7: the previous version fired W5 (inheritance unresolved). With the
    # inheritance chain introduced, defaultTextStyle 18pt is resolved, so
    # neither W5 nor a false-fire of the size gate should occur.
    assert not any(si == 7 for (si, m, d) in by_code(warns, "W5"))
    assert not any(si == 7 for (si, m, d) in by_code(errors, "E3"))
    assert any(si == 8 for (si, m, d) in by_code(warns, "W8"))


def test_e1_nofont_empty_theme_slot(tmp_path):
    """Font unspecified + empty theme ea = Malgun fallback E1. Since 0.2.1, the
    +mn-lt token in defaultTextStyle is resolved through the inheritance chain
    (effective latin=Calibri), so it's caught via a more accurate branch; pin it
    to the code rather than the message."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "폰트 미지정 한글", size=12)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    assert any(c == "E1" for (_si, c, m, _d) in errors), errors


def test_e1_render_model_slots(tmp_path):
    """Measured render model regression (2026-07-10 COM probe):
    p1 run ea is latin-only -> E1. p2 with an empty theme, latin is latin-only
    (Inter, previously a missed detection) -> E1.
    p3 with an empty theme, latin is a Hangul font -> that font actually renders
    the Hangul (a legitimate pattern) -> passes.
    p4 if run ea is a Hangul font, ea wins even when latin is latin-only -> passes."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이에이 슬롯 모노", font="Arial", ea="IBM Plex Mono", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "인터 폴백 한글", font="Inter", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "원티드 산스 한글", font="Wanted Sans", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이에이 승리 한글", font="IBM Plex Mono", ea="Wanted Sans", size=12)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    e1_pages = [si for (si, m, d) in by_code(errors, "E1")]
    assert 1 in e1_pages and 2 in e1_pages
    assert 3 not in e1_pages and 4 not in e1_pages


def test_e1_theme_ea_korean_suppresses_latin(tmp_path):
    """If the theme ea is a Hangul font, the theme ea handles the render even when
    run latin is latin-only (measured): pins the regression for the E1 false
    positive that the old ea-or-latin substitute used to emit."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "테마가 받는 한글", font="IBM Plex Mono", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "폰트 미지정 한글", size=12)
    path = patch_theme_ea(save(p, tmp_path, "fx.pptx"), "Malgun Gothic")
    errors, _w = lint_full(path)
    assert not by_code(errors, "E1"), errors


def test_e1_theme_ea_latin_only_flags(tmp_path):
    """If the theme ea itself is latin-only, font-unspecified Hangul is E1. Since
    0.2.1, the +mn-ea token in defaultTextStyle is resolved through the chain and
    caught as effective ea=Consolas (branch-independent, pinned by font name)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "폰트 미지정 한글", size=12)
    path = patch_theme_ea(save(p, tmp_path, "fx.pptx"), "Consolas")
    errors, _w = lint_full(path)
    assert any(c == "E1" and "Consolas" in d for (_si, c, m, d) in errors), errors


def test_theme_ea_by_master_relationship(tmp_path):
    """Whether theme resolution comes from the master->theme relationship rather
    than iter_parts order (key = master partname)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "관계 해석", font="Wanted Sans", size=12)
    prs = Presentation(save(p, tmp_path, "fx.pptx"))
    ea_map = jl.theme_ea_by_master(prs)
    assert len(ea_map) == 1
    key = next(iter(ea_map))
    assert "slideMaster" in key
    assert ea_map[key] == ""   # default template = empty slot


def test_e1_unit_matrix():
    """e1_violation decision matrix (a unit regression that doubles as documentation
    of the measured render model)."""
    v = jl.e1_violation
    assert v("한글", {"ea": "IBM Plex Mono"}, "") is not None          # run ea latin-only
    assert v("한글", {"ea": "맑은 고딕", "latin": "Consolas"}, "") is None   # run ea Hangul
    assert v("한글", {"latin": "IBM Plex Mono"}, "맑은 고딕") is None   # theme ea handles the render
    assert v("한글", {}, "Consolas") is not None                        # theme ea latin-only
    assert v("한글", {"latin": "Wanted Sans"}, "") is None              # empty theme + latin Hangul font
    assert v("한글", {"latin": "Inter"}, "") is not None                # empty theme + latin latin-only (previously a missed detection)
    assert v("한글", {}, "") is not None                                # nothing at all = Malgun confirmed


def test_latin_only_font_boundaries():
    """Blocklist prefix-match boundary: a Hangul-complete variant must not get
    swept up by a prefix match."""
    f = jl.is_latin_only_font
    assert f("Inter") and f("Arial") and f("Calibri") and f("IBM Plex Sans") and f("Noto Sans")
    assert not f("NanumGothicCoding")     # Hangul-complete monospace: pins the regression for an old misclassification
    assert not f("IBM Plex Sans KR")
    assert not f("Noto Sans KR") and not f("Noto Serif KR")
    assert not f("Arial Unicode MS")
    assert not f("Wanted Sans") and not f("Pretendard") and not f("Malgun Gothic")
    assert not f("") and not f(None)


def test_e2_numeric_context(tmp_path):
    """E2 context exception: a numeric-range en dash and a minus sign U+2212 pass in
    default mode; an em dash is always blocked."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "FY2020" + EN_DASH + "2024 실적", font="Wanted Sans", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, MINUS + "3.2% 하락", font="Wanted Sans", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이건" + EM_DASH + "안 됨", font="Wanted Sans", size=12)
    path = save(p, tmp_path, "fx.pptx")
    errors, _w = lint_full(path)
    e2_pages = [si for (si, m, d) in by_code(errors, "E2")]
    assert e2_pages == [3], errors
    # strict: exceptions lifted, all three pages blocked
    errors_s, _w = lint_full(path, strict=True)
    assert sorted(si for (si, m, d) in by_code(errors_s, "E2")) == [1, 2, 3]


def test_dash_violations_unit():
    dv = jl.dash_violations
    assert dv("2020" + EN_DASH + "2024") == []
    assert dv("2020" + EN_DASH + "2024", strict=True) == [EN_DASH]
    assert dv(MINUS + "3.2%") == []
    assert dv("A" + EN_DASH + "B") == [EN_DASH]        # not a numeric context -> stays blocked
    assert dv("끝" + MINUS) == [MINUS]                  # no digit follows -> stays blocked
    assert dv("a" + EM_DASH + "b") == [EM_DASH]


def test_e4_universal_measure_spc(tmp_path):
    """spc's universal measure ('1.5pt') and a garbage value: without crashing, the
    former should be E4 and the latter should be ignored."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "자간 벌어진 한글", font="Wanted Sans", size=12, spc="1.5pt")
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "정상 자간 한글", font="Wanted Sans", size=12, spc="garbage")
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    e4_pages = [si for (si, m, d) in by_code(errors, "E4")]
    assert e4_pages == [1], errors


def test_text_point_attr_unit():
    tpa = jl._text_point_attr
    assert tpa("150") == 150
    assert tpa("1.5pt") == 150
    assert tpa("0.5in") == 3600
    assert tpa("abc") is None
    assert tpa(None) is None


def test_partial_salvage_on_bad_frame(tmp_path):
    """Guard granularity regression (third adversarial panel confirmed, 2026-07-10):
    even if run1 in the same frame dies on a garbage attribute, run2's genuine E1
    violation must still be caught (per-run guard), and the swallowed section must
    be surfaced to JSON as W18. In the days of frame-level guards, run2's violation
    used to vanish entirely too, producing a false clean."""
    p = new_prs()
    s = add_slide(p)   # p1: one frame with a garbage run + a genuine E1 run
    box = tb(s, 1, 1, 5, 0.5, "쓰레기 크기 한글", font="Wanted Sans", no_size=True)
    r1 = box.text_frame.paragraphs[0].runs[0]
    r1._r.get_or_add_rPr().set("sz", "notanumber")
    r2 = box.text_frame.paragraphs[0].add_run()
    r2.text = "모노 폴백 한글"
    r2.font.name = "IBM Plex Mono"
    from pptx.util import Pt as _Pt
    r2.font.size = _Pt(12)
    s = add_slide(p)   # p2: a normal E1 defect (also pins survival across pages)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    e1_pages = [si for (si, m, d) in by_code(errors, "E1")]
    assert 1 in e1_pages, errors          # the neighboring run's violation in the same frame survives
    assert 2 in e1_pages, errors
    w18 = by_code(warns, "W18")
    assert any(si == 1 for (si, m, d) in w18), warns   # the swallowed section is surfaced in the output contract


def test_w18_geometry_decoupling(tmp_path):
    """A tboxes calculation failure (a corrupted anchor) must not also silence an
    unrelated picture's W16 (pins the regression for 2 cases reproduced measured
    by the adversarial panel). W16 survives and W18 reports the incompleteness."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 1, 1, 4, 0.5, "손상 앵커 텍스트", font="Wanted Sans", size=14)
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "bogusvalue")   # a value with no MSO_VERTICAL_ANCHOR mapping -> tboxes exception
    s.shapes.add_picture(png(tmp_path, "over.png"), Inches(12.8), Inches(2.0),
                         Inches(2.0), Inches(1.5))   # overflows 1.47in on the right
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert any("그림" in d for (_si, m, d) in by_code(warns, "W16")), warns
    assert by_code(warns, "W18"), warns


def test_w6_page_numbers_survive_token_failure(tmp_path, monkeypatch):
    """sigs double-append regression (adversarial panel confirmed): even if
    _fill_tokens dies, the sigs positions don't shift, so W6 page numbers point to
    the actual slide number."""
    def boom(slide, sw, sh):
        raise RuntimeError("token fail")
    monkeypatch.setattr(jl, "_fill_tokens", boom)
    p = new_prs()
    for i in range(6):
        s = add_slide(p)
        tb(s, 1, 0.8, 8, 0.6, "Title block %d" % i, font="Wanted Sans", size=24)
        tb(s, 1, 2.0, 6, 2.5, "Body block", font="Wanted Sans", size=12)
        tb(s, 8, 2.0, 4, 2.5, "Side block", font="Wanted Sans", size=12)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w6 = by_code(warns, "W6")
    assert w6, warns
    pages = [int(m) for m in re.findall(r"p(\d+)", w6[0][2])]
    assert pages and max(pages) <= 6, w6   # in the days of double-append, p7-p11 used to show up
    # 0.2.1: token-collection failure is also surfaced as W18 (guard wiring extended everywhere)
    assert any("w10_tokens" in d for (_si, m, d) in by_code(warns, "W18")), warns


def test_e1_theme_token_resolution(tmp_path):
    """Resolve OOXML theme tokens (the +mn-lt kind) to actual fonts (adversarial
    panel confirmed: matching the token against the blocklist literally causes an
    E1 missed detection). The default template's theme minor latin=Calibri
    (latin-only), so a Hangul run with latin="+mn-lt" must be E1."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 1, 1, 5, 0.5, "토큰 한글", size=12)
    rPr = box.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn("a:latin"), {"typeface": "+mn-lt"}))
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    assert any("Calibri" in d for (_si, m, d) in by_code(errors, "E1")), errors


def test_resolve_font_tokens_unit():
    rft = jl.resolve_font_tokens
    thm = {"mn-lt": "Calibri", "mn-ea": "맑은 고딕", "mj-lt": "Georgia", "mj-ea": ""}
    assert rft({"latin": "+mn-lt"}, thm) == {"latin": "Calibri"}
    assert rft({"ea": "+mn-ea"}, thm) == {"ea": "맑은 고딕"}
    assert rft({"latin": "+mj-lt", "ea": "Wanted Sans"}, thm) == {"latin": "Georgia", "ea": "Wanted Sans"}
    assert rft({"ea": "+mj-ea"}, thm) == {}          # if resolved to an empty value, remove the slot (leave it to the chain)
    assert rft({"latin": "+mn-lt"}, None) == {}       # if theme parsing fails, remove the slot
    assert rft({"latin": "Batang"}, thm) == {"latin": "Batang"}   # leave as-is if not a token


def test_e2_run_boundary_split(tmp_path):
    """False-positive regression for a numeric range split across a run boundary
    (adversarial panel confirmed, high): even when '2020'/'-2021' are different
    runs, the exception applies via the paragraph context. An em dash stays
    blocked even when split."""
    p = new_prs()
    s = add_slide(p)   # p1: en dash range split at a run boundary
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.5))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run(); r1.text = "2020"; r1.font.size = Pt(14)
    r2 = para.add_run(); r2.text = EN_DASH + "2024 실적"; r2.font.size = Pt(14)
    s = add_slide(p)   # p2: em dash split (always blocked)
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.5))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run(); r1.text = "말"; r1.font.size = Pt(14)
    r2 = para.add_run(); r2.text = EM_DASH + "이음"; r2.font.size = Pt(14)
    path = save(p, tmp_path, "fx.pptx")
    errors, _w = lint_full(path)
    e2_pages = [si for (si, m, d) in by_code(errors, "E2")]
    assert e2_pages == [2], errors
    errors_s, _w = lint_full(path, strict=True)   # strict blocks everything regardless of split
    assert sorted(set(si for (si, m, d) in by_code(errors_s, "E2"))) == [1, 2]


def test_size_inheritance_placeholder_gates(tmp_path):
    """Placeholder inheritance chain activation regression: placeholder text with
    no explicit size doesn't leak into W5 but is resolved to the layout lstStyle
    size, so E3/W1 actually run (the old version disarmed everything into W5)."""
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    lay = p.slide_layouts[1]   # Title and Content
    for ph in lay.placeholders:
        sz = {0: "400", 1: "800"}.get(ph.placeholder_format.idx)   # title 4pt, body 8pt
        if sz is None:
            continue
        txBody = ph.text_frame._txBody
        lst = txBody.find(qn("a:lstStyle"))
        if lst is None:
            lst = txBody.makeelement(qn("a:lstStyle"), {})
            txBody.insert(1, lst)   # after bodyPr
        lvl1 = lst.makeelement(qn("a:lvl1pPr"), {})
        defR = lst.makeelement(qn("a:defRPr"), {"sz": sz})
        lvl1.append(defR)
        lst.append(lvl1)
    s = p.slides.add_slide(lay)
    s.shapes.title.text_frame.text = "판독 불가 제목"
    s.placeholders[1].text_frame.text = "본문이 레이아웃 상속 크기로 해석되는지 보는 사십자 이상의 충분히 긴 문장입니다"
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert not by_code(warns, "W5"), warns
    assert any("4.0pt" in m for (_si, m, _d) in by_code(errors, "E3")), errors       # title 4pt
    assert any("8.0pt" in m for (_si, m, _d) in by_code(warns, "W1")), warns          # body 8pt


def test_size_inheritance_master_txstyles(tmp_path):
    """Master txStyles path: default template resolves title=44pt, body(OBJECT)=32pt."""
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text_frame.text = "제목"
    s.placeholders[1].text_frame.text = "본문"
    path = save(p, tmp_path, "fx.pptx")
    prs = Presentation(path)
    sl = prs.slides[0]
    sizer = jl.SizeResolver(prs)
    got = {}
    for sp in sl.shapes:
        if sp.has_text_frame and sp.is_placeholder:
            got[sp.placeholder_format.idx] = sizer.resolve(sp.text_frame, sp, sl, 0)
    assert got[0] == 44.0 and got[1] == 32.0, got


def test_w5_when_chain_fully_absent(tmp_path):
    """If even defaultTextStyle is removed (simulating output from an external
    generator), an unresolved size is surfaced as W5."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "크기 미상 한글", font="Wanted Sans", no_size=True)
    pres_el = p.slides._sldIdLst.getparent()
    dts = pres_el.find(qn("p:defaultTextStyle"))
    assert dts is not None
    pres_el.remove(dts)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert by_code(warns, "W5"), warns


def test_table_cell_e1(tmp_path):
    """Native table-cell path (a headline feature that had zero coverage): also
    catches font defects on CJK inside a cell."""
    p = new_prs()
    s = add_slide(p)
    gf = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(8), Inches(2))
    cell = gf.table.cell(0, 0)
    cell.text = "표 안 폴백 한글"
    r = cell.text_frame.paragraphs[0].runs[0]
    r.font.name = "Consolas"
    r.font.size = Pt(12)
    ok = gf.table.cell(1, 1)
    ok.text = "정상 셀 한글"
    r2 = ok.text_frame.paragraphs[0].runs[0]
    r2.font.name = "Wanted Sans"
    r2.font.size = Pt(12)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    e1 = by_code(errors, "E1")
    assert len(e1) == 1 and "Consolas" in e1[0][2], errors


# ---------------------------------------------------------------- layout-level gates
def test_w6_layout_clones(tmp_path):
    p = new_prs()
    for i in range(6):
        s = add_slide(p)
        tb(s, 1, 0.8, 8, 0.6, "Title block %d" % i, font="Wanted Sans", size=24)
        tb(s, 1, 2.0, 6, 2.5, "Body block", font="Wanted Sans", size=12)
        tb(s, 8, 2.0, 4, 2.5, "Side block", font="Wanted Sans", size=12)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W6" in codes(warns)


def test_w6_tunables_suppress(tmp_path):
    """W6 tunable: raising the cluster threshold makes the same deck pass (an
    intentional relief knob for template houses)."""
    p = new_prs()
    for i in range(6):
        s = add_slide(p)
        tb(s, 1, 0.8, 8, 0.6, "Title block %d" % i, font="Wanted Sans", size=24)
        tb(s, 1, 2.0, 6, 2.5, "Body block", font="Wanted Sans", size=12)
        tb(s, 8, 2.0, 4, 2.5, "Side block", font="Wanted Sans", size=12)
    path = save(p, tmp_path, "fx.pptx")
    _e, warns = lint_full(path, w6_min_cluster=10)
    assert "W6" not in codes(warns)
    # an exact clone has cosine exactly 1.00, so the sim threshold must be set to 1.0 to exclude it (confirms the parameter passes through)
    _e, warns = lint_full(path, w6_sim=1.0)
    assert "W6" not in codes(warns)


def test_w7_low_contrast_needs_render(tmp_path):
    from PIL import Image
    img = os.path.join(str(tmp_path), "bg.jpg")
    Image.new("RGB", (1536, 864), (240, 240, 240)).save(img)
    p = new_prs()
    s = add_slide(p)
    s.shapes.add_picture(img, 0, 0, Inches(13.333), Inches(7.5))
    tb(s, 2, 3, 9, 1, "white on white low contrast", font="Wanted Sans", size=40, color="F5F5F5")
    path = save(p, tmp_path, "fx.pptx")
    pages = os.path.join(str(tmp_path), "pages")
    os.makedirs(pages)
    Image.new("RGB", (1600, 900), (240, 240, 240)).save(os.path.join(pages, "p01.png"))
    _e, warns = lint_full(path, render_dir=pages)
    assert "W7" in codes(warns)


def test_render_on_textonly_deck_no_crash(tmp_path):
    """Linting a text-only deck with not a single image via --render must not crash.
    Regression: back when glob existed only as a local import inside
    contrast_check, glob.glob in the render_png_hits==0 branch died with a
    NameError, which got mislabeled as "could not open pptx" (2026-07-04)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 0.8, 9, 0.8, "제목만 있는 페이지", font="Wanted Sans", size=28)
    tb(s, 1, 2.2, 10, 3, "본문 텍스트, 이미지는 없다.", font="Wanted Sans", size=13)
    pages = os.path.join(str(tmp_path), "pages")
    os.makedirs(pages)
    from PIL import Image
    Image.new("RGB", (1600, 900), (240, 240, 240)).save(os.path.join(pages, "slide-1.png"))  # a name that breaks the convention
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"), render_dir=pages)  # fails if NameError occurs
    assert isinstance(errors, list) and isinstance(warns, list)


def test_w9_accent_vbars(tmp_path):
    p = new_prs()
    s = add_slide(p)
    for i in range(4):
        y = 1.5 + i * 1.1
        vconn(s, 5.0, y, y + 0.65, "E6A94E")
        tb(s, 5.2, y, 5.5, 0.6, "callout item %d body" % i, font="Wanted Sans", size=12)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W9" in codes(warns)


def test_w10_diagram_clone_marks(tmp_path):
    p = new_prs()
    for _ in range(2):
        s = add_slide(p)
        rect(s, 2.0, 2.0, 8.0, 3.0, "2E3540")
        for i in range(9):
            rect(s, 2.5 + i * 0.8, 3.0, 0.18, 0.18, "5B6470")
        tb(s, 1, 0.8, 8, 0.6, "diagram page", font="Wanted Sans", size=20)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W10" in codes(warns)


def test_w11_ai_copy(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 10, 0.8, "오늘날 급변하는 시장 환경에서", font="Wanted Sans", size=24)
    s = add_slide(p)
    tb(s, 1, 1, 10, 0.8, "두 사업의 시너지 효과를 극대화", font="Wanted Sans", size=14)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w11 = by_code(warns, "W11")
    assert any("오프닝" in m for (_si, m, _d) in w11)
    assert any("버즈워드" in m for (_si, m, _d) in w11)


def test_w12_footer_baseline(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 3, 8, 1, "Cover page", font="Wanted Sans", size=30)
    for i in range(6):
        s = add_slide(p)
        tb(s, 1, 1, 8, 0.6, "Content %d" % i, font="Wanted Sans", size=20)
        fy = 7.05 if i == 3 else 6.9
        tb(s, 1, fy, 5, 0.3, "footer text", font="Wanted Sans", size=8)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W12" in codes(warns)


def test_w13_native_effects(tmp_path):
    p = new_prs()
    s = add_slide(p)
    for i in range(2):
        sp = rect(s, 1 + i * 3, 2, 2, 1.2, "2E3540")
        spPr = sp._element.spPr
        eff = spPr.makeelement(qn("a:effectLst"), {})
        sh = spPr.makeelement(qn("a:outerShdw"), {"blurRad": "40000", "dist": "20000"})
        eff.append(sh)
        spPr.append(eff)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W13" in codes(warns)


def test_w14_nominal_titles_and_ghost(tmp_path):
    p = new_prs()
    for t in ("시장 현황", "경쟁 구도 분석", "제품 라인업 개요", "사업 확장 전략",
              "재무 운용 계획", "향후 추진 방안"):
        s = add_slide(p)
        tb(s, 1, 0.8, 9, 0.8, t, font="Wanted Sans", size=26)
        tb(s, 1, 2.2, 10, 3, "본문 내용", font="Wanted Sans", size=12)
    ghost = []
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"), ghost=ghost)
    assert "W14" in codes(warns)
    assert len(ghost) == 6


def test_w14_numeric_claim_titles_pass(tmp_path):
    """A number+unit title is a claim even when it ends in a noun: W14 must not
    false-fire (mitigates genre bias)."""
    p = new_prs()
    for t in ("매출 3배 성장", "점유율 42% 달성", "비용 120억 절감", "가입자 500만 돌파",
              "마진 8pp 개선", "리텐션 2x 향상"):
        s = add_slide(p)
        tb(s, 1, 0.8, 9, 0.8, t, font="Wanted Sans", size=26)
        tb(s, 1, 2.2, 10, 3, "본문 내용", font="Wanted Sans", size=12)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert "W14" not in codes(warns), by_code(warns, "W14")


# ---------------------------------------------------------------- geometry gates
def test_w15_overlap_and_intentional_layers(tmp_path):
    p = new_prs()
    s = add_slide(p)  # p1: genuine overlap (label over a value)
    tb(s, 1, 2.0, 4, 1.2, "12,400원", font="Wanted Sans", size=44)
    tb(s, 1.2, 2.15, 3, 0.5, "+11.7% 전년 대비", font="Wanted Sans", size=14)
    s = add_slide(p)  # p2: drop cap = intentional
    tb(s, 1.0, 1.0, 1.2, 1.1, "피", font="Wanted Sans", size=60)
    tb(s, 1.0, 1.05, 5, 0.5, "지컬 AI", font="Wanted Sans", size=13)
    s = add_slide(p)  # p3: identical text echo = intentional
    tb(s, 2, 2, 4, 1, "떠난 뒤에야", font="Wanted Sans", size=30)
    tb(s, 2.05, 2.05, 4, 1, "떠난 뒤에야", font="Wanted Sans", size=30)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w15 = by_code(warns, "W15")
    assert any(si == 1 for (si, m, d) in w15)
    assert not any(si == 2 for (si, m, d) in w15)
    assert not any(si == 3 for (si, m, d) in w15)


def test_w16_overflow_and_negatives(tmp_path):
    p = new_prs()
    s = add_slide(p)  # p1: text overflows right (wrap=none, single line)
    tb(s, 11.0, 1.0, 4.0, 0.6, "Quarterly revenue grew twenty two percent",
       font="Wanted Sans", size=18)
    s = add_slide(p)  # p2: long Hangul in a word_wrap=True frame overflows the bottom
    bx = tb(s, 1.0, 7.0, 3.0, 0.4,
            "긴 본문이 프레임 계산보다 길어져 바닥을 뚫고 내려가는 경우를 재현한다",
            font="Wanted Sans", size=16)
    bx.text_frame.word_wrap = True
    s = add_slide(p)  # p3: picture clipped on the right
    s.shapes.add_picture(png(tmp_path, "op.png"), Inches(12.5), Inches(2.0), Inches(2.0), Inches(1.5))
    s = add_slide(p)  # p4 (negative): full bleed + a chart boundary straddling transparent margin
    s.shapes.add_picture(png(tmp_path, "bleed.png"), Inches(0), Inches(0), Inches(13.4), Inches(7.55))
    s.shapes.add_picture(png(tmp_path, "chart.png", opaque_box=(0.1, 0.1, 0.5, 0.5)),
                         Inches(9.0), Inches(4.0), Inches(5.0), Inches(3.6))
    tb(s, 1, 1, 6, 0.6, "음성 페이지 본문", font="Wanted Sans", size=14)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w16 = by_code(warns, "W16")
    assert any(si == 1 for (si, m, d) in w16)
    assert any(si == 2 for (si, m, d) in w16)
    assert any(si == 3 for (si, m, d) in w16)
    assert not any(si == 4 for (si, m, d) in w16)


def test_w17_straddle_and_negatives(tmp_path):
    p = new_prs()
    s = add_slide(p)  # p1: straddling
    s.shapes.add_picture(png(tmp_path, "p1.png"), Inches(4.0), Inches(2.0), Inches(4.0), Inches(3.0))
    tb(s, 7.0, 3.0, 4.0, 0.4, "잘린 캡션 텍스트 예시", font="Wanted Sans", size=14)
    s = add_slide(p)  # p2 (negative): full overlay
    s.shapes.add_picture(png(tmp_path, "p2.png"), Inches(4.0), Inches(2.0), Inches(4.0), Inches(3.0))
    tb(s, 4.4, 3.0, 3.0, 0.4, "오버레이 캡션", font="Wanted Sans", size=14)
    s = add_slide(p)  # p3 (negative): straddling over a transparent margin (alpha-trimmed)
    s.shapes.add_picture(png(tmp_path, "p3.png", opaque_box=(0.0, 0.0, 0.45, 1.0)),
                         Inches(4.0), Inches(2.0), Inches(4.0), Inches(3.0))
    tb(s, 6.5, 3.0, 3.0, 0.4, "투명부 위 캡션 예시", font="Wanted Sans", size=14)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w17 = by_code(warns, "W17")
    assert any(si == 1 for (si, m, d) in w17)
    assert not any(si == 2 for (si, m, d) in w17)
    assert not any(si == 3 for (si, m, d) in w17)


def test_geo_robustness_no_false_positives(tmp_path):
    """6 false-positive scenarios reproduced measured by adversarial verification:
    all must flag 0."""
    from PIL import Image, ImageDraw
    p = new_prs()
    s = add_slide(p)  # p1: wrap=none phantom wrap
    tb(s, 1.0, 1.0, 1.0, 0.4, "Quarterly revenue grew 18% on cloud momentum",
       font="Wanted Sans", size=12)
    tb(s, 1.0, 1.4, 2.0, 0.3, "Source: 10-K", font="Wanted Sans", size=9)
    s = add_slide(p)  # p2: group-move desync
    grp = s.shapes.add_group_shape()
    box = grp.shapes.add_textbox(Inches(12.5), Inches(3.0), Inches(1.5), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "group label"
    r.font.size = Pt(12)
    grp.left, grp.top = Inches(7.5), Inches(3.0)
    s = add_slide(p)  # p3: rotated picture (bbox before rotation is outside, render is inside)
    pic = s.shapes.add_picture(png(tmp_path, "rot.png", size=(400, 50)),
                               Inches(11.0), Inches(3.5), Inches(4.0), Inches(0.5))
    pic.rotation = 90
    s = add_slide(p)  # p4: flipH (mirrored ink lands on screen)
    pic = s.shapes.add_picture(png(tmp_path, "flip.png", opaque_box=(0.0, 0.0, 0.5, 1.0)),
                               Inches(-0.8), Inches(3.0), Inches(2.0), Inches(1.0))
    pic._element.spPr.find(qn("a:xfrm")).set("flipH", "1")
    s = add_slide(p)  # p5: P-mode + tRNS transparency
    im = Image.new("RGBA", (400, 400), (0, 0, 0, 0))
    ImageDraw.Draw(im).rectangle([150, 150, 250, 250], fill=(40, 60, 90, 255))
    pmode = os.path.join(str(tmp_path), "pmode.png")
    im.convert("P").save(pmode, transparency=0)
    s.shapes.add_picture(pmode, Inches(10.33), Inches(2.0), Inches(4.0), Inches(4.0))
    s = add_slide(p)  # p6: caption over a solid card over a photo
    s.shapes.add_picture(png(tmp_path, "photo.png"), Inches(4.0), Inches(0.5), Inches(6.0), Inches(6.5))
    rect(s, 6.0, 2.5, 4.0, 1.5, "FFFFFF")
    tb(s, 6.3, 3.0, 3.4, 0.4, "카드 위 캡션 문장", font="Wanted Sans", size=14)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    geo = [x for x in warns if x[1] in ("W15", "W16", "W17")]
    assert not geo, geo


def test_autofit_parser_forms(tmp_path):
    p = new_prs()
    s = add_slide(p)
    b = tb(s, 1, 1, 4, 1, "autofit test", font="Wanted Sans", size=12)
    bodyPr = b.text_frame._txBody.find(qn("a:bodyPr"))
    na = bodyPr.makeelement(qn("a:normAutofit"),
                            {"fontScale": "62.5%", "lnSpcReduction": "20%"})
    bodyPr.append(na)
    fs, lr = jl.frame_autofit(b.text_frame)
    assert abs(fs - 0.625) < 1e-6 and abs(lr - 0.20) < 1e-6
    na.set("fontScale", "62500")
    na.set("lnSpcReduction", "20000")
    fs, lr = jl.frame_autofit(b.text_frame)
    assert abs(fs - 0.625) < 1e-6 and abs(lr - 0.20) < 1e-6


# ---------------------------------------------------------------- negative + CLI
def test_clean_deck_no_flags(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 0.8, 9, 0.8, "Clean title page", font="Wanted Sans", size=30)
    tb(s, 1, 2.2, 10, 3, "Readable body text at normal size.", font="Wanted Sans", size=13)
    s = add_slide(p)
    tb(s, 1, 0.8, 6, 0.6, "Second layout", font="Wanted Sans", size=22)
    rect(s, 7.5, 1.0, 4.5, 5.0, "20242B")
    tb(s, 1, 3.5, 5.5, 2.5, "Different grid here.", font="Wanted Sans", size=13)
    s = add_slide(p)
    tb(s, 4, 3.2, 6, 0.8, "Closing page", font="Wanted Sans", size=26)
    errors, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert not errors and not warns, (codes(errors), codes(warns))


def test_cli_json_and_exit_codes(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    bad = save(p, tmp_path, "bad.pptx")
    r = run_cli([bad, "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 1
    assert doc["summary"]["error_count"] >= 1 and not doc["summary"]["pass"]
    assert any(e["code"] == "E1" for e in doc["errors"])
    assert "ghost" in doc and isinstance(doc["ghost"], list)   # documented JSON contract

    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.6, "clean page", font="Wanted Sans", size=20)
    good = save(p, tmp_path, "good.pptx")
    r = run_cli([good, "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 0 and doc["summary"]["pass"]


def test_cli_exit2_missing_and_corrupt(tmp_path):
    """exit 2 path regression (untested in the old version): missing file and a
    corrupt pptx."""
    r = run_cli([os.path.join(str(tmp_path), "no_such.pptx")])
    assert r.returncode == 2
    assert "찾을 수 없습니다" in r.stderr
    corrupt = os.path.join(str(tmp_path), "corrupt.pptx")
    with open(corrupt, "wb") as f:
        f.write(b"this is not a pptx file at all")
    r = run_cli([corrupt])
    assert r.returncode == 2
    assert "열 수 없습니다" in r.stderr


def test_cli_text_output_and_ghost(tmp_path):
    """Human-readable text output + the --ghost section (untested in the old version)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 0.8, 9, 0.8, "고스트 타이틀 페이지", font="Wanted Sans", size=26)
    tb(s, 1, 2, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    path = save(p, tmp_path, "fx.pptx")
    r = run_cli([path, "--ghost"])
    assert r.returncode == 1
    assert "=== ARCHFORGE LINT" in r.stdout
    assert "ghost deck" in r.stdout and "고스트 타이틀" in r.stdout
    assert "ERROR p01 [E1]" in r.stdout
    assert re.search(r"--- ERROR \d+, WARN \d+ ---", r.stdout)


def test_cli_strict_gates(tmp_path):
    """--strict: a deck with only WARNs switches to exit 1, and the E2 numeric
    context exception is lifted."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 11, 0.6, "0123456789" * 5, font="Wanted Sans", size=8.5)   # W1 only
    warn_only = save(p, tmp_path, "warn.pptx")
    assert run_cli([warn_only]).returncode == 0
    assert run_cli([warn_only, "--strict"]).returncode == 1

    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 6, 0.5, "FY2020" + EN_DASH + "2024 실적", font="Wanted Sans", size=14)
    ranged = save(p, tmp_path, "range.pptx")
    assert run_cli([ranged, "--profile", "full"]).returncode == 0
    r = run_cli([ranged, "--strict", "--profile", "full", "--json"])
    doc = json.loads(r.stdout)
    assert r.returncode == 1
    assert any(e["code"] == "E2" for e in doc["errors"])


def test_cli_skip_codes(tmp_path):
    """--skip: selectively suppress genre-independent warnings (W14 etc.)."""
    p = new_prs()
    for t in ("시장 현황", "경쟁 구도 분석", "제품 라인업 개요", "사업 확장 전략",
              "재무 운용 계획", "향후 추진 방안"):
        s = add_slide(p)
        tb(s, 1, 0.8, 9, 0.8, t, font="Wanted Sans", size=26)
        tb(s, 1, 2.2, 10, 3, "본문 내용", font="Wanted Sans", size=12)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json", "--profile", "full"]).stdout)
    assert any(w["code"] == "W14" for w in doc["warnings"])
    doc = json.loads(run_cli([path, "--json", "--profile", "full", "--skip", "W14"]).stdout)
    assert not any(w["code"] == "W14" for w in doc["warnings"])


# ---------------------------------------------------------------- 0.2.1 script layer + fixes
def patch_theme_fonts(path, major_ea=None, minor_ea=None):
    """Patch the theme's major/minor ea empty slots to different values respectively
    (rewrites the zip)."""
    tmp = path + ".patched.pptx"
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if re.search(r"theme/theme\d+\.xml$", item.filename):
                if major_ea is not None:
                    m = re.search(rb"<a:majorFont>.*?</a:majorFont>", data, re.S)
                    seg = m.group(0).replace(b'<a:ea typeface=""/>',
                                             ('<a:ea typeface="%s"/>' % major_ea).encode("utf-8"))
                    data = data[:m.start()] + seg + data[m.end():]
                if minor_ea is not None:
                    m = re.search(rb"<a:minorFont>.*?</a:minorFont>", data, re.S)
                    seg = m.group(0).replace(b'<a:ea typeface=""/>',
                                             ('<a:ea typeface="%s"/>' % minor_ea).encode("utf-8"))
                    data = data[:m.start()] + seg + data[m.end():]
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


def test_title_flood_regression(tmp_path):
    """Pins the regression where the defaultTextStyle fallback 18pt got collected as
    a title candidate, flooding ghost/W14 (second external re-review confirmed).
    A deck with only size-unspecified body text must have an empty ghost."""
    p = new_prs()
    s = add_slide(p)
    for i, txt in enumerate(("본문 첫 문장입니다", "본문 둘째 문장입니다", "표 안 텍스트 셋째")):
        tb(s, 1, 1 + i, 8, 0.5, txt, font="Wanted Sans", no_size=True)
    ghost = []
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"), ghost=ghost)
    assert ghost == [], ghost
    assert "W14" not in codes(warns)


def test_title_collection_keeps_placeholder_sizes(tmp_path):
    """Placeholder inherited size (master titleStyle 44pt) is an intentional title,
    so it must still be collected."""
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text_frame.text = "진짜 제목"
    s.placeholders[1].text_frame.text = "본문"
    ghost = []
    lint_full(save(p, tmp_path, "fx.pptx"), ghost=ghost)
    assert any("진짜 제목" in t for _si, t in ghost), ghost


def test_script_layer_two_tier(tmp_path):
    """The script layer's two-tier structure (third panel: resolves both the Hanja
    regression and the JP font false positive at once).
    p1 kana+Noto Sans JP -> silent (font has kana). p2 Hangul+Noto Sans JP ->
    E1 (no Hangul). p3 kana+positive tracking -> no E4 (Japanese convention).
    p4 Hanja-only+Georgia -> E1 (font has no CJK at all). p5 Hanja-only+positive
    tracking -> no E4 since 0.6.1 (tracked hanzi is legitimate Chinese typography;
    the run must contain Hangul). p6 kana+IBM Plex Mono -> E1. p7 mixed
    Hangul+Hanja+tracking -> E4 (Korean names/legal terms keep coverage)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "こんにちは", font="Noto Sans JP", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "한글 텍스트", font="Noto Sans JP", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "あいうえお", font="Noto Sans JP", size=12, spc=100)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "漢字専用", font="Georgia", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "大韓民國", font="맑은 고딕", size=12, spc=100)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "フリーレン", font="IBM Plex Mono", size=12)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "金九 선생", font="맑은 고딕", size=12, spc=100)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    e1_pages = sorted(si for (si, m, d) in by_code(errors, "E1"))
    assert e1_pages == [2, 4, 6], errors
    e4_pages = sorted(si for (si, m, d) in by_code(errors, "E4"))
    assert e4_pages == [7], errors


def test_hangul_range_extensions():
    """Closes the missed detection for halfwidth Hangul and the Jamo Extension
    blocks (third panel)."""
    assert jl.is_hangul(chr(0xFFA1))    # halfwidth giyeok
    assert jl.is_hangul(chr(0xA960))    # Jamo Extended-A
    assert jl.is_hangul(chr(0xD7B0))    # Jamo Extended-B
    assert jl._geometry_unsupported(chr(0x0F40))   # Tibetan
    assert jl._geometry_unsupported(chr(0x1000))   # Myanmar
    assert jl._geometry_unsupported(chr(0x1780))   # Khmer


def test_title_own_lststyle_not_flooded(tmp_path):
    """Body prose in a shape's own lstStyle 20pt is not swept into ghost (closes a
    case reproduced measured by the third panel). Title eligibility = an explicit
    size or a title-family placeholder."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 1, 1, 9, 0.6, "이것은 본문 카피 문장입니다 절대 제목이 아닙니다",
             font="Wanted Sans", no_size=True)
    txBody = box.text_frame._txBody
    lst = txBody.makeelement(qn("a:lstStyle"), {})
    lvl1 = lst.makeelement(qn("a:lvl1pPr"), {})
    defR = lst.makeelement(qn("a:defRPr"), {"sz": "2000"})
    lvl1.append(defR)
    lst.append(lvl1)
    txBody.insert(1, lst)
    ghost = []
    lint_full(save(p, tmp_path, "fx.pptx"), ghost=ghost)
    assert ghost == [], ghost


def test_cli_lang_edge_cases(tmp_path):
    """--lang CLI edge cases (third panel): repeated flags let the last one win, and
    it composes safely with the skill subcommand."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json", "--lang", "en", "--lang", "ko"]).stdout)
    assert doc["lang"] == "ko"
    r = run_cli(["skill", "--lang", "ko", "--path"])
    assert r.returncode == 0 and "SKILL.md" in r.stdout
    r = run_cli(["--lang", "ko", "skill", "--path"])   # leading flag + subcommand
    assert r.returncode == 0 and "SKILL.md" in r.stdout


def test_summary_incomplete_flag(tmp_path):
    """summary.incomplete: a machine-readable signal for whether checking was
    impossible (W18) (third panel: corrects documentation overclaiming)."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 1, 1, 5, 0.5, "쓰레기 크기 한글", font="Wanted Sans", no_size=True)
    box.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr().set("sz", "notanumber")
    doc = json.loads(run_cli([save(p, tmp_path, "bad.pptx"), "--json"]).stdout)
    assert doc["summary"]["incomplete"] is True
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    doc = json.loads(run_cli([save(p, tmp_path, "good.pptx"), "--json"]).stdout)
    assert doc["summary"]["incomplete"] is False


def test_latin_only_font_additions():
    f = jl.is_latin_only_font
    assert f("Aptos") and f("Courier") and f("PT Sans") and f("PT Serif")
    assert not f("Noto Sans Mono CJK KR")   # contains cjk = Hangul-complete (second re-review fixed a false positive)
    assert not f("Source Han Sans CJK KR")


def test_e2_v2_adversarial_edges():
    """Pins the regression for 4 cases found by the E2 v2 adversarial panel
    (2026-07-10, third round)."""
    dv = jl.dash_violations
    MINUS = chr(0x2212)
    assert dv("전일대비 " + MINUS + " 3.2%") == []            # spaced minus sign (financial notation)
    assert dv("끝 " + MINUS) == [MINUS]                        # still blocked when no digit follows
    assert dv("2020" + EN_DASH + EN_DASH + "2024") == [EN_DASH, EN_DASH]   # closes the missed detection for 2 consecutive dashes
    assert dv("결론2024" + EN_DASH + "우리는") == [EN_DASH]    # closes the bypass via a word+digit mixed token
    assert dv("매출" + chr(0x00B9) + EN_DASH + "영업이익 증가") == [EN_DASH]   # superscript footnote
    assert dv("2020" + EN_DASH + "현재") == []                 # keeps the digit-starting-token range rule
    assert dv("Q1" + EN_DASH + "Q3") == []                     # keeps the both-sides-numeric rule


def test_e2_v2_range_forms(tmp_path):
    """E2 v2: passes the remaining 4 false-positive forms + keeps blocking a
    spaced one-sided-digit case (an AI-inserted-clause pattern)."""
    dv = jl.dash_violations
    assert dv("2020 " + EN_DASH + " 2024") == []          # spaced numeric range
    assert dv("Q1" + EN_DASH + "Q3") == []                 # alphanumeric token
    assert dv("5%" + EN_DASH + "10%") == []                # percent range
    assert dv("2020" + EN_DASH + "현재") == []             # attached one-sided digit
    assert dv("성장 " + EN_DASH + " 2024년에는") == [EN_DASH]   # a spaced inserted clause stays blocked
    assert dv("서울" + EN_DASH + "부산") == [EN_DASH]       # word-to-word connection stays conservatively blocked
    assert dv("2020" + EN_DASH + "현재", strict=True) == [EN_DASH]
    # integration: also pinned via a deck fixture
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 6, 0.5, "매출 5%" + EN_DASH + "10% 구간", font="Wanted Sans", size=14)
    s = add_slide(p)
    tb(s, 1, 1, 6, 0.5, "성장 " + EN_DASH + " 2024년에는 확대", font="Wanted Sans", size=14)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    assert [si for (si, m, d) in by_code(errors, "E2")] == [2], errors


def test_skip_rejects_error_codes(tmp_path):
    """--skip is WARN-only: an E code causes exit 2 rejection, and applied skips
    are recorded in the JSON (fixes a footgun)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    path = save(p, tmp_path, "fx.pptx")
    r = run_cli([path, "--skip", "E1"])
    assert r.returncode == 2 and "WARN" in r.stderr
    doc = json.loads(run_cli([path, "--json", "--profile", "full", "--skip", "W14"]).stdout)
    assert doc["summary"]["skipped_codes"] == ["W14"]
    assert any(e["code"] == "E1" for e in doc["errors"])   # E1 is still alive


def test_e1_master_lststyle_font_inheritance(tmp_path):
    """Reflects Probe 6's measured finding: the master ph lstStyle's a:ea is the
    effective render font. Latin-only means E1, a Hangul font means it passes
    (fixes the "Malgun fallback confirmed" false positive)."""
    def build(ea_font):
        p = Presentation()
        master = p.slide_masters[0]
        for ph in master.placeholders:
            if ph.placeholder_format.idx == 1:
                txBody = ph.text_frame._txBody
                lst = txBody.find(qn("a:lstStyle"))
                if lst is None:
                    lst = txBody.makeelement(qn("a:lstStyle"), {})
                    txBody.insert(1, lst)
                lvl1 = lst.makeelement(qn("a:lvl1pPr"), {})
                defR = lst.makeelement(qn("a:defRPr"), {})
                defR.append(defR.makeelement(qn("a:ea"), {"typeface": ea_font}))
                lvl1.append(defR)
                lst.append(lvl1)
        s = p.slides.add_slide(p.slide_layouts[1])
        s.shapes.title.text_frame.text = ""
        body = s.placeholders[1]
        body.text_frame.text = "마스터 상속 한글"
        return save(p, tmp_path, "fx_%s.pptx" % ("bad" if "Mono" in ea_font else "ok"))

    errors, _w = lint_full(build("IBM Plex Mono"))
    assert any("IBM Plex Mono" in d for (_si, m, d) in by_code(errors, "E1")), errors
    errors, _w = lint_full(build("맑은 고딕"))
    assert not by_code(errors, "E1"), errors


def test_e1_title_uses_major_font(tmp_path):
    """Reflects Probe 6 Q1/Q2's measured finding: the title ph rides the theme's
    majorFont ea, and the body rides the minorFont ea."""
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text_frame.text = "제목 한글"
    s.placeholders[1].text_frame.text = "본문 한글"
    path = save(p, tmp_path, "fx.pptx")
    patch_theme_fonts(path, major_ea="Consolas", minor_ea="맑은 고딕")
    errors, _w = lint_full(path)
    e1 = by_code(errors, "E1")
    assert any("제목" in d for (_si, m, d) in e1), errors     # major=Consolas -> E1
    assert not any("본문" in d for (_si, m, d) in e1), errors  # minor=Malgun Gothic -> passes


def test_vertical_and_complex_script_geometry_skip(tmp_path):
    """Vertical writing (bodyPr@vert) and complex-typesetting scripts can't be
    geometrically estimated: skip them but surface it via W18."""
    p = new_prs()
    s = add_slide(p)   # p1: a vertical-writing frame goes off-screen -> must have no W16 FP
    box = tb(s, 12.0, 1.0, 4.0, 0.6, "vertical long text overflowing edge",
             font="Wanted Sans", size=18)
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.set("vert", "eaVert")
    s = add_slide(p)   # p2: Arabic text -> geometry skip + W18
    tb(s, 11.0, 1.0, 4.0, 0.6, "مرحبا بالعالم",
       font="Arial", size=18)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    w16_text = [d for (_si, m, d) in by_code(warns, "W16") if "텍스트" in d]
    assert not w16_text, warns
    w18 = by_code(warns, "W18")
    assert any(si == 1 and "vertical_text" in d for (si, m, d) in w18), warns
    assert any(si == 2 and "complex_script" in d for (si, m, d) in w18), warns


# ---------------------------------------------------------------- 0.3.0 i18n + profiles
def test_lang_english_output(tmp_path):
    """0.3.0 i18n: in an English environment, messages are in English and lang is
    recorded in the JSON. Codes are language-independent."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json"], lang="en").stdout)
    assert doc["lang"] == "en"
    e1 = [e for e in doc["errors"] if e["code"] == "E1"][0]
    assert "Hangul" in e1["message"] and "Malgun" in e1["message"]
    doc_ko = json.loads(run_cli([path, "--json"], lang="ko").stdout)
    assert doc_ko["lang"] == "ko" and "한글" in [e for e in doc_ko["errors"] if e["code"] == "E1"][0]["message"]
    # the --lang flag beats the environment variable
    doc_flag = json.loads(run_cli([path, "--json", "--lang", "en"], lang="ko").stdout)
    assert doc_flag["lang"] == "en"


def test_lang_catalog_consistency():
    """Catalog discipline: every entry has both ko/en, and the % format-specifier
    order is identical."""
    fmt = re.compile(r"%[-#0-9.]*[sdfr%]")
    for mid, entry in jmsg.MESSAGES.items():
        assert set(entry) == {"ko", "en"}, mid
        assert fmt.findall(entry["ko"]) == fmt.findall(entry["en"]), mid


def test_profile_core_drops_style_rules(tmp_path):
    """The core profile = objective defects only: even E2 (a stylistic ERROR) is
    excluded, but the choice is left in the JSON."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이건" + EM_DASH + "차단", font="Wanted Sans", size=12)   # E2 only
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json", "--profile", "full"]).stdout)
    assert any(e["code"] == "E2" for e in doc["errors"])           # blocked when full is specified explicitly (0.4.0: default is core)
    doc = json.loads(run_cli([path, "--json", "--profile", "core"]).stdout)
    assert not doc["errors"] and doc["summary"]["pass"]
    assert doc["summary"]["profile"] == "core"
    assert "E2" in doc["summary"]["skipped_codes"]                  # not a silent bypass

    # objective defects (E1) still get blocked even under core
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    doc = json.loads(run_cli([save(p, tmp_path, "fx2.pptx"), "--json", "--profile", "core"]).stdout)
    assert any(e["code"] == "E1" for e in doc["errors"])


def test_profile_editorial_drops_w14(tmp_path):
    p = new_prs()
    for t in ("시장 현황", "경쟁 구도 분석", "제품 라인업 개요", "사업 확장 전략",
              "재무 운용 계획", "향후 추진 방안"):
        s = add_slide(p)
        tb(s, 1, 0.8, 9, 0.8, t, font="Wanted Sans", size=26)
        tb(s, 1, 2.2, 10, 3, "본문 내용", font="Wanted Sans", size=12)
    path = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([path, "--json", "--profile", "editorial"]).stdout)
    assert not any(w["code"] == "W14" for w in doc["warnings"])
    assert doc["summary"]["profile"] == "editorial"


# ---------------------------------------------------------------- 0.3.1 (third external review P0/P1)
def test_e1_para_defrpr_inheritance(tmp_path):
    """P0-1 (Probe 7 measured): the paragraph pPr/defRPr font is inherited at the
    rank right after run, and wins over lstStyle.
    Case A: paragraph ea=Hangul font -> there must be no E1 false positive.
    Case B: master lstStyle ea=Hangul font but paragraph ea=Consolas -> the
    paragraph wins, so there must be no E1 missed detection."""
    def add_para_ea(para, ea):
        pPr = para._p.get_or_add_pPr()
        defR = pPr.makeelement(qn("a:defRPr"), {})
        defR.append(defR.makeelement(qn("a:ea"), {"typeface": ea}))
        pPr.append(defR)

    # Case A: verify no false positive
    p = new_prs()
    s = add_slide(p)
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.6))
    para = box.text_frame.paragraphs[0]
    add_para_ea(para, "맑은 고딕")
    r = para.add_run(); r.text = "문단 상속 한글"; r.font.size = Pt(14)
    errors, _w = lint_full(save(p, tmp_path, "a.pptx"))
    assert not by_code(errors, "E1"), errors

    # Case B: verify no missed detection (master lstStyle Hangul font + paragraph Consolas)
    p = Presentation()
    for ph in p.slide_masters[0].placeholders:
        if ph.placeholder_format.idx == 1:
            txBody = ph.text_frame._txBody
            lst = txBody.find(qn("a:lstStyle"))
            if lst is None:
                lst = txBody.makeelement(qn("a:lstStyle"), {})
                txBody.insert(1, lst)
            lvl1 = lst.makeelement(qn("a:lvl1pPr"), {})
            defR = lst.makeelement(qn("a:defRPr"), {})
            defR.append(defR.makeelement(qn("a:ea"), {"typeface": "맑은 고딕"}))
            lvl1.append(defR)
            lst.append(lvl1)
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text_frame.text = ""
    body = s.placeholders[1]
    para = body.text_frame.paragraphs[0]
    add_para_ea(para, "Consolas")
    r = para.add_run(); r.text = "문단이 이기는 한글"; r.font.size = Pt(14)
    errors, _w = lint_full(save(p, tmp_path, "b.pptx"))
    assert any("Consolas" in d for (_si, m, d) in by_code(errors, "E1")), errors


def test_geometry_uses_style_resolver(tmp_path):
    """P0-2: geometry checks use the same inherited size as E3. If layout lstStyle
    40pt text overflows on the right, it's W16 (the old version computed with a
    default 12pt and missed it)."""
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    lay = p.slide_layouts[1]
    for ph in lay.placeholders:
        if ph.placeholder_format.idx == 1:
            txBody = ph.text_frame._txBody
            lst = txBody.find(qn("a:lstStyle"))
            if lst is None:
                lst = txBody.makeelement(qn("a:lstStyle"), {})
                txBody.insert(1, lst)
            lvl1 = lst.makeelement(qn("a:lvl1pPr"), {})
            defR = lst.makeelement(qn("a:defRPr"), {"sz": "4000"})
            lvl1.append(defR)
            lst.append(lvl1)
    s = p.slides.add_slide(lay)
    s.shapes.title.text_frame.text = ""
    body = s.placeholders[1]
    body.left, body.top = Inches(12.0), Inches(2.0)
    body.width, body.height = Inches(1.2), Inches(0.8)
    body.text_frame.word_wrap = False
    body.text_frame.text = "Revenue ABC"   # at 12pt it's 0.95in (passes), at 40pt it's 3.2in (overflows)
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert any("텍스트" in d for (_si, m, d) in by_code(warns, "W16")), warns


def test_table_cell_geometry(tmp_path):
    """P0-3: native table-cell text is included in geometry checks (a table straying
    off the right -> W16)."""
    p = new_prs()
    s = add_slide(p)
    gf = s.shapes.add_table(1, 2, Inches(11.0), Inches(2.0), Inches(4.0), Inches(1.0))
    cell = gf.table.cell(0, 1)   # the right column starts around x=13in -> off canvas
    cell.text = "overflowing table cell text"
    r = cell.text_frame.paragraphs[0].runs[0]
    r.font.size = Pt(24)
    cell.text_frame.word_wrap = False
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert any("텍스트" in d for (_si, m, d) in by_code(warns, "W16")), warns


def test_render_dir_contract(tmp_path):
    """P0-4: a missing --render folder or missing page renders are surfaced as
    incomplete; a deck with no pictures is unaffected."""
    p = new_prs()
    s = add_slide(p)
    s.shapes.add_picture(png(tmp_path, "pic.png"), Inches(2), Inches(2), Inches(3), Inches(2))
    tb(s, 1, 1, 5, 0.5, "clean text", font="Wanted Sans", size=20)
    deck = save(p, tmp_path, "fx.pptx")
    doc = json.loads(run_cli([deck, "--json", "--render",
                              os.path.join(str(tmp_path), "no_such_dir")]).stdout)
    assert doc["summary"]["incomplete"] is True, doc["summary"]
    empty = os.path.join(str(tmp_path), "empty_pages")
    os.makedirs(empty)
    doc = json.loads(run_cli([deck, "--json", "--render", empty]).stdout)
    assert doc["summary"]["incomplete"] is True, doc["summary"]   # missing render for a page that has a picture
    # a deck with no pictures is complete even with no render (there's no W7 target to check at all)
    p2 = new_prs()
    s2 = add_slide(p2)
    tb(s2, 1, 1, 5, 0.5, "text only", font="Wanted Sans", size=20)
    doc = json.loads(run_cli([save(p2, tmp_path, "t.pptx"), "--json", "--render", empty]).stdout)
    assert doc["summary"]["incomplete"] is False, doc["summary"]


def test_profile_is_engine_policy(tmp_path, monkeypatch):
    """P0-5: a profile is an execution policy. It can be used as a library via
    lint(profile=), and since an excluded rule simply doesn't run, its internal
    failure doesn't leak into W18."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이건" + EM_DASH + "차단", font="Wanted Sans", size=12)
    path = save(p, tmp_path, "fx.pptx")
    errors, _w = lint_full(path, profile="core")
    assert not by_code(errors, "E2"), errors            # profile works through the library API
    errors, _w = lint_full(path)
    assert by_code(errors, "E2"), errors                # default full stays blocking

    def boom(*a, **kw):
        raise RuntimeError("w9 fail")
    monkeypatch.setattr(jl, "accent_vbars_check", boom)
    _e, warns = lint_full(path, profile="core")           # W9 excluded -> doesn't run -> no W18 leak
    assert not any("w9" in d for (_si, m, d) in by_code(warns, "W18")), warns
    _e, warns = lint_full(path)                           # full -> runs -> guard -> W18
    assert any("w9" in d for (_si, m, d) in by_code(warns, "W18")), warns


def test_skip_validation_strengthened(tmp_path):
    """P1: a nonexistent W code or suppressing W18 causes exit 2 rejection."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    path = save(p, tmp_path, "fx.pptx")
    assert run_cli([path, "--skip", "W51"]).returncode == 2
    assert run_cli([path, "--skip", "W18"]).returncode == 2
    assert run_cli([path, "--skip", "W14"]).returncode == 0


def test_e4_hanja_message(tmp_path):
    """E4 fires on mixed Hangul+Hanja runs (Korean names/legal terms) and the message
    names the judged scope; a Hanja-only run is exempt since 0.6.1 (tracked hanzi is
    legitimate Chinese typography, external review)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "大韓民國 헌법", font="맑은 고딕", size=12, spc=100)
    errors, _w = lint_full(save(p, tmp_path, "fx.pptx"))
    e4 = by_code(errors, "E4")
    assert e4 and "한자" in e4[0][1], e4
    p2 = new_prs()
    s2 = add_slide(p2)
    tb(s2, 1, 1, 5, 0.5, "中华人民共和国", font="맑은 고딕", size=12, spc=200)
    errors2, _w2 = lint_full(save(p2, tmp_path, "fx2.pptx"))
    assert not by_code(errors2, "E4"), errors2


def test_w8_extended_hangul(tmp_path):
    """P1: with is_cjk unified, halfwidth Hangul is also caught by the W8
    small-CJK judgment."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 2, 0.4, chr(0xFFA1) * 6, font="Wanted Sans", size=6)   # halfwidth giyeok x6
    _e, warns = lint_full(save(p, tmp_path, "fx.pptx"))
    assert by_code(warns, "W8"), warns


def test_ghost_prefers_title_placeholder(tmp_path):
    """When a title placeholder exists, a bigger KPI big-number must not push the
    ghost title out (third review)."""
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text_frame.text = "시장 규모는 빠르게 확대된다"
    for r in s.shapes.title.text_frame.paragraphs[0].runs:
        r.font.size = Pt(26)
    s.placeholders[1].text_frame.text = ""
    box = s.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1.5))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "300억"
    r.font.size = Pt(60)
    ghost = []
    lint_full(save(p, tmp_path, "fx.pptx"), ghost=ghost)
    assert ghost and "시장 규모" in ghost[0][1], ghost


def test_json_schema_contract(tmp_path):
    """Start of the JSON version contract (third review): schema_version, tool,
    target_renderer."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    doc = json.loads(run_cli([save(p, tmp_path, "fx.pptx"), "--json"]).stdout)
    assert doc["schema_version"] == "1.0"
    assert doc["tool"]["name"] == "archforge" and doc["tool"]["version"]
    assert doc["target_renderer"] == "powerpoint-windows"


def test_w6_detail_english(tmp_path):
    """P1: i18n extends even to detail (in English mode, W6's detail starts with "e.g.")."""
    p = new_prs()
    for i in range(6):
        s = add_slide(p)
        tb(s, 1, 0.8, 8, 0.6, "Title block %d" % i, font="Wanted Sans", size=24)
        tb(s, 1, 2.0, 6, 2.5, "Body block", font="Wanted Sans", size=12)
        tb(s, 8, 2.0, 4, 2.5, "Side block", font="Wanted Sans", size=12)
    doc = json.loads(run_cli([save(p, tmp_path, "fx.pptx"), "--json", "--profile", "full"], lang="en").stdout)
    w6 = [w for w in doc["warnings"] if w["code"] == "W6"]
    assert w6 and w6[0]["detail"].startswith("e.g."), w6


# ---------------------------------------------------------------- 0.4.0 structural overhaul
def test_default_profile_is_core(tmp_path):
    """0.4.0 breaking change: with no options, the default is core (objective
    defects only). The style rule E2 is opt-in; the objective defect E1 stays
    blocked even under the default."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이건" + EM_DASH + "정상 통과", font="Wanted Sans", size=12)
    dash_deck = save(p, tmp_path, "dash.pptx")
    errors, _w = jl.lint(dash_deck)              # library default
    assert not by_code(errors, "E2"), errors
    r = run_cli([dash_deck, "--json"])           # CLI default
    doc = json.loads(r.stdout)
    assert r.returncode == 0 and doc["summary"]["profile"] == "core"
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    assert run_cli([save(p, tmp_path, "e1.pptx")]).returncode == 1   # objective defects stay blocked


def test_config_file(tmp_path):
    """.archforge.json: auto-discovered in the deck folder; a CLI flag beats the
    config."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "이건" + EM_DASH + "차단", font="Wanted Sans", size=12)
    deck = save(p, tmp_path, "fx.pptx")
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({"profile": "full"}, f)
    doc = json.loads(run_cli([deck, "--json"]).stdout)
    assert any(e["code"] == "E2" for e in doc["errors"])    # the config's full is applied
    doc = json.loads(run_cli([deck, "--json", "--profile", "core"]).stdout)
    assert not doc["errors"]                                 # CLI beats the config
    # fail-safe (fourth review): an unknown key (a typo) causes exit 2 rather than
    # being ignored. Prevents the accident where 'profle: full' silently runs as
    # the default core
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({"profile": "full", "no_such_key": 1}, f)
    r = run_cli([deck, "--json"])
    assert r.returncode == 2 and "no_such_key" in r.stderr
    # --no-config: ignore the config of an untrusted deck folder (a trust boundary)
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({"profile": "full"}, f)
    doc = json.loads(run_cli([deck, "--json", "--no-config"]).stdout)
    assert not doc["errors"] and doc["summary"]["config"] is None
    doc = json.loads(run_cli([deck, "--json"]).stdout)
    assert doc["summary"]["config"] and doc["summary"]["config"].endswith(".archforge.json")


def test_config_value_validation(tmp_path):
    """A config value's type/range gets a clean exit 2 rather than a traceback
    (fourth review). Same for CLI range (closes the old X1 bypass where
    --hard-min 0 silently turned off E3)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    deck = save(p, tmp_path, "fx.pptx")
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({"hard_min": "abc"}, f)
    r = run_cli([deck, "--json"])
    assert r.returncode == 2 and "hard_min" in r.stderr and "Traceback" not in r.stderr
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({"lang": "fr"}, f)
    assert run_cli([deck, "--json"]).returncode == 2
    with open(os.path.join(str(tmp_path), ".archforge.json"), "w", encoding="utf-8") as f:
        json.dump({}, f)
    assert run_cli([deck, "--hard-min", "0"]).returncode == 2   # CLI range validation


def test_baseline_v2_language_and_count(tmp_path):
    """Fingerprint v2 (fourth review HIGH fix): a baseline made under ko is also
    valid under an en run, preserves the meaning of the occurrence count, and is
    page-independent so it survives a slide insertion."""
    p = new_prs()
    for i in range(6):   # a W6 clone deck (the representative rule whose detail used to be a locale string)
        s = add_slide(p)
        tb(s, 1, 0.8, 8, 0.6, "Title block %d" % i, font="Wanted Sans", size=24)
        tb(s, 1, 2.0, 6, 2.5, "Body block", font="Wanted Sans", size=12)
        tb(s, 8, 2.0, 4, 2.5, "Side block", font="Wanted Sans", size=12)
    deck = save(p, tmp_path, "w6.pptx")
    bl = os.path.join(str(tmp_path), "bl.json")
    run_cli([deck, "--profile", "full", "--write-baseline", bl], lang="ko")
    doc = json.loads(run_cli([deck, "--profile", "full", "--json", "--baseline", bl],
                             lang="en").stdout)
    assert not any(w["code"] == "W6" for w in doc["warnings"]), doc["warnings"]

    # count semantics: the same fingerprint occurs 2 times, baseline has 1 -> only 1 is suppressed
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    one = save(p, tmp_path, "one.pptx")
    run_cli([one, "--write-baseline", bl])
    s = add_slide(p)   # same text, same font = the same fingerprint on two pages
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    two = save(p, tmp_path, "two.pptx")
    doc = json.loads(run_cli([two, "--json", "--baseline", bl]).stdout)
    assert doc["summary"]["baseline_suppressed"] == 1
    assert doc["summary"]["error_count"] == 1
    # page-independent: suppression persists even in a deck with a slide inserted before it
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "새로 삽입된 표지", font="Wanted Sans", size=20)
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    shifted = save(p, tmp_path, "shifted.pptx")
    doc = json.loads(run_cli([shifted, "--json", "--baseline", bl]).stdout)
    assert doc["summary"]["baseline_suppressed"] == 1 and doc["summary"]["error_count"] == 0
    # text-mode visibility: suppression shows as a footnote (fixes the misreading of invisible-as-clean)
    r = run_cli([shifted, "--baseline", bl])
    assert "baseline" in r.stdout


def test_lint_rejects_unknown_profile(tmp_path):
    """Library API: fixes a typo'd profile that used to silently behave as full
    (fourth review)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "clean", font="Wanted Sans", size=20)
    deck = save(p, tmp_path, "fx.pptx")
    with pytest.raises(ValueError):
        jl.lint(deck, profile="ful")


def test_baseline_flow(tmp_path):
    """baseline: accept existing violations, then report only new ones. The
    suppressed count is recorded in summary."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    deck = save(p, tmp_path, "fx.pptx")
    bl = os.path.join(str(tmp_path), "baseline.json")
    r = run_cli([deck, "--write-baseline", bl])
    assert r.returncode == 0 and os.path.exists(bl)
    doc = json.loads(run_cli([deck, "--json", "--baseline", bl]).stdout)
    assert doc["summary"]["pass"] and doc["summary"]["baseline_suppressed"] == 1
    # add a new defect -> only the new one is reported
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "새 결함 한글", font="Consolas", size=12)
    deck2 = save(p, tmp_path, "fx2.pptx")
    doc = json.loads(run_cli([deck2, "--json", "--baseline", bl]).stdout)
    assert doc["summary"]["error_count"] == 1
    assert "Consolas" in doc["errors"][0]["detail"]


def test_sarif_output(tmp_path):
    """SARIF 2.1.0 minimal contract: version, rules, results, ruleId, level."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    deck = save(p, tmp_path, "fx.pptx")
    out = os.path.join(str(tmp_path), "out.sarif")
    run_cli([deck, "--json", "--sarif", out])
    with open(out, encoding="utf-8") as f:
        doc = json.load(f)
    assert doc["version"] == "2.1.0"
    run0 = doc["runs"][0]
    assert run0["tool"]["driver"]["name"] == "archforge"
    res = run0["results"]
    assert res and res[0]["ruleId"] == "E1" and res[0]["level"] == "error"
    assert any(r["id"] == "E1" for r in run0["tool"]["driver"]["rules"])


def test_finding_location_payload(tmp_path):
    """Structured location (third review): shape_id, bbox, part, paragraph, run are
    carried in the JSON location."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 5, 0.5, "모노 폴백 한글", font="IBM Plex Mono", size=12)
    doc = json.loads(run_cli([save(p, tmp_path, "fx.pptx"), "--json"]).stdout)
    e1 = [e for e in doc["errors"] if e["code"] == "E1"][0]
    loc = e1["location"]
    assert isinstance(loc["shape_id"], int)
    assert loc["part"].endswith("slide1.xml")
    assert len(loc["bbox"]) == 4 and abs(loc["bbox"][0] - 1.0) < 0.01
    assert loc["paragraph"] == 0 and loc["run"] == 0


def test_finding_tuple_compat_and_lazy_message():
    """Finding's 4-tuple backward compatibility and locale neutrality (rendering the
    same finding in two languages)."""
    from archforge.findings import Finding
    f = Finding(3, "E2", "e2", (), "cp=U+2014")
    page, code, msg, detail = f
    assert (page, code, detail) == (3, "E2", "cp=U+2014") and f[1] == "E2" and len(f) == 4
    jmsg.set_lang("ko")
    ko = f.message
    jmsg.set_lang("en")
    en = f.message
    jmsg.set_lang("ko")
    assert ko != en and "Dash" in en


# ---------------------------------------------------------------- skill packaging
def _repo_root():
    return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def test_skill_pack_sync():
    """The repo-root skills/ (for discovery) and the package-bundled copy (the
    canonical one) must not drift apart."""
    root = _repo_root()
    a = os.path.join(root, "skills", "archforge-pptx-lint", "SKILL.md")
    b = os.path.join(root, "src", "archforge", "skills", "archforge-pptx-lint", "SKILL.md")
    with open(a, "rb") as fa, open(b, "rb") as fb:
        assert fa.read() == fb.read(), "루트 skills/ 사본과 패키지 정본이 다릅니다"


def test_skill_frontmatter_name_matches_dir():
    """Agent Skills spec: frontmatter name == the skill directory name (regression
    for an external review finding)."""
    root = _repo_root()
    path = os.path.join(root, "skills", "archforge-pptx-lint", "SKILL.md")
    with open(path, encoding="utf-8") as f:
        head = f.read(500)
    m = re.search(r"^name:\s*(\S+)", head, re.M)
    assert m and m.group(1) == "archforge-pptx-lint"


def test_cli_skill_subcommand(tmp_path):
    """archforge skill: both the output and install paths match the package-bundled
    copy."""
    r = run_cli(["skill"])
    assert r.returncode == 0
    assert "name: archforge-pptx-lint" in r.stdout
    r = run_cli(["skill", "--install", str(tmp_path)])
    assert r.returncode == 0
    dst = os.path.join(str(tmp_path), "archforge-pptx-lint", "SKILL.md")
    assert os.path.exists(dst)
    with open(dst, encoding="utf-8") as f:
        installed = f.read()
    assert "name: archforge-pptx-lint" in installed


# ---------------------------------------------------------------- 0.5.0: loc reinforcement + a:fld/a:br
def _by(items, code):
    return [f for f in items if f.code == code]


def _add_fld(para, text, sz=None, latin=None, spc=None):
    """Plant an a:fld (auto field) directly into the paragraph. python-pptx runs
    skip this element, so it's done via raw XML."""
    fld = para._p.makeelement(qn("a:fld"), {
        "id": "{93C0AD05-B0B5-4E56-9A2A-9F2113D1B94A}", "type": "slidenum"})
    attrs = {}
    if sz is not None:
        attrs["sz"] = str(sz)
    if spc is not None:
        attrs["spc"] = str(spc)
    rPr = fld.makeelement(qn("a:rPr"), attrs)
    if latin:
        rPr.append(rPr.makeelement(qn("a:latin"), {"typeface": latin}))
    fld.append(rPr)
    t = fld.makeelement(qn("a:t"), {})
    t.text = text
    fld.append(t)
    para._p.append(fld)
    return fld


def test_fld_field_runs_gated(tmp_path):
    """a:fld also renders with the same rPr as a normal run, so it's subject to
    the same gates (carried over from the fourth review): an Arial Hangul field
    is E1, a positive-tracking field is E4. loc carries a field marker, with no
    run index."""
    p = new_prs()
    s = add_slide(p)
    box = s.shapes.add_textbox(Inches(1), Inches(6.8), Inches(3), Inches(0.4))
    _add_fld(box.text_frame.paragraphs[0], "3 페이지", sz=1200, latin="Arial", spc=200)
    errors, _w = lint_full(save(p, tmp_path, "fld.pptx"))
    e1 = _by(errors, "E1")
    assert e1, codes(errors)
    assert e1[0].loc.get("field") is True and "run" not in e1[0].loc
    assert _by(errors, "E4"), codes(errors)


def test_br_offsets_and_fld_not_title(tmp_path):
    """a:br contributes only a single line-break character in the E2 context
    (offset preserved: an em dash in the run after br is still caught). A large
    auto field (a 40pt page number) is not collected as a ghost title."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 1, 1, 8, 1, "첫 줄", size=14)
    para = box.text_frame.paragraphs[0]
    para._p.append(para._p.makeelement(qn("a:br"), {}))
    r2 = para.add_run()
    r2.text = "결론" + EM_DASH + "정리"
    r2.font.size = Pt(14)
    box2 = s.shapes.add_textbox(Inches(11), Inches(6.5), Inches(2), Inches(0.8))
    _add_fld(box2.text_frame.paragraphs[0], "01", sz=4000)
    ghost = []
    errors, _w = lint_full(save(p, tmp_path, "br.pptx"), ghost=ghost)
    assert "E2" in codes(errors), codes(errors)
    assert not any("01" in str(g) for g in ghost), ghost


def test_table_cell_loc(tmp_path):
    """A table-cell finding's loc.cell == [row, col], 0-based (carried over from
    the fourth review)."""
    p = new_prs()
    s = add_slide(p)
    gfx = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    cell = gfx.table.cell(1, 0)
    cell.text = "한글"
    r = cell.text_frame.paragraphs[0].runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(12)
    errors, _w = lint_full(save(p, tmp_path, "tbl_loc.pptx"))
    e1 = _by(errors, "E1")
    assert e1, codes(errors)
    assert e1[0].loc.get("cell") == [1, 0], e1[0].loc


def test_group_child_loc_bbox_absolute(tmp_path):
    """Under group-move desync (off!=chOff), a run-level finding's loc bbox is in
    slide absolute coordinates, not the group's chOff coordinate system (carried
    over from the fourth review). Child raw x=12.5in, group moved to 7.5in ->
    absolute x is 7.5in."""
    p = new_prs()
    s = add_slide(p)
    grp = s.shapes.add_group_shape()
    box = grp.shapes.add_textbox(Inches(12.5), Inches(3.0), Inches(1.5), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "깨알 주석"
    r.font.size = Pt(3)   # E3
    grp.left, grp.top = Inches(7.5), Inches(3.0)
    errors, _w = lint_full(save(p, tmp_path, "grp_loc.pptx"))
    e3 = _by(errors, "E3")
    assert e3, codes(errors)
    bbox = e3[0].loc.get("bbox")
    assert bbox and abs(bbox[0] - 7.5) < 0.05, e3[0].loc


def test_w15_w16_locations(tmp_path):
    """W15/W16 carry loc (the effective glyph's absolute bbox), and W15 identifies
    the counterpart frame via related (carried over from the fourth review: W15-17
    lacked location)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1.0, 1.0, 5.0, 1.0, "겹침 판정 대상 문장 하나", size=24)
    tb(s, 1.1, 1.05, 5.0, 1.0, "겹침 판정 대상 문장 둘", size=24)
    tb(s, 12.8, 3.0, 3.0, 0.5, "화면 밖으로 넘어가는 문장", size=18)
    _e, warns = lint_full(save(p, tmp_path, "w15loc.pptx"))
    w15 = _by(warns, "W15")
    assert w15, codes(warns)
    assert w15[0].loc and "bbox" in w15[0].loc and "related" in w15[0].loc, w15[0].loc
    assert "bbox" in w15[0].loc["related"]
    w16 = _by(warns, "W16")
    assert w16, codes(warns)
    assert w16[0].loc and "bbox" in w16[0].loc, w16[0].loc


# ---------------------------------------------------------------- 0.5.0: scan/demo subcommands
def test_cli_demo_and_scan(tmp_path):
    """demo: generates broken (fires 6 kinds of defects) + fixed (clean) and lints
    them on the spot, exit 0.
    scan: recursive directory aggregate JSON, exit 1 if even one fails, exit 2 if
    there are 0 matches."""
    d = os.path.join(str(tmp_path), "demo")
    r = run_cli(["demo", "--dir", d])
    assert r.returncode == 0, r.stderr
    assert os.path.exists(os.path.join(d, "broken.pptx"))
    assert os.path.exists(os.path.join(d, "fixed.pptx"))
    assert "ERROR 0, WARN 0" in r.stdout   # fixed must be clean (demo contract)

    r = run_cli(["scan", d, "--profile", "full", "--json"])
    assert r.returncode == 1, r.stderr
    doc = json.loads(r.stdout)
    assert doc["summary"]["file_count"] == 2
    assert doc["summary"]["failed_files"] == 1
    assert doc["summary"]["pass"] is False
    names = {os.path.basename(f["file"]) for f in doc["files"]}
    assert names == {"broken.pptx", "fixed.pptx"}
    for fdoc in doc["files"]:
        assert fdoc["summary"]["profile"] == "full"

    # fixed only: exit 0 + a text summary line
    r = run_cli(["scan", os.path.join(d, "fixed.pptx"), "--profile", "full"])
    assert r.returncode == 0, r.stdout + r.stderr
    assert "1" in r.stdout and "0" in r.stdout   # scan summary (1 file, 0 failed)

    # 0 glob matches = exit 2, not a silent pass (prevents a CI footgun)
    r = run_cli(["scan", os.path.join(str(tmp_path), "none_*.pptx")])
    assert r.returncode == 2


def test_examples_contract(tmp_path):
    """Pins the documented contract for examples/: style_warnings is clean under
    core, and W13/W14 under full."""
    from archforge import demo as jdemo
    p = os.path.join(str(tmp_path), "warn.pptx")
    jdemo.build_warnings(p)
    e, w = jl.lint(p, profile="core")
    assert not e and not w, (codes(e), codes(w))
    e, w = lint_full(p)
    assert not e, codes(e)
    assert {"W13", "W14"} <= set(codes(w)), codes(w)


def test_demo_en_variant(tmp_path):
    """English-edition demo deck contract: monolingual English (user decision), so it
    seeds only the script-independent defects. broken_en is exactly E2+E3 plus W15/W16;
    the Hangul-only defects (E1/E4) are the Korean deck's job. fixed_en is clean under
    full. The README(en) assets are this deck's actual renders."""
    from archforge import demo as jdemo
    b = os.path.join(str(tmp_path), "b.pptx")
    fx = os.path.join(str(tmp_path), "f.pptx")
    jdemo.build_broken(b, lang="en")
    jdemo.build_fixed(fx, lang="en")
    e, w = lint_full(b)
    assert set(codes(e)) == {"E2", "E3"}, codes(e)
    assert {"W15", "W16"} <= set(codes(w)), codes(w)
    e, w = lint_full(fx)
    assert not e and not w, (codes(e), codes(w))


def test_cli_scan_sarif_multi(tmp_path):
    """scan --sarif: multiple files are merged into one SARIF run, each with its
    own per-file artifactLocation."""
    d = os.path.join(str(tmp_path), "demo")
    r = run_cli(["demo", "--dir", d])
    assert r.returncode == 0, r.stderr
    sarif_path = os.path.join(str(tmp_path), "out.sarif")
    r = run_cli(["scan", d, "--profile", "full", "--sarif", sarif_path])
    assert r.returncode == 1
    with open(sarif_path, encoding="utf-8") as f:
        doc = json.load(f)
    uris = {res["locations"][0]["physicalLocation"]["artifactLocation"]["uri"]
            for res in doc["runs"][0]["results"]}
    assert any(u.endswith("broken.pptx") for u in uris), uris


# ---------------------------------------------------------------- 0.6.0: hardening batch
def test_nan_thresholds_rejected(tmp_path):
    """NaN slipped through every range comparison (NaN <= 0 is False) and json.load even
    accepts a bare NaN literal, silently disabling E3 via CLI or an attacker-controlled
    config. Both paths must exit 2 now (external verification finding)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "tiny", size=3)
    deck = save(p, tmp_path, "nan.pptx")
    r = run_cli([deck, "--hard-min", "nan"])
    assert r.returncode == 2, (r.returncode, r.stderr)
    cfgdir = os.path.join(str(tmp_path), "cfg")
    os.makedirs(cfgdir)
    shutil.copy(deck, os.path.join(cfgdir, "deck.pptx"))
    with open(os.path.join(cfgdir, ".archforge.json"), "w", encoding="utf-8") as f:
        f.write('{"hard_min": NaN}')
    r = run_cli([os.path.join(cfgdir, "deck.pptx")])
    assert r.returncode == 2, (r.returncode, r.stderr)


def test_scan_isolates_broken_files(tmp_path):
    """One corrupt deck must not abort the batch (external review P0): it becomes a
    per-file error entry, the rest is still checked, and the aggregate JSON survives."""
    d = os.path.join(str(tmp_path), "batch")
    os.makedirs(d)
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "clean deck", size=14)
    save(p, d, "ok.pptx")
    with open(os.path.join(d, "corrupt.pptx"), "w") as f:
        f.write("not a zip")
    r = run_cli(["scan", d, "--json"])
    assert r.returncode == 1, r.stderr
    doc = json.loads(r.stdout)
    assert doc["summary"]["file_count"] == 2
    assert doc["summary"]["error_files"] == 1
    statuses = {os.path.basename(f["file"]): f["status"] for f in doc["files"]}
    assert statuses["corrupt.pptx"] == "error"
    assert statuses["ok.pptx"] == "pass"


def test_scan_rejects_shared_baseline(tmp_path):
    """A single CLI baseline across multiple decks would suppress findings across
    unrelated files (fingerprints carry no file identity): exit 2."""
    d = os.path.join(str(tmp_path), "two")
    os.makedirs(d)
    for name in ("a.pptx", "b.pptx"):
        p = new_prs()
        s = add_slide(p)
        tb(s, 1, 1, 4, 1, "x", size=14)
        save(p, d, name)
    bl = os.path.join(str(tmp_path), "bl.json")
    with open(bl, "w") as f:
        f.write('{"schema_version": "2", "findings": []}')
    r = run_cli(["scan", d, "--baseline", bl])
    assert r.returncode == 2, (r.returncode, r.stdout, r.stderr)


def test_scan_config_lang_does_not_leak(tmp_path):
    """A deck folder config's lang must not change the language of later files or the
    aggregate report (external verification finding: lazy rendering made the last
    file's language dominate everything)."""
    d1 = os.path.join(str(tmp_path), "en_deck")
    d2 = os.path.join(str(tmp_path), "plain")
    os.makedirs(d1)
    os.makedirs(d2)
    for d in (d1, d2):
        p = new_prs()
        s = add_slide(p)
        tb(s, 1, 1, 4, 1, "x", size=14)
        save(p, d, "deck.pptx")
    with open(os.path.join(d1, ".archforge.json"), "w") as f:
        f.write('{"lang": "en"}')
    r = run_cli(["scan", os.path.join(d1, "deck.pptx"), os.path.join(d2, "deck.pptx"),
                 "--json"], lang="ko")
    assert r.returncode == 0, r.stderr
    doc = json.loads(r.stdout)
    assert doc["lang"] == "ko", doc["lang"]


def test_strict_split_flags(tmp_path):
    """--fail-on-warning fails a WARN-only deck; the default does not. --strict stays
    the union alias."""
    from archforge import demo as jdemo
    p = os.path.join(str(tmp_path), "warn.pptx")
    jdemo.build_warnings(p)
    r = run_cli([p, "--profile", "full"])
    assert r.returncode == 0, r.stdout + r.stderr
    r = run_cli([p, "--profile", "full", "--fail-on-warning"])
    assert r.returncode == 1
    r = run_cli([p, "--profile", "full", "--strict"])
    assert r.returncode == 1
    doc = json.loads(run_cli([p, "--profile", "full", "--fail-on-warning", "--json"]).stdout)
    assert doc["summary"]["pass"] is False


def test_write_baseline_validates_flags_first(tmp_path):
    """A typo'd --skip used to record a baseline as if nothing were wrong: validation
    now precedes recording (external review)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "v.pptx")
    bl = os.path.join(str(tmp_path), "out_bl.json")
    r = run_cli([deck, "--write-baseline", bl, "--skip", "E1"])
    assert r.returncode == 2
    assert not os.path.exists(bl)
    r = run_cli([deck, "--write-baseline", bl, "--skip", "W999"])
    assert r.returncode == 2
    assert not os.path.exists(bl)


def test_config_no_config_conflict(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "c.pptx")
    cfg = os.path.join(str(tmp_path), "x.json")
    with open(cfg, "w") as f:
        f.write("{}")
    r = run_cli([deck, "--config", cfg, "--no-config"])
    assert r.returncode == 2


def test_version_flag():
    r = run_cli(["--version"])
    assert r.returncode == 0
    assert "archforge" in r.stdout


def test_br_starts_new_visual_line_in_geometry(tmp_path):
    """An explicit a:br splits the visual line: without the split, this no-wrap line
    would run past the canvas edge and fire W16 (external review: br was in E2 context
    but lost in geometry)."""
    long_half = "wide segment text " * 2
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 8.0, 3.0, 2.0, 1.0, long_half, size=18)
    para = box.text_frame.paragraphs[0]
    para._p.append(para._p.makeelement(qn("a:br"), {}))
    r2 = para.add_run()
    r2.text = long_half
    r2.font.size = Pt(18)
    _e, warns = lint_full(save(p, tmp_path, "br_geo.pptx"))
    assert "W16" not in codes(warns), codes(warns)
    # Control: the same text without the break does overflow
    p2 = new_prs()
    s2 = add_slide(p2)
    tb(s2, 8.0, 3.0, 2.0, 1.0, long_half + " " + long_half, size=18)
    _e2, warns2 = lint_full(save(p2, tmp_path, "br_geo_ctrl.pptx"))
    assert "W16" in codes(warns2), codes(warns2)


def test_insets_shift_glyph_origin(tmp_path):
    """A huge lIns pushes the glyph start right of the frame edge: only inset-aware
    geometry sees this short text cross the canvas (external review: insets ignored)."""
    p = new_prs()
    s = add_slide(p)
    box = tb(s, 11.0, 3.0, 2.3, 0.6, "clipped by inset", size=16)
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(1.5 * 914400)))
    _e, warns = lint_full(save(p, tmp_path, "inset.pptx"))
    assert "W16" in codes(warns), codes(warns)


def test_merged_cells_single_count(tmp_path):
    """Continuation cells of a merged region mirror the origin's text: walking them
    double-counted findings (external review)."""
    p = new_prs()
    s = add_slide(p)
    gfx = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    tbl = gfx.table
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    tbl.cell(0, 0).text = "한글"
    r = tbl.cell(0, 0).text_frame.paragraphs[0].runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(12)
    errors, _w = lint_full(save(p, tmp_path, "merge.pptx"))
    e1 = [f for f in errors if f.code == "E1"]
    assert len(e1) == 1, codes(errors)


def test_zip_preflight_blocks_bomb(tmp_path):
    """A pptx-shaped decompression bomb is a usage error with a stated reason, not a
    hang (external review: no package-level budgets existed)."""
    import zipfile
    bomb = os.path.join(str(tmp_path), "bomb.pptx")
    with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("ppt/huge.bin", b"\x00" * 30_000_000)
    r = run_cli([bomb])
    assert r.returncode == 2
    assert "preflight" in (r.stderr or "")


def test_sarif_rule_metadata_static(tmp_path):
    """SARIF rule titles must be static (no %-placeholders) and results carry
    partialFingerprints for cross-run tracking (external review)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "tiny", size=3)
    deck = save(p, tmp_path, "s.pptx")
    out = os.path.join(str(tmp_path), "o.sarif")
    r = run_cli([deck, "--sarif", out])
    assert r.returncode == 1
    with open(out, encoding="utf-8") as f:
        doc = json.load(f)
    run0 = doc["runs"][0]
    for rule in run0["tool"]["driver"]["rules"]:
        assert "%" not in rule["shortDescription"]["text"], rule
        assert rule["helpUri"].startswith("https://")
    assert all("partialFingerprints" in res for res in run0["results"])


# ---------------------------------------------------------------- 0.6.1: contract batch
def test_scan_empty_pattern_fails(tmp_path):
    """One input matching nothing must not hide behind another that matched: exit 2
    unless --allow-empty-pattern (external review P0)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "ok.pptx")
    ghost_glob = os.path.join(str(tmp_path), "nope", "**", "*.pptx")
    r = run_cli(["scan", deck, ghost_glob])
    assert r.returncode == 2, (r.returncode, r.stderr)
    r = run_cli(["scan", deck, ghost_glob, "--allow-empty-pattern"])
    assert r.returncode == 0, (r.returncode, r.stderr)
    r = run_cli(["scan", deck, "--json"])
    doc = json.loads(r.stdout)
    assert doc["scan"]["inputs"][0]["matches"] == 1


def test_scan_global_usage_error_exits_2(tmp_path):
    """A bad CLI flag is a global usage error (exit 2), not N per-file error entries
    (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "g.pptx")
    for extra in (["--skip", "W999"], ["--skip", "E1"],
                  ["--config", os.path.join(str(tmp_path), "missing.json")],
                  ["--hard-min", "nan"]):
        r = run_cli(["scan", deck] + extra)
        assert r.returncode == 2, (extra, r.returncode, r.stdout, r.stderr)


def test_summary_policy_recorded(tmp_path):
    """The active failure policy travels with the verdict so JSON consumers can tell
    why identical counts pass or fail (external review P0)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "pol.pptx")
    doc = json.loads(run_cli([deck, "--json", "--fail-on-warning"]).stdout)
    assert doc["summary"]["policy"] == {"fail_on_warning": True,
                                        "fail_incomplete": False,
                                        "e2_no_exemptions": False}
    doc = json.loads(run_cli([deck, "--json", "--strict"]).stdout)
    assert all(doc["summary"]["policy"].values())


def test_field_only_complex_script_marks_incomplete(tmp_path):
    """Arabic living entirely inside an a:fld used to bypass the complex-script screen
    (built from para.runs) while still entering width math: now the screen sees field
    text and the frame is skipped into W18 (external review P0)."""
    p = new_prs()
    s = add_slide(p)
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.6))
    _add_fld(box.text_frame.paragraphs[0], "مرحبا بك",
             sz=1800)
    _e, warns = lint_full(save(p, tmp_path, "rtl_fld.pptx"))
    assert "W18" in codes(warns), codes(warns)


def test_w7_unknown_explicit_color_abstains(tmp_path):
    """An explicit color the decoder cannot resolve (hslClr) must stop resolution and
    abstain, not fall through to an inherited color (external review P1: a white hslClr
    run was judged with inherited black = false positive)."""
    from PIL import Image
    p = new_prs()
    s = add_slide(p)
    d = os.path.join(str(tmp_path), "pages")
    os.makedirs(d)
    Image.new("RGB", (1333, 750), (10, 10, 10)).save(os.path.join(d, "p01.png"))
    s.shapes.add_picture(png(tmp_path, "bg.png", size=(800, 500)),
                         Inches(1), Inches(1), Inches(8), Inches(4))
    box = tb(s, 2, 2, 5, 1, "white text via hslClr", size=20)
    r = box.text_frame.paragraphs[0].runs[0]
    rPr = r._r.get_or_add_rPr()
    fill = rPr.makeelement(qn("a:solidFill"), {})
    hsl = rPr.makeelement(qn("a:hslClr"), {"hue": "0", "sat": "0%", "lum": "100%"})
    fill.append(hsl)
    rPr.insert(0, fill)
    _e, warns = lint_full(save(p, tmp_path, "hsl.pptx"), render_dir=d)
    assert "W7" not in codes(warns), codes(warns)
    assert "W18" in codes(warns), codes(warns)


def test_config_rejects_bool_and_fractional_cluster(tmp_path):
    """float(True) == 1.0 slipped through as a threshold, and w6_cluster 1.9 silently
    truncated to 1 (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    for body in ('{"hard_min": true}', '{"w6_cluster": 1.9}'):
        d = os.path.join(str(tmp_path), "c%d" % (abs(hash(body)) % 97))
        os.makedirs(d, exist_ok=True)
        deck = save(p, d, "deck.pptx")
        with open(os.path.join(d, ".archforge.json"), "w") as f:
            f.write(body)
        r = run_cli([deck])
        assert r.returncode == 2, (body, r.returncode, r.stderr)


def test_rules_and_explain_subcommands():
    r = run_cli(["rules"])
    assert r.returncode == 0 and "E1" in r.stdout and "W18" in r.stdout
    r = run_cli(["explain", "w15", "--lang", "en"])
    assert r.returncode == 0 and "overlap" in r.stdout.lower()
    r = run_cli(["explain", "W99"])
    assert r.returncode == 2


def test_module_entry_no_runpy_warning(tmp_path):
    """`python -m archforge` must not emit the runpy RuntimeWarning that
    `python -m archforge.lint` triggers (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "m.pptx")
    env = dict(os.environ)
    env["ARCHFORGE_LANG"] = "en"
    r = subprocess.run([sys.executable, "-m", "archforge", deck], capture_output=True,
                       text=True, encoding="utf-8", stdin=subprocess.DEVNULL, env=env)
    assert r.returncode == 0, r.stderr
    assert "RuntimeWarning" not in (r.stderr or "")


def test_lint_subcommand_alias(tmp_path):
    p = new_prs()
    s = add_slide(p)
    tb(s, 1, 1, 4, 1, "x", size=14)
    deck = save(p, tmp_path, "alias.pptx")
    r = run_cli(["lint", deck, "--json"])
    assert r.returncode == 0, r.stderr
    assert json.loads(r.stdout)["summary"]["pass"] is True


def test_action_runner_rejects_bool_typo(tmp_path):
    """A typo'd boolean input must fail the job, not silently disable a safety default
    (external review P0)."""
    runner = os.path.join(_repo_root(), "action_runner.py")
    env = dict(os.environ)
    env.update({"AF_FILES": "whatever.pptx", "AF_FAIL_INCOMPLETE": "ture"})
    r = subprocess.run([sys.executable, runner], capture_output=True, text=True,
                       encoding="utf-8", env=env, stdin=subprocess.DEVNULL)
    assert r.returncode == 2
    assert "fail-incomplete" in (r.stderr or "")


def test_w15_loc_carries_paragraph_and_cell_fields(tmp_path):
    """Geometry locations gained paragraph indexes and (for fld-only paragraphs) the
    field marker (external review P1)."""
    p = new_prs()
    s = add_slide(p)
    tb(s, 1.0, 1.0, 5.0, 1.0, "overlap sentence one here", size=24)
    tb(s, 1.1, 1.05, 5.0, 1.0, "overlap sentence two here", size=24)
    _e, warns = lint_full(save(p, tmp_path, "w15p.pptx"))
    w15 = [f for f in warns if f.code == "W15"]
    assert w15 and "paragraph" in w15[0].loc, w15[0].loc if w15 else warns


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-q"]))
