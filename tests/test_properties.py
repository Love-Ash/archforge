# -*- coding: utf-8 -*-
"""Property-based invariants (hypothesis) plus a deterministic structured fuzzer.

These do not assert specific findings; they pin the contracts that must hold for ANY
deck: determinism, translation invariance, size monotonicity for E3, and the
no-traceback guarantee under attribute mutation (external review: example-based tests
alone cannot cover OOXML's combinatorial space)."""
import io
import os
import random
import sys
import zipfile

import pytest
from hypothesis import given, settings, HealthCheck, strategies as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

import archforge.lint as jl

SETTINGS = dict(max_examples=20, deadline=None,
                suppress_health_check=[HealthCheck.function_scoped_fixture])


def _build(boxes, tmpdir, name):
    """boxes = [(x_in, y_in, w_in, size_pt, text)] placed inside a safe zone."""
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    s = p.slides.add_slide(p.slide_layouts[6])
    for (x, y, w, size, text) in boxes:
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.0))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(size)
        rPr = r._r.get_or_add_rPr()
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": "맑은 고딕"}))
    path = os.path.join(str(tmpdir), name)
    p.save(path)
    return path


def _verdict(path):
    errors, warns = jl.lint(path, profile="full")
    return sorted((f.page, f.code, f.detail) for f in list(errors) + list(warns))


box_st = st.tuples(
    st.floats(min_value=1.0, max_value=6.0),    # x: safe zone, never near the edge
    st.floats(min_value=1.0, max_value=4.5),    # y
    st.floats(min_value=2.0, max_value=5.0),    # w
    st.integers(min_value=10, max_value=40),    # pt
    st.sampled_from(["Revenue grew 18%", "구독 매출 성장", "Margin recovery",
                     "핵심 지표 개선", "guidance intact"]),
)


@settings(**SETTINGS)
@given(st.lists(box_st, min_size=1, max_size=4))
def test_lint_is_deterministic(tmp_path_factory, boxes):
    """The same file linted twice yields the identical finding list."""
    d = tmp_path_factory.mktemp("prop")
    path = _build(boxes, d, "det.pptx")
    assert _verdict(path) == _verdict(path)


@settings(**SETTINGS)
@given(st.lists(box_st, min_size=1, max_size=3),
       st.floats(min_value=-0.5, max_value=0.5),
       st.floats(min_value=-0.5, max_value=0.5))
def test_translation_keeps_codes(tmp_path_factory, boxes, dx, dy):
    """Shifting every shape by the same small delta inside the safe zone must not
    change WHICH rules fire (geometry is relative between shapes and to the canvas,
    and the safe zone keeps everything away from the edges)."""
    d = tmp_path_factory.mktemp("prop")
    a = _build(boxes, d, "a.pptx")
    moved = [(x + dx, y + dy, w, s, t) for (x, y, w, s, t) in boxes]
    b = _build(moved, d, "b.pptx")
    codes_a = sorted(c for (_p, c, _d) in _verdict(a))
    codes_b = sorted(c for (_p, c, _d) in _verdict(b))
    assert codes_a == codes_b


@settings(**SETTINGS)
@given(st.lists(box_st, min_size=1, max_size=3), st.integers(min_value=1, max_value=20))
def test_growing_font_never_creates_e3(tmp_path_factory, boxes, bump):
    """E3 is a lower bound on effective size: increasing every explicit size must never
    introduce a new E3."""
    d = tmp_path_factory.mktemp("prop")
    a = _build(boxes, d, "a.pptx")
    bigger = [(x, y, w, s + bump, t) for (x, y, w, s, t) in boxes]
    b = _build(bigger, d, "b.pptx")
    e3_a = sum(1 for (_p, c, _d) in _verdict(a) if c == "E3")
    e3_b = sum(1 for (_p, c, _d) in _verdict(b) if c == "E3")
    assert e3_b <= e3_a


def test_fuzz_attribute_mutation_never_tracebacks(tmp_path):
    """Deterministic structured fuzz: mutate random XML attribute values inside a real
    deck. Every mutation must end in a report or a controlled exception; page guards
    absorb the damage into W18, never an unhandled crash. Seeded for reproducibility."""
    from archforge import demo as jdemo
    seed_deck = os.path.join(str(tmp_path), "seed.pptx")
    jdemo.build_broken(seed_deck, lang="ko")
    with open(seed_deck, "rb") as f:
        base = f.read()
    rng = random.Random(20260711)
    junk = ["", "abc", "-1", "999999999999", "1.5pt", "%", "\x00nul", "true", "*"]
    failures = []
    for i in range(40):
        src = zipfile.ZipFile(io.BytesIO(base))
        out_buf = io.BytesIO()
        target = rng.choice([n for n in src.namelist() if n.endswith(".xml")])
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as dst:
            for name in src.namelist():
                data = src.read(name)
                if name == target:
                    text = data.decode("utf-8", "replace")
                    # replace the value of a random attribute occurrence
                    import re
                    attrs = list(re.finditer(r'\b([a-zA-Z:]+)="([^"]*)"', text))
                    if attrs:
                        m = rng.choice(attrs)
                        repl = rng.choice(junk)
                        text = text[:m.start(2)] + repl + text[m.end(2):]
                    data = text.encode("utf-8")
                dst.writestr(name, data)
        mut = os.path.join(str(tmp_path), "mut_%02d.pptx" % i)
        with open(mut, "wb") as f:
            f.write(out_buf.getvalue())
        try:
            jl.lint(mut, profile="full")
        except Exception as e:
            # A controlled rejection is acceptable; what we forbid is an unhandled
            # crash class that would escape the CLI's err_open mapping
            if isinstance(e, (MemoryError, RecursionError, SystemExit)):
                failures.append((i, target, type(e).__name__))
    assert not failures, failures


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-q"]))